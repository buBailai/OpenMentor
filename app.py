import os
import json
import re
import requests
import threading
import time
import sqlite3
import shutil
from functools import wraps
from jinja2 import Environment, FileSystemLoader, select_autoescape
import urllib.parse
from flask import Flask, render_template, redirect, url_for, request, flash, jsonify, make_response, send_file, send_from_directory, session, Response, abort
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, Date, Boolean, ForeignKey, UniqueConstraint, LargeBinary, inspect, text, func
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, relationship
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
import uuid
import random
import string
from datetime import datetime, timedelta
import io
import base64
from dotenv import load_dotenv
import logging
# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()  # 输出到控制台
    ]
)

logger = logging.getLogger(__name__)


# 加载环境变量
load_dotenv()

# 创建上传文件目录
UPLOAD_FOLDER = 'static/uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# 允许的文件扩展名
ALLOWED_EXTENSIONS = {'pdf', 'html', 'htm', 'jpg', 'zip'}
HTML_EDITOR_EXTENSIONS = {'html', 'htm'}
HTML_EDITOR_MAX_BYTES = 5 * 1024 * 1024
PROMPT_MATRIX_CACHE_PATH = os.path.join('data', 'prompt_matrix_cache.json')
PROMPT_MATRIX_CACHE_TTL_SECONDS = 60
_PROMPT_MATRIX_MEMORY_CACHE = {'ts': 0, 'items': None, 'status': None}
CHECKIN_CACHE_TTL_SECONDS = 60
_CHECKIN_MEMORY_CACHE = {'ts': 0, 'counts': None}
PUBLIC_PLAZA_CACHE_TTL_SECONDS = 60
_PUBLIC_PLAZA_MEMORY_CACHE = {'ts': 0, 'submissions': None}
CHECKIN_ALLOWED_PROVINCES = {
    '北京', '天津', '河北', '山西', '内蒙古',
    '辽宁', '吉林', '黑龙江',
    '上海', '江苏', '浙江', '安徽', '福建', '江西', '山东',
    '河南', '湖北', '湖南',
    '广东', '广西', '海南',
    '重庆', '四川', '贵州', '云南', '西藏',
    '陕西', '甘肃', '青海', '宁夏', '新疆',
    '香港', '澳门', '台湾',
    '海外',
}

app = Flask(__name__)
app.config['TEMPLATES_AUTO_RELOAD'] = True  # waitress 下也按 mtime 重读模板（覆盖模板文件即生效，无需重启）
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'dev_secret_key')
DATABASE_URL = 'sqlite:///openmentor.db'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB 全局上限（OpenMentor 单文件限 10MB 在端点处校验）
app.config['JSON_AS_ASCII'] = False  # 确保JSON响应中的中文正确显示，不转义为Unicode

APP_NAME = 'OpenMentor'
APP_EDITION = (os.getenv('OPENMENTOR_EDITION') or 'community').strip().lower()
if APP_EDITION not in {'community', 'campus'}:
    logger.warning(f"未知 OPENMENTOR_EDITION={APP_EDITION!r}，已回退为 community")
    APP_EDITION = 'community'


def is_campus_edition():
    return APP_EDITION == 'campus'


def _sqlite_db_path_from_url(database_url):
    if database_url.startswith('sqlite:///'):
        return database_url.replace('sqlite:///', '', 1)
    return None


DB_FILE_PATH = _sqlite_db_path_from_url(DATABASE_URL)
DB_EXISTED_BEFORE_START = bool(DB_FILE_PATH and os.path.exists(DB_FILE_PATH))


def _has_non_placeholder_files(path):
    if not os.path.exists(path):
        return False
    for root, dirs, files in os.walk(path):
        for name in files:
            if name not in {'.DS_Store', '.gitkeep'}:
                return True
    return False


def _guard_against_orphan_runtime_data():
    """防止 DB 缺失但附件/报告还在时静默创建空库，避免用户误以为数据被清空。"""
    if DB_EXISTED_BEFORE_START or not DB_FILE_PATH:
        return
    runtime_paths = ['static/uploads', 'static/reports']
    if any(_has_non_placeholder_files(p) for p in runtime_paths):
        raise RuntimeError(
            '检测到 openmentor.db 不存在，但 static/uploads 或 static/reports 中已有数据。'
            '为避免误创建空数据库，请先检查数据库文件或从备份恢复。'
        )


_guard_against_orphan_runtime_data()

# 初始化SQLAlchemy
engine = create_engine(DATABASE_URL, connect_args={'check_same_thread': False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# 注册模板全局变量
@app.template_global()
def get_app_name():
    return APP_NAME


_OM_VERSION_CACHE = {'mtime': 0, 'value': '1.2.1'}

def _parse_latest_version_from_changelog():
    """从 CHANGELOG.md 解析最新发布版本号（跳过 [Unreleased]），失败回退到上次缓存值。
    用作 footer / 全局版本号显示的 source of truth：每次发版只需新增一段 ## [X.Y.Z]。
    按 mtime 缓存，避免每次请求都做无用 IO；CHANGELOG 一改，下次请求即生效，无需重启服务。"""
    cl_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'CHANGELOG.md')
    try:
        m_time = os.path.getmtime(cl_path)
        if m_time == _OM_VERSION_CACHE['mtime']:
            return _OM_VERSION_CACHE['value']
        with open(cl_path, 'r', encoding='utf-8') as f:
            for line in f:
                m = re.match(r'^##\s*\[(\d+\.\d+\.\d+)\]', line.strip())
                if m:
                    _OM_VERSION_CACHE['mtime'] = m_time
                    _OM_VERSION_CACHE['value'] = m.group(1)
                    return m.group(1)
        # 文件读完没找到 → 缓存 mtime 防反复扫
        _OM_VERSION_CACHE['mtime'] = m_time
    except Exception:
        pass
    return _OM_VERSION_CACHE['value']


@app.context_processor
def _inject_om_version():
    """让所有模板都能用 {{ om_version }} 拿到最新版本号；动态读 CHANGELOG，改后无需重启"""
    return {'om_version': _parse_latest_version_from_changelog()}

# 用于存储分析任务进度的字典（在生产环境中应使用Redis等）
analysis_progress = {}
analysis_results = {}
# 用于跟踪已成功生成报告的任务ID，避免重复生成
completed_reports = set()
# 线程锁，确保对共享数据的安全访问
progress_lock = threading.Lock()

# 工具函数
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def save_uploaded_file(file):
    try:
        if file and allowed_file(file.filename):
            unique_filename = str(uuid.uuid4()) + '_' + file.filename
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            file.save(filepath)
            # 确保路径使用正斜杠，以便在URL中正确使用
            filepath = filepath.replace('\\', '/')
            return unique_filename, filepath
    except Exception as e:
        logger.error(f"保存文件失败: {str(e)}")
    return None, None

def _file_extension(filename):
    return filename.rsplit('.', 1)[1].lower() if filename and '.' in filename else ''

def _is_html_attachment(attachment):
    ext = _file_extension(attachment.file_name) or _file_extension(attachment.file_path)
    return ext in HTML_EDITOR_EXTENSIONS

def _resolve_upload_file_path(stored_path):
    """把附件路径限制在 static/uploads 下，避免在线编辑器读写任意文件。"""
    if not stored_path:
        return None

    upload_root = os.path.realpath(os.path.abspath(os.path.join(app.root_path, app.config['UPLOAD_FOLDER'])))
    if os.path.isabs(stored_path):
        candidate = os.path.realpath(stored_path)
    else:
        normalized = stored_path.replace('\\', '/').lstrip('/')
        if normalized.startswith('uploads/'):
            normalized = 'static/' + normalized
        candidate = os.path.realpath(os.path.abspath(os.path.join(app.root_path, normalized)))

    if candidate == upload_root or candidate.startswith(upload_root + os.sep):
        return candidate
    return None

def _read_editable_html_file(file_path):
    with open(file_path, 'rb') as f:
        raw = f.read(HTML_EDITOR_MAX_BYTES + 1)
    if len(raw) > HTML_EDITOR_MAX_BYTES:
        raise ValueError('HTML 文件超过在线编辑大小限制（5MB）')

    for encoding in ('utf-8-sig', 'utf-8', 'gb18030'):
        try:
            return raw.decode(encoding), encoding
        except UnicodeDecodeError:
            continue
    return raw.decode('utf-8', errors='replace'), 'utf-8'

def _builtin_prompt_matrix_records():
    return [
        {
            'enabled': True, 'sort': 10, 'category': '数据采集', 'key': 'simple',
            'title': '基础提交页面',
            'prompt': "创建完应用后，以POST方式向 '{{api_url}}' 发送数据",
            'note': '用于生成最基础的数据采集页面，只需要把页面里的表单数据提交到当前任务的 API 地址。',
            'bv': 'BV1WgGq65EHa', 'video_title': '基础提交页面制作视频演示', 'edition': '全部',
        },
        {
            'enabled': True, 'sort': 20, 'category': '数据采集', 'key': 'attachment',
            'title': '附件采集功能',
            'prompt': '创建带附件上传的 HTML 表单 → form 设 method="POST" 和 enctype="multipart/form-data"，附件用 <input type="file" name="字段名">；用 fetch + new FormData(form) 提交到 "{{api_url}}"，不要手动设 Content-Type',
            'note': '用于生成支持附件上传的采集页面。重点是使用 multipart/form-data、FormData，并且不要手动设置 Content-Type。',
            'bv': 'BV1ZrGi6PEiu', 'video_title': '附件采集功能制作视频演示', 'edition': '全部',
        },
        {
            'enabled': True, 'sort': 30, 'category': '数据采集', 'key': 'assistantWidget',
            'title': 'AI 导师悬浮球功能',
            'prompt': '请在这个 HTML 页面右下角增加一个「AI 导师」悬浮球：右下角固定显示圆形按钮，点击后弹出小窗；小窗顶部显示「AI 导师」和关闭按钮，主体用 iframe 嵌入 AI 导师学生访问链接（把链接替换成 "{{chat_url_example}}"）；桌面端小窗约 380px × 560px，手机端接近全屏宽度、高度约 75vh；不要破坏原页面表单、按钮、样式和 JavaScript；所有 CSS 和 JavaScript 写在当前 HTML 文件里；请输出完整 HTML 代码',
            'note': '用于在采集页面右下角加入 AI 导师悬浮球，把数据采集和 AI 导师答疑放到同一个学生页面里。',
            'bv': 'BV1N4Gi6yEgx', 'video_title': 'AI 导师悬浮球制作视频演示', 'edition': '全部',
        },
        {
            'enabled': True, 'sort': 40, 'category': '数据分析', 'key': 'apiDashboard',
            'title': '基础数据大屏',
            'prompt': '请通过数据接口（API 链接）读取数据，制作一个数据大屏 HTML 页面。\n\n数据接口地址：{{api_url}}\n\n请根据我提供的采集页面 HTML 文件了解各字段含义、字段名称和展示方式；页面打开后用 fetch(GET) 请求上面的 API 地址，读取返回数据并渲染统计卡片、图表和明细表。请处理空数据、字段缺失、接口请求失败，以及附件字段（如 _attachments）可能存在的情况。\n\n完整数据格式参考如下，请严格按照这个 JSON 结构读取字段：\n「完整数据格式替换」\n\n请输出一个可直接保存运行的完整 HTML 文件，CSS 和 JavaScript 写在同一个 HTML 文件里，不需要改动后端。',
            'note': '建议先提交 1-2 条测试数据，再打开数据接口链接，把页面里的完整数据内容复制出来，替换提示词里的「完整数据格式替换」。',
            'bv': 'BV1wYGi66EF5', 'video_title': '基础数据大屏制作视频演示', 'edition': '全部',
        },
        {
            'enabled': True, 'sort': 50, 'category': '数据分析', 'key': 'reveal',
            'title': '本地附件查看功能',
            'prompt': '想给附件加「📁 打开服务器文件夹」按钮（自部署场景免下载直接看）→ 发 POST 到 "{{attachment_reveal_url}}"，body 用 FormData 设 path=附件 url 去掉开头的 /（例：static/uploads/quickform/.../xxx.png）；会在跑 OpenMentor 那台电脑上弹出 Finder / 资源管理器并高亮该文件',
            'note': '用于自部署场景的数据大屏，让老师在运行 OpenMentor 的电脑上直接打开附件所在目录。',
            'bv': 'BV1zzGi6vEgQ', 'video_title': '本地附件查看功能制作视频演示', 'edition': '全部',
        },
    ]

def _prompt_matrix_cell_text(value):
    if value is None:
        return ''
    if isinstance(value, bool):
        return '是' if value else ''
    if isinstance(value, (int, float)):
        return str(value)
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, list):
        parts = []
        for item in value:
            if isinstance(item, dict):
                parts.append(str(item.get('text') or item.get('name') or item.get('value') or '').strip())
            else:
                parts.append(str(item).strip())
        return '，'.join([p for p in parts if p])
    if isinstance(value, dict):
        return str(value.get('text') or value.get('name') or value.get('value') or '').strip()
    return str(value).strip()

def _prompt_matrix_enabled(value):
    if isinstance(value, bool):
        return value
    text = _prompt_matrix_cell_text(value).lower()
    return text in {'1', 'true', 'yes', 'y', '是', '启用', '已启用', '勾选'}

def _normalize_prompt_matrix_record(fields):
    def pick(*names):
        for name in names:
            if name in fields:
                return fields.get(name)
        return None

    try:
        sort_value = int(float(_prompt_matrix_cell_text(pick('序号', '排序')) or '999'))
    except ValueError:
        sort_value = 999

    enabled_value = pick('启用', '是否启用')

    return {
        'enabled': True if enabled_value is None else _prompt_matrix_enabled(enabled_value),
        'sort': sort_value,
        'category': _prompt_matrix_cell_text(pick('分类')),
        'key': _prompt_matrix_cell_text(pick('场景Key', '场景 key', '场景KEY', 'key')),
        'title': _prompt_matrix_cell_text(pick('场景名称', '名称', '标题')),
        'prompt': _prompt_matrix_cell_text(pick('提示词内容', '提示词')),
        'note': _prompt_matrix_cell_text(pick('说明文字', '说明')),
        'bv': _prompt_matrix_cell_text(pick('B站BV号', 'BV号', 'B站 BV 号')),
        'video_title': _prompt_matrix_cell_text(pick('视频标题')),
        'edition': _prompt_matrix_cell_text(pick('适用版本')) or '全部',
    }

def _lark_prompt_matrix_configured():
    return all((os.getenv(k) or '').strip() for k in ('LARK_APP_ID', 'LARK_APP_SECRET', 'LARK_BASE_TOKEN', 'LARK_TABLE_ID'))

def _get_lark_tenant_access_token():
    app_id = (os.getenv('LARK_APP_ID') or '').strip()
    app_secret = (os.getenv('LARK_APP_SECRET') or '').strip()
    if not app_id or not app_secret:
        raise RuntimeError('同步应用凭证未配置')
    token_resp = requests.post(
        'https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal',
        json={'app_id': app_id, 'app_secret': app_secret},
        timeout=10,
    )
    token_resp.raise_for_status()
    token_data = token_resp.json()
    tenant_token = token_data.get('tenant_access_token')
    if not tenant_token:
        raise RuntimeError(token_data.get('msg') or '无法获取同步访问令牌')
    return tenant_token

def _fetch_lark_prompt_matrix_records():
    base_token = (os.getenv('LARK_BASE_TOKEN') or '').strip()
    table_id = (os.getenv('LARK_TABLE_ID') or '').strip()
    view_id = (os.getenv('LARK_VIEW_ID') or '').strip()

    tenant_token = _get_lark_tenant_access_token()
    params = {'page_size': 100}
    if view_id:
        params['view_id'] = view_id
    headers = {'Authorization': f'Bearer {tenant_token}'}
    url = f'https://open.feishu.cn/open-apis/bitable/v1/apps/{base_token}/tables/{table_id}/records'
    records = []
    page_token = ''
    while True:
        if page_token:
            params['page_token'] = page_token
        resp = requests.get(url, headers=headers, params=params, timeout=10)
        resp.raise_for_status()
        payload = resp.json()
        if payload.get('code') not in (0, None):
            raise RuntimeError(payload.get('msg') or '读取数据使用矩阵失败')
        data = payload.get('data') or {}
        for item in data.get('items') or []:
            records.append(_normalize_prompt_matrix_record(item.get('fields') or {}))
        if not data.get('has_more'):
            break
        page_token = data.get('page_token') or ''
        if not page_token:
            break
    return records

def _read_prompt_matrix_cache():
    if not os.path.exists(PROMPT_MATRIX_CACHE_PATH):
        return None
    try:
        with open(PROMPT_MATRIX_CACHE_PATH, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return data.get('items') if isinstance(data, dict) else None
    except Exception as e:
        logger.warning(f'[提示词矩阵] 读取缓存失败: {e}')
        return None

def _write_prompt_matrix_cache(items):
    try:
        os.makedirs(os.path.dirname(PROMPT_MATRIX_CACHE_PATH), exist_ok=True)
        with open(PROMPT_MATRIX_CACHE_PATH, 'w', encoding='utf-8') as f:
            json.dump({'updated_at': datetime.now().isoformat(), 'items': items}, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.warning(f'[提示词矩阵] 写入缓存失败: {e}')

def _load_prompt_matrix_records():
    now = time.time()
    if _PROMPT_MATRIX_MEMORY_CACHE['items'] and now - _PROMPT_MATRIX_MEMORY_CACHE['ts'] < PROMPT_MATRIX_CACHE_TTL_SECONDS:
        return _PROMPT_MATRIX_MEMORY_CACHE['items'], _PROMPT_MATRIX_MEMORY_CACHE['status']

    status = '内置预览数据'
    items = None
    if _lark_prompt_matrix_configured():
        try:
            items = _fetch_lark_prompt_matrix_records()
            status = f'已同步 · {datetime.now().strftime("%H:%M")}'
            _write_prompt_matrix_cache(items)
        except Exception as e:
            logger.warning(f'[提示词矩阵] 同步失败: {e}')
            cached = _read_prompt_matrix_cache()
            if cached:
                items = cached
                status = f'同步失败，已使用本地缓存：{e}'
            else:
                status = f'同步失败，已使用内置预览：{e}'

    if items is None:
        cached = _read_prompt_matrix_cache()
        if cached:
            items = cached
            status = '本地缓存数据'
        else:
            items = _builtin_prompt_matrix_records()
            if not _lark_prompt_matrix_configured():
                status = '内置预览数据（同步配置未完成）'

    _PROMPT_MATRIX_MEMORY_CACHE.update({'ts': now, 'items': items, 'status': status})
    return items, status

def _sync_prompt_matrix_from_lark():
    if not _lark_prompt_matrix_configured():
        raise RuntimeError('同步配置未完成')
    items = _fetch_lark_prompt_matrix_records()
    _write_prompt_matrix_cache(items)
    _PROMPT_MATRIX_MEMORY_CACHE.update({
        'ts': time.time(),
        'items': items,
        'status': f'手动同步 · {datetime.now().strftime("%H:%M")}',
    })
    return items

def _lark_checkin_configured():
    return all((os.getenv(k) or '').strip() for k in ('LARK_APP_ID', 'LARK_APP_SECRET', 'LARK_CHECKIN_BASE_TOKEN', 'LARK_CHECKIN_TABLE_ID'))

OM_LEGACY_CHECKIN_API_URL = 'https://quickform.cn/api/xrtzrtf5vq'

def _checkin_config():
    return {
        'base_token': (os.getenv('LARK_CHECKIN_BASE_TOKEN') or '').strip(),
        'table_id': (os.getenv('LARK_CHECKIN_TABLE_ID') or '').strip(),
        'view_id': (os.getenv('LARK_CHECKIN_VIEW_ID') or '').strip(),
        'field': (os.getenv('LARK_CHECKIN_FIELD') or '省份').strip() or '省份',
        'count_field': (os.getenv('LARK_CHECKIN_COUNT_FIELD') or '人数').strip() or '人数',
    }

def _checkin_records_url():
    cfg = _checkin_config()
    return f'https://open.feishu.cn/open-apis/bitable/v1/apps/{cfg["base_token"]}/tables/{cfg["table_id"]}/records'

def _normalize_checkin_province(value):
    province = _prompt_matrix_cell_text(value)
    return province if province in CHECKIN_ALLOWED_PROVINCES else ''

def _extract_checkin_province(record):
    if not isinstance(record, dict):
        return ''
    for key in ('province', '省份'):
        province = _normalize_checkin_province(record.get(key))
        if province:
            return province
    for key in ('data', 'fields', 'payload', 'form_data'):
        nested = record.get(key)
        if isinstance(nested, dict):
            province = _extract_checkin_province(nested)
            if province:
                return province
    return ''

def _create_legacy_checkin_record(province):
    resp = requests.post(OM_LEGACY_CHECKIN_API_URL, json={'province': province}, timeout=10)
    if not resp.ok:
        raise RuntimeError(resp.text[:200] or f'HTTP {resp.status_code}')
    _CHECKIN_MEMORY_CACHE.update({'ts': 0, 'counts': None})
    return True


def _checkin_parse_count(value):
    """人数字段兼容数字 / 文本 / 富文本数组三种存法，统一解析成 int。"""
    if isinstance(value, (int, float)):
        return int(value)
    text = _prompt_matrix_cell_text(value)
    try:
        return int(float(text)) if text else 0
    except ValueError:
        return 0


def _fetch_lark_checkin_rows():
    """读取聚合表全部行（每省一行）。返回 {省份: {'record_id', 'count', 'raw_count'}}。"""
    cfg = _checkin_config()
    headers = {'Authorization': f'Bearer {_get_lark_tenant_access_token()}'}
    params = {'page_size': 100}
    if cfg['view_id']:
        params['view_id'] = cfg['view_id']
    rows = {}
    page_token = ''
    while True:
        if page_token:
            params['page_token'] = page_token
        resp = requests.get(_checkin_records_url(), headers=headers, params=params, timeout=10)
        try:
            payload = resp.json()
        except (ValueError, json.JSONDecodeError):
            payload = {}
        if not resp.ok:
            raise RuntimeError(payload.get('msg') or payload.get('message') or resp.text[:200] or f'HTTP {resp.status_code}')
        if payload.get('code') not in (0, None):
            raise RuntimeError(payload.get('msg') or '读取打卡数据失败')
        data = payload.get('data') or {}
        for item in data.get('items') or []:
            fields = item.get('fields') or {}
            province = _normalize_checkin_province(fields.get(cfg['field']))
            if not province:
                continue
            raw = fields.get(cfg['count_field'])
            rows[province] = {
                'record_id': item.get('record_id') or item.get('id'),
                'count': _checkin_parse_count(raw),
                'raw_count': raw,
            }
        if not data.get('has_more'):
            break
        page_token = data.get('page_token') or ''
        if not page_token:
            break
    return rows


def _create_lark_checkin_record(province):
    """打卡 = 给对应省份行的「人数」+1（编辑行而非新增行，保证表恒定 ~35 行）。
    省份行不存在时兜底新建一行（人数=1）。"""
    if not _lark_checkin_configured():
        raise RuntimeError('打卡同步配置未完成')
    cfg = _checkin_config()
    headers = {
        'Authorization': f'Bearer {_get_lark_tenant_access_token()}',
        'Content-Type': 'application/json',
    }
    rows = _fetch_lark_checkin_rows()
    row = rows.get(province)
    if row and row.get('record_id'):
        new_count = row['count'] + 1
        # 数字字段 API 读出来可能是字符串，但写入必须是数字；文本字段则相反。
        # 先按数字写，遇到字段类型转换错误再回退成字符串。
        record_url = f"{_checkin_records_url()}/{row['record_id']}"
        resp = requests.put(
            record_url,
            headers=headers,
            json={'fields': {cfg['count_field']: new_count}},
            timeout=10,
        )
        if not resp.ok and 'FieldConvFail' in (resp.text or ''):
            resp = requests.put(
                record_url,
                headers=headers,
                json={'fields': {cfg['count_field']: str(new_count)}},
                timeout=10,
            )
    else:
        resp = requests.post(
            _checkin_records_url(),
            headers=headers,
            json={'fields': {cfg['field']: province, cfg['count_field']: 1}},
            timeout=10,
        )
        if not resp.ok and 'FieldConvFail' in (resp.text or ''):
            resp = requests.post(
                _checkin_records_url(),
                headers=headers,
                json={'fields': {cfg['field']: province, cfg['count_field']: '1'}},
                timeout=10,
            )
    try:
        payload = resp.json()
    except (ValueError, json.JSONDecodeError):
        payload = {}
    if not resp.ok:
        raise RuntimeError(payload.get('msg') or payload.get('message') or resp.text[:200] or f'HTTP {resp.status_code}')
    if payload.get('code') not in (0, None):
        raise RuntimeError(payload.get('msg') or '写入打卡数据失败')
    _CHECKIN_MEMORY_CACHE.update({'ts': 0, 'counts': None})
    return payload.get('data') or {}

def _create_checkin_record(province):
    if _lark_checkin_configured():
        return _create_lark_checkin_record(province)
    return _create_legacy_checkin_record(province)

def _fetch_legacy_checkin_counts():
    """旧 quickform 接口仍是一行一条记录，拉回后在服务端聚合成 {省份: 人数}。"""
    resp = requests.get(f'{OM_LEGACY_CHECKIN_API_URL}/all', timeout=10)
    try:
        payload = resp.json()
    except (ValueError, json.JSONDecodeError):
        payload = {}
    if not resp.ok:
        raise RuntimeError(payload.get('message') or resp.text[:200] or f'HTTP {resp.status_code}')
    records = payload.get('submissions') or payload.get('data') or payload.get('items') or []
    counts = {}
    for record in records:
        province = _extract_checkin_province(record)
        if province:
            counts[province] = counts.get(province, 0) + 1
    return counts


def _fetch_lark_checkin_counts():
    """聚合表模式：~35 行（每省一行 + 人数列），一次请求读完。"""
    if not _lark_checkin_configured():
        raise RuntimeError('打卡同步配置未完成')
    rows = _fetch_lark_checkin_rows()
    return {province: row['count'] for province, row in rows.items()}


def _load_checkin_counts():
    now = time.time()
    if _CHECKIN_MEMORY_CACHE['counts'] is not None and now - _CHECKIN_MEMORY_CACHE['ts'] < CHECKIN_CACHE_TTL_SECONDS:
        return _CHECKIN_MEMORY_CACHE['counts']
    counts = _fetch_lark_checkin_counts() if _lark_checkin_configured() else _fetch_legacy_checkin_counts()
    _CHECKIN_MEMORY_CACHE.update({'ts': now, 'counts': counts})
    return counts

def _render_prompt_matrix_placeholders(text, task):
    host_url = request.host_url
    host_root = request.host_url.rstrip('/')
    replacements = {
        '{{api_url}}': f'{host_root}/api/{task.task_id}',
        '{{host_url}}': host_url,
        '{{chat_url_example}}': f'{host_root}/chat/某个AI导师ID',
        '{{attachment_reveal_url}}': f'{host_root}{url_for("attachment_reveal")}',
    }
    for key, value in replacements.items():
        text = (text or '').replace(key, value)
    return text

def _prepare_prompt_matrix_for_task(task):
    raw_items, status = _load_prompt_matrix_records()
    edition_aliases = {
        'community': {'全部', '个人版', '开源版', 'community'},
        'campus': {'全部', '校园版', 'campus'},
    }.get(APP_EDITION, {'全部'})
    items = []
    used_keys = set()
    for idx, raw in enumerate(raw_items, 1):
        edition = (raw.get('edition') or '全部').strip()
        edition_parts = {part.strip() for part in re.split(r'[,，/、\s]+', edition) if part.strip()}
        if not raw.get('enabled') or (edition_parts and not edition_parts.intersection(edition_aliases)):
            continue
        key = re.sub(r'[^A-Za-z0-9_]+', '_', raw.get('key') or f'item_{idx}').strip('_') or f'item_{idx}'
        js_key = f'remote_{key}'
        while js_key in used_keys:
            js_key = f'{js_key}_{idx}'
        used_keys.add(js_key)
        item = {
            'js_key': js_key,
            'textarea_id': f'{js_key}_instruction',
            'note_id': f'{js_key}_note',
            'category': raw.get('category') or '其他',
            'title': raw.get('title') or raw.get('key') or '未命名场景',
            'prompt': _render_prompt_matrix_placeholders(raw.get('prompt') or '', task),
            'note': _render_prompt_matrix_placeholders(raw.get('note') or '', task),
            'bv': raw.get('bv') or '',
            'video_title': raw.get('video_title') or f'{raw.get("title") or "提示词"}视频演示',
            'sort': raw.get('sort', 999),
        }
        items.append(item)

    def sort_matrix_items(matrix_items):
        return sorted(matrix_items, key=lambda x: (x.get('sort', 999), x.get('title', '')))

    grouped = []
    for category in ['数据采集', '数据分析']:
        cat_items = sort_matrix_items([x for x in items if x['category'] == category])
        if cat_items:
            grouped.append({'name': category, 'items': cat_items})
    other_items = sort_matrix_items([x for x in items if x['category'] not in {'数据采集', '数据分析'}])
    if other_items:
        grouped.append({'name': '其他', 'items': other_items})

    page_items = [item for group in grouped for item in group['items']]
    scenarios = {x['js_key']: {'title': x['title'], 'instructionId': x['textarea_id'], 'noteId': x['note_id']} for x in page_items}
    videos = {x['js_key']: {'bv': x['bv'], 'title': x['video_title']} for x in page_items}
    return {'groups': grouped, 'items': page_items, 'scenarios': scenarios, 'videos': videos, 'status': status}

def generate_custom_id():
    """
    生成11位自定义ID：9位数字和字母组合 + 2位大写字母
    例如：oU59mLzPJPU
    """
    chars = string.ascii_letters + string.digits
    prefix = ''.join(random.choices(chars, k=9))
    suffix = ''.join(random.choices(string.ascii_uppercase, k=2))
    return prefix + suffix

# 数据库模型
class User(UserMixin, Base):
    __tablename__ = 'user'
    id = Column(Integer, primary_key=True)
    username = Column(String(50), unique=True, nullable=False)
    email = Column(String(100), unique=True, nullable=False)
    password = Column(String(200), nullable=False)
    role = Column(String(20), default='teacher', nullable=False)  # admin / teacher
    is_active_flag = Column(Boolean, default=True, nullable=False)
    display_name = Column(String(100))
    force_password_change = Column(Boolean, default=False, nullable=False)
    last_login_at = Column(DateTime)
    created_at = Column(DateTime, default=datetime.now)
    tasks = relationship('Task', back_populates='author')
    ai_config = relationship('AIConfig', back_populates='user', uselist=False)

    @property
    def is_active(self):
        return bool(self.is_active_flag)

class Task(Base):
    __tablename__ = 'task'
    id = Column(Integer, primary_key=True)
    title = Column(String(200), nullable=False)
    description = Column(Text)
    created_at = Column(DateTime, default=datetime.now)
    user_id = Column(Integer, ForeignKey('user.id'))
    author = relationship('User', back_populates='tasks')
    submission = relationship('Submission', back_populates='task', cascade='all, delete-orphan')
    attachments = relationship('Attachment', back_populates='task', cascade='all, delete-orphan')
    task_id = Column(String(11), unique=True, default=generate_custom_id)
    analysis_report = Column(Text)
    report_file_path = Column(String(500))
    report_generated_at = Column(DateTime)

    # ==================== OpenMentor 扩展字段（V1）====================
    # 兼容性原则：所有新增字段均带默认值，不影响现有 QuickForm 任务的语义
    # is_assistant=False 表示老 QuickForm 任务；True 表示 AI 导师（OpenMentor 新增）
    is_assistant = Column(Boolean, default=False, nullable=False)
    system_prompt = Column(Text)                           # AI 导师的系统提示词
    welcome_message = Column(Text)                         # 学生进入对话的欢迎语
    selected_model_name = Column(String(50))               # 该助学使用的模型（覆盖用户全局设置）
    allow_image_input = Column(Boolean, default=True)      # 允许学生上传图片
    allow_file_upload = Column(Boolean, default=True)      # 允许学生上传文件
    allow_voice_input = Column(Boolean, default=True)      # 允许学生语音输入（需老师配置硅基流动 key 做转写）
    allow_tts = Column(Boolean, default=False)             # 允许学生朗读 AI 回复（默认关；Edge TTS 免费引擎）
    audio_mode = Column(String(20), default='transcribe')  # 语音处理：transcribe=转写成文字 / direct=原声直达模型（音频理解）
    allow_image_generation = Column(Boolean, default=False)
    allow_video_generation = Column(Boolean, default=False)  # 允许 AI 生成视频（豆包 seedance，异步任务）
    allow_call_teacher = Column(Boolean, default=False)     # 学生「找老师」按钮（掌上讲台 App 配套，默认关） # 允许 AI 生成图片
    share_mode = Column(String(20), default='independent')  # independent / group(V2)
    group_max_size = Column(Integer, default=6)            # 小组人数上限（V2）
    max_messages_per_student_daily = Column(Integer, default=50)  # 每生每日消息上限
    link_expires_at = Column(DateTime)                     # 链接有效期截止（max +7 天）
    blocked_keywords = Column(Text)                        # 关键词黑名单（JSON）
    prompt_drafts = Column(Text)                           # AI 辅助生成的提示词历史草稿（JSON, 最多 5 条）
    avatar_path = Column(String(500))                      # AI 导师头像（相对 static 路径，可选）
    roster_mode = Column(String(20), default='off')        # off=自由填写 / strict=必须从花名册中选择
    status = Column(String(20), default='active')          # active / disabled / expired
    template_key = Column(String(50))                      # 内置模板任务标识（如 artwork_gallery），用于去重/识别
    pinned = Column(Boolean, default=False)                # 置顶（数据任务列表排前）
    conversations = relationship('Conversation', back_populates='assistant', cascade='all, delete-orphan')

class Attachment(Base):
    __tablename__ = 'attachment'
    id = Column(Integer, primary_key=True)
    task_id = Column(Integer, ForeignKey('task.id'), nullable=False)
    task = relationship('Task', back_populates='attachments')
    file_name = Column(String(200), nullable=False)
    file_path = Column(String(500), nullable=False)
    created_at = Column(DateTime, default=datetime.now)

class Submission(Base):
    __tablename__ = 'submission'
    id = Column(Integer, primary_key=True)
    task_id = Column(Integer, ForeignKey('task.id'))
    task = relationship('Task', back_populates='submission')
    data = Column(Text, nullable=False)
    submitted_at = Column(DateTime, default=datetime.now)

class AIConfig(Base):
    __tablename__ = 'ai_config'
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey('user.id'), unique=True)
    user = relationship('User', back_populates='ai_config')
    selected_model = Column(String(50), default='deepseek')
    model_configs = relationship('AIModelConfig', back_populates='ai_config', cascade='all, delete-orphan')

class AIModelConfig(Base):
    __tablename__ = 'ai_model_config'
    __table_args__ = (
        UniqueConstraint('ai_config_id', 'model_name', name='uq_ai_model_config_scope'),
    )
    id = Column(Integer, primary_key=True)
    ai_config_id = Column(Integer, ForeignKey('ai_config.id'))
    ai_config = relationship('AIConfig', back_populates='model_configs')
    model_name = Column(String(50))           # 提供商代号：deepseek/doubao/qwen/glm/siliconflow/ollama
    api_key = Column(String(200))
    api_url = Column(String(200))
    extra_settings = Column(Text)             # 对话模型名（如 qwen-vl-plus / glm-4v）
    image_gen_model = Column(String(100))     # OpenMentor: 图片生成模型名（如 doubao-seedream-3-0-t2i-250415）
    video_gen_model = Column(String(100))     # OpenMentor: 视频生成模型名（如 doubao-seedance-1-5-pro-251215，目前仅豆包）

class QFConfig(Base):
    __tablename__ = 'qf_config'
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey('user.id'), unique=True)
    user = relationship('User', back_populates='qf_config')
    username = Column(String(100))
    password = Column(String(200))

User.qf_config = relationship('QFConfig', back_populates='user', uselist=False)

# ==================== OpenMentor 新增数据表 ====================

class Conversation(Base):
    """OpenMentor: 学生与 AI 导师的会话（一个学生一个会话）"""
    __tablename__ = 'conversation'
    id = Column(Integer, primary_key=True)
    assistant_id = Column(Integer, ForeignKey('task.id'), nullable=False)
    assistant = relationship('Task', back_populates='conversations')
    student_class = Column(String(100), nullable=False)
    student_name = Column(String(100), nullable=False)
    student_token = Column(String(64), nullable=False, index=True)  # hash(class+name)，用于 localStorage 识别
    started_at = Column(DateTime, default=datetime.now)
    last_active_at = Column(DateTime, default=datetime.now)
    daily_message_count = Column(Integer, default=0)                # 今日已发条数
    daily_reset_date = Column(Date)                                  # 上次重置日期
    total_tokens = Column(Integer, default=0)                        # 累计消耗 tokens
    is_blocked = Column(Boolean, default=False)                      # 是否被老师手动封禁
    messages = relationship(
        'Message',
        back_populates='conversation',
        cascade='all, delete-orphan',
        order_by='Message.created_at'
    )


class RosterEntry(Base):
    """OpenMentor: 班级花名册（属于某位老师，跨 AI 导师共享）"""
    __tablename__ = 'roster_entry'
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey('user.id'), nullable=False, index=True)
    class_name = Column(String(100), nullable=False)
    student_name = Column(String(100), nullable=False)
    student_no = Column(String(50))                  # 学号（可选）
    notes = Column(String(200))                      # 备注（可选）
    created_at = Column(DateTime, default=datetime.now)
    __table_args__ = (
        UniqueConstraint('user_id', 'class_name', 'student_name', name='uq_roster_user_class_name'),
    )


class StudentReport(Base):
    """B14: AI 生成的学情报告（单学生 / 全班）持久化。"""
    __tablename__ = 'student_report'
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey('user.id'), nullable=False, index=True)
    assistant_id = Column(Integer, ForeignKey('task.id'), nullable=False, index=True)
    report_type = Column(String(20), nullable=False)        # 'student' / 'class'
    student_class = Column(String(100))                     # 单学生：班级；班级报告：班级筛选（可空）
    student_name = Column(String(100))                      # 单学生：姓名；班级报告：空
    scope_label = Column(String(200))                       # 显示标签（如「四年2班 · 张三」「全部学生」）
    content = Column(Text, nullable=False)                  # Markdown 全文
    model_used = Column(String(50))
    student_count = Column(Integer)
    message_count = Column(Integer)
    omitted = Column(Integer)
    created_at = Column(DateTime, default=datetime.now, index=True)


class ClassOption(Base):
    """B9: 自由填写模式下的"年级 / 班级"下拉选项（仅 roster_mode=off 使用）
    每位老师维护一份年级+班级清单，避免学生手填的班级名五花八门，便于聚合。
    grade 可空（不区分年级时只用班级一层）。
    """
    __tablename__ = 'class_option'
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey('user.id'), nullable=False, index=True)
    grade = Column(String(50))                       # 年级（可空，例：高一 / 三年级）
    class_name = Column(String(100), nullable=False) # 班级（例：1班 / (3)班）
    sort_order = Column(Integer, default=0)
    created_at = Column(DateTime, default=datetime.now)
    __table_args__ = (
        UniqueConstraint('user_id', 'grade', 'class_name', name='uq_class_option_user_grade_class'),
    )


class CampusPlazaTemplate(Base):
    """校园版：校内 AI 导师广场模板。只保存配置，不保存 API Key 或学生数据。"""
    __tablename__ = 'campus_plaza_template'
    id = Column(Integer, primary_key=True)
    source_assistant_id = Column(Integer, ForeignKey('task.id'))
    owner_user_id = Column(Integer, ForeignKey('user.id'), nullable=False, index=True)
    title = Column(String(200), nullable=False)
    subject = Column(String(30), nullable=False)
    grade = Column(String(20), nullable=False)
    purpose = Column(String(60), nullable=False)
    description = Column(Text)
    nickname = Column(String(30))
    school = Column(String(60))
    config_json = Column(Text, nullable=False)
    status = Column(String(20), default='active', nullable=False)
    created_at = Column(DateTime, default=datetime.now, index=True)
    updated_at = Column(DateTime, default=datetime.now)


class Message(Base):
    """OpenMentor: 对话中的一条消息"""
    __tablename__ = 'message'
    id = Column(Integer, primary_key=True)
    conversation_id = Column(Integer, ForeignKey('conversation.id'), nullable=False)
    conversation = relationship('Conversation', back_populates='messages')
    role = Column(String(20), nullable=False)                # user / assistant / system
    content = Column(Text, nullable=False)                    # 文本内容
    image_paths = Column(Text)                                # 附带图片路径（JSON 数组）
    file_paths = Column(Text)                                 # 附带文件路径（JSON 数组）
    audio_paths = Column(Text)                                # 附带语音/音频路径（JSON 数组）
    generated_image_path = Column(String(500))                # AI 生成的图片路径
    generated_video_path = Column(String(500))                # AI 生成的视频路径
    tokens_used = Column(Integer, default=0)                  # 本条消息消耗的 tokens
    triggered_keyword = Column(String(200))                   # 命中的黑名单关键词（如有）
    rating = Column(Integer)                                  # 学生评分：1=👍 / -1=👎 / NULL=未评
    created_at = Column(DateTime, default=datetime.now, index=True)


class KnowledgeDoc(Base):
    """OpenMentor: AI 导师知识库文档（RAG）"""
    __tablename__ = 'knowledge_doc'
    id = Column(Integer, primary_key=True)
    assistant_id = Column(Integer, ForeignKey('task.id'), nullable=False, index=True)
    file_name = Column(String(300), nullable=False)           # 原始文件名（展示用）
    file_path = Column(String(500))                           # 落盘路径（保留原件供下载/重建）
    status = Column(String(20), default='ready')              # ready / failed
    error_msg = Column(String(500))
    chunk_count = Column(Integer, default=0)
    created_at = Column(DateTime, default=datetime.now)
    chunks = relationship('KnowledgeChunk', back_populates='doc', cascade='all, delete-orphan')


class KnowledgeChunk(Base):
    """OpenMentor: 知识库切片 + 向量（float32 BLOB，检索用暴力余弦）"""
    __tablename__ = 'knowledge_chunk'
    id = Column(Integer, primary_key=True)
    doc_id = Column(Integer, ForeignKey('knowledge_doc.id'), nullable=False, index=True)
    doc = relationship('KnowledgeDoc', back_populates='chunks')
    assistant_id = Column(Integer, nullable=False, index=True)  # 冗余存一份，检索时免 join
    chunk_index = Column(Integer, default=0)
    content = Column(Text, nullable=False)
    embedding = Column(LargeBinary)                            # float32 数组打包（array('f').tobytes()）


class AppDevice(Base):
    """OpenMentor: 掌上讲台 App 配对设备（长效令牌 + 推送 token）"""
    __tablename__ = 'app_device'
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey('user.id'), nullable=False, index=True)
    device_name = Column(String(100))
    platform = Column(String(20))                 # ios / android
    push_token = Column(String(200))              # APNs/厂商推送 token（推送功能上线后使用）
    token_hash = Column(String(64), index=True)   # 长效令牌的 sha256（不存明文）
    revoked = Column(Boolean, default=False)
    created_at = Column(DateTime, default=datetime.now)
    last_seen_at = Column(DateTime, default=datetime.now)


# 创建数据库表
Base.metadata.create_all(engine)


def _quote_ident(name):
    return '"' + name.replace('"', '""') + '"'


def _migration_lock_path():
    os.makedirs('backups', exist_ok=True)
    return os.path.join('backups', '.openmentor_migration.lock')


def _acquire_migration_lock():
    """简易跨平台迁移锁，防止两个启动器同时改库。"""
    lock_path = _migration_lock_path()
    now = time.time()
    try:
        if os.path.exists(lock_path):
            age = now - os.path.getmtime(lock_path)
            if age > 30 * 60:
                logger.warning('[OpenMentor 迁移] 发现超过 30 分钟的旧锁，自动移除')
                os.remove(lock_path)
        fd = os.open(lock_path, os.O_CREAT | os.O_EXCL | os.O_WRONLY)
        os.write(fd, f'pid={os.getpid()} time={datetime.now().isoformat()}\n'.encode('utf-8'))
        return fd, lock_path
    except FileExistsError:
        raise RuntimeError('数据库迁移锁已存在，可能有另一个 OpenMentor 正在启动或迁移。请稍后再试。')


def _release_migration_lock(lock_info):
    fd, lock_path = lock_info
    try:
        os.close(fd)
    except Exception:
        pass
    try:
        if os.path.exists(lock_path):
            os.remove(lock_path)
    except Exception as e:
        logger.warning(f'[OpenMentor 迁移] 删除迁移锁失败: {e}')


def _backup_sqlite_database(reason='migration'):
    if not DB_FILE_PATH or not os.path.exists(DB_FILE_PATH):
        return None
    os.makedirs('backups', exist_ok=True)
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    backup_path = os.path.join('backups', f'openmentor_{ts}_before_{reason}.db')
    src = sqlite3.connect(DB_FILE_PATH)
    dst = sqlite3.connect(backup_path)
    try:
        src.backup(dst)
    finally:
        dst.close()
        src.close()
    logger.info(f'[OpenMentor 迁移] 已备份数据库到 {backup_path}')
    return backup_path


def _ensure_schema_migration_table(conn):
    conn.execute(text("""
        CREATE TABLE IF NOT EXISTS schema_migration (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            namespace VARCHAR(50) NOT NULL,
            version VARCHAR(100) NOT NULL,
            applied_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            notes TEXT,
            UNIQUE(namespace, version)
        )
    """))


def _mark_migration_applied(conn, namespace, version, notes=''):
    conn.execute(
        text("""
            INSERT OR IGNORE INTO schema_migration (namespace, version, notes)
            VALUES (:namespace, :version, :notes)
        """),
        {'namespace': namespace, 'version': version, 'notes': notes},
    )


def _sqlite_index_columns(conn, index_name):
    rows = conn.exec_driver_sql(f"PRAGMA index_info({_quote_ident(index_name)})").fetchall()
    return [row[2] for row in rows]


def _ai_model_config_needs_scope_migration(conn, table_names=None):
    table_names = table_names or inspect(engine).get_table_names()
    if 'ai_model_config' not in table_names:
        return False

    table_sql_row = conn.execute(text("""
        SELECT sql FROM sqlite_master
        WHERE type = 'table' AND name = 'ai_model_config'
    """)).fetchone()
    table_sql = table_sql_row[0] if table_sql_row else ''

    has_global_model_unique = 'UNIQUE (model_name)' in table_sql
    has_scoped_unique = False
    for row in conn.exec_driver_sql("PRAGMA index_list('ai_model_config')").fetchall():
        # PRAGMA index_list: seq, name, unique, origin, partial
        if not row[2]:
            continue
        cols = _sqlite_index_columns(conn, row[1])
        if cols == ['model_name']:
            has_global_model_unique = True
        if cols == ['ai_config_id', 'model_name']:
            has_scoped_unique = True

    return has_global_model_unique or not has_scoped_unique


def _migrate_ai_model_config_scope(conn, table_names):
    """把模型配置唯一性从全局 model_name 改为 ai_config_id + model_name。"""
    if not _ai_model_config_needs_scope_migration(conn, table_names):
        return False

    conn.execute(text("""
        CREATE TABLE ai_model_config_new (
            id INTEGER NOT NULL PRIMARY KEY,
            ai_config_id INTEGER,
            model_name VARCHAR(50),
            api_key VARCHAR(200),
            api_url VARCHAR(200),
            extra_settings TEXT,
            image_gen_model VARCHAR(100),
            FOREIGN KEY(ai_config_id) REFERENCES ai_config (id)
        )
    """))
    conn.execute(text("""
        INSERT INTO ai_model_config_new
            (id, ai_config_id, model_name, api_key, api_url, extra_settings, image_gen_model)
        SELECT
            id, ai_config_id, model_name, api_key, api_url, extra_settings, image_gen_model
        FROM ai_model_config
    """))
    conn.execute(text("DROP TABLE ai_model_config"))
    conn.execute(text("ALTER TABLE ai_model_config_new RENAME TO ai_model_config"))
    conn.execute(text("""
        CREATE UNIQUE INDEX IF NOT EXISTS uq_ai_model_config_scope
        ON ai_model_config (ai_config_id, model_name)
    """))
    return True


def migrate_database():
    """OpenMentor 非破坏性数据库迁移。

    规则：
    - 只追加字段/表，不覆盖 openmentor.db
    - 真正执行结构变更前自动备份 SQLite
    - migration 需要幂等，可重复启动
    """
    inspector = inspect(engine)
    table_names = inspector.get_table_names()

    plans = []
    if 'user' in table_names:
        plans.append(('user', [
            ('role', "VARCHAR(20) DEFAULT 'teacher' NOT NULL"),
            ('is_active_flag', "BOOLEAN DEFAULT 1 NOT NULL"),
            ('display_name', "VARCHAR(100)"),
            ('force_password_change', "BOOLEAN DEFAULT 0 NOT NULL"),
            ('last_login_at', "DATETIME"),
        ]))
    if 'task' in table_names:
        plans.append(('task', [
            ('is_assistant', "BOOLEAN DEFAULT 0"),
            ('system_prompt', "TEXT"),
            ('welcome_message', "TEXT"),
            ('selected_model_name', "VARCHAR(50)"),
            ('allow_image_input', "BOOLEAN DEFAULT 1"),
            ('allow_file_upload', "BOOLEAN DEFAULT 1"),
            ('allow_image_generation', "BOOLEAN DEFAULT 0"),
            ('share_mode', "VARCHAR(20) DEFAULT 'independent'"),
            ('group_max_size', "INTEGER DEFAULT 6"),
            ('max_messages_per_student_daily', "INTEGER DEFAULT 50"),
            ('link_expires_at', "DATETIME"),
            ('blocked_keywords', "TEXT"),
            ('prompt_drafts', "TEXT"),
            ('avatar_path', "VARCHAR(500)"),
            ('roster_mode', "VARCHAR(20) DEFAULT 'off'"),
            ('status', "VARCHAR(20) DEFAULT 'active'"),
            ('allow_voice_input', "BOOLEAN DEFAULT 1"),
            ('allow_tts', "BOOLEAN DEFAULT 0"),
            ('audio_mode', "VARCHAR(20) DEFAULT 'transcribe'"),
            ('allow_video_generation', "BOOLEAN DEFAULT 0"),
            ('allow_call_teacher', "BOOLEAN DEFAULT 0"),
            ('template_key', "VARCHAR(50)"),
            ('pinned', "BOOLEAN DEFAULT 0"),
        ]))
    if 'ai_model_config' in table_names:
        plans.append(('ai_model_config', [
            ('image_gen_model', "VARCHAR(100)"),
            ('video_gen_model', "VARCHAR(100)"),
        ]))
    if 'message' in table_names:
        plans.append(('message', [
            ('rating', "INTEGER"),
            ('audio_paths', "TEXT"),
            ('generated_video_path', "VARCHAR(500)"),
        ]))

    pending_cols = []
    for table_name, new_cols in plans:
        existing_cols = {c['name'] for c in inspector.get_columns(table_name)}
        for col_name, col_def in new_cols:
            if col_name not in existing_cols:
                pending_cols.append((table_name, col_name, col_def))

    schema_table_missing = 'schema_migration' not in table_names
    with engine.connect() as conn:
        ai_model_scope_migration_needed = _ai_model_config_needs_scope_migration(conn, table_names)

    if not pending_cols and not schema_table_missing and not ai_model_scope_migration_needed:
        logger.info("[OpenMentor 迁移] 表结构已是最新，无需迁移")
        return

    lock_info = _acquire_migration_lock()
    try:
        backup_path = _backup_sqlite_database('migration') if DB_EXISTED_BEFORE_START else None
        total_added = []
        with engine.begin() as conn:
            _ensure_schema_migration_table(conn)
            for table_name, col_name, col_def in pending_cols:
                conn.execute(text(
                    f"ALTER TABLE {_quote_ident(table_name)} ADD COLUMN {_quote_ident(col_name)} {col_def}"
                ))
                total_added.append(f'{table_name}.{col_name}')

            if _migrate_ai_model_config_scope(conn, table_names):
                total_added.append('ai_model_config scoped unique index')
                _mark_migration_applied(
                    conn,
                    'core',
                    '20260522-ai-model-config-scope',
                    'AI model config uniqueness changed from model_name to ai_config_id + model_name',
                )

            if 'user' in table_names:
                # 旧个人版升级：旧用户默认启用；如果没有 admin，最早创建的用户成为 admin/owner，避免升级后锁死。
                conn.execute(text("UPDATE \"user\" SET role = 'teacher' WHERE role IS NULL OR role = ''"))
                conn.execute(text("UPDATE \"user\" SET is_active_flag = 1 WHERE is_active_flag IS NULL"))
                conn.execute(text("UPDATE \"user\" SET force_password_change = 0 WHERE force_password_change IS NULL"))
                admin_count = conn.execute(text("SELECT COUNT(*) FROM \"user\" WHERE role = 'admin'")).scalar() or 0
                if admin_count == 0:
                    first_id = conn.execute(text("SELECT id FROM \"user\" ORDER BY id ASC LIMIT 1")).scalar()
                    if first_id:
                        conn.execute(text("UPDATE \"user\" SET role = 'admin' WHERE id = :id"), {'id': first_id})

            _mark_migration_applied(conn, 'core', '2026-05-22-user-role-and-safe-migration', f'backup={backup_path or ""}')

        if total_added:
            logger.info(f"[OpenMentor 迁移] 新增 {len(total_added)} 个列: {', '.join(total_added)}")
        else:
            logger.info("[OpenMentor 迁移] 已初始化 schema_migration，无新增列")
    except Exception:
        logger.exception('[OpenMentor 迁移] 迁移失败，服务将停止启动')
        raise
    finally:
        _release_migration_lock(lock_info)


migrate_database()


def _ensure_default_admin_user():
    """新安装或纯代码 clone 后，如果库里没有用户，自动创建默认 admin。"""
    db = SessionLocal()
    try:
        user_count = db.query(User).count()
        if user_count > 0:
            return
        admin = User(
            username='admin',
            email='admin@openmentor.local',
            password=generate_password_hash('openmentor'),
            role='admin',
            display_name='管理员',
            is_active_flag=True,
            force_password_change=True,
        )
        db.add(admin)
        db.commit()
        logger.info("[OpenMentor 初始化] 已创建默认管理员账号 admin / openmentor")
    except Exception:
        db.rollback()
        logger.exception('[OpenMentor 初始化] 创建默认管理员账号失败')
        raise
    finally:
        db.close()


_ensure_default_admin_user()


# ============ 内置模板任务（随安装/升级自动下发，置顶数据任务） ============
# 模板 HTML 存放于 static/task_templates/<key>/，含占位符 __OM_TASK_API_ID__；
# 下发时为每位用户新建一个任务，把占位符替换成该任务自己的 API 号，附件复制到 uploads。
# HTML 内 API 基址用 window.location.origin，任何 IP/域名安装后零配置自动连本机接口。
TEMPLATE_TASKS = [
    {
        'key': 'artwork_gallery',
        'title': '作品采集与展览',
        'description': '一站式作品流程：学生扫码上传作品 → 教师数据大屏实时统计 → 作品展览画廊展出。'
                       '挂载的三个页面均自适应手机/平板，接口自动匹配本机地址，开箱即用。',
        'dir': 'artwork_gallery',
        'files': [
            ('01_collect.html', '作品实时采集.html'),
            ('02_dashboard.html', '作品采集数据大屏.html'),
            ('03_gallery.html', '作品展览.html'),
        ],
    },
]


def _provision_template_task(db, user_id, tpl):
    """为某用户下发一个模板任务（幂等：已存在同 template_key 则跳过）。返回 True=新建。"""
    exists = db.query(Task).filter_by(user_id=user_id, template_key=tpl['key']).first()
    if exists:
        return False
    seed_dir = os.path.join(app.root_path, 'static', 'task_templates', tpl['dir'])
    if not os.path.isdir(seed_dir):
        logger.warning(f"[模板任务] 种子目录缺失，跳过：{seed_dir}")
        return False
    task = Task(
        title=tpl['title'],
        description=tpl['description'],
        user_id=user_id,
        is_assistant=False,
        template_key=tpl['key'],
        pinned=True,
    )
    db.add(task)
    db.flush()  # 取 task.task_id（API 号）
    for seed_name, display_name in tpl['files']:
        seed_path = os.path.join(seed_dir, seed_name)
        if not os.path.exists(seed_path):
            continue
        with open(seed_path, encoding='utf-8') as f:
            html = f.read().replace('__OM_TASK_API_ID__', task.task_id)
        unique = f"{uuid.uuid4()}_{display_name}"
        out_rel = f"static/uploads/{unique}"
        out_abs = os.path.join(app.root_path, out_rel)
        with open(out_abs, 'w', encoding='utf-8') as f:
            f.write(html)
        db.add(Attachment(task_id=task.id, file_name=display_name, file_path=out_rel))
    db.commit()
    logger.info(f"[模板任务] 已为用户 {user_id} 下发「{tpl['title']}」(task {task.task_id})")
    return True


def _ensure_template_tasks_for_user(user_id):
    """确保某用户拥有全部内置模板任务（懒下发，列表查看时调用）。"""
    db = SessionLocal()
    try:
        for tpl in TEMPLATE_TASKS:
            try:
                _provision_template_task(db, user_id, tpl)
            except Exception:
                db.rollback()
                logger.exception(f"[模板任务] 下发失败 user={user_id} key={tpl.get('key')}")
    finally:
        db.close()


def _ensure_template_tasks_all_users():
    """启动时为所有现有用户补齐模板任务（安装/升级自动拥有）。"""
    db = SessionLocal()
    try:
        user_ids = [u.id for u in db.query(User.id).all()]
    finally:
        db.close()
    for uid in user_ids:
        _ensure_template_tasks_for_user(uid)


try:
    _ensure_template_tasks_all_users()
except Exception:
    logger.exception('[模板任务] 启动下发失败（已忽略，不影响启动）')


# OpenMentor 调试：捕获所有 400/413 异常并打印详情
from werkzeug.exceptions import HTTPException, BadRequest, RequestEntityTooLarge

@app.errorhandler(BadRequest)
def _handle_bad_request(e):
    logger.warning(f'[400 拦截] path={request.path}, method={request.method}, content_type={request.content_type!r}, content_length={request.content_length}, headers={dict(request.headers)}, description={e.description!r}')
    return jsonify({'code': 400, 'message': f'BadRequest: {e.description}'}), 400

@app.errorhandler(RequestEntityTooLarge)
def _handle_too_large(e):
    logger.warning(f'[413 拦截] path={request.path}, content_length={request.content_length}, MAX={app.config.get("MAX_CONTENT_LENGTH")}')
    return jsonify({'code': 413, 'message': f'文件过大（请求体 {request.content_length} 字节超过 Flask 全局上限）'}), 413




# 初始化Flask-Login
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

# 使用werkzeug.security进行密码加密（无需初始化）

def is_admin_user(user=None):
    user = user or current_user
    return bool(getattr(user, 'is_authenticated', False) and getattr(user, 'role', '') == 'admin')


def _is_safe_next_url(target):
    """只允许登录后跳回本站内部地址，避免开放重定向。"""
    if not target:
        return False
    ref = urllib.parse.urlparse(request.host_url)
    test = urllib.parse.urlparse(urllib.parse.urljoin(request.host_url, target))
    return test.scheme in ('http', 'https') and ref.netloc == test.netloc


def _safe_next_url(target, default_endpoint='dashboard'):
    return target if _is_safe_next_url(target) else url_for(default_endpoint)


def campus_admin_required(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        if not is_campus_edition():
            abort(404)
        if not current_user.is_authenticated:
            return login_manager.unauthorized()
        if not is_admin_user(current_user):
            flash('需要校园版管理员权限', 'danger')
            return redirect(url_for('dashboard'))
        return func(*args, **kwargs)
    return wrapper


@app.context_processor
def _inject_edition_flags():
    return {
        'om_edition': APP_EDITION,
        'om_is_campus': is_campus_edition(),
        'om_is_campus_admin': is_campus_edition() and is_admin_user(current_user),
    }


@login_manager.user_loader
def load_user(user_id):
    db = SessionLocal()
    try:
        return db.query(User).get(int(user_id))
    finally:
        db.close()


@app.before_request
def _enforce_account_state():
    if not current_user.is_authenticated:
        return None

    endpoint = request.endpoint or ''
    allowed = {'logout', 'profile', 'static', 'pwa_manifest', 'pwa_service_worker'}

    if not current_user.is_active:
        logout_user()
        flash('该账号已被禁用，请联系管理员', 'danger')
        return redirect(url_for('login'))

    if getattr(current_user, 'force_password_change', False) and endpoint not in allowed:
        next_target = request.full_path if request.query_string else request.path
        flash('请先修改临时密码', 'warning')
        return redirect(url_for('profile', active_tab='password', next=next_target))

    return None

def read_file_content(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                return f"二进制文件 (大小: {len(content)} 字节)"
        except Exception as e:
            logger.error(f"读取文件内容失败: {str(e)}")
            return f"无法读取文件内容: {str(e)}"
    except Exception as e:
        logger.error(f"读取文件内容失败: {str(e)}")
        return f"无法读取文件内容: {str(e)}"

def generate_analysis_prompt(task, submission=None, file_content=None):
    """
    根据任务信息生成分析提示词
    """
    # 获取提交数据
    if not submission:
        db = SessionLocal()
        try:
            submission = db.query(Submission).filter_by(task_id=task.id).all()
        finally:
            db.close()
    
    # 构建提示词
    prompt = f"""你是一个数据分析专家，请基于以下表单数据提供详细的分析报告：

任务标题：{task.title}
任务描述：{task.description or '无'}

提交数据摘要：
"""
    
    # 添加提交数据摘要
    if submission:
        prompt += f"共有 {len(submission)} 条提交记录\n"
        
        # 分析前3条提交数据作为示例
        for i, sub in enumerate(submission[:3]):
            try:
                data = json.loads(sub.data)
                prompt += f"\n提交 #{i+1}:\n"
                for key, value in data.items():
                    prompt += f"  - {key}: {value}\n"
            except:
                prompt += f"\n提交 #{i+1}: {sub.data[:100]}...\n"
    else:
        prompt += "暂无提交数据\n"
    
    # 添加文件信息
    if file_content:
        prompt += f"\n附件内容摘要：\n{file_content[:500]}...\n" if len(file_content) > 500 else f"\n附件内容：\n{file_content}\n"
    
    # 添加分析要求
    prompt += """

请提供一个全面的数据分析报告，包括但不限于：
1. 数据概览：总提交量、关键数据分布等
2. 主要发现：数据中的趋势、模式和异常
3. 深入分析：基于数据的详细洞察
4. 建议和结论：基于分析结果的实用建议

请以中文撰写报告，使用Markdown格式，包括适当的标题、列表和表格来增强可读性。
"""
    
    return prompt

# Cherry Studio 企业版：OpenAI 兼容网关，服务端地址可自定义（云端 / 学校自部署）
CHERRY_DEFAULT_API_URL = 'https://express-ent-admin.cherryin.ai/v1'
CHERRY_DEFAULT_MODEL = 'deepseek/deepseek-v4-flash'


def _cherry_api_base(api_url):
    """把老师填的服务端地址规范成以 /v1 结尾的 base。
    兼容三种填法：https://host / https://host/v1 / http://ip:8088"""
    base = (api_url or '').strip() or CHERRY_DEFAULT_API_URL
    if not base.startswith('http'):
        base = 'http://' + base
    base = base.rstrip('/')
    if not base.endswith('/v1'):
        base += '/v1'
    return base


def _cherry_chat_url(api_url):
    return _cherry_api_base(api_url) + '/chat/completions'


def call_ai_model(prompt, ai_config):
    """
    调用AI模型生成分析报告
    """
    def get_model_config(model_name):
        for mc in ai_config.model_configs:
            if mc.model_name == model_name:
                return mc
        return None

    if ai_config.selected_model == 'deepseek':
        model_cfg = get_model_config('deepseek')
        api_key = model_cfg.api_key if model_cfg else ''
        custom_model = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        url = "https://api.deepseek.com/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        data = {
            "model": custom_model or "deepseek-chat",
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=60)
            response.raise_for_status()
            result = response.json()
            return result["choices"][0]["message"]["content"]
        except Exception as e:
            logger.error(f"DeepSeek API调用失败: {str(e)}")
            raise Exception(f"DeepSeek API调用失败: {str(e)}")

    elif ai_config.selected_model == 'doubao':
        model_cfg = get_model_config('doubao')
        api_key = model_cfg.api_key if model_cfg else ''
        custom_model = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        if not custom_model:
            raise Exception('豆包必须在「个人设置」中填写「模型名称」（视觉模型推荐 doubao-seed-1-6-vision-250815；或填自创建的 endpoint id）')
        url = "https://ark.cn-beijing.volces.com/api/v3/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        data = {
            "model": custom_model,
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            response = requests.post(url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            result = response.json()
            return result["choices"][0]["message"]["content"]
        except Exception as e:
            logger.error(f"豆包API调用失败: {str(e)}")
            raise Exception(f"豆包API调用失败: {str(e)}")

    elif ai_config.selected_model == 'qwen':
        model_cfg = get_model_config('qwen')
        api_key = model_cfg.api_key if model_cfg else ''
        custom_model = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        # OpenMentor: 改用阿里云的 OpenAI 兼容端点，与流式接口口径统一
        url = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        data = {
            "model": custom_model or "qwen-plus",
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            logger.info(f"调用阿里云百炼API，模型: {data['model']}")
            logger.info(f"请求URL: {url}")
            logger.info(f"请求头: {headers}")
            logger.info(f"请求数据: {json.dumps(data, ensure_ascii=False)[:200]}...")

            response = requests.post(url, headers=headers, json=data, timeout=120)

            logger.info(f"阿里云百炼API响应状态码: {response.status_code}")
            logger.info(f"阿里云百炼API响应头: {dict(response.headers)}")
            logger.info(f"阿里云百炼API响应内容: {response.text[:500]}...")

            if response.status_code != 200:
                raise Exception(f"阿里云百炼API调用失败，状态码: {response.status_code}，响应: {response.text[:200]}")

            if not response.text:
                raise Exception("阿里云百炼API返回空响应")

            try:
                result = response.json()
                logger.info(f"阿里云百炼API响应JSON结构: {list(result.keys()) if isinstance(result, dict) else '非字典结构'}")
            except ValueError as ve:
                raise Exception(f"阿里云百炼API返回非JSON响应: {response.text[:200]}")

            if isinstance(result, dict) and "code" in result and result["code"] != "200":
                raise Exception(f"阿里云百炼API调用失败: {result.get('message', '未知错误')} (错误码: {result.get('code')})")

            if isinstance(result, dict):
                if "output" in result and "text" in result["output"]:
                    return result["output"]["text"]
                elif "choices" in result and len(result["choices"]) > 0:
                    choice = result["choices"][0]
                    if "message" in choice and "content" in choice["message"]:
                        return choice["message"]["content"]
                    elif "text" in choice:
                        return choice["text"]
                elif "data" in result and "choices" in result["data"] and len(result["data"]["choices"]) > 0:
                    choice = result["data"]["choices"][0]
                    if "message" in choice and "content" in choice["message"]:
                        return choice["message"]["content"]

            raise Exception(f"阿里云百炼API返回未知格式的响应: {str(result)[:200]}")
        except requests.exceptions.RequestException as re:
            logger.error(f"阿里云百炼API网络请求异常: {str(re)}")
            raise Exception(f"阿里云百炼API网络请求异常: {str(re)}")
        except Exception as e:
            logger.error(f"阿里云百炼API调用失败: {str(e)}")
            raise Exception(f"阿里云百炼API调用失败: {str(e)}")

    elif ai_config.selected_model == 'glm':
        model_cfg = get_model_config('glm')
        api_key = model_cfg.api_key if model_cfg else ''
        custom_model = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        url = "https://open.bigmodel.cn/api/paas/v4/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        data = {
            "model": custom_model or "glm-4",
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            logger.info(f"调用GLM API，URL: {url}")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            logger.info(f"GLM API响应状态码: {response.status_code}")
            logger.info(f"GLM API响应内容: {response.text[:500]}")
            response.raise_for_status()
            result = response.json()
            return result["choices"][0]["message"]["content"]
        except Exception as e:
            logger.error(f"GLM API调用失败: {str(e)}")
            raise Exception(f"GLM API调用失败: {str(e)}")

    elif ai_config.selected_model == 'siliconflow':
        model_cfg = get_model_config('siliconflow')
        api_key = model_cfg.api_key if model_cfg else ''
        model_name = model_cfg.extra_settings if model_cfg and model_cfg.extra_settings else 'Qwen/Qwen2.5-72B-Instruct'
        url = "https://api.siliconflow.cn/v1/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        data = {
            "model": model_name,
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            logger.info(f"调用硅基流动API，URL: {url}")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            logger.info(f"硅基流动API响应状态码: {response.status_code}")
            logger.info(f"硅基流动API响应内容: {response.text[:500]}")
            response.raise_for_status()
            result = response.json()
            return result["choices"][0]["message"]["content"]
        except Exception as e:
            logger.error(f"硅基流动API调用失败: {str(e)}")
            raise Exception(f"硅基流动API调用失败: {str(e)}")

    elif ai_config.selected_model == 'cherry':
        model_cfg = get_model_config('cherry')
        api_key = model_cfg.api_key if model_cfg else ''
        model_name = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        url = _cherry_chat_url(model_cfg.api_url if model_cfg else '')
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        data = {
            "model": model_name or CHERRY_DEFAULT_MODEL,
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            logger.info(f"调用 Cherry Studio 企业版 API，URL: {url}, 模型: {data['model']}")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            if response.status_code != 200:
                raise Exception(f"返回 {response.status_code}: {response.text[:200]}")
            result = response.json()
            return result["choices"][0]["message"]["content"]
        except Exception as e:
            logger.error(f"Cherry Studio 企业版 API 调用失败: {str(e)}")
            raise Exception(f"Cherry Studio 企业版 API 调用失败: {str(e)}")

    elif ai_config.selected_model == 'ollama':
        model_cfg = get_model_config('ollama')
        api_url = model_cfg.api_url if model_cfg else 'http://localhost:11434'
        extra_settings = model_cfg.extra_settings if model_cfg else ''
        ollama_model = extra_settings if extra_settings else 'llama3.2'
        if not api_url.startswith('http'):
            api_url = 'http://' + api_url
        url = f"{api_url.rstrip('/')}/v1/chat/completions"
        headers = {
            "Content-Type": "application/json"
        }
        data = {
            "model": ollama_model,
            "messages": [
                {"role": "system", "content": "你是一个专业的数据分析助手。请基于用户提供的数据，生成一份详细、专业、有洞察力的分析报告。"},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 4000
        }

        try:
            logger.info(f"调用Ollama API，URL: {url}，模型: {ollama_model}")
            response = requests.post(url, headers=headers, json=data, timeout=180)
            logger.info(f"Ollama API响应状态码: {response.status_code}")
            response.raise_for_status()
            result = response.json()
            return result["choices"][0]["message"]["content"]
        except Exception as e:
            logger.error(f"Ollama API调用失败: {str(e)}")
            raise Exception(f"Ollama API调用失败: {str(e)}")

    else:
        raise Exception(f"不支持的AI模型: {ai_config.selected_model}")

# 创建Jinja2环境用于后台线程渲染模板
_template_env = None

def get_template_env():
    """获取Jinja2模板环境（延迟初始化）"""
    global _template_env
    if _template_env is None:
        template_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'templates')
        _template_env = Environment(
            loader=FileSystemLoader(template_dir),
            autoescape=select_autoescape(['html', 'xml'])
        )
    return _template_env

def save_analysis_report(task_id, report_content):
    """
    保存分析报告到文件系统和数据库
    """
    db = SessionLocal()
    try:
        task = db.query(Task).filter_by(id=task_id).first()
        if task:
            # 检查报告内容是否为空或只包含空白字符
            if not report_content or not report_content.strip():
                # 如果报告内容为空，生成友好的提示内容
                report_content = "本次分析未能生成有效内容。可能是由于以下原因：\n\n- 提交的数据量不足\n- 数据质量问题\n- AI模型处理异常\n\n请尝试提交更多数据或修改提示词后重新分析。"
            
            # 使用模板生成HTML报告内容
            template = get_template_env().get_template('simple_report.html')
            html_report = template.render(
                task_title=task.title,
                report_time=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                report_content=report_content
            )
            
            # 保存HTML报告到文件
            report_dir = 'static/reports'
            if not os.path.exists(report_dir):
                os.makedirs(report_dir)
            
            report_filename = f"report_{task_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            report_path = os.path.join(report_dir, report_filename)
            
            with open(report_path, 'w', encoding='utf-8') as f:
                f.write(html_report)
            
            # 更新数据库中的报告信息
            task.analysis_report = report_content
            task.report_file_path = report_path
            task.report_generated_at = datetime.now()
            db.commit()
            
            # 添加到已完成报告集合
            with progress_lock:
                completed_reports.add(task_id)
            
            logger.info(f"任务 {task_id} 的分析报告已保存")
    except Exception as e:
        logger.error(f"保存分析报告失败: {str(e)}")
    finally:
        db.close()

def timeout(seconds, error_message="函数执行超时"):
    """
    超时装饰器（使用线程实现，避免信号处理问题）
    """
    import threading
    from functools import wraps
    
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            # 结果容器
            result = [None]
            exception = [None]
            
            # 目标函数
            def target():
                try:
                    result[0] = func(*args, **kwargs)
                except Exception as e:
                    exception[0] = e
            
            # 创建并启动线程
            thread = threading.Thread(target=target)
            thread.daemon = True
            thread.start()
            thread.join(seconds)
            
            # 检查线程是否仍在运行
            if thread.is_alive():
                # 线程超时，抛出异常
                raise TimeoutError(error_message)
            elif exception[0]:
                # 函数执行中出现异常
                raise exception[0]
            else:
                # 正常返回结果
                return result[0]
        
        return wrapper
    
    return decorator

def perform_analysis_with_custom_prompt(task_id, user_id, ai_config_id, custom_prompt):
    """
    使用自定义提示词执行分析任务
    """
    db = SessionLocal()
    try:
        # 获取任务信息
        task = db.query(Task).filter_by(id=task_id, user_id=user_id).first()
        if not task:
            with progress_lock:
                analysis_progress[task_id] = {
                    'status': 'error',
                    'message': '任务不存在'
                }
            return
        
        # 获取提交数据
        submission = db.query(Submission).filter_by(task_id=task_id).all()
        
        # 读取附件内容（如果有）
        file_content = None
        if task.attachments:
            # 读取第一个附件的内容
            first_attachment = task.attachments[0]
            if os.path.exists(first_attachment.file_path):
                file_content = read_file_content(first_attachment.file_path)
        
        # 获取AI配置
        ai_config = db.query(AIConfig).filter_by(id=ai_config_id).first()
        if not ai_config:
            with progress_lock:
                analysis_progress[task_id] = {
                    'status': 'error',
                    'message': 'AI配置不存在'
                }
            return
        
        # 验证AI配置是否正确
        if ai_config.selected_model == 'deepseek' and not ai_config.deepseek_api_key:
            with progress_lock:
                analysis_progress[task_id] = {
                    'status': 'error',
                    'message': 'DeepSeek API密钥未配置'
                }
            logging.error(f"任务 {task_id}：DeepSeek API密钥未配置")
            return
        elif ai_config.selected_model == 'doubao' and not ai_config.doubao_api_key:
            with progress_lock:
                analysis_progress[task_id] = {
                    'status': 'error',
                    'message': '豆包API密钥未配置完整'
                }
            logging.error(f"任务 {task_id}：豆包API密钥未配置完整")
            return
        
        logging.info(f"任务 {task_id}：使用模型 {ai_config.selected_model}")
        
        # 进度1：正在生成提示词
        with progress_lock:
            analysis_progress[task_id] = {
                'status': 'in_progress',
                'progress': 0,
                'message': '正在生成提示词...'
            }
        
        # 生成分析提示词
        prompt = custom_prompt
        
        # 进度2：大模型分析中
        with progress_lock:
            analysis_progress[task_id] = {
                'status': 'in_progress',
                'progress': 1,
                'message': '大模型分析中，这可能需要几分钟时间...'
            }
        logging.info(f"任务 {task_id}：调用AI模型进行分析")
        
        # 设置AI调用的超时时间，根据模型类型调整
        timeout_seconds = 120 if ai_config.selected_model == 'deepseek' else (120 if ai_config.selected_model == 'qwen' else 90)
        
        # 带超时的AI模型调用
        @timeout(seconds=timeout_seconds, error_message=f"调用{ai_config.selected_model}模型超时（{timeout_seconds}秒）")
        def call_ai_with_timeout(prompt, config):
            logging.info(f"开始调用 {config.selected_model} API，提示词长度: {len(prompt)} 字符，超时设置: {timeout_seconds}秒")
            return call_ai_model(prompt, config)
        
        # 调用AI模型
        try:
            analysis_report = call_ai_with_timeout(prompt, ai_config)
            logging.info(f"成功获取 {ai_config.selected_model} API 响应，报告长度: {len(analysis_report)} 字符")
        except TimeoutError as timeout_error:
            # 处理超时错误
            error_msg = str(timeout_error)
            logging.error(f"任务 {task_id}：{error_msg}")
            with progress_lock:
                analysis_progress[task_id] = {
                    'status': 'error',
                    'message': f"分析超时：{error_msg}，请检查网络连接或稍后重试"
                }
            return
        except Exception as api_error:
            logging.error(f"任务 {task_id}：AI模型调用失败: {str(api_error)}")
            logging.error(f"详细错误堆栈: {traceback.format_exc()}")
            with progress_lock:
                analysis_progress[task_id] = {
                    'status': 'error',
                    'message': f'API调用失败: {str(api_error)}'
                }
            return
        
        # 检查是否是错误消息
        if analysis_report.startswith("错误：") or \
           (analysis_report.startswith("DeepSeek API调用") and "失败" in analysis_report) or \
           (analysis_report.startswith("豆包API调用") and "失败" in analysis_report):
            logging.error(f"任务 {task_id}：AI模型返回错误: {analysis_report}")
            raise Exception(analysis_report)
        
        # 保存结果到文件和数据库
        with progress_lock:
            save_analysis_report(task_id, analysis_report)
            analysis_results[task_id] = analysis_report
            analysis_progress[task_id] = {
                'status': 'completed',
                'progress': 3,
                'message': '分析完成，请查看报告'
            }
            
    except Exception as e:
        # 处理错误
        with progress_lock:
            analysis_progress[task_id] = {
                'status': 'error',
                'message': f'分析过程中出错: {str(e)}'
            }
    finally:
        db.close()

# 路由函数
@app.route('/login', methods=['GET', 'POST'])
def login():
    next_page = request.values.get('next', '').strip()

    if current_user.is_authenticated and request.method == 'GET':
        if getattr(current_user, 'force_password_change', False):
            return redirect(url_for('profile', active_tab='password', next=next_page))
        return redirect(_safe_next_url(next_page)) if next_page else redirect(url_for('dashboard'))

    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        db = SessionLocal()
        try:
            user = db.query(User).filter_by(username=username).first()
            
            if user and check_password_hash(user.password, password):
                if not user.is_active:
                    flash('该账号已被禁用，请联系管理员', 'danger')
                    return render_template('login.html')
                user.last_login_at = datetime.now()
                db.commit()
                login_user(user)
                
                if getattr(user, 'force_password_change', False) or password == 'openmentor':
                    flash('请修改您的默认密码', 'warning')
                    profile_args = {'active_tab': 'password'}
                    if next_page and _is_safe_next_url(next_page):
                        profile_args['next'] = next_page
                    return redirect(url_for('profile', **profile_args))
                
                return redirect(_safe_next_url(next_page)) if next_page else redirect(url_for('dashboard'))
            else:
                flash('用户名或密码错误', 'danger')
        finally:
            db.close()
    
    return render_template('login.html', next_page=next_page)

@app.route('/logout')
def logout():
    logout_user()
    return redirect(url_for('login'))

# ==================== 简体 / 繁体切换 ====================
# 实现方式：模板全部保持简体；用户切到繁体后（cookie），在响应出口把整页 HTML
# 用 OpenCC s2twp 词表转换（词组级，"软件→軟體 / 网络连接→網路連線"）。
# 只转 text/html；JSON、SSE 流不动（学生/AI 的对话内容属于用户数据，不做转换，
# AI 回复改由 system prompt 要求其直接用繁体输出）。

LANG_COOKIE = 'om_lang'  # zh-Hans（默认，不写 cookie 也算）/ zh-Hant
_opencc_s2t = None
_opencc_failed = False


def _get_s2t_converter():
    """懒加载 OpenCC；环境缺包时降级为 None（繁体开关静默失效，不影响运行）。"""
    global _opencc_s2t, _opencc_failed
    if _opencc_s2t is None and not _opencc_failed:
        try:
            from opencc import OpenCC
            _opencc_s2t = OpenCC('s2twp')
            logger.info('[简繁] OpenCC s2twp 转换器已加载')
        except Exception as e:
            _opencc_failed = True
            logger.warning(f'[简繁] OpenCC 不可用，繁体模式停用（pip install opencc-python-reimplemented 可启用）: {e}')
    return _opencc_s2t


def _lang_is_hant():
    try:
        return request.cookies.get(LANG_COOKIE) == 'zh-Hant'
    except RuntimeError:
        return False


@app.after_request
def _serve_traditional_chinese(response):
    """繁体模式下在出口把 HTML 页面整体简转繁。失败时按原文返回，绝不影响功能。"""
    try:
        if not _lang_is_hant():
            return response
        if not (response.content_type or '').startswith('text/html'):
            return response
        if response.direct_passthrough or not (200 <= response.status_code < 300):
            return response
        cc = _get_s2t_converter()
        if not cc:
            return response
        response.set_data(cc.convert(response.get_data(as_text=True)))
    except Exception as e:
        logger.warning(f'[简繁] 页面转换失败，按简体返回: {e}')
    return response


@app.route('/lang/<string:lang_code>')
def set_language(lang_code):
    """切换界面语言（cookie 一年有效），跳回来源页。"""
    if lang_code not in ('zh-Hans', 'zh-Hant'):
        lang_code = 'zh-Hans'
    target = request.referrer
    if not target or not target.startswith(request.host_url):
        target = url_for('index')
    resp = redirect(target)
    resp.set_cookie(LANG_COOKIE, lang_code, max_age=365 * 24 * 3600, samesite='Lax')
    return resp


@app.route('/')
def index():
    return render_template('home.html')


@app.route('/privacy')
def privacy_policy():
    return render_template('privacy.html')

@app.route('/api/checkin', methods=['POST'])
def api_checkin_submit():
    body = request.get_json(silent=True) or {}
    province = (body.get('province') or '').strip()
    if province not in CHECKIN_ALLOWED_PROVINCES:
        return jsonify({'success': False, 'message': '无效的地区'}), 400
    try:
        _create_checkin_record(province)
        return jsonify({'success': True, 'province': province})
    except Exception as e:
        logger.warning(f'[首页打卡] 提交失败: {e}')
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/checkin/stats')
def api_checkin_stats():
    try:
        counts = _load_checkin_counts()
        return jsonify({
            'success': True,
            'counts': counts,
            'total': sum(counts.values()),
        })
    except Exception as e:
        logger.warning(f'[首页打卡] 读取统计失败: {e}')
        return jsonify({'success': False, 'message': str(e), 'counts': {}, 'total': 0}), 500


# ==================== PWA：根路径暴露 manifest 和 service worker（B6）====================
@app.route('/manifest.json')
def pwa_manifest():
    resp = send_from_directory(app.static_folder, 'manifest.json')
    resp.headers['Content-Type'] = 'application/manifest+json; charset=utf-8'
    resp.headers['Cache-Control'] = 'public, max-age=3600'
    return resp


@app.route('/sw.js')
def pwa_service_worker():
    resp = send_from_directory(app.static_folder, 'sw.js')
    resp.headers['Content-Type'] = 'application/javascript; charset=utf-8'
    resp.headers['Service-Worker-Allowed'] = '/'
    resp.headers['Cache-Control'] = 'no-cache'  # 永远拉最新版 sw
    return resp

@app.route('/dashboard')
@login_required
def dashboard():
    _ensure_template_tasks_for_user(current_user.id)  # 懒下发内置模板任务
    db = SessionLocal()
    try:
        # 仅展示数据任务（QuickForm 原有功能），不显示 AI 导师（后者在 /assistant/list 中管理）
        # 置顶优先，其次按创建时间倒序
        tasks = (
            db.query(Task)
            .filter_by(user_id=current_user.id, is_assistant=False)
            .order_by(Task.pinned.desc(), Task.created_at.desc())
            .all()
        )
        pinned_count = sum(1 for t in tasks if t.pinned)
        return render_template('dashboard.html', tasks=tasks, pinned_count=pinned_count, max_pinned=MAX_PINNED_TASKS)
    finally:
        db.close()


MAX_PINNED_TASKS = 3


@app.route('/task/<int:task_id>/pin', methods=['POST'])
@login_required
def task_toggle_pin(task_id):
    """数据任务置顶/取消置顶（最多置顶 MAX_PINNED_TASKS 个）。"""
    db = SessionLocal()
    try:
        task = db.query(Task).filter_by(id=task_id, user_id=current_user.id, is_assistant=False).first()
        if not task:
            return jsonify({'success': False, 'message': '任务不存在'}), 404
        if task.pinned:
            task.pinned = False
            db.commit()
            return jsonify({'success': True, 'pinned': False, 'message': '已取消置顶'})
        cur = db.query(Task).filter_by(user_id=current_user.id, is_assistant=False, pinned=True).count()
        if cur >= MAX_PINNED_TASKS:
            return jsonify({'success': False, 'message': f'最多置顶 {MAX_PINNED_TASKS} 个，请先取消其他置顶'}), 400
        task.pinned = True
        db.commit()
        return jsonify({'success': True, 'pinned': True, 'message': '已置顶'})
    finally:
        db.close()

@app.route('/generate_report/<int:task_id>', methods=['GET', 'POST'])
@login_required
def generate_report(task_id):
    """
    在新页面中生成分析报告
    """
    # 添加详细的请求日志
    logger.info(f"收到生成报告请求 - Task ID: {task_id}, Method: {request.method}")
    logger.info(f"请求URL: {request.url}")
    logger.info(f"请求参数: {dict(request.args)}")
    logger.info(f"表单数据: {dict(request.form)}")
    logger.info(f"请求头: {dict(request.headers)}")
    
    db = SessionLocal()
    try:
        # 检查任务权限
        task = db.query(Task).filter_by(id=task_id, user_id=current_user.id).first()
        if not task:
            logger.warning(f"任务不存在或无权访问 - Task ID: {task_id}, User ID: {current_user.id}")
            return render_template('generate_report.html', error='任务不存在或无权访问')
        
        # 获取AI配置
        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
        if not ai_config or not ai_config.selected_model:
            flash('请先在配置页面设置AI模型和API密钥', 'warning')
            return redirect(url_for('profile'))
        
        # 针对不同模型验证必需的API密钥
        if ai_config.selected_model == 'deepseek' and not ai_config.deepseek_api_key:
            flash('请先配置DeepSeek API密钥', 'warning')
            return redirect(url_for('profile'))
        elif ai_config.selected_model == 'doubao' and not ai_config.doubao_api_key:
            flash('请先配置豆包API密钥', 'warning')
            return redirect(url_for('profile'))
        elif ai_config.selected_model == 'qwen' and not ai_config.qwen_api_key:
            flash('请先配置阿里云百炼API密钥', 'warning')
            return redirect(url_for('profile'))
        
        # 获取提示词
        custom_prompt = None
        if request.method == 'GET' and 'prompt' in request.args:
            custom_prompt = request.args.get('prompt')
            logger.info(f"从GET参数获取提示词，长度: {len(custom_prompt) if custom_prompt else 0}")
        elif request.method == 'POST' and 'custom_prompt' in request.form:
            custom_prompt = request.form.get('custom_prompt')
            logger.info(f"从POST表单获取提示词，长度: {len(custom_prompt) if custom_prompt else 0}")
        
        # 如果没有提示词，生成默认提示词
        if not custom_prompt:
            logger.info("未提供自定义提示词，生成默认提示词")
            submission = db.query(Submission).filter_by(task_id=task_id).all()
            file_content = None
            if task.attachments:
                # 读取第一个附件的内容
                first_attachment = task.attachments[0]
                if os.path.exists(first_attachment.file_path):
                    file_content = read_file_content(first_attachment.file_path)
            custom_prompt = generate_analysis_prompt(task, submission, file_content)
            logger.info(f"生成默认提示词，长度: {len(custom_prompt) if custom_prompt else 0}")
        else:
            logger.info(f"使用自定义提示词，长度: {len(custom_prompt)}")
        
        # 验证提示词不为空
        if not custom_prompt or not custom_prompt.strip():
            logger.warning("提示词为空或只包含空白字符")
            return render_template('generate_report.html', task=task, error="提示词不能为空", ai_config=ai_config)
        
        logger.info(f"开始生成报告任务 {task_id}，使用模型 {ai_config.selected_model}")
        
        # 执行分析
        try:
            # 进度显示
            progress_message = "正在使用AI模型分析数据..."
            
            # 设置超时时间
            timeout_seconds = 120 if ai_config.selected_model == 'deepseek' else 90
            
            # 调用AI模型
            @timeout(seconds=timeout_seconds, error_message=f"调用{ai_config.selected_model}模型超时（{timeout_seconds}秒）")
            def call_ai_with_timeout(prompt, config):
                return call_ai_model(prompt, config)
            
            # 执行分析
            analysis_report = call_ai_with_timeout(custom_prompt, ai_config)
            
            # 保存报告
            save_analysis_report(task_id, analysis_report)
            
            # 成功显示报告
            return render_template('generate_report.html', 
                                 task=task, 
                                 report=analysis_report,
                                 preview_prompt=custom_prompt,
                                 ai_config=ai_config)
            
        except Exception as e:
            logger.error(f"生成报告失败: {str(e)}")
            return render_template('generate_report.html', 
                                 task=task, 
                                 error=f'生成报告失败: {str(e)}',
                                 preview_prompt=custom_prompt,
                                 ai_config=ai_config)
            
    except Exception as e:
        logger.error(f"访问生成报告页面失败: {str(e)}")
        flash('生成报告时出现错误', 'danger')
        return redirect(url_for('task_detail', task_id=task_id))
    finally:
        db.close()

@app.route('/create_task', methods=['GET', 'POST'])
@login_required
def create_task():
    if request.method == 'POST':
        title = request.form.get('title')
        description = request.form.get('description')
        
        db = SessionLocal()
        try:
            task = Task(title=title, description=description, user_id=current_user.id)
            db.add(task)
            db.commit()
            
            # 处理多附件上传
            # 支持file、file_2、file_3等多个文件字段
            file_fields = ['file', 'file_2', 'file_3']
            for field_name in file_fields:
                if field_name in request.files and request.files[field_name].filename != '':
                    file = request.files[field_name]
                    unique_filename, filepath = save_uploaded_file(file)
                    if unique_filename:
                        attachment = Attachment(
                            task_id=task.id,
                            file_name=file.filename,
                            file_path=filepath
                        )
                        db.add(attachment)
            
            db.commit()
            
            flash('数据任务创建成功', 'success')
            return redirect(url_for('task_detail', task_id=task.id))
        finally:
            db.close()
    return render_template('create_task.html')

@app.route('/import_task', methods=['GET', 'POST'])
@login_required
def import_task():
    tasks = []
    error = None

    host = request.host.lower()
    if '127.0.0.1' in host or 'localhost' in host:
        flash('导入任务不能使用127.0.0.1的方式访问网站', 'danger')
        return render_template('import_task.html', tasks=[], error=None)

    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')

        try:
            session['quickform_username'] = username
            session['quickform_password'] = password

            url = 'https://quickform.cn/cli/list'
            data = {
                'username': username,
                'password': password
            }
            response = requests.post(url, data=data)
            response.raise_for_status()

            result = response.json()
            if result.get('success'):
                tasks = result.get('tasks', [])
            else:
                error = result.get('message', '获取任务列表失败')
        except Exception as e:
            error = f'请求失败: {str(e)}'
    else:
        tasks_param = request.args.get('tasks')
        if tasks_param:
            try:
                tasks = json.loads(tasks_param)
            except:
                tasks = []

    return render_template('import_task.html', tasks=tasks, error=error)

@app.route('/import_task_action/<string:apiid>')
@login_required
def import_task_action(apiid):
    import requests
    import re
    import os
    import uuid
    
    task_name = request.args.get('task_name', '导入的任务')
    
    host = request.host.lower()
    if '127.0.0.1' in host or 'localhost' in host:
        flash('导入任务不能使用127.0.0.1的方式访问网站', 'danger')
        return redirect(url_for('import_task'))
    
    db = SessionLocal()
    try:
        quickform_username = session.get('quickform_username')
        quickform_password = session.get('quickform_password')
        
        if not quickform_username or not quickform_password:
            qf_config = db.query(QFConfig).filter_by(user_id=current_user.id).first()
            if qf_config and qf_config.username and qf_config.password:
                quickform_username = qf_config.username
                quickform_password = qf_config.password
            else:
                flash('请先获取任务列表以验证quickform.cn账号', 'danger')
                return redirect(url_for('import_task'))
        
        quickform_url = 'https://quickform.cn'
        show_data = {
            'username': quickform_username,
            'password': quickform_password,
            'apiid': apiid
        }
        
        try:
            response = requests.post(
                f'{quickform_url}/cli/show',
                data=show_data,
                timeout=30,
                allow_redirects=True
            )
            import logging
            logger = logging.getLogger(__name__)
            logger.info(f"CLI show response status: {response.status_code}")
            logger.info(f"CLI show response headers: {dict(response.headers)}")
            logger.info(f"CLI show response text: {response.text[:500] if response.text else 'empty'}")
            
            if response.status_code != 200:
                flash(f'获取任务信息失败: HTTP {response.status_code}', 'danger')
                return redirect(url_for('import_task'))
            
            task_info = response.json()
        except json.JSONDecodeError as e:
            flash(f'获取任务信息失败: 响应格式错误', 'danger')
            return redirect(url_for('import_task'))
        except Exception as e:
            flash(f'获取任务信息失败: {str(e)}', 'danger')
            return redirect(url_for('import_task'))
        
        if not task_info.get('success'):
            flash(f'获取任务信息失败: {task_info.get("message", "未知错误")}', 'danger')
            return redirect(url_for('import_task'))
        
        task_title = task_info.get('name', task_name)
        task_intro = task_info.get('intro', '')
        tutorial_link = task_info.get('tutorial', '')
        share_url = task_info.get('share_url', '')
        attachments_info = task_info.get('attachments', [])
        
        existing_task = db.query(Task).filter_by(task_id=apiid).first()
        if existing_task:
            new_api_id = generate_custom_id()
            flash(f'API {apiid} 已存在，已生成新API: {new_api_id}', 'info')
        else:
            new_api_id = apiid
        
        new_task = Task(
            title=task_title,
            description=task_intro,
            user_id=current_user.id,
            task_id=new_api_id
        )
        db.add(new_task)
        db.flush()
        
        for attachment in attachments_info:
            attachment_name = attachment.get('name', '')
            attachment_url = attachment.get('url', '')
            
            if not attachment_url or not attachment_name.endswith('.html'):
                continue
            
            try:
                html_response = requests.get(attachment_url, timeout=30)
                html_content = html_response.text
                
                pattern = rf'https?://quickform\.cn/api/([a-zA-Z0-9]+)'
                new_api_pattern = request.host_url.rstrip('/') + '/api/' + new_api_id
                modified_html = re.sub(pattern, new_api_pattern, html_content)
                
                unique_filename = f"{uuid.uuid4().hex}_{attachment_name}"
                uploads_dir = os.path.join(app.root_path, 'static', 'uploads')
                os.makedirs(uploads_dir, exist_ok=True)
                file_path = os.path.join(uploads_dir, unique_filename)
                
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write(modified_html)
                
                relative_path = f'uploads/{unique_filename}'
                
                db_attachment = Attachment(
                    task_id=new_task.id,
                    file_name=attachment_name,
                    file_path=relative_path
                )
                db.add(db_attachment)
            except Exception as e:
                flash(f'下载附件 {attachment_name} 失败: {str(e)}', 'warning')
        
        db.commit()
        flash(f'任务"{task_title}"导入成功，API ID: {new_api_id}', 'success')
        return redirect(url_for('task_detail', task_id=new_task.id))
    except Exception as e:
        flash(f'任务导入失败: {str(e)}', 'danger')
        return redirect(url_for('import_task'))
    finally:
        db.close()

@app.route('/import_task_by_url')
@login_required
def import_task_by_url():
    host = request.host.lower()
    if '127.0.0.1' in host or 'localhost' in host:
        flash('导入任务不能使用127.0.0.1的方式访问网站', 'danger')
        return redirect(url_for('import_task'))
    
    task_url = request.args.get('url', '')
    
    match = re.search(r'/api/([a-zA-Z0-9]+)', task_url)
    if not match:
        flash('无效的任务URL格式', 'danger')
        return redirect(url_for('import_task'))
    
    apiid = match.group(1)
    
    return redirect(url_for('import_task_action', apiid=apiid, task_name=f'任务{apiid}'))

@app.route('/import_task_from_file', methods=['POST'])
@login_required
def import_task_from_file():
    import zipfile
    import io
    import re
    import os
    import logging
    
    logger = logging.getLogger(__name__)
    logger.info(f"Request files: {request.files}")
    logger.info(f"Request form: {request.form}")
    
    if 'task_file' not in request.files:
        flash('没有文件上传', 'danger')
        return redirect(url_for('import_task'))
    
    file = request.files['task_file']
    if file.filename == '':
        flash('没有选择文件', 'danger')
        return redirect(url_for('import_task'))
    
    host = request.host.lower()
    if '127.0.0.1' in host or 'localhost' in host:
        flash('导入任务不能使用127.0.0.1的方式访问网站', 'danger')
        return redirect(url_for('import_task'))
    
    try:
        zip_bytes = file.read()
        zip_file = zipfile.ZipFile(io.BytesIO(zip_bytes))
        
        json_content = zip_file.read('quickform-task-migration.json').decode('utf-8')
        import json
        task_data = json.loads(json_content)
        
        original_api_id = task_data.get('api_id', '')
        title = task_data.get('title', '未命名任务')
        description = task_data.get('description', '')
        html_files = task_data.get('html_files', [])
        export_api_base = task_data.get('export_api_base', 'https://quickform.cn')
        
        db = SessionLocal()
        try:
            existing_task = db.query(Task).filter_by(task_id=original_api_id).first()
            
            if existing_task:
                new_api_id = generate_custom_id()
                flash(f'API {original_api_id} 已存在，已生成新API: {new_api_id}', 'info')
            else:
                new_api_id = original_api_id
            
            new_task = Task(
                title=title,
                description=description,
                user_id=current_user.id,
                task_id=new_api_id
            )
            db.add(new_task)
            db.flush()
            
            for html_file_info in html_files:
                archive_name = html_file_info.get('archive_name', '')
                original_name = html_file_info.get('original_name', '')
                
                if archive_name and archive_name in zip_file.namelist():
                    html_content = zip_file.read(archive_name).decode('utf-8')
                    
                    export_api_base = task_data.get('export_api_base', 'https://quickform.cn').rstrip('/')
                    new_api_pattern = request.host_url.rstrip('/') + '/api/' + new_api_id
                    pattern = rf'https?://quickform\.cn/api/([a-zA-Z0-9]+)'
                    modified_html = re.sub(pattern, new_api_pattern, html_content)
                    
                    unique_filename = f"{uuid.uuid4().hex}_{original_name}"
                    uploads_dir = os.path.join(app.root_path, 'static', 'uploads')
                    os.makedirs(uploads_dir, exist_ok=True)
                    file_path = os.path.join(uploads_dir, unique_filename)
                    
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(modified_html)
                    
                    relative_path = f'uploads/{unique_filename}'
                    
                    attachment = Attachment(
                        task_id=new_task.id,
                        file_name=original_name,
                        file_path=relative_path
                    )
                    db.add(attachment)
            
            db.commit()
            flash(f'任务"{title}"导入成功，API ID: {new_api_id}', 'success')
            return redirect(url_for('task_detail', task_id=new_task.id))
        finally:
            db.close()
    except zipfile.BadZipFile:
        flash('无效的压缩包文件', 'danger')
    except KeyError as e:
        flash(f'压缩包内缺少必要文件: {str(e)}', 'danger')
    except Exception as e:
        flash(f'导入失败: {str(e)}', 'danger')
    
    return redirect(url_for('import_task'))

@app.route('/task/<int:task_id>/upload', methods=['POST'])
@login_required
def upload_task_attachment(task_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            return jsonify({'success': False, 'message': '任务不存在'})
        if task.user_id != current_user.id:
            return jsonify({'success': False, 'message': '无权访问此任务'})

        if 'file' not in request.files:
            return jsonify({'success': False, 'message': '没有文件'})
        file = request.files['file']
        if file.filename == '':
            return jsonify({'success': False, 'message': '没有选择文件'})

        unique_filename, filepath = save_uploaded_file(file)
        if unique_filename:
            attachment = Attachment(
                task_id=task.id,
                file_name=file.filename,
                file_path=filepath
            )
            db.add(attachment)
            db.commit()
            return jsonify({'success': True, 'message': '文件上传成功'})
        else:
            return jsonify({'success': False, 'message': '文件保存失败'})
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)})
    finally:
        db.close()

@app.route('/task/<int:task_id>/attachment/<int:attachment_id>/edit', methods=['GET', 'POST'])
@login_required
def edit_task_html_attachment(task_id, attachment_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权编辑此任务', 'danger')
            return redirect(url_for('dashboard'))

        attachment = db.query(Attachment).get(attachment_id)
        if not attachment or attachment.task_id != task.id:
            flash('附件不存在或不属于当前任务', 'danger')
            return redirect(url_for('task_detail', task_id=task.id))
        if not _is_html_attachment(attachment):
            flash('只能在线编辑 HTML / HTM 文件', 'warning')
            return redirect(url_for('task_detail', task_id=task.id))

        file_path = _resolve_upload_file_path(attachment.file_path)
        if not file_path or not os.path.isfile(file_path):
            flash('附件文件不存在或路径异常', 'danger')
            return redirect(url_for('task_detail', task_id=task.id))

        current_api_url = request.host_url.rstrip('/') + url_for('submit_form', task_id=task.task_id)
        preview_url = url_for('static', filename='uploads/' + os.path.basename(file_path))
        content = None
        encoding = request.form.get('encoding') or 'utf-8'

        if request.method == 'POST':
            content = request.form.get('content', '')
            if len(content.encode('utf-8')) > HTML_EDITOR_MAX_BYTES:
                flash('HTML 文件超过在线编辑大小限制（5MB）', 'danger')
            else:
                if encoding not in {'utf-8', 'utf-8-sig', 'gb18030'}:
                    encoding = 'utf-8'
                with open(file_path, 'w', encoding=encoding, newline='') as f:
                    f.write(content)
                flash('HTML 文件已保存', 'success')
                return redirect(url_for('edit_task_html_attachment', task_id=task.id, attachment_id=attachment.id))

        if content is None:
            try:
                content, encoding = _read_editable_html_file(file_path)
            except ValueError as e:
                flash(str(e), 'danger')
                return redirect(url_for('task_detail', task_id=task.id))

        assistant_options = []
        assistants = (
            db.query(Task)
            .filter_by(user_id=current_user.id, is_assistant=True)
            .order_by(Task.created_at.desc())
            .all()
        )
        for assistant in assistants:
            assistant_options.append({
                'id': assistant.id,
                'title': assistant.title,
                'task_id': assistant.task_id,
                'status': assistant.status or 'active',
                'url': request.host_url.rstrip('/') + url_for('chat_entry', task_id=assistant.task_id),
            })

        return render_template(
            'html_attachment_editor.html',
            task=task,
            attachment=attachment,
            content=content,
            encoding=encoding,
            current_api_url=current_api_url,
            assistant_options=assistant_options,
            preview_url=preview_url,
            auto_detect=request.args.get('action') == 'repair',
        )
    finally:
        db.close()

@app.route('/task/<int:task_id>')
@login_required
def task_detail(task_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权访问此任务', 'danger')
            return redirect(url_for('dashboard'))
        
        submission = db.query(Submission).filter_by(task_id=task.id).order_by(Submission.submitted_at.desc()).all()

        # B16 续：把每条预解析一下，分离附件，方便模板渲染"查看附件"按钮 + 模态框
        items = []
        for sub in submission:
            attachments = []
            display_data = sub.data
            try:
                parsed = json.loads(sub.data)
                if isinstance(parsed, dict) and '_attachments' in parsed:
                    attachments = parsed.get('_attachments') or []
                    main_dict = {k: v for k, v in parsed.items() if k != '_attachments'}
                    display_data = json.dumps(main_dict, ensure_ascii=False, indent=2)
            except (json.JSONDecodeError, TypeError, ValueError):
                pass
            # 复用一个简单包装对象，保留 sub.id / sub.submitted_at 给模板
            sub.display_data = display_data
            sub.attachments = attachments
            items.append(sub)

        remote_prompt_matrix = _prepare_prompt_matrix_for_task(task)
        return render_template(
            'task_detail.html',
            task=task,
            submission=items,
            remote_prompt_matrix=remote_prompt_matrix,
        )
    finally:
        db.close()

@app.route('/task/<int:task_id>/prompt_matrix/sync', methods=['POST'])
@login_required
def sync_task_prompt_matrix(task_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权访问此任务', 'danger')
            return redirect(url_for('dashboard'))

        try:
            items = _sync_prompt_matrix_from_lark()
            flash(f'数据使用矩阵已同步，共 {len(items)} 条。', 'success')
        except Exception as e:
            logger.warning(f'[提示词矩阵] 手动同步失败: {e}')
            flash(f'数据使用矩阵同步失败：{e}', 'danger')
        return redirect(url_for('task_detail', task_id=task.id))
    finally:
        db.close()

@app.route('/task/<int:task_id>/data')
@login_required
def task_data_view(task_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权访问此任务', 'danger')
            return redirect(url_for('dashboard'))
        
        submission = db.query(Submission).filter_by(task_id=task.id).order_by(Submission.submitted_at.desc()).all()

        # B16: 把每条记录预解析一下，把附件单独抽出来供模板渲染缩略图 / 下载链接
        items = []
        for sub in submission:
            attachments = []
            display_data = sub.data
            try:
                parsed = json.loads(sub.data)
                if isinstance(parsed, dict) and '_attachments' in parsed:
                    attachments = parsed.get('_attachments') or []
                    main_dict = {k: v for k, v in parsed.items() if k != '_attachments'}
                    display_data = json.dumps(main_dict, ensure_ascii=False, indent=2)
            except (json.JSONDecodeError, TypeError, ValueError):
                pass
            items.append({
                'id': sub.id,
                'submitted_at': sub.submitted_at,
                'display_data': display_data,
                'attachments': attachments,
            })

        class SimplePagination:
            def __init__(self, page, per_page, total):
                self.page = page
                self.per_page = per_page
                self.total = total
                self.pages = 1

        total_submissions = len(submission)
        pagination = SimplePagination(1, 10, total_submissions)

        return render_template('task_data_view.html', task=task, submissions=items, total_submissions=total_submissions, pagination=pagination)
    finally:
        db.close()

@app.route('/edit_task/<int:task_id>', methods=['GET', 'POST'])
@login_required
def edit_task(task_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权编辑此任务', 'danger')
            return redirect(url_for('dashboard'))
        
        if request.method == 'POST':
            title = request.form.get('title')
            description = request.form.get('description')
            
            # 更新任务信息
            task.title = title
            task.description = description
            
            # 处理删除附件
            remove_attachments = request.form.getlist('remove_attachments')
            for attachment_id in remove_attachments:
                attachment = db.query(Attachment).get(int(attachment_id))
                if attachment and attachment.task_id == task.id:
                    # 删除物理文件
                    if os.path.exists(attachment.file_path):
                        os.remove(attachment.file_path)
                    db.delete(attachment)
            
            # 处理新附件上传
            # 支持file、file_2、file_3等多个文件字段
            file_fields = ['file', 'file_2', 'file_3']
            for field_name in file_fields:
                if field_name in request.files and request.files[field_name].filename != '':
                    file = request.files[field_name]
                    unique_filename, filepath = save_uploaded_file(file)
                    if unique_filename:
                        attachment = Attachment(
                            task_id=task.id,
                            file_name=file.filename,
                            file_path=filepath
                        )
                        db.add(attachment)
            
            db.commit()
            flash('任务更新成功', 'success')
            return redirect(url_for('task_detail', task_id=task.id))
        
        return render_template('edit_task.html', task=task)
    finally:
        db.close()

@app.route('/delete_submission/<int:submission_id>', methods=['POST', 'GET'])
@login_required
def delete_submission(submission_id):
    """删除单个提交数据"""
    db = SessionLocal()
    try:
        submission = db.query(Submission).get(submission_id)
        if not submission:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'success': False, 'message': '提交数据不存在'})
            flash('提交数据不存在', 'danger')
            return redirect(url_for('dashboard'))
        
        task = db.query(Task).get(submission.task_id)
        if not task:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'success': False, 'message': '任务不存在'})
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        
        if task.user_id != current_user.id:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'success': False, 'message': '无权删除此提交数据'})
            flash('无权删除此提交数据', 'danger')
            return redirect(url_for('dashboard'))
        
        db.delete(submission)
        db.commit()
        
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return jsonify({'success': True, 'message': '提交数据已删除'})
        
        flash('提交数据已删除', 'success')
        return redirect(url_for('task_data_view', task_id=task.id))
    finally:
        db.close()

@app.route('/clear_all_submissions/<int:task_id>', methods=['GET'])
@login_required
def clear_all_submissions(task_id):
    """清空所有提交数据"""
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'success': False, 'message': '任务不存在'})
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        
        if task.user_id != current_user.id:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'success': False, 'message': '无权删除此任务数据'})
            flash('无权删除此任务数据', 'danger')
            return redirect(url_for('dashboard'))
        
        db.query(Submission).filter_by(task_id=task.id).delete()
        db.commit()
        
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return jsonify({'success': True, 'message': '已清空所有提交数据'})
        
        flash('已清空所有提交数据', 'success')
        return redirect(url_for('task_data_view', task_id=task.id))
    finally:
        db.close()

@app.route('/delete_multiple_submissions/<int:task_id>', methods=['POST'])
@login_required
def delete_multiple_submissions(task_id):
    """批量删除提交数据"""
    db = SessionLocal()
    try:
        # 查询任务
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        
        # 检查用户权限
        if task.user_id != current_user.id:
            flash('无权删除此任务的提交数据', 'danger')
            return redirect(url_for('dashboard'))
        
        # 获取要删除的提交数据ID列表
        submission_ids = request.form.getlist('submission_ids')
        if not submission_ids:
            flash('请选择要删除的提交数据', 'warning')
            return redirect(url_for('task_detail', task_id=task_id))
        
        # 转换为整数并过滤
        submission_ids = [int(sid) for sid in submission_ids if sid.isdigit()]
        
        # 查询这些提交数据
        submissions = db.query(Submission).filter(
            Submission.id.in_(submission_ids),
            Submission.task_id == task_id
        ).all()
        
        # 检查是否所有提交数据都属于当前用户
        for sub in submissions:
            if sub.task.user_id != current_user.id:
                flash('无权删除部分提交数据', 'danger')
                return redirect(url_for('task_detail', task_id=task_id))
        
        # 删除提交数据
        for submission in submissions:
            db.delete(submission)
        
        db.commit()
        flash(f'已删除 {len(submissions)} 条提交数据', 'success')
        return redirect(url_for('task_detail', task_id=task_id))
    except ValueError:
        flash('无效的提交数据ID', 'danger')
        return redirect(url_for('task_detail', task_id=task_id))
    finally:
        db.close()

@app.route('/delete_attachment/<int:attachment_id>', methods=['POST'])
@login_required
def delete_attachment(attachment_id):
    import os
    db = SessionLocal()
    try:
        attachment = db.query(Attachment).get(attachment_id)
        if not attachment:
            return jsonify({'success': False, 'message': '附件不存在'})
        
        task = db.query(Task).get(attachment.task_id)
        if not task or task.user_id != current_user.id:
            return jsonify({'success': False, 'message': '无权删除此附件'})
        
        file_path = attachment.file_path
        if os.path.exists(file_path):
            os.remove(file_path)
        
        db.delete(attachment)
        db.commit()
        return jsonify({'success': True, 'message': '删除成功'})
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)})
    finally:
        db.close()

@app.route('/delete_task/<int:task_id>', methods=['POST'])
@login_required
def delete_task(task_id):
    import os
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权删除此任务', 'danger')
            return redirect(url_for('dashboard'))
        
        attachments = db.query(Attachment).filter_by(task_id=task.id).all()
        for attachment in attachments:
            file_path = os.path.join(app.root_path, 'static', attachment.file_path)
            if os.path.exists(file_path):
                os.remove(file_path)
        
        db.delete(task)
        db.commit()
        flash('任务已删除', 'success')
        return redirect(url_for('dashboard'))
    finally:
        db.close()

@app.route('/test_api_key', methods=['POST', 'OPTIONS'])
@login_required
def test_api_key():
    if request.method == 'OPTIONS':
        response = make_response()
        response.headers['Access-Control-Allow-Origin'] = '*'
        response.headers['Access-Control-Allow-Methods'] = 'POST, OPTIONS'
        response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
        return response

    try:
        data = request.get_json()
        model = data.get('model')
        api_key = data.get('api_key', '')
        api_url = data.get('api_url', '')
        model_name = data.get('model_name', '')

        if not model:
            return jsonify({'success': False, 'error': '缺少必要参数'}), 400

        class TestModelConfig:
            def __init__(self, model_name, api_key, api_url, extra_settings):
                self.model_name = model_name
                self.api_key = api_key
                self.api_url = api_url
                self.extra_settings = extra_settings

        class TestAIConfig:
            def __init__(self, model, api_key, api_url, model_name):
                self.selected_model = model
                self.model_configs = [TestModelConfig(model, api_key, api_url, model_name or ('llama3.2' if model == 'ollama' else ''))]

        test_config = TestAIConfig(model, api_key, api_url, model_name)

        test_prompt = '这是一个API密钥测试，请回复"测试成功"'

        result = call_ai_model(test_prompt, test_config)

        if result and ('测试成功' in result or 'success' in result.lower()):
            return jsonify({'success': True, 'message': 'API密钥有效'}), 200
        else:
            return jsonify({'success': True, 'message': 'API密钥有效，但返回内容不符合预期'}), 200

    except Exception as e:
        logger.error(f"API密钥测试失败: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/<string:task_id>/all', methods=['GET', 'OPTIONS'])
def get_all_submissions(task_id):
    # 永久重定向到/api/<string:task_id>路由
    return redirect(url_for('submit_form', task_id=task_id), code=301)

@app.route('/api/<string:task_id>', methods=['GET', 'POST', 'OPTIONS'])
def submit_form(task_id):
    # 处理预检请求
    if request.method == 'OPTIONS':
        response = make_response()
        response.headers['Access-Control-Allow-Origin'] = '*'
        response.headers['Access-Control-Allow-Methods'] = 'GET, POST, OPTIONS'
        response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
        return response
        
    db = SessionLocal()
    try:
        task = db.query(Task).filter_by(task_id=task_id).first()
        if not task:
            response = jsonify({'error': '任务不存在'})
            response.headers['Access-Control-Allow-Origin'] = '*'
            return response, 404
        
        if request.method == 'GET':
            # 返回所有回收的数据
            submissions = db.query(Submission).filter_by(task_id=task.id).all()
            all_data = []
            for sub in submissions:
                try:
                    data = json.loads(sub.data)
                except:
                    data = sub.data
                all_data.append({
                    'data': data,
                    'id': sub.id,
                    'submitted_at': sub.submitted_at.strftime('%Y-%m-%d %H:%M:%S')
                })
            # 直接构建JSON字符串以确保键的顺序（json 已在模块顶部 import）
            submission_count = len(all_data)
            response_data = {
                'note': f'Total {submission_count} submission(s).',
                'submissions': all_data,
                'task_id': task_id,
                'task_title': task.title,
                'total_submissions': submission_count
            }
            # 使用json.dumps确保键的顺序
            json_response = json.dumps(response_data, ensure_ascii=False, sort_keys=False)
            response = make_response(json_response)
            response.headers['Content-Type'] = 'application/json; charset=utf-8'
            response.headers['Access-Control-Allow-Origin'] = '*'
            return response, 200
        
        # 处理POST请求 - 回收数据
        form_data = {}

        # 检查Content-Type并选择合适的数据获取方式
        if request.is_json:
            # 如果是JSON请求，尝试获取JSON数据
            try:
                form_data = request.get_json() or {}
            except Exception as e:
                logger.error(f"解析JSON数据失败: {str(e)}")
                form_data = {}
        else:
            # 如果不是JSON请求，获取表单数据
            form_data = request.form.to_dict()

        # 如果表单数据仍然为空，尝试从请求体获取原始数据
        if not form_data:
            try:
                form_data = request.get_data(as_text=True)
            except Exception as e:
                logger.error(f"获取请求体数据失败: {str(e)}")
                form_data = {}

        # B16: 处理附件上传（multipart with file fields）
        attachments_info = []
        upload_errors = []
        if request.files:
            # 只在 form_data 是 dict 时才能附加 _attachments；如果是字符串（raw body）就不动
            attachments_info, upload_errors = _save_quickform_uploads(request.files, task.task_id)
            if attachments_info and isinstance(form_data, dict):
                form_data['_attachments'] = attachments_info

        # 序列化：如果带附件 → JSON（保证 _attachments 可解析）；否则保持原 str(form_data) 行为不破坏旧消费者
        if attachments_info and isinstance(form_data, dict):
            stored_data = json.dumps(form_data, ensure_ascii=False)
        else:
            stored_data = str(form_data)
        submission = Submission(task_id=task.id, data=stored_data)
        db.add(submission)
        db.commit()

        response_payload = {'message': '提交成功'}
        if attachments_info:
            response_payload['attachments'] = [
                {'name': a['name'], 'url': a['url'], 'size': a['size']} for a in attachments_info
            ]
        if upload_errors:
            response_payload['upload_warnings'] = upload_errors
        response = jsonify(response_payload)
        response.headers['Access-Control-Allow-Origin'] = '*'
        return response, 200
    finally:
        db.close()

@app.route('/export/<int:task_id>')
@login_required
def export_data(task_id):
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task or task.user_id != current_user.id:
            flash('无权访问此数据', 'danger')
            return redirect(url_for('dashboard'))
        
        submission = db.query(Submission).filter_by(task_id=task.id).all()
        
        if not submission:
            flash('没有可导出的数据', 'info')
            return redirect(url_for('task_detail', task_id=task_id))
        
        # 尝试解析提交数据并转换为DataFrame
        data_list = []
        for sub in submission:
            try:
                data = json.loads(sub.data)
                data['submitted_at'] = sub.submitted_at.strftime('%Y-%m-%d %H:%M:%S')
                data_list.append(data)
            except:
                # 如果解析失败，添加原始数据
                data_list.append({
                    'submitted_at': sub.submitted_at.strftime('%Y-%m-%d %H:%M:%S'),
                    'raw_data': sub.data
                })
        
        import pandas as pd
        df = pd.DataFrame(data_list)
        
        # 创建CSV文件
        output = io.BytesIO()
        df.to_csv(output, index=False, encoding='utf-8-sig')
        output.seek(0)
        
        # 发送文件（兼容不同版本的Flask）
        filename = f"{task.title}_数据导出_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        try:
            # 尝试使用新版本Flask的参数
            return send_file(output, download_name=filename, as_attachment=True, mimetype='text/csv; charset=utf-8')
        except TypeError:
            # 如果新参数不被支持，回退到旧版本的参数
            return send_file(output, attachment_filename=filename, as_attachment=True, mimetype='text/csv; charset=utf-8')
    except Exception as e:
        flash(f'导出数据时出错: {str(e)}', 'danger')
        return redirect(url_for('task_detail', task_id=task_id))
    finally:
        db.close()

@app.route('/export_json/<int:task_id>')
@login_required
def export_json(task_id):
    """
    导出任务提交数据为JSON格式
    """
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task or task.user_id != current_user.id:
            flash('无权访问此数据', 'danger')
            return redirect(url_for('dashboard'))
        
        submission = db.query(Submission).filter_by(task_id=task.id).all()
        
        if not submission:
            flash('没有可导出的数据', 'info')
            return redirect(url_for('task_detail', task_id=task_id))
        
        # 构建JSON数据
        data_list = []
        for sub in submission:
            try:
                data = json.loads(sub.data)
                data['_submitted_at'] = sub.submitted_at.strftime('%Y-%m-%d %H:%M:%S')
                data['_submission_id'] = sub.id
                data_list.append(data)
            except:
                # 如果解析失败，添加原始数据
                data_list.append({
                    '_submitted_at': sub.submitted_at.strftime('%Y-%m-%d %H:%M:%S'),
                    '_submission_id': sub.id,
                    '_raw_data': sub.data
                })
        
        # 创建JSON输出
        output = io.BytesIO()
        json_data = {
            'task_title': task.title,
            'task_id': task.id,
            'export_time': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'total_records': len(data_list),
            'data': data_list
        }
        output.write(json.dumps(json_data, ensure_ascii=False, indent=2).encode('utf-8'))
        output.seek(0)
        
        # 发送文件（兼容不同版本的Flask）
        filename = f"{task.title}_数据导出_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        try:
            # 尝试使用新版本Flask的参数
            return send_file(output, download_name=filename, as_attachment=True, mimetype='application/json; charset=utf-8')
        except TypeError:
            # 如果新参数不被支持，回退到旧版本的参数
            return send_file(output, attachment_filename=filename, as_attachment=True, mimetype='application/json; charset=utf-8')
    except Exception as e:
        flash(f'导出数据时出错: {str(e)}', 'danger')
        return redirect(url_for('task_detail', task_id=task_id))
    finally:
        db.close()

@app.route('/api/qf/test_connection', methods=['POST'])
@login_required
def test_qf_connection():
    import requests
    db = SessionLocal()
    try:
        qf_config = db.query(QFConfig).filter_by(user_id=current_user.id).first()
        if not qf_config or not qf_config.username or not qf_config.password:
            return jsonify({'success': False, 'message': '请先保存用户名和密码'})

        try:
            response = requests.post(
                'https://quickform.cn/cli/list',
                json={'username': qf_config.username, 'password': qf_config.password},
                timeout=10
            )
            result = response.json()
            if result.get('success'):
                return jsonify({'success': True, 'message': '连接成功', 'tasks': result.get('tasks', [])})
            else:
                return jsonify({'success': False, 'message': result.get('message', '认证失败')})
        except Exception as e:
            return jsonify({'success': False, 'message': f'连接失败: {str(e)}'})
    finally:
        db.close()

@app.route('/api/qf/list', methods=['GET'])
@login_required
def get_qf_task_list():
    import requests
    db = SessionLocal()
    try:
        qf_config = db.query(QFConfig).filter_by(user_id=current_user.id).first()
        if not qf_config or not qf_config.username or not qf_config.password:
            return jsonify({'success': False, 'message': '请先在设置中配置QF数据互联'})

        try:
            response = requests.post(
                'https://quickform.cn/cli/list',
                json={'username': qf_config.username, 'password': qf_config.password},
                timeout=10
            )
            result = response.json()
            if result.get('success'):
                return jsonify({'success': True, 'tasks': result.get('tasks', [])})
            else:
                return jsonify({'success': False, 'message': result.get('message', '认证失败')})
        except Exception as e:
            return jsonify({'success': False, 'message': f'连接失败: {str(e)}'})
    finally:
        db.close()

@app.route('/api/system/init', methods=['POST'])
@login_required
def system_init():
    if is_campus_edition():
        return jsonify({'success': False, 'message': '校园版不支持从个人设置页执行系统初始化'}), 403
    import os
    db = SessionLocal()
    try:
        try:
            ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
            if ai_config:
                db.query(AIModelConfig).filter(
                    AIModelConfig.ai_config_id == ai_config.id,
                    AIModelConfig.model_name != 'ollama'
                ).delete(synchronize_session=False)
                ollama_cfg = db.query(AIModelConfig).filter(
                    AIModelConfig.ai_config_id == ai_config.id,
                    AIModelConfig.model_name == 'ollama'
                ).first()
                if ollama_cfg:
                    ollama_cfg.api_key = ''
                if not ollama_cfg:
                    ollama_cfg = AIModelConfig(
                        ai_config_id=ai_config.id,
                        model_name='ollama',
                        api_key='',
                        api_url='http://localhost:11434',
                        extra_settings='llama3'
                    )
                    db.add(ollama_cfg)
            
            qf_configs = db.query(QFConfig).filter_by(user_id=current_user.id).all()
            for qf in qf_configs:
                db.delete(qf)
            
            user = db.query(User).filter_by(id=current_user.id).first()
            if user:
                user.username = 'admin'
                user.password = generate_password_hash('openmentor')
            
            all_tasks = db.query(Task).filter_by(user_id=current_user.id).order_by(Task.id).all()
            tasks_to_delete = all_tasks[3:] if len(all_tasks) > 3 else []
            
            for task in tasks_to_delete:
                attachments = db.query(Attachment).filter_by(task_id=task.id).all()
                for att in attachments:
                    if os.path.exists(att.file_path):
                        os.remove(att.file_path)
                    db.delete(att)
                db.delete(task)
            
            db.commit()
            return jsonify({'success': True, 'message': '系统初始化成功'})
        except Exception as e:
            db.rollback()
            return jsonify({'success': False, 'message': f'初始化失败: {str(e)}'})
    finally:
        db.close()

@app.route('/profile', methods=['GET', 'POST'])
@login_required
def profile():
    db = SessionLocal()
    try:
        next_page = request.values.get('next', '').strip()
        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()

        if request.method == 'POST':
            password_changed = False
            if 'selected_model' in request.form:
                selected_model = request.form.get('selected_model')

                if not ai_config:
                    ai_config = AIConfig(user_id=current_user.id, selected_model=selected_model)
                    db.add(ai_config)
                    db.flush()
                else:
                    ai_config.selected_model = selected_model

                # OpenMentor: 改用 UPSERT 语义（不再先 DELETE 再重建）
                # 防御性：若表单中某字段为空字符串，保留 DB 中已有值
                #   - 这样即使部分提交（比如脚本/测试只填一部分），不会误清空其他模型的 Key
                #   - 用户清空字段后保存 → 仍然保留旧值（推荐通过专门的「清除」按钮触发清除）
                # 用户可主动设置某字段为空（高级场景）：通过 image_gen_model 这种"非必填"字段，
                #   form 提交空字符串 + DB 已有非空 → 保留 DB；DB 为空 → 保持空。
                existing = {mc.model_name: mc for mc in ai_config.model_configs}

                def _pick(form_val, existing_val):
                    """空字符串 → 沿用旧值；非空 → 用新值"""
                    return form_val if (form_val or '').strip() else (existing_val or '')

                model_configs = [
                    ('deepseek',    request.form.get('deepseek_api_key', ''),    '', request.form.get('deepseek_model', '').strip(),       '', ''),
                    ('doubao',      request.form.get('doubao_api_key', ''),      '', request.form.get('doubao_model', '').strip(),         request.form.get('doubao_image_gen_model', '').strip(), request.form.get('doubao_video_gen_model', '').strip()),
                    ('qwen',        request.form.get('qwen_api_key', ''),        '', request.form.get('qwen_model', '').strip(),           request.form.get('qwen_image_gen_model', '').strip(), ''),
                    ('glm',         request.form.get('glm_api_key', ''),         '', request.form.get('glm_model', '').strip(),            request.form.get('glm_image_gen_model', '').strip(), ''),
                    ('siliconflow', request.form.get('siliconflow_api_key', ''), '', request.form.get('siliconflow_model', 'Qwen/Qwen2.5-72B-Instruct').strip(), request.form.get('siliconflow_image_gen_model', '').strip(), ''),
                    ('cherry',      request.form.get('cherry_api_key', ''),      request.form.get('cherry_api_url', '').strip(), request.form.get('cherry_model', '').strip(), request.form.get('cherry_image_gen_model', '').strip(), ''),
                    ('ollama',      '', request.form.get('ollama_api_url', 'http://localhost:11434'), request.form.get('ollama_model', 'llama3.2').strip(), '', ''),
                ]

                for model_name, api_key_in, api_url_in, extra_in, image_gen_in, video_gen_in in model_configs:
                    cfg = existing.get(model_name)
                    if cfg is None:
                        # 没有则新建（仅当至少有一项非空）
                        if not (api_key_in or api_url_in or extra_in or image_gen_in or video_gen_in):
                            continue
                        cfg = AIModelConfig(
                            ai_config_id=ai_config.id,
                            model_name=model_name,
                            api_key=api_key_in,
                            api_url=api_url_in,
                            extra_settings=extra_in,
                            image_gen_model=image_gen_in or None,
                            video_gen_model=video_gen_in or None,
                        )
                        db.add(cfg)
                    else:
                        # 存在则更新：空值沿用旧值（防丢 key）
                        cfg.api_key = _pick(api_key_in, cfg.api_key)
                        cfg.api_url = _pick(api_url_in, cfg.api_url)
                        cfg.extra_settings = _pick(extra_in, cfg.extra_settings)
                        cfg.image_gen_model = _pick(image_gen_in, cfg.image_gen_model) or None
                        cfg.video_gen_model = _pick(video_gen_in, cfg.video_gen_model) or None

                db.commit()
                flash('AI配置更新成功', 'success')

            elif 'update_qf_config' in request.form:
                qf_username = request.form.get('qf_username', '').strip()
                qf_password = request.form.get('qf_password', '').strip()
                
                qf_config = db.query(QFConfig).filter_by(user_id=current_user.id).first()
                if not qf_config:
                    qf_config = QFConfig(user_id=current_user.id, username=qf_username, password=qf_password)
                    db.add(qf_config)
                else:
                    qf_config.username = qf_username
                    qf_config.password = qf_password
                
                db.commit()
                flash('QF配置更新成功', 'success')

            elif 'change_username' in request.form:
                new_username = request.form.get('username', '').strip()
                user = db.query(User).filter_by(id=current_user.id).first()
                exists = db.query(User).filter(User.username == new_username, User.id != current_user.id).first() if new_username else None
                if exists:
                    flash('用户名已存在', 'danger')
                elif user and new_username:
                    user.username = new_username
                    db.commit()
                    flash('用户名修改成功', 'success')
                else:
                    flash('用户名修改失败', 'danger')

            elif 'change_password' in request.form:
                current_password = request.form.get('current_password')
                new_password = (request.form.get('new_password') or '').strip()
                confirm_password = (request.form.get('confirm_password') or '').strip()

                user = db.query(User).filter_by(id=current_user.id).first()
                if len(new_password) < 6:
                    flash('新密码长度至少为 6 个字符', 'danger')
                elif new_password != confirm_password:
                    flash('两次输入的新密码不一致', 'danger')
                elif user and check_password_hash(user.password, current_password):
                    user.password = generate_password_hash(new_password)
                    user.force_password_change = False
                    db.commit()
                    password_changed = True
                    flash('密码修改成功', 'success')
                else:
                    flash('当前密码错误', 'danger')

            active_tab = request.form.get('active_tab', 'config')
            if password_changed and next_page and _is_safe_next_url(next_page):
                return redirect(next_page)
            profile_args = {'active_tab': active_tab}
            if next_page and _is_safe_next_url(next_page):
                profile_args['next'] = next_page
            return redirect(url_for('profile', **profile_args))

        model_configs_dict = {}
        if ai_config:
            for mc in ai_config.model_configs:
                model_configs_dict[mc.model_name] = mc

        qf_config = db.query(QFConfig).filter_by(user_id=current_user.id).first()

        return render_template(
            'profile.html',
            user=current_user,
            ai_config=ai_config,
            model_configs_dict=model_configs_dict,
            qf_config=qf_config,
            next_page=next_page if _is_safe_next_url(next_page) else '',
        )
    finally:
        db.close()


def _generate_temp_password(length=10):
    alphabet = string.ascii_letters + string.digits
    return ''.join(random.choices(alphabet, k=length))


def _legacy_teacher_email(username):
    """兼容旧 user.email 非空唯一约束；校园版界面不再使用邮箱登录或展示。"""
    safe = ''.join(ch if ch.isalnum() or ch in ('-', '_', '.') else '_' for ch in username)
    return f'{safe or "teacher"}@openmentor.local'


def _parse_teacher_import_file(file_storage):
    """解析教师账号导入文件，表头支持：手机号/账号/用户名 + 姓名/显示名。"""
    filename = (file_storage.filename or '').lower()
    if not filename:
        return None, '请选择要上传的 Excel 文件'

    try:
        if filename.endswith('.xlsx'):
            from openpyxl import load_workbook
            wb = load_workbook(file_storage, read_only=True, data_only=True)
            ws = wb.active
            rows = [list(r) for r in ws.iter_rows(values_only=True)]
        elif filename.endswith('.csv'):
            import csv
            import io
            raw = file_storage.read()
            text_data = None
            for enc in ('utf-8-sig', 'utf-8', 'gbk', 'gb18030'):
                try:
                    text_data = raw.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if text_data is None:
                return None, 'CSV 编码无法识别，请另存为 UTF-8 后再试'
            rows = list(csv.reader(io.StringIO(text_data)))
        else:
            return None, '请上传 .xlsx 或 .csv 文件'
    except Exception as e:
        logger.exception('[校园版] 教师账号导入文件解析失败')
        return None, f'解析失败：{e}'

    rows = [r for r in rows if any(_normalize_text(c) for c in r)]
    if len(rows) < 2:
        return None, '文件内容为空或仅有表头'

    headers = [_normalize_text(c).lower() for c in rows[0]]
    username_aliases = {'手机号', '手机', '账号', '用户名', '登录账号', '联系电话', 'phone', 'mobile', 'username'}
    name_aliases = {'姓名', '教师姓名', '老师姓名', '显示名', '名称', 'name', 'displayname', 'display_name'}
    username_idx = next((i for i, h in enumerate(headers) if h in username_aliases), None)
    name_idx = next((i for i, h in enumerate(headers) if h in name_aliases), None)

    if username_idx is None or name_idx is None:
        return None, '未识别到「手机号/账号」和「姓名」两列，请使用模板或检查表头'

    entries = []
    seen = set()
    for row_no, row in enumerate(rows[1:], start=2):
        username = _normalize_text(row[username_idx] if username_idx < len(row) else '')
        display_name = _normalize_text(row[name_idx] if name_idx < len(row) else '')
        if not username and not display_name:
            continue
        if not username or not display_name:
            return None, f'第 {row_no} 行缺少手机号/账号或姓名'
        if username in seen:
            return None, f'导入文件中账号重复：{username}'
        seen.add(username)
        entries.append({'username': username, 'display_name': display_name})

    if not entries:
        return None, '没有可导入的教师账号'
    return entries, None


@app.route('/campus/teachers/template.xlsx')
@login_required
@campus_admin_required
def campus_teachers_template():
    from openpyxl import Workbook
    import io
    wb = Workbook()
    ws = wb.active
    ws.title = '教师账号'
    ws.append(['手机号', '姓名'])
    ws.append(['13800000001', '张老师'])
    ws.append(['13800000002', '李老师'])
    ws.column_dimensions['A'].width = 18
    ws.column_dimensions['B'].width = 14
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return send_file(
        buf,
        as_attachment=True,
        download_name='openmentor_teacher_accounts_template.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    )


@app.route('/campus/teachers', methods=['GET', 'POST'])
@login_required
@campus_admin_required
def campus_teachers():
    db = SessionLocal()
    try:
        generated_password = None
        if request.method == 'POST':
            username = (request.form.get('username') or '').strip()
            display_name = (request.form.get('display_name') or '').strip()
            temp_password = (request.form.get('temp_password') or '').strip() or 'openmentor'

            if not username:
                flash('请填写手机号/账号', 'danger')
            elif len(temp_password) < 6:
                flash('临时密码长度至少为 6 个字符', 'danger')
            elif db.query(User).filter_by(username=username).first():
                flash('账号已存在', 'danger')
            else:
                teacher = User(
                    username=username,
                    email=_legacy_teacher_email(username),
                    display_name=display_name or username,
                    password=generate_password_hash(temp_password),
                    role='teacher',
                    is_active_flag=True,
                    force_password_change=True,
                )
                db.add(teacher)
                db.commit()
                generated_password = temp_password
                flash(f'教师账号「{username}」已创建，临时密码：{temp_password}', 'success')

        search_query = (request.args.get('q') or '').strip()
        try:
            page = int(request.args.get('page', 1))
        except (TypeError, ValueError):
            page = 1
        try:
            per_page = int(request.args.get('per_page', 10))
        except (TypeError, ValueError):
            per_page = 10
        if per_page not in (10, 20, 50):
            per_page = 10
        if page < 1:
            page = 1

        user_query = db.query(User)
        if search_query:
            like = f'%{search_query}%'
            user_query = user_query.filter(
                (User.display_name.ilike(like)) |
                (User.username.ilike(like))
            )

        total_users = user_query.count()
        total_pages = max((total_users + per_page - 1) // per_page, 1)
        if page > total_pages:
            page = total_pages

        users = (
            user_query
            .order_by(User.role, User.id)
            .offset((page - 1) * per_page)
            .limit(per_page)
            .all()
        )
        start_index = (page - 1) * per_page + 1 if total_users else 0
        end_index = min(page * per_page, total_users)

        stats = {}
        for user in users:
            task_count = db.query(Task).filter_by(user_id=user.id, is_assistant=False).count()
            assistant_count = db.query(Task).filter_by(user_id=user.id, is_assistant=True).count()
            roster_count = db.query(RosterEntry).filter_by(user_id=user.id).count()
            stats[user.id] = {
                'task_count': task_count,
                'assistant_count': assistant_count,
                'roster_count': roster_count,
            }

        return render_template(
            'campus_teachers.html',
            users=users,
            stats=stats,
            generated_password=generated_password,
            search_query=search_query,
            page=page,
            per_page=per_page,
            per_page_options=(10, 20, 50),
            total_users=total_users,
            total_pages=total_pages,
            start_index=start_index,
            end_index=end_index,
        )
    finally:
        db.close()


@app.route('/campus/teachers/import', methods=['POST'])
@login_required
@campus_admin_required
def campus_teachers_import():
    file_storage = request.files.get('teacher_file')
    entries, error = _parse_teacher_import_file(file_storage) if file_storage else (None, '请选择要上传的 Excel 文件')
    if error:
        flash(error, 'danger')
        return redirect(url_for('campus_teachers'))

    db = SessionLocal()
    try:
        created = 0
        skipped = []
        for entry in entries:
            username = entry['username']
            if db.query(User).filter_by(username=username).first():
                skipped.append(username)
                continue
            teacher = User(
                username=username,
                email=_legacy_teacher_email(username),
                display_name=entry['display_name'],
                password=generate_password_hash('openmentor'),
                role='teacher',
                is_active_flag=True,
                force_password_change=True,
            )
            db.add(teacher)
            created += 1
        db.commit()

        msg = f'已导入 {created} 个教师账号，默认密码：openmentor，首次登录需修改密码。'
        if skipped:
            msg += f' 已跳过 {len(skipped)} 个已存在账号：{", ".join(skipped[:8])}'
            if len(skipped) > 8:
                msg += ' 等'
        flash(msg, 'success' if created else 'warning')
        return redirect(url_for('campus_teachers'))
    except Exception as e:
        db.rollback()
        logger.exception('[校园版] 批量导入教师账号失败')
        flash(f'导入失败：{e}', 'danger')
        return redirect(url_for('campus_teachers'))
    finally:
        db.close()


def _campus_period_since(period):
    now = datetime.now()
    if period == 'today':
        return datetime(now.year, now.month, now.day)
    if period == '7d':
        return now - timedelta(days=7)
    if period == '90d':
        return now - timedelta(days=90)
    if period == 'all':
        return None
    return now - timedelta(days=30)


def _campus_apply_since(query, column, since):
    return query.filter(column >= since) if since else query


def _campus_attachment_type(path):
    ext = os.path.splitext(path or '')[1].lower().lstrip('.')
    if ext in {'jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp', 'svg'}:
        return '图片'
    if ext in {'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx', 'pdf', 'txt', 'md', 'csv'}:
        return '文档'
    if ext in {'mp3', 'wav', 'm4a', 'aac', 'flac', 'mp4', 'mov', 'avi', 'mkv', 'webm'}:
        return '音视频'
    return '其他'


@app.route('/campus/usage')
@login_required
@campus_admin_required
def campus_usage_dashboard():
    period_options = [
        ('30d', '近30天'),
        ('7d', '近7天'),
        ('90d', '近90天'),
        ('today', '今天'),
        ('all', '全部'),
    ]
    period_values = {k for k, _ in period_options}
    period = request.args.get('period', '30d')
    if period not in period_values:
        period = '30d'
    since = _campus_period_since(period)

    try:
        teacher_id = int(request.args.get('teacher_id', 0) or 0)
    except (TypeError, ValueError):
        teacher_id = 0

    db = SessionLocal()
    try:
        teachers = (
            db.query(User)
            .filter(User.role.in_(('admin', 'teacher')))
            .order_by(User.role, User.display_name, User.username)
            .all()
        )
        teacher_ids = {u.id for u in teachers}
        if teacher_id not in teacher_ids:
            teacher_id = 0

        def task_scope(query):
            return query.filter(Task.user_id == teacher_id) if teacher_id else query

        def submission_scope(query):
            query = query.join(Task, Submission.task_id == Task.id)
            return query.filter(Task.user_id == teacher_id) if teacher_id else query

        def attachment_scope(query):
            query = query.join(Task, Attachment.task_id == Task.id)
            return query.filter(Task.user_id == teacher_id) if teacher_id else query

        def conversation_scope(query):
            query = query.join(Task, Conversation.assistant_id == Task.id)
            return query.filter(Task.user_id == teacher_id) if teacher_id else query

        def message_scope(query):
            query = query.join(Conversation, Message.conversation_id == Conversation.id).join(Task, Conversation.assistant_id == Task.id)
            return query.filter(Task.user_id == teacher_id) if teacher_id else query

        def user_scope(query):
            return query.filter(User.id == teacher_id) if teacher_id else query

        def campus_plaza_scope(query):
            return query.filter(CampusPlazaTemplate.owner_user_id == teacher_id) if teacher_id else query

        users_query = user_scope(db.query(User).filter(User.role.in_(('admin', 'teacher'))))
        total_users = users_query.count()
        active_users = users_query.filter(User.is_active_flag == True).count()
        forced_password_users = users_query.filter(User.force_password_change == True).count()
        new_users = _campus_apply_since(users_query, User.created_at, since).count()

        ai_assistants = _campus_apply_since(
            task_scope(db.query(Task)).filter(Task.is_assistant == True),
            Task.created_at,
            since,
        ).count()
        data_tasks = _campus_apply_since(
            task_scope(db.query(Task)).filter(Task.is_assistant == False),
            Task.created_at,
            since,
        ).count()
        conversations = _campus_apply_since(
            conversation_scope(db.query(Conversation)),
            Conversation.started_at,
            since,
        ).count()
        submissions = _campus_apply_since(
            submission_scope(db.query(Submission)),
            Submission.submitted_at,
            since,
        ).count()
        attachments = _campus_apply_since(
            attachment_scope(db.query(Attachment)),
            Attachment.created_at,
            since,
        ).count()
        messages_total = _campus_apply_since(
            message_scope(db.query(Message)),
            Message.created_at,
            since,
        ).count()
        student_messages = _campus_apply_since(
            message_scope(db.query(Message)).filter(Message.role == 'user'),
            Message.created_at,
            since,
        ).count()
        assistant_messages = _campus_apply_since(
            message_scope(db.query(Message)).filter(Message.role == 'assistant'),
            Message.created_at,
            since,
        ).count()
        safety_hits = _campus_apply_since(
            message_scope(db.query(Message)).filter(Message.triggered_keyword.isnot(None)),
            Message.created_at,
            since,
        ).count()
        image_generations = _campus_apply_since(
            message_scope(db.query(Message)).filter(Message.generated_image_path.isnot(None)),
            Message.created_at,
            since,
        ).count()

        campus_shared_query = _campus_apply_since(
            campus_plaza_scope(db.query(CampusPlazaTemplate)),
            CampusPlazaTemplate.created_at,
            since,
        )
        campus_shared_templates = campus_shared_query.count()
        campus_shared_active = campus_shared_query.filter(CampusPlazaTemplate.status == 'active').count()
        campus_shared_disabled = campus_shared_query.filter(CampusPlazaTemplate.status == 'disabled').count()

        roster_query = db.query(RosterEntry)
        if teacher_id:
            roster_query = roster_query.filter(RosterEntry.user_id == teacher_id)
        roster_entries = _campus_apply_since(roster_query, RosterEntry.created_at, since).count()

        report_query = db.query(StudentReport)
        if teacher_id:
            report_query = report_query.filter(StudentReport.user_id == teacher_id)
        student_reports = _campus_apply_since(report_query, StudentReport.created_at, since).count()

        usage_count = student_messages + submissions
        primary_metrics = [
            {'label': '用户数', 'value': total_users, 'sub': f'启用 {active_users} · 待改密 {forced_password_users}', 'icon': 'bi-people', 'tone': 'ink'},
            {'label': '使用次数', 'value': usage_count, 'sub': f'学生消息 {student_messages} · 数据提交 {submissions}', 'icon': 'bi-lightning-charge', 'tone': 'green'},
            {'label': 'AI导师数', 'value': ai_assistants, 'sub': f'对话 {conversations} · 回复 {assistant_messages}', 'icon': 'bi-robot', 'tone': 'blue'},
            {'label': '数据任务', 'value': data_tasks, 'sub': f'提交 {submissions} · 附件 {attachments}', 'icon': 'bi-list-task', 'tone': 'amber'},
            {'label': '对话数量', 'value': conversations, 'sub': f'消息 {messages_total} · 安全拦截 {safety_hits}', 'icon': 'bi-chat-dots', 'tone': 'red'},
            {'label': '附件数量', 'value': attachments, 'sub': f'报告 {student_reports} · 花名册 {roster_entries}', 'icon': 'bi-paperclip', 'tone': 'violet'},
            {'label': '校内共享', 'value': campus_shared_templates, 'sub': f'启用 {campus_shared_active} · 禁用 {campus_shared_disabled}', 'icon': 'bi-share', 'tone': 'green'},
        ]

        extra_metrics = [
            {'label': '新增用户', 'value': new_users},
            {'label': '学生消息', 'value': student_messages},
            {'label': 'AI回复', 'value': assistant_messages},
            {'label': '图片生成', 'value': image_generations},
            {'label': '安全拦截', 'value': safety_hits},
            {'label': '学情报告', 'value': student_reports},
            {'label': '花名册', 'value': roster_entries},
            {'label': '共享导师', 'value': campus_shared_templates},
        ]

        attachment_rows = _campus_apply_since(
            attachment_scope(db.query(Attachment.file_name, Attachment.file_path, Attachment.created_at)),
            Attachment.created_at,
            since,
        ).all()
        attachment_type_counts = {'图片': 0, '文档': 0, '音视频': 0, '其他': 0}
        for row in attachment_rows:
            attachment_type_counts[_campus_attachment_type(row.file_name or row.file_path)] += 1
        attachment_type_items = [{'label': k, 'value': v} for k, v in attachment_type_counts.items()]

        trend_start = since
        if not trend_start:
            trend_start = datetime.now() - timedelta(days=13)
        max_days = 30
        day_count = max(1, min(max_days, (datetime.now().date() - trend_start.date()).days + 1))
        trend_start_date = datetime.now().date() - timedelta(days=day_count - 1)
        trend = {
            (trend_start_date + timedelta(days=i)).isoformat(): {
                'label': (trend_start_date + timedelta(days=i)).strftime('%m-%d'),
                'messages': 0,
                'submissions': 0,
            }
            for i in range(day_count)
        }

        msg_rows = message_scope(
            db.query(Message.created_at).filter(
                Message.role == 'user',
                Message.created_at >= datetime.combine(trend_start_date, datetime.min.time())
            )
        ).all()
        for (created_at,) in msg_rows:
            if created_at:
                key = created_at.date().isoformat()
                if key in trend:
                    trend[key]['messages'] += 1

        sub_rows = submission_scope(
            db.query(Submission.submitted_at).filter(
                Submission.submitted_at >= datetime.combine(trend_start_date, datetime.min.time())
            )
        ).all()
        for (submitted_at,) in sub_rows:
            if submitted_at:
                key = submitted_at.date().isoformat()
                if key in trend:
                    trend[key]['submissions'] += 1

        trend_items = list(trend.values())
        trend_max = max([item['messages'] + item['submissions'] for item in trend_items] + [1])

        scoped_teacher_ids = [teacher_id] if teacher_id else list(teacher_ids)
        if scoped_teacher_ids:
            assistant_by_teacher = {
                uid: count for uid, count in db.query(Task.user_id, func.count(Task.id))
                .filter(Task.user_id.in_(scoped_teacher_ids), Task.is_assistant == True)
                .group_by(Task.user_id)
                .all()
            }
            data_task_by_teacher = {
                uid: count for uid, count in db.query(Task.user_id, func.count(Task.id))
                .filter(Task.user_id.in_(scoped_teacher_ids), Task.is_assistant == False)
                .group_by(Task.user_id)
                .all()
            }
            submission_by_teacher = {
                uid: count for uid, count in submission_scope(db.query(Task.user_id, func.count(Submission.id)))
                .group_by(Task.user_id)
                .all()
            }
            conversation_by_teacher = {
                uid: count for uid, count in conversation_scope(db.query(Task.user_id, func.count(Conversation.id)))
                .group_by(Task.user_id)
                .all()
            }
            attachment_by_teacher = {
                uid: count for uid, count in attachment_scope(db.query(Task.user_id, func.count(Attachment.id)))
                .group_by(Task.user_id)
                .all()
            }
        else:
            assistant_by_teacher = {}
            data_task_by_teacher = {}
            submission_by_teacher = {}
            conversation_by_teacher = {}
            attachment_by_teacher = {}

        teacher_lookup = {u.id: u for u in teachers}
        teacher_rows = []
        for uid in scoped_teacher_ids:
            user = teacher_lookup.get(uid)
            if not user:
                continue
            row_usage = (submission_by_teacher.get(uid, 0) or 0) + (conversation_by_teacher.get(uid, 0) or 0)
            teacher_rows.append({
                'name': user.display_name or user.username,
                'username': user.username,
                'is_active': user.is_active,
                'usage': row_usage,
                'assistants': assistant_by_teacher.get(uid, 0) or 0,
                'data_tasks': data_task_by_teacher.get(uid, 0) or 0,
                'submissions': submission_by_teacher.get(uid, 0) or 0,
                'conversations': conversation_by_teacher.get(uid, 0) or 0,
                'attachments': attachment_by_teacher.get(uid, 0) or 0,
            })
        teacher_rows = sorted(teacher_rows, key=lambda r: (r['usage'], r['assistants'], r['data_tasks']), reverse=True)[:10]

        recent_conversations = (
            conversation_scope(db.query(Conversation, Task, User))
            .join(User, Task.user_id == User.id)
            .order_by(Conversation.last_active_at.desc())
            .limit(8)
            .all()
        )
        recent_items = [
            {
                'teacher': user.display_name or user.username,
                'assistant': task.title,
                'student': f'{conv.student_class} · {conv.student_name}',
                'time': conv.last_active_at.strftime('%m-%d %H:%M') if conv.last_active_at else '-',
            }
            for conv, task, user in recent_conversations
        ]

        selected_teacher = teacher_lookup.get(teacher_id)
        return render_template(
            'campus_usage.html',
            period=period,
            period_options=period_options,
            teacher_id=teacher_id,
            teachers=teachers,
            selected_teacher=selected_teacher,
            since=since,
            primary_metrics=primary_metrics,
            extra_metrics=extra_metrics,
            attachment_type_items=attachment_type_items,
            trend_items=trend_items,
            trend_max=trend_max,
            teacher_rows=teacher_rows,
            recent_items=recent_items,
        )
    finally:
        db.close()


@app.route('/campus/teachers/<int:user_id>/toggle', methods=['POST'])
@login_required
@campus_admin_required
def campus_teacher_toggle(user_id):
    if user_id == current_user.id:
        flash('不能禁用当前登录的管理员账号', 'danger')
        return redirect(url_for('campus_teachers'))

    db = SessionLocal()
    try:
        user = db.query(User).get(user_id)
        if not user:
            flash('账号不存在', 'danger')
        elif user.role == 'admin':
            flash('暂不允许在这里禁用管理员账号', 'danger')
        else:
            user.is_active_flag = not user.is_active_flag
            db.commit()
            flash(f'已{"启用" if user.is_active_flag else "禁用"}教师账号「{user.username}」', 'success')
        return redirect(url_for('campus_teachers'))
    finally:
        db.close()


@app.route('/campus/teachers/<int:user_id>/reset_password', methods=['POST'])
@login_required
@campus_admin_required
def campus_teacher_reset_password(user_id):
    temp_password = (request.form.get('temp_password') or '').strip() or _generate_temp_password()
    if len(temp_password) < 6:
        flash('临时密码长度至少为 6 个字符', 'danger')
        return redirect(url_for('campus_teachers'))

    db = SessionLocal()
    try:
        user = db.query(User).get(user_id)
        if not user:
            flash('账号不存在', 'danger')
        elif user.role == 'admin' and user.id != current_user.id:
            flash('暂不允许重置其他管理员密码', 'danger')
        else:
            user.password = generate_password_hash(temp_password)
            user.force_password_change = True
            db.commit()
            flash(f'账号「{user.username}」密码已重置，临时密码：{temp_password}', 'success')
        return redirect(url_for('campus_teachers'))
    finally:
        db.close()

@app.route('/analyze/<int:task_id>/smart_analyze', methods=['GET'])
@login_required
def smart_analyze(task_id):
    """
    智能分析页面 - 显示分析选项和数据统计
    """
    db = SessionLocal()
    try:
        # 检查用户是否拥有该任务
        task = db.query(Task).filter_by(id=task_id, user_id=current_user.id).first()
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        
        # 获取提交数据数量和列表
        submission = db.query(Submission).filter_by(task_id=task_id).all()
        submission_count = len(submission)
        
        # 检查是否有AI配置
        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
        
        # 检查是否有APIKEY
        has_api_key = False
        if ai_config and ai_config.selected_model:
            model_cfg = db.query(AIModelConfig).filter_by(
                ai_config_id=ai_config.id,
                model_name=ai_config.selected_model
            ).first()
            if model_cfg:
                if model_cfg.api_key:
                    has_api_key = True
                elif model_cfg.api_url and ai_config.selected_model == 'ollama':
                    has_api_key = True
        
        # 读取附件内容（如果有）
        file_content = None
        if task.attachments:
            # 读取第一个附件的内容
            first_attachment = task.attachments[0]
            if os.path.exists(first_attachment.file_path):
                file_content = read_file_content(first_attachment.file_path)
        
        # 生成预览提示词
        preview_prompt = generate_analysis_prompt(task, submission, file_content)
        
        # 获取报告内容（如果存在）
        report = task.analysis_report if task and task.analysis_report else None
        
        return render_template('smart_analyze.html', 
                             task=task, 
                             report=report,
                             preview_prompt=preview_prompt,
                             submission_count=submission_count,
                             has_api_key=has_api_key,
                             now=datetime.now())
    finally:
        db.close()

@app.route('/download_report/<int:task_id>')
@login_required
def download_report(task_id):
    """
    下载分析报告
    """
    db = SessionLocal()
    try:
        task = db.query(Task).get(task_id)
        if not task:
            flash('任务不存在', 'danger')
            return redirect(url_for('dashboard'))
        if task.user_id != current_user.id:
            flash('无权访问此任务', 'danger')
            return redirect(url_for('dashboard'))
        
        # 保存任务信息
        report_file_path = task.report_file_path
        task_title = task.title
        report_content = task.analysis_report
        
        # 如果有报告文件且存在，直接发送
        if report_file_path and os.path.exists(report_file_path):
            db.close()
            import re
            safe_title = re.sub(r'[^a-zA-Z0-9_\u4e00-\u9fa5]', '_', task_title)
            safe_filename = f"{safe_title}_分析报告.html"
            
            try:
                return send_file(
                    report_file_path,
                    as_attachment=True,
                    download_name=safe_filename,
                    mimetype='text/html; charset=utf-8'
                )
            except TypeError:
                return send_file(
                    report_file_path,
                    as_attachment=True,
                    attachment_filename=safe_filename,
                    mimetype='text/html; charset=utf-8'
                )
        
        # 如果没有报告文件，但有数据库中的报告内容，直接生成HTML并下载
        if report_content and report_content.strip():
            import re
            from io import BytesIO
            safe_title = re.sub(r'[^a-zA-Z0-9_\u4e00-\u9fa5]', '_', task_title)
            safe_filename = f"{safe_title}_分析报告.html"
            
            # 使用模板渲染HTML报告
            report_time = task.report_generated_at.strftime('%Y-%m-%d %H:%M:%S') if task.report_generated_at else '未知'
            html_content = render_template('simple_report.html', 
                                         task_title=task_title, 
                                         report_time=report_time, 
                                         report_content=report_content)
            
            db.close()
            
            # 直接返回HTML内容作为下载
            html_bytes = html_content.encode('utf-8')
            return send_file(
                BytesIO(html_bytes),
                as_attachment=True,
                download_name=safe_filename,
                mimetype='text/html; charset=utf-8'
            )
        
        db.close()
        # 没有报告内容
        flash('该任务尚未生成分析报告，请先进行智能分析', 'info')
        return redirect(url_for('smart_analyze', task_id=task_id))
        
    except Exception as e:
        flash(f'下载报告时出错: {str(e)}', 'danger')
        return redirect(url_for('dashboard'))
    finally:
        if 'db' in locals() and db:
            db.close()

# ============================================================
# OpenMentor: AI 导师管理路由（V1 / Day 3-4）
# ============================================================
# 所有 AI 导师走 /assistant/* 路由，与 /task/* 完全隔离，互不影响。
# Task 表中 is_assistant=True 的记录由本组路由管理。

# 默认黑名单关键词（老师可在编辑页修改）
DEFAULT_BLOCKED_KEYWORDS = {
    "violence": {
        "label": "暴力自残",
        "enabled": True,
        "keywords": ["自杀", "自残", "伤害他人", "爆炸", "武器制作"]
    },
    "porn": {
        "label": "色情低俗",
        "enabled": True,
        "keywords": ["色情", "裸体"]
    },
    "jailbreak": {
        "label": "越狱提示词",
        "enabled": True,
        "keywords": ["忽略前面", "DAN", "无限制模式", "扮演无道德", "ignore previous"]
    },
    "cheat": {
        "label": "学习作弊（可选）",
        "enabled": False,
        "keywords": ["完整答案", "直接帮我做", "考试题"]
    }
}

# 多模态能力对照表（按模型自动决定哪些功能可启用）
# image_input：是否原生支持视觉理解（取决于对话模型本身）
# image_generation：是否支持文生图（取决于老师是否在 image_gen_model 字段配置了图生模型）
MODEL_CAPABILITIES = {
    'deepseek':    {'image_input': False, 'image_generation': False, 'label': 'DeepSeek'},
    'doubao':      {'image_input': True,  'image_generation': True,  'video_generation': True, 'label': '豆包（字节）'},
    'qwen':        {'image_input': True,  'image_generation': True,  'label': '阿里通义'},
    'glm':         {'image_input': True,  'image_generation': True,  'label': '智谱 GLM'},
    'siliconflow': {'image_input': True,  'image_generation': True,  'label': '硅基流动'},
    'cherry':      {'image_input': True,  'image_generation': True,  'label': 'Cherry Studio 企业版'},
    'ollama':      {'image_input': False, 'image_generation': False, 'label': 'Ollama 本地'},
}

# 图片生成 API 配置：
# - OpenAI 兼容（同步）：豆包 / GLM / SiliconFlow → POST /v1/images/generations
# - DashScope 异步（提交 → 轮询）：阿里通义 → 走 _generate_image_qwen_async
IMAGE_GEN_PROVIDERS = {
    'doubao': {
        'mode': 'openai',
        'url': 'https://ark.cn-beijing.volces.com/api/v3/images/generations',
        'default_model': 'doubao-seedream-3-0-t2i-250415',
    },
    'glm': {
        'mode': 'openai',
        'url': 'https://open.bigmodel.cn/api/paas/v4/images/generations',
        'default_model': 'cogview-3-flash',
    },
    'siliconflow': {
        'mode': 'openai',
        'url': 'https://api.siliconflow.cn/v1/images/generations',
        'default_model': 'black-forest-labs/FLUX.1-schnell',
    },
    'cherry': {
        # URL 动态：按老师填的服务端地址（_cherry_api_base）+ /images/generations
        # 网关后面模型五花八门（gpt-image / gemini-image 等），不传 size/n，由网关用默认值
        'mode': 'openai',
        'url': None,
        'default_model': 'openai/gpt-image-2',
    },
    'qwen': {
        # 阿里通义图生主推 qwen-image 系列（multimodal-generation 同步接口）
        # wanx* 系列保留向后兼容（通过 text2image 异步接口）但不再推荐
        'mode': 'qwen_dispatch',
        'task_url_template': 'https://dashscope.aliyuncs.com/api/v1/tasks/{task_id}',
        'wanx_submit_url': 'https://dashscope.aliyuncs.com/api/v1/services/aigc/text2image/image-synthesis',
        'qwen_image_submit_url': 'https://dashscope.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
        'default_model': 'qwen-image',
    },
}

# 视频生成（异步任务）：目前仅豆包 seedance；学生可选文生视频 / 图生视频两种模式
# 流程：创建任务 → 前端轮询 → 成功后服务器下载视频落盘（豆包链接 24h 过期）→ 存为 assistant 消息
VIDEO_GEN_PROVIDERS = {
    'doubao': {
        'create_url': 'https://ark.cn-beijing.volces.com/api/v3/contents/generations/tasks',
        'task_url_template': 'https://ark.cn-beijing.volces.com/api/v3/contents/generations/tasks/{task_id}',
        # 默认参数：5 秒 / 720p / 无水印（控制成本与等待时间，实测 ~21 秒生成）
        'default_params': '--resolution 720p --duration 5 --watermark false',
    },
}

# 进行中的视频任务登记：ark_task_id -> {'conv_id', 'prompt', 'done', 'result'}
# 内存态（服务重启后旧任务失效，前端会收到“任务已过期请重新生成”）
_VIDEO_TASKS = {}


def _parse_link_expires(expires_hours_str, default_hours=48):
    """解析链接有效期（小时数 → DateTime）。最大不超过 7 天（168 小时）。"""
    try:
        h = int(expires_hours_str) if expires_hours_str else default_hours
    except (TypeError, ValueError):
        h = default_hours
    h = max(1, min(168, h))  # 1 ~ 168 小时
    from datetime import timedelta
    return datetime.now() + timedelta(hours=h)


def _serialize_blocked_keywords(form):
    """从表单提取黑名单 JSON。表单提供每类的 enabled checkbox + 自定义关键词 textarea"""
    result = {}
    for key, default in DEFAULT_BLOCKED_KEYWORDS.items():
        enabled = form.get(f'bk_{key}_enabled') == 'on'
        custom_text = form.get(f'bk_{key}_keywords', '').strip()
        keywords = [k.strip() for k in re.split(r'[,，\n]+', custom_text) if k.strip()] if custom_text else default['keywords']
        result[key] = {
            'label': default['label'],
            'enabled': enabled,
            'keywords': keywords,
        }
    return json.dumps(result, ensure_ascii=False)


def _load_blocked_keywords(json_str):
    """加载黑名单 JSON；返回 dict（缺失字段以默认值补全）"""
    try:
        loaded = json.loads(json_str) if json_str else {}
    except (json.JSONDecodeError, TypeError):
        loaded = {}
    merged = {}
    for key, default in DEFAULT_BLOCKED_KEYWORDS.items():
        merged[key] = {**default, **loaded.get(key, {})}
    return merged


def _dt_text(value, with_seconds=False):
    """把 DateTime 转成前端友好的字符串。"""
    if not value:
        return None
    return value.strftime('%Y-%m-%d %H:%M:%S' if with_seconds else '%Y-%m-%d %H:%M')


def _truthy(value):
    """兼容 JSON / 表单传来的布尔值。"""
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in ('1', 'true', 'yes', 'on')
    return bool(value)


def _split_mobile_keywords(value, fallback=None):
    if isinstance(value, list):
        keywords = [str(x).strip() for x in value if str(x).strip()]
    else:
        text = str(value or '').strip()
        keywords = [x.strip() for x in re.split(r'[,，\n]+', text) if x.strip()] if text else []
    return keywords or list(fallback or [])


def _mobile_blocked_keywords_json(current_json, incoming):
    """移动端黑名单编辑器的 JSON 标准化。"""
    current = _load_blocked_keywords(current_json)
    incoming = incoming or {}
    result = {}
    for key, default in DEFAULT_BLOCKED_KEYWORDS.items():
        old = current.get(key, default)
        new = incoming.get(key, {}) if isinstance(incoming, dict) else {}
        if not isinstance(new, dict):
            new = {}
        keywords_value = new.get('keywords_text', new.get('keywords', old.get('keywords', default['keywords'])))
        result[key] = {
            'label': default['label'],
            'enabled': _truthy(new.get('enabled', old.get('enabled', default['enabled']))),
            'keywords': _split_mobile_keywords(keywords_value, old.get('keywords', default['keywords'])),
        }
    return json.dumps(result, ensure_ascii=False)


def _mobile_assistant_summary(db, assistant):
    """移动控制台列表卡片统计。"""
    convs = db.query(Conversation).filter_by(assistant_id=assistant.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS).all()
    conv_ids = [c.id for c in convs]
    students = {}
    for c in convs:
        key = (c.student_class or '', c.student_name or '')
        item = students.setdefault(key, {'blocked': False, 'last_active_at': None})
        item['blocked'] = item['blocked'] or bool(c.is_blocked)
        if c.last_active_at and (not item['last_active_at'] or c.last_active_at > item['last_active_at']):
            item['last_active_at'] = c.last_active_at

    msg_count = 0
    alert_count = 0
    latest_alert = None
    if conv_ids:
        msg_count = db.query(Message).filter(Message.conversation_id.in_(conv_ids)).count()
        alert_q = (
            db.query(Message)
            .filter(Message.conversation_id.in_(conv_ids))
            .filter(Message.triggered_keyword.isnot(None))
            .filter(Message.triggered_keyword != '')
        )
        alert_count = alert_q.count()
        latest_alert = alert_q.order_by(Message.created_at.desc()).first()

    active_cutoff = datetime.now() - timedelta(minutes=10)
    active_students = sum(1 for s in students.values() if s['last_active_at'] and s['last_active_at'] >= active_cutoff)
    is_expired = bool(assistant.link_expires_at and assistant.link_expires_at < datetime.now())
    return {
        'id': assistant.id,
        'title': assistant.title,
        'description': assistant.description or '',
        'status': assistant.status or 'active',
        'effective_status': 'expired' if is_expired else (assistant.status or 'active'),
        'student_url': f"{request.host_url.rstrip('/')}/chat/{assistant.task_id}",
        'selected_model_name': assistant.selected_model_name or '默认',
        'allow_image_input': bool(assistant.allow_image_input),
        'allow_file_upload': bool(assistant.allow_file_upload),
        'allow_image_generation': bool(assistant.allow_image_generation),
        'daily_limit': assistant.max_messages_per_student_daily or 50,
        'student_count': len(students),
        'active_students': active_students,
        'blocked_students': sum(1 for s in students.values() if s['blocked']),
        'session_count': len(convs),
        'message_count': msg_count,
        'alert_count': alert_count,
        'latest_alert_at': _dt_text(latest_alert.created_at) if latest_alert else None,
        'link_expires_at': _dt_text(assistant.link_expires_at),
    }


@app.route('/assistant/list')
@login_required
def assistant_list():
    """AI 导师列表"""
    db = SessionLocal()
    try:
        assistants = (
            db.query(Task)
            .filter_by(user_id=current_user.id, is_assistant=True)
            .order_by(Task.created_at.desc())
            .all()
        )
        # 计算每个助学的会话数
        items = []
        now = datetime.now()
        for a in assistants:
            conv_count = db.query(Conversation).filter_by(assistant_id=a.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS).count()
            msg_count = (
                db.query(Message)
                .join(Conversation, Message.conversation_id == Conversation.id)
                .filter(Conversation.assistant_id == a.id)
                .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
                .count()
            )
            is_expired = a.link_expires_at and a.link_expires_at < now
            items.append({
                'task': a,
                'conv_count': conv_count,
                'msg_count': msg_count,
                'is_expired': is_expired,
                'effective_status': 'expired' if is_expired else (a.status or 'active'),
            })
        return render_template('assistant_list.html', items=items)
    finally:
        db.close()


AVATAR_DIR = os.path.join('static', 'uploads', 'avatars')
ALLOWED_AVATAR_EXTS = {'png', 'jpg', 'jpeg', 'webp', 'gif'}


def _save_assistant_avatar(file, assistant_id):
    """保存 AI 导师头像。失败返回 (None, error_message)，成功返回 (rel_path, None)"""
    if not file or not file.filename:
        return None, None
    ext = _ext_of(file.filename)
    if ext not in ALLOWED_AVATAR_EXTS:
        return None, f'头像格式不支持（请用 PNG/JPG/WebP/GIF）'
    file.seek(0, os.SEEK_END)
    size = file.tell()
    file.seek(0)
    if size > 3 * 1024 * 1024:
        return None, '头像不能大于 3 MB'
    if size <= 0:
        return None, '空文件'

    if not os.path.exists(AVATAR_DIR):
        os.makedirs(AVATAR_DIR, exist_ok=True)

    save_ext = 'jpg' if ext in ('heic', 'heif') else ext
    unique = f'{assistant_id}_{uuid.uuid4().hex[:6]}.{save_ext}'
    save_path = os.path.join(AVATAR_DIR, unique)

    if ext in ('heic', 'heif'):
        try:
            from PIL import Image
            from io import BytesIO
            img = Image.open(BytesIO(file.read()))
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            img.save(save_path, 'JPEG', quality=92)
        except Exception as e:
            return None, f'HEIC 转换失败: {e}'
    else:
        file.save(save_path)

    return save_path.replace('\\', '/'), None


def _delete_old_avatar(rel_path):
    """删除旧头像文件"""
    if not rel_path:
        return
    try:
        if os.path.exists(rel_path) and rel_path.startswith(AVATAR_DIR):
            os.remove(rel_path)
    except OSError:
        pass


@app.route('/assistant/create', methods=['GET', 'POST'])
@login_required
def assistant_create():
    """创建 AI 导师"""
    db = SessionLocal()
    try:
        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()

        if request.method == 'POST':
            title = request.form.get('title', '').strip()
            if not title:
                flash('请填写助学名称', 'warning')
                return redirect(url_for('assistant_create'))

            assistant = Task(
                title=title,
                description=request.form.get('description', '').strip() or None,
                user_id=current_user.id,
                is_assistant=True,
                system_prompt=request.form.get('system_prompt', '').strip() or None,
                welcome_message=request.form.get('welcome_message', '').strip() or None,
                selected_model_name=request.form.get('selected_model_name') or (ai_config.selected_model if ai_config else 'deepseek'),
                allow_image_input=(request.form.get('allow_image_input') == 'on'),
                allow_file_upload=(request.form.get('allow_file_upload') == 'on'),
                allow_image_generation=(request.form.get('allow_image_generation') == 'on'),
                allow_video_generation=(request.form.get('allow_video_generation') == 'on'),
                allow_call_teacher=(request.form.get('allow_call_teacher') == 'on'),
                allow_voice_input=(request.form.get('allow_voice_input') == 'on'),
                allow_tts=(request.form.get('allow_tts') == 'on'),
                audio_mode=('direct' if request.form.get('audio_mode') == 'direct' else 'transcribe'),
                share_mode=('group' if request.form.get('share_mode') == 'group' else 'independent'),
                roster_mode=('strict' if request.form.get('roster_mode') == 'strict' else 'off'),
                max_messages_per_student_daily=int(request.form.get('max_messages_per_student_daily') or 50),
                link_expires_at=_parse_link_expires(request.form.get('link_expires_hours')),
                blocked_keywords=_serialize_blocked_keywords(request.form),
                status='active',
            )
            db.add(assistant)
            db.commit()
            db.refresh(assistant)

            # 头像（可选）
            avatar_file = request.files.get('avatar')
            if avatar_file and avatar_file.filename:
                rel_path, err = _save_assistant_avatar(avatar_file, assistant.id)
                if err:
                    flash(f'头像保存失败：{err}', 'warning')
                else:
                    assistant.avatar_path = rel_path
                    db.commit()

            flash(f'AI 导师「{assistant.title}」创建成功，请把链接分享给学生。', 'success')
            return redirect(url_for('assistant_detail', assistant_id=assistant.id))

        # GET：展示空表单
        return render_template(
            'assistant_form.html',
            mode='create',
            assistant=None,
            ai_config=ai_config,
            model_capabilities=MODEL_CAPABILITIES,
            blocked_keywords=DEFAULT_BLOCKED_KEYWORDS,
            default_link_hours=48,
        )
    finally:
        db.close()


@app.route('/assistant/edit/<int:assistant_id>', methods=['GET', 'POST'])
@login_required
def assistant_edit(assistant_id):
    """编辑 AI 导师"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))

        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()

        if request.method == 'POST':
            title = request.form.get('title', '').strip()
            if not title:
                flash('请填写助学名称', 'warning')
                return redirect(url_for('assistant_edit', assistant_id=assistant_id))

            assistant.title = title
            assistant.description = request.form.get('description', '').strip() or None
            assistant.system_prompt = request.form.get('system_prompt', '').strip() or None
            assistant.welcome_message = request.form.get('welcome_message', '').strip() or None
            assistant.selected_model_name = request.form.get('selected_model_name') or (ai_config.selected_model if ai_config else 'deepseek')
            assistant.allow_image_input = (request.form.get('allow_image_input') == 'on')
            assistant.allow_file_upload = (request.form.get('allow_file_upload') == 'on')
            assistant.allow_image_generation = (request.form.get('allow_image_generation') == 'on')
            assistant.allow_video_generation = (request.form.get('allow_video_generation') == 'on')
            assistant.allow_call_teacher = (request.form.get('allow_call_teacher') == 'on')
            assistant.allow_voice_input = (request.form.get('allow_voice_input') == 'on')
            assistant.allow_tts = (request.form.get('allow_tts') == 'on')
            assistant.audio_mode = 'direct' if request.form.get('audio_mode') == 'direct' else 'transcribe'
            assistant.roster_mode = 'strict' if request.form.get('roster_mode') == 'strict' else 'off'
            assistant.share_mode = 'group' if request.form.get('share_mode') == 'group' else 'independent'
            assistant.max_messages_per_student_daily = int(request.form.get('max_messages_per_student_daily') or 50)

            # 仅在用户主动重置时更新链接有效期
            if request.form.get('reset_expires') == 'on':
                assistant.link_expires_at = _parse_link_expires(request.form.get('link_expires_hours'))

            assistant.blocked_keywords = _serialize_blocked_keywords(request.form)

            # 头像处理
            if request.form.get('remove_avatar') == 'on':
                _delete_old_avatar(assistant.avatar_path)
                assistant.avatar_path = None
            else:
                avatar_file = request.files.get('avatar')
                if avatar_file and avatar_file.filename:
                    rel_path, err = _save_assistant_avatar(avatar_file, assistant.id)
                    if err:
                        flash(f'头像保存失败：{err}', 'warning')
                    else:
                        _delete_old_avatar(assistant.avatar_path)
                        assistant.avatar_path = rel_path

            db.commit()
            flash(f'AI 导师「{assistant.title}」已更新', 'success')
            return redirect(url_for('assistant_detail', assistant_id=assistant.id))

        # GET：渲染表单
        return render_template(
            'assistant_form.html',
            mode='edit',
            assistant=assistant,
            ai_config=ai_config,
            model_capabilities=MODEL_CAPABILITIES,
            blocked_keywords=_load_blocked_keywords(assistant.blocked_keywords),
            default_link_hours=48,
        )
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>')
@login_required
def assistant_detail(assistant_id):
    """AI 导师详情：显示学生链接、二维码、统计"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))

        conv_count = db.query(Conversation).filter_by(assistant_id=assistant.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS).count()
        msg_count = (
            db.query(Message)
            .join(Conversation, Message.conversation_id == Conversation.id)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
            .count()
        )
        now = datetime.now()
        is_expired = assistant.link_expires_at and assistant.link_expires_at < now
        student_url = f"{request.host_url.rstrip('/')}/chat/{assistant.task_id}"

        kb_docs = (
            db.query(KnowledgeDoc)
            .filter_by(assistant_id=assistant.id)
            .order_by(KnowledgeDoc.created_at.desc())
            .all()
        )
        kb_enabled = bool(_get_siliconflow_key(db, assistant))

        return render_template(
            'assistant_detail.html',
            assistant=assistant,
            conv_count=conv_count,
            msg_count=msg_count,
            is_expired=is_expired,
            student_url=student_url,
            model_capabilities=MODEL_CAPABILITIES,
            kb_docs=kb_docs,
            kb_enabled=kb_enabled,
        )
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit')
@login_required
def assistant_audit(assistant_id):
    """审计大屏：列出该助学的所有学生会话 + 实时统计"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        return render_template('assistant_audit.html', assistant=assistant)
    finally:
        db.close()


@app.route('/api/assistant/<int:assistant_id>/audit/conversations')
@login_required
def api_assistant_audit_conversations(assistant_id):
    """审计概览数据 API（前端轮询自动刷新）"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': '不存在'}), 404

        convs = db.query(Conversation).filter_by(assistant_id=assistant.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS).order_by(Conversation.last_active_at.desc()).all()
        items = []
        total_msgs = 0
        total_blocked_hits = 0
        total_thumbs_up = 0
        total_thumbs_down = 0
        for c in convs:
            msg_count = len(c.messages)
            user_msg_count = sum(1 for x in c.messages if x.role == 'user')
            blocked_hits = sum(1 for x in c.messages if x.triggered_keyword)
            up_count = sum(1 for x in c.messages if x.rating == 1)
            down_count = sum(1 for x in c.messages if x.rating == -1)
            total_msgs += msg_count
            total_blocked_hits += blocked_hits
            total_thumbs_up += up_count
            total_thumbs_down += down_count
            last = c.messages[-1] if c.messages else None
            items.append({
                'id': c.id,
                'student_class': c.student_class,
                'student_name': c.student_name,
                'started_at': c.started_at.strftime('%Y-%m-%d %H:%M') if c.started_at else None,
                'last_active_at': c.last_active_at.strftime('%Y-%m-%d %H:%M') if c.last_active_at else None,
                'last_active_ts': c.last_active_at.timestamp() if c.last_active_at else 0,
                'msg_count': msg_count,
                'user_msg_count': user_msg_count,
                'daily_message_count': c.daily_message_count or 0,
                'daily_limit': assistant.max_messages_per_student_daily or 50,
                'blocked_hits': blocked_hits,
                'thumbs_up': up_count,
                'thumbs_down': down_count,
                'is_blocked': bool(c.is_blocked),
                'last_role': last.role if last else None,
                'last_content': (last.content[:80] + '...') if last and last.content and len(last.content) > 80 else (last.content if last else None),
            })

        # B8: 按 (班级, 姓名) 聚合，把同一身份的多设备会话合成一条
        students_map = {}  # (class, name) → aggregated dict
        for it in items:
            key = (it['student_class'], it['student_name'])
            if key not in students_map:
                students_map[key] = {
                    'student_class': it['student_class'],
                    'student_name': it['student_name'],
                    'msg_count': 0,
                    'user_msg_count': 0,
                    'daily_message_count': 0,
                    'daily_limit': it['daily_limit'],
                    'blocked_hits': 0,
                    'thumbs_up': 0,
                    'thumbs_down': 0,
                    'is_blocked_any': False,
                    'last_active_ts': 0,
                    'last_active_at': None,
                    'last_role': None,
                    'last_content': None,
                    'sessions': [],
                }
            agg = students_map[key]
            agg['msg_count'] += it['msg_count']
            agg['user_msg_count'] += it['user_msg_count']
            agg['daily_message_count'] += it['daily_message_count']
            agg['blocked_hits'] += it['blocked_hits']
            agg['thumbs_up'] += it['thumbs_up']
            agg['thumbs_down'] += it['thumbs_down']
            agg['is_blocked_any'] = agg['is_blocked_any'] or it['is_blocked']
            if it['last_active_ts'] > agg['last_active_ts']:
                agg['last_active_ts'] = it['last_active_ts']
                agg['last_active_at'] = it['last_active_at']
                agg['last_role'] = it['last_role']
                agg['last_content'] = it['last_content']
            agg['sessions'].append({
                'id': it['id'],
                'started_at': it['started_at'],
                'last_active_at': it['last_active_at'],
                'msg_count': it['msg_count'],
                'daily_message_count': it['daily_message_count'],
                'is_blocked': it['is_blocked'],
            })
        students = sorted(students_map.values(), key=lambda s: -s['last_active_ts'])
        for s in students:
            s['session_count'] = len(s['sessions'])
            s['primary_session_id'] = s['sessions'][0]['id'] if s['sessions'] else None
            s['student_id'] = f"{s['student_class']}|{s['student_name']}"

        rated = total_thumbs_up + total_thumbs_down
        satisfaction_pct = round(total_thumbs_up * 100 / rated) if rated > 0 else None

        return jsonify({
            'code': 200,
            'assistant': {
                'id': assistant.id,
                'title': assistant.title,
                'status': assistant.status,
                'selected_model_name': assistant.selected_model_name,
                'link_expires_at': assistant.link_expires_at.strftime('%Y-%m-%d %H:%M') if assistant.link_expires_at else None,
            },
            'stats': {
                'student_count': len(students_map),       # B8: 不重复学生数（按 class+name 去重）
                'session_count': len(convs),              # 设备会话数（同一学生多设备会有多个）
                'conversation_count': len(students_map),  # 兼容旧字段：旧前端读这个
                'total_messages': total_msgs,
                'blocked_hits': total_blocked_hits,
                'thumbs_up': total_thumbs_up,
                'thumbs_down': total_thumbs_down,
                'satisfaction_pct': satisfaction_pct,
            },
            'students': students,            # B8: 按学生聚合（推荐使用）
            'conversations': items,          # 原始 flat list（兼容）
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/<int:conv_id>')
@login_required
def assistant_audit_conv(assistant_id, conv_id):
    """单个学生会话详情"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        conv = db.query(Conversation).filter_by(id=conv_id, assistant_id=assistant.id).first()
        if not conv:
            flash('会话不存在', 'danger')
            return redirect(url_for('assistant_audit', assistant_id=assistant.id))
        return render_template('assistant_audit_conv.html', assistant=assistant, conversation=conv)
    finally:
        db.close()


@app.route('/api/assistant/<int:assistant_id>/audit/<int:conv_id>/messages')
@login_required
def api_assistant_audit_messages(assistant_id, conv_id):
    """会话消息 API（前端轮询）"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': '不存在'}), 404
        conv = db.query(Conversation).filter_by(id=conv_id, assistant_id=assistant.id).first()
        if not conv:
            return jsonify({'code': 404, 'message': '会话不存在'}), 404

        def _display_filename(path):
            base = os.path.basename(path or '')
            if '_' in base:
                prefix, _, rest = base.partition('_')
                if len(prefix) == 8 and all(c in '0123456789abcdef' for c in prefix.lower()):
                    return rest
            return base

        msgs = []
        for m in conv.messages:
            try:
                imgs = json.loads(m.image_paths) if m.image_paths else []
            except (json.JSONDecodeError, TypeError):
                imgs = []
            try:
                files = json.loads(m.file_paths) if m.file_paths else []
            except (json.JSONDecodeError, TypeError):
                files = []
            try:
                audios = json.loads(m.audio_paths) if m.audio_paths else []
            except (json.JSONDecodeError, TypeError):
                audios = []
            msgs.append({
                'id': m.id,
                'role': m.role,
                'content': m.content,
                'triggered_keyword': m.triggered_keyword,
                'rating': m.rating,
                'image_urls': ['/' + p for p in imgs],
                'files': [{'name': _display_filename(p), 'url': '/' + p} for p in files],
                'audio_urls': ['/' + p for p in audios],
                'generated_image_url': ('/' + m.generated_image_path) if m.generated_image_path else None,
                'generated_video_url': ('/' + m.generated_video_path) if m.generated_video_path else None,
                'tokens_used': m.tokens_used or 0,
                'created_at': m.created_at.strftime('%H:%M:%S') if m.created_at else None,
                'created_at_full': m.created_at.strftime('%Y-%m-%d %H:%M:%S') if m.created_at else None,
            })

        # B8: 同一学生（class+name）在该 AI 导师下的所有会话（多设备）
        sibling_convs = db.query(Conversation).filter_by(
            assistant_id=assistant.id,
            student_class=conv.student_class,
            student_name=conv.student_name,
        ).order_by(Conversation.last_active_at.desc()).all()
        sibling_sessions = []
        for sc in sibling_convs:
            msg_count = len(sc.messages)
            sibling_sessions.append({
                'id': sc.id,
                'started_at': sc.started_at.strftime('%Y-%m-%d %H:%M') if sc.started_at else None,
                'last_active_at': sc.last_active_at.strftime('%Y-%m-%d %H:%M') if sc.last_active_at else None,
                'msg_count': msg_count,
                'is_current': sc.id == conv.id,
                'is_blocked': bool(sc.is_blocked),
            })

        return jsonify({
            'code': 200,
            'conversation': {
                'id': conv.id,
                'student_class': conv.student_class,
                'student_name': conv.student_name,
                'started_at': conv.started_at.strftime('%Y-%m-%d %H:%M') if conv.started_at else None,
                'last_active_at': conv.last_active_at.strftime('%Y-%m-%d %H:%M') if conv.last_active_at else None,
                'daily_used': conv.daily_message_count or 0,
                'daily_limit': assistant.max_messages_per_student_daily or 50,
                'is_blocked': bool(conv.is_blocked),
            },
            'sibling_sessions': sibling_sessions,
            'messages': msgs,
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/<int:conv_id>/block', methods=['POST'])
@login_required
def assistant_block_student(assistant_id, conv_id):
    """封禁/解禁某个学生会话"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': '不存在'}), 404
        conv = db.query(Conversation).filter_by(id=conv_id, assistant_id=assistant.id).first()
        if not conv:
            return jsonify({'success': False, 'message': '会话不存在'}), 404
        conv.is_blocked = not conv.is_blocked
        db.commit()
        return jsonify({'success': True, 'is_blocked': conv.is_blocked,
                        'message': '已封禁此学生' if conv.is_blocked else '已解禁此学生'})
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/student/block', methods=['POST'])
@login_required
def assistant_block_student_by_name(assistant_id):
    """B8: 按 (class, name) 批量封禁/解禁学生在所有设备上的会话"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': '不存在'}), 404
        class_name = (request.form.get('class_name') or '').strip()
        student_name = (request.form.get('student_name') or '').strip()
        action = (request.form.get('action') or 'block').strip()
        if not class_name or not student_name:
            return jsonify({'success': False, 'message': '缺少 class_name / student_name'}), 400
        target = action == 'block'
        convs = db.query(Conversation).filter_by(
            assistant_id=assistant.id,
            student_class=class_name,
            student_name=student_name,
        ).all()
        if not convs:
            return jsonify({'success': False, 'message': '没有匹配的会话'}), 404
        affected = 0
        for c in convs:
            if c.is_blocked != target:
                c.is_blocked = target
                affected += 1
        db.commit()
        return jsonify({
            'success': True,
            'affected': affected,
            'is_blocked': target,
            'message': ('已封禁' if target else '已解禁') + f' {class_name} · {student_name} 的 {len(convs)} 个会话',
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/export')
@login_required
def assistant_audit_export(assistant_id):
    """导出该助学的所有对话为 Excel"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))

        convs = db.query(Conversation).filter_by(assistant_id=assistant.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS).order_by(Conversation.id).all()
        rows = []
        for c in convs:
            for m in c.messages:
                rows.append({
                    '会话ID': c.id,
                    '班级': c.student_class,
                    '姓名': c.student_name,
                    '消息时间': m.created_at.strftime('%Y-%m-%d %H:%M:%S') if m.created_at else '',
                    '角色': {'user': '学生', 'assistant': 'AI', 'system': '系统'}.get(m.role, m.role),
                    '内容': (m.content or '')[:1000],
                    '命中关键词': m.triggered_keyword or '',
                    '附图数': len(json.loads(m.image_paths or '[]')),
                    '附件数': len(json.loads(m.file_paths or '[]')),
                    'AI生成图': '是' if m.generated_image_path else '',
                })
        if not rows:
            flash('该 AI 导师暂无对话记录', 'info')
            return redirect(url_for('assistant_audit', assistant_id=assistant.id))

        import pandas as pd
        df = pd.DataFrame(rows)
        from io import BytesIO
        bio = BytesIO()
        with pd.ExcelWriter(bio, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='对话记录', index=False)
        bio.seek(0)
        safe_title = re.sub(r'[\\/:*?"<>|]', '_', assistant.title)[:60]
        filename = f"{safe_title}_对话记录_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
        return send_file(bio, as_attachment=True, download_name=filename,
                         mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    finally:
        db.close()


# ============ B12 / B13: AI 学情报告 ============

# 单条对话条目格式化的字符长度上限（避免 token 爆炸）
_REPORT_MAX_CONTENT_PER_MSG = 400
_REPORT_MAX_TOTAL_CHARS = 30000  # 喂给 AI 的总字符数上限（约 ~12K tokens）
_REPORT_MAX_CUSTOM_INSTRUCTION_CHARS = 3000


def _clean_report_instruction(value):
    value = (value or '').strip()
    if len(value) > _REPORT_MAX_CUSTOM_INSTRUCTION_CHARS:
        value = value[:_REPORT_MAX_CUSTOM_INSTRUCTION_CHARS]
    return value


def _report_instruction_block(value):
    instruction = _clean_report_instruction(value)
    if not instruction:
        return ''
    return f"""

【老师本次报告要求】
{instruction}

请优先满足老师本次报告要求；如果老师没有指定报告结构，再参考下面的默认结构。"""


def _format_messages_for_report(conv, max_per_msg=_REPORT_MAX_CONTENT_PER_MSG):
    """把一段会话格式化成 AI 可读的对话片段。"""
    lines = []
    for m in conv.messages:
        role_label = {'user': '学生', 'assistant': 'AI', 'system': '系统'}.get(m.role, m.role)
        text = (m.content or '').strip().replace('\n', ' ')
        if len(text) > max_per_msg:
            text = text[:max_per_msg] + '...'
        if m.triggered_keyword:
            text += f' [命中关键词: {m.triggered_keyword}]'
        if m.rating == 1:
            text += ' [学生👍]'
        elif m.rating == -1:
            text += ' [学生👎]'
        lines.append(f'  {role_label}: {text}')
    return '\n'.join(lines)


def _truncate_to_total(blocks, total_limit):
    """按总字符数截断 blocks，超过的部分丢弃并加注省略说明。"""
    out = []
    used = 0
    omitted = 0
    for b in blocks:
        if used + len(b) > total_limit:
            omitted += 1
            continue
        out.append(b)
        used += len(b)
    return out, omitted


def _generate_report_via_ai(meta_prompt, ai_config):
    """复用 call_ai_model；返回 (text, error)"""
    try:
        text = call_ai_model(meta_prompt, ai_config)
    except Exception as e:
        logger.exception('生成 AI 报告失败')
        return None, str(e)
    if not text or not text.strip():
        return None, 'AI 返回空内容，请重试'
    return text.strip(), None


def _stream_report_sse(meta_prompt, ai_config, meta_event_payload, persist_factory):
    """B19 报告流式输出 generator。
    yield SSE 字符串：event: meta / event: delta / event: done / event: error。
    persist_factory(full_text) → 返回包含 report_id 的 dict（成功后调用，写库）。
    """
    # 立刻发一个 meta 事件，让前端立即关掉 spinner 显示骨架
    yield _sse_event('meta', meta_event_payload)

    # 找出当前选中的模型
    model_key = (ai_config.selected_model or '').strip() if ai_config else ''
    model_cfg = None
    if ai_config:
        model_cfg = next((mc for mc in ai_config.model_configs if mc.model_name == model_key), None)

    full_chunks = []
    try:
        for chunk in _stream_assistant_response(
            model_key, model_cfg,
            [{'role': 'user', 'content': meta_prompt}]
        ):
            if chunk:
                full_chunks.append(chunk)
                yield _sse_event('delta', {'content': chunk})
    except Exception as e:
        logger.exception('流式生成报告失败')
        yield _sse_event('error', {'message': str(e)})
        return

    full_text = (''.join(full_chunks)).strip()
    if not full_text:
        yield _sse_event('error', {'message': 'AI 返回空内容'})
        return

    # 持久化
    try:
        rec_payload = persist_factory(full_text)
        yield _sse_event('done', rec_payload)
    except Exception as e:
        logger.exception('保存报告失败')
        yield _sse_event('error', {'message': '保存失败：' + str(e)})


@app.route('/assistant/<int:assistant_id>/audit/student_report', methods=['POST'])
@login_required
def assistant_student_report(assistant_id):
    """B12：单学生 AI 学情报告。
    POST 参数：student_class, student_name
    返回：{ code, success, report (markdown), session_count, message_count }
    """
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在或无权访问'}), 404

        student_class = (request.form.get('student_class') or '').strip()
        student_name = (request.form.get('student_name') or '').strip()
        report_instruction = _clean_report_instruction(request.form.get('report_instruction'))
        if not student_class or not student_name:
            return jsonify({'code': 400, 'message': '缺少 student_class / student_name'}), 400

        convs = db.query(Conversation).filter_by(
            assistant_id=assistant.id,
            student_class=student_class,
            student_name=student_name,
        ).order_by(Conversation.started_at).all()
        total_msgs = sum(len(c.messages) for c in convs)
        if not convs or total_msgs == 0:
            return jsonify({'code': 400, 'message': '该学生暂无对话记录'}), 400

        # 拼接对话
        blocks = []
        for i, c in enumerate(convs, 1):
            session_label = f'### 设备会话 {i}（{c.started_at.strftime("%Y-%m-%d %H:%M") if c.started_at else "?"}）\n'
            blocks.append(session_label + _format_messages_for_report(c) + '\n')
        truncated_blocks, omitted = _truncate_to_total(blocks, _REPORT_MAX_TOTAL_CHARS)
        if not truncated_blocks:
            return jsonify({'code': 400, 'message': '对话内容超出可分析长度'}), 400
        chat_text = '\n'.join(truncated_blocks)

        meta_prompt = f"""你是一位有教育心理学背景的学情分析师。请基于以下"学生与 AI 助教"的对话记录，为老师生成一份学情报告。

【AI 助教信息】
- 助教名称：{assistant.title}
- 助教定位：{(assistant.description or '（未设置）')[:200]}

【学生信息】
- 班级：{student_class}
- 姓名：{student_name}
- 共 {len(convs)} 个设备会话，{total_msgs} 条消息（{'已截取重要片段，省略 ' + str(omitted) + ' 段' if omitted else '完整记录'}）

【对话记录】
{chat_text}
{_report_instruction_block(report_instruction)}

请输出一份**结构化、可执行**的 Markdown 报告，包含以下模块（每个模块 2-4 句话即可，避免空话）：

## 一、学习画像
（提问类型分布、关注的学科 / 知识点、学习风格初判）

## 二、知识盲点
（从提问中能看出来的薄弱点、概念误区，列 2-5 条）

## 三、AI 互动效果
（学生评分情况、是否真在思考还是只想要答案、AI 回应是否到位）

## 四、安全与态度
（是否触发黑名单、整体提问态度是否积极，有无异常）

## 五、给老师的建议
（针对这位学生的下一步辅导建议，列 2-3 条具体可操作的）

要求：
- 中文输出，事实就事实，不要空泛赞美
- 涉及具体提问时可以引用学生原话作为论据（用「」括起来）
- 不要重复"学生"两个字开头每段，直接陈述"知识画像显示..."这样
- 总长度控制在 800-1200 字
- 直接输出 Markdown 正文，不要包代码块"""

        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
        if not ai_config:
            return jsonify({'code': 500, 'message': '请先在「AI 配置」中选择并配置 AI 模型'}), 500
        cur_cfg = next((mc for mc in ai_config.model_configs if mc.model_name == ai_config.selected_model), None)
        if (not cur_cfg or not (cur_cfg.api_key or '').strip()) and ai_config.selected_model != 'ollama':
            return jsonify({'code': 500, 'message': f'当前选中的 {ai_config.selected_model} 未配置 API Key'}), 500

        # B19: 流式模式
        if (request.form.get('stream') or request.args.get('stream') or '').strip() == '1':
            from flask import stream_with_context
            assistant_title = assistant.title
            model_used = ai_config.selected_model
            uid = current_user.id
            aid = assistant.id

            def _persist_student(full_text):
                _db = SessionLocal()
                try:
                    rec = StudentReport(
                        user_id=uid, assistant_id=aid, report_type='student',
                        student_class=student_class, student_name=student_name,
                        scope_label=f'{student_class} · {student_name}',
                        content=full_text, model_used=model_used,
                        student_count=1, message_count=total_msgs, omitted=omitted,
                    )
                    _db.add(rec); _db.commit(); _db.refresh(rec)
                    return {
                        'success': True, 'report_id': rec.id,
                        'session_count': len(convs), 'message_count': total_msgs,
                        'omitted_blocks': omitted, 'model_used': model_used,
                        'student_class': student_class, 'student_name': student_name,
                        'assistant_title': assistant_title,
                        'created_at': rec.created_at.strftime('%Y-%m-%d %H:%M'),
                    }
                finally:
                    _db.close()

            meta_payload = {
                'session_count': len(convs), 'message_count': total_msgs,
                'omitted_blocks': omitted, 'model_used': model_used,
                'student_class': student_class, 'student_name': student_name,
                'assistant_title': assistant_title,
                'scope_label': f'{student_class} · {student_name}',
            }

            return Response(
                stream_with_context(_stream_report_sse(meta_prompt, ai_config, meta_payload, _persist_student)),
                mimetype='text/event-stream',
                headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'},
            )

        report, err = _generate_report_via_ai(meta_prompt, ai_config)
        if err:
            return jsonify({'code': 500, 'message': f'生成失败：{err}'}), 500

        # B14: 持久化到 student_report
        rec = StudentReport(
            user_id=current_user.id,
            assistant_id=assistant.id,
            report_type='student',
            student_class=student_class,
            student_name=student_name,
            scope_label=f'{student_class} · {student_name}',
            content=report,
            model_used=ai_config.selected_model,
            student_count=1,
            message_count=total_msgs,
            omitted=omitted,
        )
        db.add(rec)
        db.commit()
        db.refresh(rec)

        return jsonify({
            'code': 200,
            'success': True,
            'report_id': rec.id,
            'report': report,
            'session_count': len(convs),
            'message_count': total_msgs,
            'omitted_blocks': omitted,
            'model_used': ai_config.selected_model,
            'student_class': student_class,
            'student_name': student_name,
            'assistant_title': assistant.title,
            'created_at': rec.created_at.strftime('%Y-%m-%d %H:%M'),
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/class_report', methods=['POST'])
@login_required
def assistant_class_report(assistant_id):
    """B13：按班级（或全部）的全班学情报告。
    POST 参数：class_name（可空 = 全部学生）
    """
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在或无权访问'}), 404

        class_filter = (request.form.get('class_name') or '').strip()
        report_instruction = _clean_report_instruction(request.form.get('report_instruction'))
        selected_raw = (request.form.get('selected_students') or '').strip()
        selected_lookup = set()
        if selected_raw:
            try:
                selected_items = json.loads(selected_raw)
            except json.JSONDecodeError:
                return jsonify({'code': 400, 'message': '选中学生参数格式错误'}), 400
            if not isinstance(selected_items, list):
                return jsonify({'code': 400, 'message': '选中学生参数格式错误'}), 400
            for item in selected_items:
                if not isinstance(item, dict):
                    continue
                cls = (item.get('student_class') or item.get('class_name') or '').strip()
                name = (item.get('student_name') or item.get('name') or '').strip()
                if cls and name:
                    selected_lookup.add((cls, name))
            if not selected_lookup:
                return jsonify({'code': 400, 'message': '没有有效的选中学生'}), 400

        q = db.query(Conversation).filter_by(assistant_id=assistant.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
        if class_filter and not selected_lookup:
            q = q.filter(Conversation.student_class == class_filter)
        convs = q.order_by(Conversation.student_class, Conversation.student_name, Conversation.started_at).all()
        if selected_lookup:
            convs = [c for c in convs if (c.student_class, c.student_name) in selected_lookup]
        if not convs:
            return jsonify({'code': 400, 'message': '没有匹配的对话记录'}), 400

        # 按 (class, name) 聚合
        grouped = {}
        for c in convs:
            key = (c.student_class, c.student_name)
            grouped.setdefault(key, []).append(c)

        student_count = len(grouped)
        total_msgs = sum(len(c.messages) for c in convs)
        if total_msgs == 0:
            return jsonify({'code': 400, 'message': '没有可分析的消息'}), 400

        blocks = []
        for (cls, name), conv_list in grouped.items():
            stu_msgs = sum(len(c.messages) for c in conv_list)
            header = f'### {cls} · {name}（{len(conv_list)} 个会话，共 {stu_msgs} 条消息）\n'
            stu_lines = []
            for c in conv_list:
                stu_lines.append(_format_messages_for_report(c))
            blocks.append(header + '\n'.join(stu_lines) + '\n')
        truncated_blocks, omitted = _truncate_to_total(blocks, _REPORT_MAX_TOTAL_CHARS)
        if not truncated_blocks:
            return jsonify({'code': 400, 'message': '对话内容过多，请按班级缩小范围后再生成'}), 400
        chat_text = '\n'.join(truncated_blocks)

        scope_label = f'已选 {len(selected_lookup)} 名学生' if selected_lookup else (f'班级「{class_filter}」' if class_filter else '全部学生')
        meta_prompt = f"""你是一位有教育心理学背景的学情分析师。请基于以下"AI 助教 {scope_label}"的全部对话记录，为老师生成一份**班级级别**的学情报告。

【AI 助教信息】
- 助教名称：{assistant.title}
- 助教定位：{(assistant.description or '（未设置）')[:200]}

【数据范围】
- 范围：{scope_label}
- 学生数：{student_count}
- 消息总数：{total_msgs}
- {'数据较多已截取最多 ' + str(len(truncated_blocks)) + ' 个学生的数据，省略 ' + str(omitted) + ' 个学生' if omitted else '完整记录'}

【对话记录（按学生分组）】
{chat_text}
{_report_instruction_block(report_instruction)}

请输出一份**结构化、可执行**的 Markdown 报告：

## 一、整体使用情况
（活跃度概览：高活跃学生、沉默学生、平均提问数；无需具体名单，给量级判断即可）

## 二、共性知识盲点
（从全班提问中归纳出最频繁的疑惑点、错误概念，列 3-6 条，配学生原话做佐证）

## 三、提问类型分布
（学生主要是在求"答案"、求"思路"还是求"解释"？反映学习风格）

## 四、AI 互动质量
（满意度概况、AI 回应有无系统性偏差、有无被学生用来作弊或越狱的迹象）

## 五、值得老师关注的学生
（最多 3-5 名：行为有异常、知识盲点突出、或被 AI 帮了很多但没真正学会的）

## 六、教学改进建议
（针对这个班级 / 这次使用，给老师下一步教学的 3 条具体建议）

要求：
- 中文输出，避免空泛话
- 涉及具体学生时只在第五部分点名，前几部分用类型化描述
- 引用对话原话时用「」括起来
- 总长度控制在 1200-1800 字
- 直接输出 Markdown 正文，不要包代码块"""

        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
        if not ai_config:
            return jsonify({'code': 500, 'message': '请先在「AI 配置」中选择并配置 AI 模型'}), 500
        cur_cfg = next((mc for mc in ai_config.model_configs if mc.model_name == ai_config.selected_model), None)
        if (not cur_cfg or not (cur_cfg.api_key or '').strip()) and ai_config.selected_model != 'ollama':
            return jsonify({'code': 500, 'message': f'当前选中的 {ai_config.selected_model} 未配置 API Key'}), 500

        # B19: 流式模式
        if (request.form.get('stream') or request.args.get('stream') or '').strip() == '1':
            from flask import stream_with_context
            assistant_title = assistant.title
            model_used = ai_config.selected_model
            uid = current_user.id
            aid = assistant.id

            def _persist_class(full_text):
                _db = SessionLocal()
                try:
                    rec = StudentReport(
                        user_id=uid, assistant_id=aid, report_type='class',
                        student_class=(None if selected_lookup else (class_filter or None)), student_name=None,
                        scope_label=scope_label,
                        content=full_text, model_used=model_used,
                        student_count=student_count, message_count=total_msgs, omitted=omitted,
                    )
                    _db.add(rec); _db.commit(); _db.refresh(rec)
                    return {
                        'success': True, 'report_id': rec.id,
                        'student_count': student_count, 'message_count': total_msgs,
                        'omitted_students': omitted, 'model_used': model_used,
                        'class_filter': class_filter,
                        'assistant_title': assistant_title,
                        'scope_label': scope_label,
                        'created_at': rec.created_at.strftime('%Y-%m-%d %H:%M'),
                    }
                finally:
                    _db.close()

            meta_payload = {
                'student_count': student_count, 'message_count': total_msgs,
                'omitted_students': omitted, 'model_used': model_used,
                'class_filter': class_filter,
                'assistant_title': assistant_title,
                'scope_label': scope_label,
            }
            return Response(
                stream_with_context(_stream_report_sse(meta_prompt, ai_config, meta_payload, _persist_class)),
                mimetype='text/event-stream',
                headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'},
            )

        report, err = _generate_report_via_ai(meta_prompt, ai_config)
        if err:
            return jsonify({'code': 500, 'message': f'生成失败：{err}'}), 500

        # B14: 持久化
        rec = StudentReport(
            user_id=current_user.id,
            assistant_id=assistant.id,
            report_type='class',
            student_class=(None if selected_lookup else (class_filter or None)),
            student_name=None,
            scope_label=scope_label,
            content=report,
            model_used=ai_config.selected_model,
            student_count=student_count,
            message_count=total_msgs,
            omitted=omitted,
        )
        db.add(rec)
        db.commit()
        db.refresh(rec)

        return jsonify({
            'code': 200,
            'success': True,
            'report_id': rec.id,
            'report': report,
            'student_count': student_count,
            'message_count': total_msgs,
            'omitted_students': omitted,
            'model_used': ai_config.selected_model,
            'class_filter': class_filter,
            'assistant_title': assistant.title,
            'scope_label': scope_label,
            'created_at': rec.created_at.strftime('%Y-%m-%d %H:%M'),
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/reports')
@login_required
def assistant_reports_list(assistant_id):
    """B14: 列出该 AI 导师下当前老师所有 AI 报告（按时间倒序）。
    Query 参数:
      - student_class / student_name 选填，过滤到某学生
      - report_type 选填 ('student' / 'class')
    """
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在或无权访问'}), 404

        q = db.query(StudentReport).filter_by(
            user_id=current_user.id,
            assistant_id=assistant.id,
        )
        rt = (request.args.get('report_type') or '').strip()
        if rt:
            q = q.filter(StudentReport.report_type == rt)
        cls = (request.args.get('student_class') or '').strip()
        if cls:
            q = q.filter(StudentReport.student_class == cls)
        name = (request.args.get('student_name') or '').strip()
        if name:
            q = q.filter(StudentReport.student_name == name)

        items = q.order_by(StudentReport.created_at.desc()).limit(200).all()
        return jsonify({
            'code': 200,
            'reports': [{
                'id': r.id,
                'report_type': r.report_type,
                'scope_label': r.scope_label,
                'student_class': r.student_class,
                'student_name': r.student_name,
                'model_used': r.model_used,
                'student_count': r.student_count,
                'message_count': r.message_count,
                'omitted': r.omitted,
                'created_at': r.created_at.strftime('%Y-%m-%d %H:%M') if r.created_at else None,
                'snippet': (r.content or '')[:120],
            } for r in items],
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/reports/<int:report_id>')
@login_required
def assistant_report_detail(assistant_id, report_id):
    """B14: 拉取单条历史报告全文。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在或无权访问'}), 404
        rec = db.query(StudentReport).filter_by(
            id=report_id, user_id=current_user.id, assistant_id=assistant.id
        ).first()
        if not rec:
            return jsonify({'code': 404, 'message': '报告不存在'}), 404
        return jsonify({
            'code': 200,
            'success': True,
            'id': rec.id,
            'report_type': rec.report_type,
            'scope_label': rec.scope_label,
            'student_class': rec.student_class,
            'student_name': rec.student_name,
            'content': rec.content,
            'model_used': rec.model_used,
            'student_count': rec.student_count,
            'message_count': rec.message_count,
            'omitted': rec.omitted,
            'created_at': rec.created_at.strftime('%Y-%m-%d %H:%M:%S') if rec.created_at else None,
            'assistant_title': assistant.title,
        })
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/audit/reports/<int:report_id>/delete', methods=['POST'])
@login_required
def assistant_report_delete(assistant_id, report_id):
    """B14: 删除一条历史报告。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        rec = db.query(StudentReport).filter_by(
            id=report_id, user_id=current_user.id, assistant_id=assistant.id
        ).first()
        if not rec:
            return jsonify({'success': False, 'message': '报告不存在'}), 404
        db.delete(rec)
        db.commit()
        return jsonify({'success': True})
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/toggle', methods=['POST'])
@login_required
def assistant_toggle(assistant_id):
    """切换 AI 导师启用/禁用状态（紧急熔断）"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        if assistant.status == 'active':
            assistant.status = 'disabled'
            msg = '已禁用，学生无法访问'
        else:
            assistant.status = 'active'
            msg = '已启用'
        db.commit()
        _emit_app_event(assistant.user_id, 'assistant_toggled', {
            'assistant_id': assistant.id, 'assistant_title': assistant.title, 'status': assistant.status,
        })
        return jsonify({'success': True, 'status': assistant.status, 'message': msg})
    finally:
        db.close()


# ---- 教师移动控制台 ----

@app.route('/mobile')
@login_required
def mobile_console():
    """手机端 AI 导师课堂控制台。"""
    return render_template('mobile_console.html')


@app.route('/mobile/assistant/<int:assistant_id>')
@login_required
def mobile_assistant_console(assistant_id):
    """手机端单个 AI 导师控制页。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('mobile_console'))
        return render_template(
            'mobile_assistant.html',
            assistant=assistant,
            blocked_keywords=_load_blocked_keywords(assistant.blocked_keywords),
            summary=_mobile_assistant_summary(db, assistant),
        )
    finally:
        db.close()


@app.route('/api/mobile/assistants')
@login_required
def api_mobile_assistants():
    """移动控制台：当前教师的全部 AI 导师。"""
    db = SessionLocal()
    try:
        assistants = (
            db.query(Task)
            .filter_by(user_id=current_user.id, is_assistant=True)
            .order_by(Task.created_at.desc())
            .all()
        )
        return jsonify({
            'code': 200,
            'user': {
                'display_name': current_user.display_name or current_user.username,
                'username': current_user.username,
            },
            'items': [_mobile_assistant_summary(db, a) for a in assistants],
        })
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/status', methods=['POST'])
@login_required
def api_mobile_assistant_status(assistant_id):
    """移动控制台：紧急熔断 / 重新启用。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        body = request.get_json(silent=True) or {}
        status = (body.get('status') or '').strip()
        if status not in ('active', 'disabled'):
            status = 'disabled' if assistant.status == 'active' else 'active'
        assistant.status = status
        db.commit()
        return jsonify({
            'success': True,
            'status': assistant.status,
            'message': '已启用' if assistant.status == 'active' else '已紧急熔断，学生暂时无法访问',
        })
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/settings', methods=['POST'])
@login_required
def api_mobile_assistant_settings(assistant_id):
    """移动控制台：保存课堂常用开关。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        body = request.get_json(silent=True) or {}
        for field in ('allow_image_input', 'allow_file_upload', 'allow_image_generation', 'allow_video_generation', 'allow_voice_input', 'allow_tts', 'allow_call_teacher'):
            if field in body:
                setattr(assistant, field, _truthy(body.get(field)))
        if 'max_messages_per_student_daily' in body:
            try:
                limit = int(body.get('max_messages_per_student_daily') or 50)
            except (TypeError, ValueError):
                return jsonify({'success': False, 'message': '每日消息上限必须是数字'}), 400
            if limit < 1 or limit > 500:
                return jsonify({'success': False, 'message': '每日消息上限建议设置在 1-500 之间'}), 400
            assistant.max_messages_per_student_daily = limit
        db.commit()
        return jsonify({'success': True, 'summary': _mobile_assistant_summary(db, assistant)})
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/link', methods=['POST'])
@login_required
def api_mobile_assistant_link(assistant_id):
    """移动控制台：重置学生链接有效期。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        body = request.get_json(silent=True) or {}
        try:
            hours = int(body.get('expires_hours') or 48)
        except (TypeError, ValueError):
            return jsonify({'success': False, 'message': '有效时长必须是数字'}), 400
        if hours < 1 or hours > 168:
            return jsonify({'success': False, 'message': '学生链接有效时长需要设置在 1-168 小时之间'}), 400
        assistant.link_expires_at = datetime.now() + timedelta(hours=hours)
        db.commit()
        return jsonify({
            'success': True,
            'summary': _mobile_assistant_summary(db, assistant),
            'message': f'学生链接有效期已重置为 {hours} 小时',
        })
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/blacklist', methods=['POST'])
@login_required
def api_mobile_assistant_blacklist(assistant_id):
    """移动控制台：保存黑名单分类与关键词。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        body = request.get_json(silent=True) or {}
        assistant.blocked_keywords = _mobile_blocked_keywords_json(
            assistant.blocked_keywords,
            body.get('categories') or body.get('blocked_keywords') or {},
        )
        db.commit()
        return jsonify({
            'success': True,
            'blocked_keywords': _load_blocked_keywords(assistant.blocked_keywords),
        })
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/alerts')
@login_required
def api_mobile_assistant_alerts(assistant_id):
    """移动控制台：黑名单命中记录，供手机端轮询提醒。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        try:
            since_id = int(request.args.get('since_id') or 0)
        except ValueError:
            since_id = 0
        try:
            limit = int(request.args.get('limit') or 20)
        except ValueError:
            limit = 20
        limit = max(1, min(limit, 50))
        q = (
            db.query(Message, Conversation)
            .join(Conversation, Message.conversation_id == Conversation.id)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Message.triggered_keyword.isnot(None))
            .filter(Message.triggered_keyword != '')
        )
        if since_id > 0:
            rows = q.filter(Message.id > since_id).order_by(Message.id.desc()).limit(limit).all()
        else:
            rows = q.order_by(Message.id.desc()).limit(limit).all()
        items = []
        max_id = since_id
        for m, c in rows:
            max_id = max(max_id, m.id)
            content = (m.content or '').strip()
            items.append({
                'id': m.id,
                'conversation_id': c.id,
                'student_class': c.student_class,
                'student_name': c.student_name,
                'triggered_keyword': m.triggered_keyword,
                'content': (content[:120] + '...') if len(content) > 120 else content,
                'created_at': _dt_text(m.created_at, with_seconds=True),
                'is_blocked': bool(c.is_blocked),
            })
        return jsonify({'code': 200, 'max_id': max_id, 'items': items})
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/students')
@login_required
def api_mobile_assistant_students(assistant_id):
    """移动控制台：按学生聚合的最近使用与封禁状态。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        convs = db.query(Conversation).filter_by(assistant_id=assistant.id).filter(Conversation.student_class != TEACHER_PREVIEW_CLASS).order_by(Conversation.last_active_at.desc()).all()
        students = {}
        for c in convs:
            key = (c.student_class, c.student_name)
            if key not in students:
                students[key] = {
                    'student_class': c.student_class,
                    'student_name': c.student_name,
                    'session_count': 0,
                    'msg_count': 0,
                    'daily_message_count': 0,
                    'blocked_hits': 0,
                    'is_blocked_any': False,
                    'last_active_ts': 0,
                    'last_active_at': None,
                    'last_content': '',
                }
            agg = students[key]
            agg['session_count'] += 1
            agg['daily_message_count'] += c.daily_message_count or 0
            agg['is_blocked_any'] = agg['is_blocked_any'] or bool(c.is_blocked)
            msg_count = len(c.messages)
            agg['msg_count'] += msg_count
            agg['blocked_hits'] += sum(1 for m in c.messages if m.triggered_keyword)
            last = c.messages[-1] if c.messages else None
            ts = c.last_active_at.timestamp() if c.last_active_at else 0
            if ts > agg['last_active_ts']:
                agg['last_active_ts'] = ts
                agg['last_active_at'] = _dt_text(c.last_active_at)
                agg['last_content'] = ((last.content or '')[:80] + '...') if last and last.content and len(last.content) > 80 else ((last.content or '') if last else '')
        items = sorted(students.values(), key=lambda s: -s['last_active_ts'])
        return jsonify({
            'code': 200,
            'daily_limit': assistant.max_messages_per_student_daily or 50,
            'items': items,
        })
    finally:
        db.close()


@app.route('/api/mobile/assistant/<int:assistant_id>/student/block', methods=['POST'])
@login_required
def api_mobile_assistant_student_block(assistant_id):
    """移动控制台：按学生身份显式封禁 / 解禁。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'success': False, 'message': 'AI 导师不存在'}), 404
        body = request.get_json(silent=True) or {}
        class_name = (body.get('student_class') or '').strip()
        student_name = (body.get('student_name') or '').strip()
        target = _truthy(body.get('blocked'))
        if not class_name or not student_name:
            return jsonify({'success': False, 'message': '缺少班级或姓名'}), 400
        convs = db.query(Conversation).filter_by(
            assistant_id=assistant.id,
            student_class=class_name,
            student_name=student_name,
        ).all()
        if not convs:
            return jsonify({'success': False, 'message': '没有匹配的学生会话'}), 404
        for c in convs:
            c.is_blocked = target
        db.commit()
        return jsonify({
            'success': True,
            'is_blocked': target,
            'affected': len(convs),
            'message': ('已封禁' if target else '已解禁') + f' {class_name} · {student_name}',
        })
    finally:
        db.close()


# ---- 模板导出 / 导入 / 克隆 ----

OPENMENTOR_TEMPLATE_VERSION = '1.0'


def _assistant_to_template_dict(assistant):
    """把 AI 导师序列化为可分享的模板（剔除敏感字段：API Key/对话/学生数据）"""
    return {
        'openmentor_template_version': OPENMENTOR_TEMPLATE_VERSION,
        'title': assistant.title,
        'description': assistant.description or '',
        'system_prompt': assistant.system_prompt or '',
        'welcome_message': assistant.welcome_message or '',
        'selected_model_name': assistant.selected_model_name or 'deepseek',
        'allow_image_input': bool(assistant.allow_image_input),
        'allow_file_upload': bool(assistant.allow_file_upload),
        'allow_image_generation': bool(assistant.allow_image_generation),
        'allow_video_generation': bool(assistant.allow_video_generation),
        'allow_call_teacher': bool(assistant.allow_call_teacher),
        'allow_voice_input': bool(assistant.allow_voice_input),
        'allow_tts': bool(assistant.allow_tts),
        'audio_mode': assistant.audio_mode or 'transcribe',
        'share_mode': assistant.share_mode or 'independent',
        'group_max_size': assistant.group_max_size or 6,
        'max_messages_per_student_daily': assistant.max_messages_per_student_daily or 50,
        'blocked_keywords': json.loads(assistant.blocked_keywords) if assistant.blocked_keywords else None,
    }


def _template_dict_to_assistant_kwargs(tpl, user_id):
    """从模板 dict 构造 Assistant 创建参数（导入用）"""
    title = (tpl.get('title') or '').strip()
    if not title:
        raise ValueError('模板缺少 title')
    bk = tpl.get('blocked_keywords')
    bk_json = json.dumps(bk, ensure_ascii=False) if bk else None
    from datetime import timedelta
    return {
        'title': title,
        'description': tpl.get('description') or None,
        'user_id': user_id,
        'is_assistant': True,
        'system_prompt': tpl.get('system_prompt') or None,
        'welcome_message': tpl.get('welcome_message') or None,
        'selected_model_name': tpl.get('selected_model_name') or 'deepseek',
        'allow_image_input': bool(tpl.get('allow_image_input', True)),
        'allow_file_upload': bool(tpl.get('allow_file_upload', True)),
        'allow_image_generation': bool(tpl.get('allow_image_generation', False)),
        'allow_video_generation': bool(tpl.get('allow_video_generation', False)),
        'allow_call_teacher': bool(tpl.get('allow_call_teacher', False)),
        'allow_voice_input': bool(tpl.get('allow_voice_input', True)),
        'allow_tts': bool(tpl.get('allow_tts', False)),
        'audio_mode': ('direct' if tpl.get('audio_mode') == 'direct' else 'transcribe'),
        'share_mode': tpl.get('share_mode') or 'independent',
        'group_max_size': int(tpl.get('group_max_size') or 6),
        'max_messages_per_student_daily': int(tpl.get('max_messages_per_student_daily') or 50),
        'link_expires_at': datetime.now() + timedelta(hours=48),
        'blocked_keywords': bk_json,
        'status': 'active',
    }


@app.route('/assistant/<int:assistant_id>/export')
@login_required
def assistant_export_template(assistant_id):
    """导出 AI 导师为 JSON 模板"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        tpl = _assistant_to_template_dict(assistant)
        from io import BytesIO
        bio = BytesIO()
        bio.write(json.dumps(tpl, ensure_ascii=False, indent=2).encode('utf-8'))
        bio.seek(0)
        safe_title = re.sub(r'[\\/:*?"<>|\s]+', '_', assistant.title)[:60]
        filename = f'OpenMentor_{safe_title}.json'
        return send_file(bio, as_attachment=True, download_name=filename,
                         mimetype='application/json; charset=utf-8')
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/clone', methods=['POST'])
@login_required
def assistant_clone(assistant_id):
    """在自己账户内克隆一份 AI 导师"""
    db = SessionLocal()
    try:
        src = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not src:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        tpl = _assistant_to_template_dict(src)
        tpl['title'] = f'{src.title}（副本）'
        kwargs = _template_dict_to_assistant_kwargs(tpl, current_user.id)
        new_assistant = Task(**kwargs)
        db.add(new_assistant)
        db.commit()
        db.refresh(new_assistant)
        flash(f'已克隆为「{new_assistant.title}」', 'success')
        return redirect(url_for('assistant_detail', assistant_id=new_assistant.id))
    finally:
        db.close()


# ---- AI 导师广场（公共广场使用同步表格；旧 QuickForm 数据源保留用于回退 / 迁移）----

OM_LEGACY_PLAZA_API_URL = 'https://quickform.cn/api/pa9eaj18kz'
OM_PLAZA_PRESET_SUBJECTS = ['语文', '数学', '英语', '物理', '化学', '生物', '历史', '政治', '地理', '信息', '科学', '音乐', '体育', '美术', '心理', '劳动', '综合实践', '通用']
OM_PLAZA_PRESET_GRADES = ['小学', '初中', '高中', '大学', '培训']


def _lark_public_plaza_configured():
    return all((os.getenv(k) or '').strip() for k in ('LARK_APP_ID', 'LARK_APP_SECRET', 'LARK_PLAZA_BASE_TOKEN', 'LARK_PLAZA_TABLE_ID'))


def _public_plaza_config():
    return {
        'base_token': (os.getenv('LARK_PLAZA_BASE_TOKEN') or '').strip(),
        'table_id': (os.getenv('LARK_PLAZA_TABLE_ID') or '').strip(),
        'view_id': (os.getenv('LARK_PLAZA_VIEW_ID') or '').strip(),
        'data_field': (os.getenv('LARK_PLAZA_DATA_FIELD') or '数据').strip() or '数据',
        'enabled_field': (os.getenv('LARK_PLAZA_ENABLED_FIELD') or '启动').strip() or '启动',
    }


def _public_plaza_records_url():
    cfg = _public_plaza_config()
    return f'https://open.feishu.cn/open-apis/bitable/v1/apps/{cfg["base_token"]}/tables/{cfg["table_id"]}/records'


def _public_plaza_enabled(value):
    if value is None:
        return True
    if isinstance(value, bool):
        return value
    text_value = _prompt_matrix_cell_text(value).strip().lower()
    if not text_value:
        return True
    return text_value not in {'0', 'false', 'no', 'n', '否', '禁用', '停用', 'disabled'}


def _public_plaza_data_text(value):
    if value is None:
        return ''
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, list):
        parts = []
        for item in value:
            if isinstance(item, dict):
                parts.append(str(item.get('text') or item.get('name') or item.get('value') or ''))
            else:
                parts.append(str(item))
        return ''.join(parts).strip()
    if isinstance(value, dict):
        return str(value.get('text') or value.get('name') or value.get('value') or '').strip()
    return str(value).strip()


def _normalize_legacy_public_plaza_record(record):
    if not isinstance(record, dict):
        return None
    if isinstance(record.get('config'), dict):
        return record
    for key in ('data', 'fields', 'payload', 'form_data'):
        nested = record.get(key)
        if isinstance(nested, dict) and isinstance(nested.get('config'), dict):
            return nested
    return None


def _write_legacy_public_plaza_payload(payload):
    resp = requests.post(OM_LEGACY_PLAZA_API_URL, json=payload, timeout=10)
    if not resp.ok:
        raise RuntimeError(resp.text[:200] or f'HTTP {resp.status_code}')
    _PUBLIC_PLAZA_MEMORY_CACHE.update({'ts': 0, 'submissions': None})
    return True


def _fetch_legacy_public_plaza_submissions():
    resp = requests.get(f'{OM_LEGACY_PLAZA_API_URL}/all', timeout=10)
    try:
        data = resp.json()
    except (ValueError, json.JSONDecodeError):
        data = {}
    if not resp.ok:
        raise RuntimeError(data.get('message') or resp.text[:200] or f'HTTP {resp.status_code}')
    records = data.get('submissions') or data.get('data') or data.get('items') or []
    submissions = []
    for record in records:
        item = _normalize_legacy_public_plaza_record(record)
        if item:
            submissions.append(item)
    return submissions


def _write_public_plaza_payload(payload):
    if not _lark_public_plaza_configured():
        return _write_legacy_public_plaza_payload(payload)
    cfg = _public_plaza_config()
    headers = {
        'Authorization': f'Bearer {_get_lark_tenant_access_token()}',
        'Content-Type': 'application/json',
    }
    resp = requests.post(
        _public_plaza_records_url(),
        headers=headers,
        json={'fields': {cfg['data_field']: json.dumps(payload, ensure_ascii=False)}},
        timeout=10,
    )
    try:
        data = resp.json()
    except (ValueError, json.JSONDecodeError):
        data = {}
    if not resp.ok:
        raise RuntimeError(data.get('msg') or data.get('message') or resp.text[:200] or f'HTTP {resp.status_code}')
    if data.get('code') not in (0, None):
        raise RuntimeError(data.get('msg') or '写入公共广场失败')
    _PUBLIC_PLAZA_MEMORY_CACHE.update({'ts': 0, 'submissions': None})
    return data.get('data') or {}


def _fetch_public_plaza_submissions():
    if not _lark_public_plaza_configured():
        return _fetch_legacy_public_plaza_submissions()
    cfg = _public_plaza_config()
    headers = {'Authorization': f'Bearer {_get_lark_tenant_access_token()}'}
    params = {'page_size': 100}
    if cfg['view_id']:
        params['view_id'] = cfg['view_id']
    submissions = []
    page_token = ''
    while True:
        if page_token:
            params['page_token'] = page_token
        resp = requests.get(_public_plaza_records_url(), headers=headers, params=params, timeout=10)
        try:
            data = resp.json()
        except (ValueError, json.JSONDecodeError):
            data = {}
        if not resp.ok:
            raise RuntimeError(data.get('msg') or data.get('message') or resp.text[:200] or f'HTTP {resp.status_code}')
        if data.get('code') not in (0, None):
            raise RuntimeError(data.get('msg') or '读取公共广场失败')
        payload = data.get('data') or {}
        for item in payload.get('items') or []:
            fields = item.get('fields') or {}
            enabled_value = fields.get(cfg['enabled_field'])
            if enabled_value is None and cfg['enabled_field'] != '启动':
                enabled_value = fields.get('启动')
            if enabled_value is None and cfg['enabled_field'] != '启用':
                enabled_value = fields.get('启用')
            if not _public_plaza_enabled(enabled_value):
                continue
            raw = _public_plaza_data_text(fields.get(cfg['data_field']))
            if not raw:
                continue
            try:
                parsed = json.loads(raw)
            except (TypeError, ValueError, json.JSONDecodeError):
                logger.warning('[公共广场] 跳过无法解析的数据记录')
                continue
            if isinstance(parsed, dict) and isinstance(parsed.get('config'), dict):
                submissions.append(parsed)
        if not payload.get('has_more'):
            break
        page_token = payload.get('page_token') or ''
        if not page_token:
            break
    return submissions


def _load_public_plaza_submissions(force_refresh=False):
    now = time.time()
    if not force_refresh and _PUBLIC_PLAZA_MEMORY_CACHE['submissions'] is not None and now - _PUBLIC_PLAZA_MEMORY_CACHE['ts'] < PUBLIC_PLAZA_CACHE_TTL_SECONDS:
        return _PUBLIC_PLAZA_MEMORY_CACHE['submissions']
    submissions = _fetch_public_plaza_submissions()
    _PUBLIC_PLAZA_MEMORY_CACHE.update({'ts': now, 'submissions': submissions})
    return submissions


@app.route('/changelog')
@login_required
def changelog():
    """更新日志页面：按版本顺序展示对老师有感的功能更新"""
    return render_template('changelog.html')


@app.route('/plaza')
@login_required
def plaza():
    """AI 导师广场页面：浏览全国老师分享的 AI 导师，一键克隆"""
    return render_template(
        'plaza.html',
        plaza_api_url=url_for('api_public_plaza_templates'),
        preset_subjects=OM_PLAZA_PRESET_SUBJECTS,
        preset_grades=OM_PLAZA_PRESET_GRADES,
    )


@app.route('/api/plaza/public/templates')
@login_required
def api_public_plaza_templates():
    """公共广场列表。返回格式保持与旧 QuickForm /all 兼容。"""
    try:
        force_refresh = request.args.get('refresh') in {'1', 'true', 'yes'}
        return jsonify({'code': 200, 'success': True, 'submissions': _load_public_plaza_submissions(force_refresh=force_refresh)})
    except Exception as e:
        logger.warning(f'[公共广场] 读取失败: {e}')
        return jsonify({'code': 500, 'success': False, 'message': str(e), 'submissions': []}), 500


@app.route('/api/assistant/<int:assistant_id>/share_to_plaza', methods=['POST'])
@login_required
def assistant_share_to_plaza(assistant_id):
    """老师把自己的 AI 导师配置分享到公共广场、校内广场或两者。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(
            id=assistant_id, user_id=current_user.id, is_assistant=True
        ).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在或无权访问'}), 404

        body = request.get_json(silent=True) or {}
        subject = (body.get('subject') or '').strip()[:30]
        grade = (body.get('grade') or '').strip()[:20]
        purpose = (body.get('purpose') or '').strip()[:60]
        if not subject or not grade or not purpose:
            return jsonify({'code': 400, 'message': '学科 / 学段 / 用途都是必填'}), 400

        destinations = body.get('destinations')
        if not isinstance(destinations, list):
            destinations = ['public']
        destinations = {str(x).strip() for x in destinations if str(x).strip()}
        allowed_destinations = {'public'}
        if is_campus_edition():
            allowed_destinations.add('campus')
        destinations = destinations & allowed_destinations
        if not destinations:
            return jsonify({'code': 400, 'message': '请选择至少一种分享方式'}), 400

        nickname = (body.get('nickname') or '').strip()[:30]
        school = (body.get('school') or '').strip()[:60]
        description = (body.get('description') or '').strip()[:500]

        tpl = _assistant_to_template_dict(assistant)
        # 老师可在分享时改写一份"专门展示给广场看"的描述
        if description:
            tpl['description'] = description

        plaza_payload = {
            'subject': subject,
            'grade': grade,
            'purpose': purpose,
            'description': description or (tpl.get('description') or '')[:300],
            'nickname': nickname or '匿名',
            'school': school,
            'shared_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'om_version': _parse_latest_version_from_changelog(),
            'config': tpl,
        }

        shared_targets = []
        if 'public' in destinations:
            try:
                _write_public_plaza_payload(plaza_payload)
                shared_targets.append('公共广场')
            except Exception as e:
                logger.exception(f'分享 AI 导师到公共广场失败: {e}')
                return jsonify({'code': 500, 'message': f'连接公共广场服务失败: {e}'}), 500

        if 'campus' in destinations:
            local_tpl = CampusPlazaTemplate(
                source_assistant_id=assistant.id,
                owner_user_id=current_user.id,
                title=tpl.get('title') or assistant.title,
                subject=subject,
                grade=grade,
                purpose=purpose,
                description=plaza_payload['description'],
                nickname=nickname or (current_user.display_name or current_user.username),
                school=school,
                config_json=json.dumps(tpl, ensure_ascii=False),
                status='active',
            )
            db.add(local_tpl)
            db.commit()
            shared_targets.append('校内广场')

        return jsonify({
            'code': 200,
            'success': True,
            'message': '已成功分享到' + '、'.join(shared_targets),
        })
    finally:
        db.close()


@app.route('/api/campus/plaza/templates')
@login_required
def campus_plaza_templates():
    """校园版校内广场模板列表。普通教师只看启用项；admin 可看全部并管理状态。"""
    if not is_campus_edition():
        return jsonify({'code': 404, 'message': '仅校园版支持校内广场'}), 404

    db = SessionLocal()
    try:
        q = db.query(CampusPlazaTemplate, User).join(User, CampusPlazaTemplate.owner_user_id == User.id)
        if not is_admin_user(current_user):
            q = q.filter(CampusPlazaTemplate.status == 'active')
        rows = q.order_by(CampusPlazaTemplate.created_at.desc()).all()
        items = []
        for rec, user in rows:
            try:
                config = json.loads(rec.config_json or '{}')
            except Exception:
                config = {}
            items.append({
                'id': rec.id,
                'subject': rec.subject,
                'grade': rec.grade,
                'purpose': rec.purpose,
                'description': rec.description or '',
                'nickname': rec.nickname or (user.display_name or user.username),
                'school': rec.school or '',
                'shared_at': rec.created_at.strftime('%Y-%m-%d %H:%M:%S') if rec.created_at else '',
                'status': rec.status,
                'owner_username': user.username,
                'config': config,
            })
        return jsonify({'code': 200, 'success': True, 'submissions': items, 'can_manage': is_admin_user(current_user)})
    finally:
        db.close()


@app.route('/api/campus/plaza/templates/<int:template_id>/toggle', methods=['POST'])
@login_required
@campus_admin_required
def campus_plaza_template_toggle(template_id):
    """admin 启用/禁用校内广场模板；不做删除。"""
    db = SessionLocal()
    try:
        rec = db.query(CampusPlazaTemplate).get(template_id)
        if not rec:
            return jsonify({'code': 404, 'message': '模板不存在'}), 404
        rec.status = 'disabled' if rec.status == 'active' else 'active'
        rec.updated_at = datetime.now()
        db.commit()
        return jsonify({'code': 200, 'success': True, 'status': rec.status})
    finally:
        db.close()


@app.route('/api/plaza/clone', methods=['POST'])
@login_required
def plaza_clone():
    """从广场卡片直接克隆一份 AI 导师到当前老师的库里"""
    db = SessionLocal()
    try:
        body = request.get_json(silent=True) or {}
        config = body.get('config')
        if not isinstance(config, dict):
            return jsonify({'code': 400, 'message': '克隆数据无效（缺少 config 对象）'}), 400
        try:
            kwargs = _template_dict_to_assistant_kwargs(config, current_user.id)
        except ValueError as e:
            return jsonify({'code': 400, 'message': f'模板字段错误：{e}'}), 400

        # 标题加「广场克隆」后缀，方便老师区分来源
        kwargs['title'] = f"{kwargs['title']}（广场克隆）"

        new_assistant = Task(**kwargs)
        db.add(new_assistant)
        db.commit()
        db.refresh(new_assistant)

        return jsonify({
            'code': 200,
            'success': True,
            'assistant_id': new_assistant.id,
            'message': f'已克隆为「{new_assistant.title}」',
            'redirect_url': url_for('assistant_detail', assistant_id=new_assistant.id),
        })
    finally:
        db.close()


@app.route('/assistant/import', methods=['GET', 'POST'])
@login_required
def assistant_import():
    """从 JSON 模板导入 AI 导师"""
    if request.method == 'POST':
        file = request.files.get('template_file')
        json_text = (request.form.get('template_json') or '').strip()
        if not file and not json_text:
            flash('请上传 JSON 文件或粘贴 JSON 文本', 'warning')
            return redirect(url_for('assistant_import'))

        try:
            if file and file.filename:
                content = file.read().decode('utf-8', errors='ignore')
            else:
                content = json_text
            tpl = json.loads(content)
        except (json.JSONDecodeError, UnicodeDecodeError) as e:
            flash(f'JSON 解析失败：{e}', 'danger')
            return redirect(url_for('assistant_import'))

        if not isinstance(tpl, dict):
            flash('模板格式错误：根节点应为对象', 'danger')
            return redirect(url_for('assistant_import'))

        # 兼容性检查：版本 + 必填
        ver = tpl.get('openmentor_template_version')
        if ver and ver != OPENMENTOR_TEMPLATE_VERSION:
            flash(f'模板版本 {ver} 与当前版本 {OPENMENTOR_TEMPLATE_VERSION} 不一致，已尝试导入但可能字段不完整', 'warning')

        try:
            kwargs = _template_dict_to_assistant_kwargs(tpl, current_user.id)
        except ValueError as e:
            flash(f'模板字段错误：{e}', 'danger')
            return redirect(url_for('assistant_import'))

        db = SessionLocal()
        try:
            new_assistant = Task(**kwargs)
            db.add(new_assistant)
            db.commit()
            db.refresh(new_assistant)
            flash(f'已导入「{new_assistant.title}」。请记得在「个人设置」配置对应的 AI Key。', 'success')
            return redirect(url_for('assistant_detail', assistant_id=new_assistant.id))
        finally:
            db.close()

    return render_template('assistant_import.html')


@app.route('/assistant/<int:assistant_id>/delete', methods=['POST'])
@login_required
def assistant_delete(assistant_id):
    """删除 AI 导师（含级联删除会话和消息）"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        title = assistant.title
        # 知识库文档（独立表 + 落盘文件）一并清理
        kb_docs = db.query(KnowledgeDoc).filter_by(assistant_id=assistant.id).all()
        for d in kb_docs:
            if d.file_path and os.path.exists(d.file_path):
                try:
                    os.remove(d.file_path)
                except OSError:
                    pass
            db.delete(d)  # cascade 删除 chunks
        db.delete(assistant)  # 级联删除 conversations & messages
        db.commit()
        flash(f'AI 导师「{title}」已删除', 'success')
        return redirect(url_for('assistant_list'))
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/screen')
@login_required
def assistant_screen(assistant_id):
    """课堂投屏模式：全屏大二维码 + 实时进入/提问统计，专为投影仪设计。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        student_url = f"{request.host_url.rstrip('/')}/chat/{assistant.task_id}"
        return render_template(
            'assistant_screen.html',
            assistant=assistant,
            student_url=student_url,
            is_group=_assistant_is_group(assistant),
        )
    finally:
        db.close()


@app.route('/api/assistant/<int:assistant_id>/screen_stats')
@login_required
def api_assistant_screen_stats(assistant_id):
    """投屏页实时统计：今日进入数 + 今日提问数（排除教师试聊）。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': '不存在'}), 404
        today_start = datetime.combine(datetime.now().date(), datetime.min.time())
        joined_today = (
            db.query(Conversation)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
            .filter(Conversation.last_active_at >= today_start)
            .count()
        )
        questions_today = (
            db.query(Message)
            .join(Conversation, Message.conversation_id == Conversation.id)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
            .filter(Message.role == 'user')
            .filter(Message.created_at >= today_start)
            .count()
        )
        return jsonify({'joined': joined_today, 'questions': questions_today})
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/try')
@login_required
def assistant_try_chat(assistant_id):
    """老师试聊：以「教师试聊」班级身份进入学生聊天页。
    试聊会话不计入导师统计/审计列表/报告（按 TEACHER_PREVIEW_CLASS 排除）。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))

        teacher_label = (current_user.display_name or current_user.username or '老师')[:50]
        token = _hash_student_token(TEACHER_PREVIEW_CLASS, f'preview|{current_user.id}', None)
        conv = db.query(Conversation).filter_by(assistant_id=assistant.id, student_token=token).first()
        if not conv:
            conv = Conversation(
                assistant_id=assistant.id,
                student_class=TEACHER_PREVIEW_CLASS,
                student_name=teacher_label,
                student_token=token,
                started_at=datetime.now(),
                last_active_at=datetime.now(),
                daily_reset_date=datetime.now().date(),
            )
            db.add(conv)
            db.commit()

        resp = redirect(url_for('chat_entry', task_id=assistant.task_id))
        resp.set_cookie(_student_cookie_name(assistant.task_id), token, httponly=True, samesite='Lax')
        if _assistant_is_group(assistant):
            from urllib.parse import quote
            resp.set_cookie(_member_cookie_name(assistant.task_id), quote(teacher_label), httponly=True, samesite='Lax')
        return resp
    finally:
        db.close()


# ------------------------------------------------------------
# OpenMentor: 掌上讲台 App 服务端地基（/api/app/v1）
# 设计文档：创AI案例/OpenMentor-App/功能设计-v0.2.md 第 4 节
# 原则：独立命名空间、纯增量；不装 App 时全部沉睡。推送发送逻辑暂不实现（仅事件钩子 + 注册接口）。
# ------------------------------------------------------------

import secrets as _secrets
from collections import deque as _deque
import itertools as _itertools

APP_API_LEVEL = 1
APP_PAIR_CODE_TTL = 600          # 配对码有效期 10 分钟
_APP_PAIR_CODES = {}             # code -> {'user_id', 'expires'}
_APP_EVENTS = {}                 # user_id -> deque({'id','type','ts','data'})
_APP_EVENT_SEQ = _itertools.count(1)
_APP_LOCK = threading.Lock()
_CALL_TEACHER_LAST = {}          # (conv_id, member_name) -> last_ts（找老师限频：每生 60 秒 1 次）


def _emit_app_event(user_id, event_type, data=None):
    """向某位老师的事件流写一条事件（App 轮询 /api/app/v1/events 消费；
    将来推送功能也订阅这里）。失败静默，绝不影响业务。"""
    try:
        with _APP_LOCK:
            q = _APP_EVENTS.setdefault(int(user_id), _deque(maxlen=200))
            q.append({
                'id': next(_APP_EVENT_SEQ),
                'type': event_type,
                'ts': datetime.now().isoformat(),
                'data': data or {},
            })
    except Exception as e:
        logger.warning(f'[App事件] 发射失败（已忽略）: {e}')
    try:
        _t, _b = _push_text_for_event(event_type, data)
        if _t:
            _push_to_user(user_id, _t, _b, event_type, data=data)
    except Exception as e:
        logger.warning(f'[推送] 触发失败（已忽略）: {e}')


# ============ 掌上讲台锁屏推送（APNs） ============
# 直连模式：.env 配 OM_APNS_KEY_PATH / OM_APNS_KEY_ID（开发者自有部署）；
# 中继模式：.env 配 OM_PUSH_RELAY_URL，事件转给中继服务器代发（多租户铺开用）。
# 都没配 → 推送静默关闭，业务零影响；依赖 httpx[http2] + PyJWT，缺了也只是关闭推送。
APNS_KEY_PATH = os.environ.get('OM_APNS_KEY_PATH', '')
APNS_KEY_ID = os.environ.get('OM_APNS_KEY_ID', '')
APNS_TEAM_ID = os.environ.get('OM_APNS_TEAM_ID', 'VYD397KDCW')
APNS_TOPIC = os.environ.get('OM_APNS_TOPIC', 'top.openmentor.openmentorApp')
PUSH_RELAY_URL = os.environ.get('OM_PUSH_RELAY_URL', '')
APNS_PROXY = os.environ.get('OM_APNS_PROXY', '')
_APNS_JWT_CACHE = {'token': None, 'ts': 0.0}
_PUSH_LAST_SENT = {}        # (user_id, event_type) -> ts，同类事件 8 秒内只推 1 条防刷屏
_OFFLINE_PUSH_STATE = {'serving': False, 'sent': False}


def _push_enabled():
    if PUSH_RELAY_URL:
        return True
    if not (APNS_KEY_PATH and APNS_KEY_ID and os.path.exists(APNS_KEY_PATH)):
        return False
    try:
        import httpx, h2, jwt, cryptography  # noqa: F401  (cryptography: PyJWT 的 ES256 后端)
        return True
    except ImportError:
        return False


def _apns_jwt():
    import jwt as _pyjwt
    now = time.time()
    if _APNS_JWT_CACHE['token'] and now - _APNS_JWT_CACHE['ts'] < 2700:
        return _APNS_JWT_CACHE['token']
    with open(APNS_KEY_PATH) as f:
        key = f.read()
    tok = _pyjwt.encode({'iss': APNS_TEAM_ID, 'iat': int(now)}, key,
                        algorithm='ES256', headers={'kid': APNS_KEY_ID})
    _APNS_JWT_CACHE.update(token=tok, ts=now)
    return tok


def _apns_send(tokens, title, body, event_type='', timeout=8, data=None):
    """同步发送（调用方自行放后台线程）。返回已失效的 token 列表。"""
    dead = []
    if PUSH_RELAY_URL:
        try:
            requests.post(PUSH_RELAY_URL.rstrip('/') + '/push', json={
                'tokens': tokens, 'title': title, 'body': body, 'event': event_type,
                'data': data or {},
            }, timeout=timeout, proxies={'http': None, 'https': None})
        except Exception as e:
            logger.warning(f'[推送] 中继发送失败: {e}')
        return dead
    import httpx
    payload = {
        'aps': {'alert': {'title': title, 'body': body}, 'sound': 'default',
                'thread-id': event_type or 'openmentor'},
        'om_event': event_type,
        # 通知即操作：App 点开推送直接弹处理面板要用的上下文（assistant_id/conversation_id 等）
        'om_data': data or {},
    }
    headers = {
        'authorization': f'bearer {_apns_jwt()}',
        'apns-topic': APNS_TOPIC,
        'apns-push-type': 'alert',
        'apns-priority': '10',
    }
    # 国内网络直连 api.push.apple.com 可能被墙：优先走 OM_APNS_PROXY（如本机 Clash），
    # 失败再直连（学校等可直连的网络不依赖代理）
    candidates = []
    if APNS_PROXY:
        candidates.append({'proxy': APNS_PROXY})
    candidates.append({'trust_env': True})
    for kw in candidates:
        ok, errors = 0, 0
        try:
            with httpx.Client(http2=True, timeout=timeout, **kw) as client:
                for t in tokens:
                    try:
                        r = client.post(f'https://api.push.apple.com/3/device/{t}',
                                        json=payload, headers=headers)
                        ok += 1
                        if r.status_code == 410 or (r.status_code == 400 and 'BadDeviceToken' in r.text):
                            dead.append(t)
                        elif r.status_code != 200:
                            logger.warning(f'[推送] APNs {r.status_code}: {r.text[:200]}')
                    except Exception as e:
                        errors += 1
                        logger.warning(f'[推送] 单条发送失败({kw}): {e}')
        except Exception as e:
            logger.warning(f'[推送] APNs 连接失败({kw}): {e}')
            continue
        if ok > 0 or errors == 0:
            break  # 本通道可用（或没东西可发），不再换下一通道
        dead = []
    return dead


def _push_to_user(user_id, title, body, event_type='', data=None):
    """给某位老师的全部已登记设备发锁屏推送（后台线程，绝不阻塞业务）。"""
    if not _push_enabled():
        return
    now = time.time()
    key = (int(user_id), event_type)
    if now - _PUSH_LAST_SENT.get(key, 0) < 8:
        return
    _PUSH_LAST_SENT[key] = now

    def _worker():
        db = SessionLocal()
        try:
            devs = db.query(AppDevice).filter(
                AppDevice.user_id == int(user_id),
                AppDevice.revoked == False,  # noqa: E712
                AppDevice.push_token.isnot(None)).all()
            # 去重：同一台手机可能在多条设备记录上留有相同 token（换账号重配对的历史）
            tokens = list({d.push_token for d in devs if d.push_token})
            if not tokens:
                logger.info(f'[推送] {event_type} → 用户{user_id} 无已登记设备，跳过')
                return
            logger.info(f'[推送] {event_type} → 用户{user_id} 的 {len(tokens)} 台设备')
            dead = _apns_send(tokens, title, body, event_type, data=data)
            if dead:
                for d in devs:
                    if d.push_token in dead:
                        d.push_token = None
                db.commit()
        except Exception as e:
            logger.warning(f'[推送] 任务失败（已忽略）: {e}')
        finally:
            db.close()

    threading.Thread(target=_worker, daemon=True).start()


def _push_text_for_event(event_type, data):
    d = data or {}
    if event_type == 'call_teacher':
        return ('学生举手', f"{d.get('student_class', '')} {d.get('student_name', '')} "
                           f"在「{d.get('assistant_title', '')}」请求帮助")
    if event_type == 'blocked_keyword':
        return ('触发敏感词', f"{d.get('student_name', '')} 触发敏感词（{d.get('category', '')}）")
    if event_type == 'quota_warning':
        return ('额度预警', f"{d.get('student_name', '')} 额度已用 {d.get('used', '')}/{d.get('limit', '')}")
    if event_type == 'assistant_toggled':
        return ('导师状态', f"「{d.get('assistant_title', '')}」已"
                           f"{'启用' if d.get('status') == 'active' else '熔断'}")
    return (None, None)


def _push_server_offline():
    """服务器正常退出（关机/关进程）时，给所有设备推一条下线提醒。
    只在真正 serve 过之后触发，避免脚本 import app.py 时误报。"""
    if not _OFFLINE_PUSH_STATE['serving'] or _OFFLINE_PUSH_STATE['sent'] or not _push_enabled():
        return
    _OFFLINE_PUSH_STATE['sent'] = True
    db = SessionLocal()
    try:
        devs = db.query(AppDevice).filter(
            AppDevice.revoked == False,  # noqa: E712
            AppDevice.push_token.isnot(None)).all()
        tokens = list({d.push_token for d in devs if d.push_token})
        if tokens:
            _apns_send(tokens, '服务器已下线', 'OpenMentor 服务器刚刚关闭，学生端将无法访问',
                       'server_offline', timeout=4)
            logger.info(f'[推送] 已向 {len(tokens)} 台设备发送下线提醒')
    except Exception as e:
        logger.warning(f'[推送] 下线提醒失败: {e}')
    finally:
        db.close()


def app_token_required(func):
    """掌上讲台 App 的 Bearer 长效令牌鉴权（独立于网页 Session）。"""
    @wraps(func)
    def wrapper(*args, **kwargs):
        auth = request.headers.get('Authorization') or ''
        token = auth[7:].strip() if auth.startswith('Bearer ') else ''
        if not token:
            return jsonify({'code': 401, 'message': '缺少 App 令牌'}), 401
        import hashlib as _hl
        token_hash = _hl.sha256(token.encode('utf-8')).hexdigest()
        db = SessionLocal()
        try:
            dev = db.query(AppDevice).filter_by(token_hash=token_hash, revoked=False).first()
            if not dev:
                return jsonify({'code': 401, 'message': 'App 令牌无效或已撤销，请重新扫码配对'}), 401
            dev.last_seen_at = datetime.now()
            db.commit()
            request.om_app_user_id = dev.user_id
            request.om_app_device_id = dev.id
        finally:
            db.close()
        return func(*args, **kwargs)
    return wrapper


@app.route('/app/pairing')
@login_required
def app_pairing():
    """网页端：生成掌上讲台 App 配对二维码（地址 + 一次性配对码，10 分钟有效）。"""
    code = _secrets.token_hex(4).upper()
    now = time.time()
    with _APP_LOCK:
        # 清理过期码
        for c in [c for c, v in _APP_PAIR_CODES.items() if v['expires'] < now]:
            _APP_PAIR_CODES.pop(c, None)
        _APP_PAIR_CODES[code] = {'user_id': current_user.id, 'expires': now + APP_PAIR_CODE_TTL}
    pair_payload = json.dumps({
        'om_pair': 1,
        'addr': request.host_url.rstrip('/'),
        'code': code,
    }, ensure_ascii=False)
    return render_template('app_pairing.html', pair_payload=pair_payload, pair_code=code)


@app.route('/api/app/v1/capabilities')
def api_app_capabilities():
    """能力握手（无需令牌）：App 据此决定哪些功能可用/置灰。"""
    return jsonify({
        'app_api': APP_API_LEVEL,
        'version': _parse_latest_version_from_changelog(),
        'edition': APP_EDITION,
        'app_name': get_app_name(),
        'features': ['pair', 'events', 'assistants', 'announce', 'call_teacher', 'push_register', 'actions']
                    + (['push'] if _push_enabled() else []),
    })


@app.route('/api/app/v1/pair', methods=['POST'])
def api_app_pair():
    """App 扫码配对：一次性配对码换长效令牌。"""
    body = request.get_json(silent=True) or {}
    code = (body.get('code') or '').strip().upper()
    device_name = (body.get('device_name') or '未命名设备').strip()[:100]
    platform = (body.get('platform') or '').strip().lower()[:20]
    now = time.time()
    with _APP_LOCK:
        info = _APP_PAIR_CODES.get(code)
        if not info or info['expires'] < now:
            _APP_PAIR_CODES.pop(code, None)
            return jsonify({'code': 400, 'message': '配对码无效或已过期，请在电脑端重新生成'}), 400
        _APP_PAIR_CODES.pop(code, None)  # 一次性
        user_id = info['user_id']

    token = _secrets.token_urlsafe(32)
    import hashlib as _hl
    db = SessionLocal()
    try:
        dev = AppDevice(
            user_id=user_id,
            device_name=device_name,
            platform=platform,
            token_hash=_hl.sha256(token.encode('utf-8')).hexdigest(),
        )
        db.add(dev)
        db.commit()
        user = db.query(User).get(user_id)
        logger.info(f'[App配对] 用户 {user_id} 绑定设备「{device_name}」({platform})')
        return jsonify({
            'code': 200,
            'token': token,
            'teacher': (user.display_name or user.username) if user else '',
            'version': _parse_latest_version_from_changelog(),
            'edition': APP_EDITION,
        })
    finally:
        db.close()


@app.route('/api/app/v1/me')
@app_token_required
def api_app_me():
    db = SessionLocal()
    try:
        user = db.query(User).get(request.om_app_user_id)
        devices = db.query(AppDevice).filter_by(user_id=request.om_app_user_id, revoked=False).all()
        return jsonify({
            'code': 200,
            'teacher': (user.display_name or user.username) if user else '',
            'username': user.username if user else '',
            'devices': [{
                'id': d.id, 'name': d.device_name, 'platform': d.platform,
                'last_seen_at': d.last_seen_at.isoformat() if d.last_seen_at else None,
                'is_current': d.id == request.om_app_device_id,
            } for d in devices],
        })
    finally:
        db.close()


@app.route('/api/app/v1/push/register', methods=['POST'])
@app_token_required
def api_app_push_register():
    """登记推送 token（发送逻辑随 App 正式开发再上线，本接口先收好）。"""
    body = request.get_json(silent=True) or {}
    push_token = (body.get('push_token') or '').strip()[:200]
    db = SessionLocal()
    try:
        dev = db.query(AppDevice).get(request.om_app_device_id)
        if dev:
            dev.push_token = push_token or None
            db.commit()
        return jsonify({'code': 200})
    finally:
        db.close()


@app.route('/api/app/v1/assistants')
@app_token_required
def api_app_assistants():
    """老师的 AI 导师列表 + 课堂统计（复用移动控制台口径）。"""
    db = SessionLocal()
    try:
        assistants = (
            db.query(Task)
            .filter_by(user_id=request.om_app_user_id, is_assistant=True)
            .order_by(Task.created_at.desc())
            .all()
        )
        return jsonify({'code': 200, 'items': [_mobile_assistant_summary(db, a) for a in assistants]})
    finally:
        db.close()


@app.route('/api/app/v1/events')
@app_token_required
def api_app_events():
    """事件流轮询（敏感词命中 / 熔断 / 额度预警 / 找老师）。?after_id=N 增量拉取。"""
    after_id = request.args.get('after_id', type=int) or 0
    with _APP_LOCK:
        q = _APP_EVENTS.get(request.om_app_user_id) or []
        events = [e for e in q if e['id'] > after_id]
    return jsonify({
        'code': 200,
        'events': events,
        'last_id': events[-1]['id'] if events else after_id,
    })


@app.route('/api/app/v1/assistant/<int:assistant_id>/announce', methods=['POST'])
@app_token_required
def api_app_announce(assistant_id):
    """全班公告：向该导师今日活跃的所有学生会话插入一条系统消息。"""
    body = request.get_json(silent=True) or {}
    text = (body.get('text') or '').strip()
    if not text:
        return jsonify({'code': 400, 'message': '公告内容不能为空'}), 400
    if len(text) > 300:
        return jsonify({'code': 400, 'message': '公告过长（≤ 300 字符）'}), 400
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        today_start = datetime.combine(datetime.now().date(), datetime.min.time())
        convs = (
            db.query(Conversation)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
            .filter(Conversation.last_active_at >= today_start)
            .all()
        )
        for c in convs:
            db.add(Message(conversation_id=c.id, role='system', content=f'📢 老师公告：{text}'))
        db.commit()
        logger.info(f'[App公告] 导师 {assistant.id} → {len(convs)} 个今日会话')
        return jsonify({'code': 200, 'sent_to': len(convs)})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/toggle', methods=['POST'])
@app_token_required
def api_app_assistant_toggle(assistant_id):
    """App：熔断/恢复导师。body 带 {'status': 'active'|'disabled'} 为幂等设置，不带则切换。
    与网页版 assistant_toggle 同语义（通知即操作的核心动作）。"""
    body = request.get_json(silent=True) or {}
    target = body.get('status')
    if target not in (None, 'active', 'disabled'):
        return jsonify({'code': 400, 'message': 'status 只能是 active 或 disabled'}), 400
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        new_status = target or ('disabled' if assistant.status == 'active' else 'active')
        changed = assistant.status != new_status
        assistant.status = new_status
        db.commit()
        if changed:
            _emit_app_event(assistant.user_id, 'assistant_toggled', {
                'assistant_id': assistant.id, 'assistant_title': assistant.title, 'status': assistant.status,
            })
        logger.info(f'[App操作] 导师 {assistant.id} → {new_status}')
        return jsonify({'code': 200, 'status': assistant.status,
                        'message': '已熔断，学生无法访问' if new_status == 'disabled' else '已恢复'})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/conversation/<int:conv_id>/block', methods=['POST'])
@app_token_required
def api_app_block_student(assistant_id, conv_id):
    """App：封禁/解禁某个学生会话。body 带 {'blocked': true|false} 为幂等设置，不带则切换。"""
    body = request.get_json(silent=True) or {}
    target = body.get('blocked')
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        conv = db.query(Conversation).filter_by(id=conv_id, assistant_id=assistant.id).first()
        if not conv:
            return jsonify({'code': 404, 'message': '学生会话不存在'}), 404
        conv.is_blocked = bool(target) if target is not None else (not conv.is_blocked)
        db.commit()
        logger.info(f'[App操作] 会话 {conv.id}（{conv.student_name}）→ {"封禁" if conv.is_blocked else "解禁"}')
        return jsonify({'code': 200, 'is_blocked': bool(conv.is_blocked),
                        'student_name': conv.student_name,
                        'message': f'已封禁 {conv.student_name}' if conv.is_blocked else f'已解禁 {conv.student_name}'})
    finally:
        db.close()


_SUMMARY_STOPWORDS = {
    '什么', '怎么', '为什么', '哪个', '哪些', '可以', '帮我', '一下', '这个', '那个', '老师',
    '请问', '问题', '知道', '应该', '如果', '因为', '所以', '然后', '还有', '没有', '我们',
    '你们', '他们', '这样', '那样', '怎样', '如何', '是不是', '能不能', '告诉', '解释',
}


@app.route('/api/app/v1/assistant/<int:assistant_id>/summary')
@app_token_required
def api_app_assistant_summary(assistant_id):
    """App：下课一键总结——时间窗内的课堂数据（人数/提问/热词/风险/样例），供生成分享卡片。
    热词依赖 jieba（可选）；未安装则返回空列表，其余照常。"""
    minutes = max(5, min(1440, request.args.get('minutes', type=int) or 60))
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        cutoff = datetime.now() - timedelta(minutes=minutes)
        convs = (
            db.query(Conversation)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
            .all()
        )
        conv_by_id = {c.id: c for c in convs}
        msgs = []
        if conv_by_id:
            msgs = (
                db.query(Message)
                .filter(Message.conversation_id.in_(list(conv_by_id)))
                .filter(Message.created_at >= cutoff)
                .filter(Message.role == 'user')
                .order_by(Message.created_at.asc())
                .all()
            )

        # 学生参与度与风险明细
        per_student = {}
        risk_items = []
        for m in msgs:
            c = conv_by_id.get(m.conversation_id)
            if not c:
                continue
            name = f'{c.student_class} {c.student_name}'.strip()
            per_student[name] = per_student.get(name, 0) + 1
            if m.triggered_keyword:
                parts = (m.triggered_keyword or '').split(':', 1)
                risk_items.append({
                    'student': name,
                    'category': parts[0],
                    'keyword': parts[1] if len(parts) > 1 else '',
                    'time': m.created_at.strftime('%H:%M') if m.created_at else '',
                })
        top_students = sorted(per_student.items(), key=lambda kv: -kv[1])[:3]

        # 热词（jieba 可选）
        hot_words = []
        try:
            import jieba
            counts = {}
            for m in msgs:
                if m.triggered_keyword:
                    continue
                for w in jieba.cut(m.content or ''):
                    w = w.strip()
                    if len(w) < 2 or w in _SUMMARY_STOPWORDS or not any('一' <= ch <= '鿿' for ch in w):
                        continue
                    counts[w] = counts.get(w, 0) + 1
            hot_words = [w for w, n in sorted(counts.items(), key=lambda kv: -kv[1])[:8] if n >= 2]
        except ImportError:
            pass

        # 典型提问样例：取最长的 3 条正常提问
        normal = [m.content for m in msgs if not m.triggered_keyword and (m.content or '').strip()]
        samples = [s[:60] for s in sorted(set(normal), key=len, reverse=True)[:3]]

        return jsonify({'code': 200, 'summary': {
            'assistant_title': assistant.title,
            'window_minutes': minutes,
            'generated_at': datetime.now().strftime('%Y-%m-%d %H:%M'),
            'student_count': len(per_student),
            'question_count': len(msgs),
            'violation_count': len(risk_items),
            'top_students': [{'name': n, 'count': c} for n, c in top_students],
            'hot_words': hot_words,
            'risk_items': risk_items[:10],
            'question_samples': samples,
        }})
    finally:
        db.close()


@app.route('/api/app/v1/datatasks')
@app_token_required
def api_app_datatasks():
    """App：数据任务列表（QuickForm 老本行）。附件区分挂载 HTML 页与普通文件。"""
    _ensure_template_tasks_for_user(request.om_app_user_id)  # 懒下发内置模板任务
    db = SessionLocal()
    try:
        tasks = (
            db.query(Task)
            .filter_by(user_id=request.om_app_user_id, is_assistant=False)
            .order_by(Task.pinned.desc(), Task.created_at.desc())
            .all()
        )
        items = []
        for t in tasks:
            sub_count = db.query(Submission).filter_by(task_id=t.id).count()
            atts = []
            for a in (t.attachments or []):
                fname = (a.file_path or '').split('/')[-1].split('\\')[-1]
                ext = (a.file_name or '').rsplit('.', 1)[-1].lower() if '.' in (a.file_name or '') else ''
                atts.append({
                    'name': a.file_name,
                    'url': f'/static/uploads/{fname}',
                    'is_html': ext in ('html', 'htm'),
                })
            items.append({
                'id': t.id,
                'title': t.title,
                'submission_count': sub_count,
                'attachments': atts,
                'pinned': bool(t.pinned),
                'is_template': bool(t.template_key),
                'created_at': t.created_at.strftime('%Y-%m-%d') if t.created_at else '',
            })
        return jsonify({'code': 200, 'items': items})
    finally:
        db.close()


@app.route('/api/app/v1/datatask/<int:task_id>/submissions')
@app_token_required
def api_app_datatask_submissions(task_id):
    """App：某数据任务的提交记录（含学生上传的附件，URL 为公开静态路径）。"""
    db = SessionLocal()
    try:
        task = db.query(Task).filter_by(id=task_id, user_id=request.om_app_user_id, is_assistant=False).first()
        if not task:
            return jsonify({'code': 404, 'message': '数据任务不存在'}), 404
        subs = (
            db.query(Submission)
            .filter_by(task_id=task.id)
            .order_by(Submission.submitted_at.desc())
            .limit(200)
            .all()
        )
        items = []
        for sub in subs:
            fields = None
            raw = sub.data or ''
            attachments = []
            try:
                parsed = json.loads(sub.data)
                if isinstance(parsed, dict):
                    attachments = parsed.pop('_attachments', None) or []
                    fields = {str(k): ('' if v is None else str(v)) for k, v in parsed.items()}
                    raw = ''
            except (json.JSONDecodeError, TypeError, ValueError):
                pass
            items.append({
                'id': sub.id,
                'submitted_at': sub.submitted_at.strftime('%m-%d %H:%M') if sub.submitted_at else '',
                'fields': fields,
                'raw': raw,
                'attachments': attachments,
            })
        return jsonify({'code': 200, 'title': task.title, 'items': items})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/renew', methods=['POST'])
@app_token_required
def api_app_assistant_renew(assistant_id):
    """App：过期导师复原——重置学生链接有效时长（1-168 小时，默认 48），与网页/移动控制台同语义。"""
    body = request.get_json(silent=True) or {}
    try:
        hours = int(body.get('hours') or 48)
    except (TypeError, ValueError):
        return jsonify({'code': 400, 'message': '有效时长必须是数字'}), 400
    if hours < 1 or hours > 168:
        return jsonify({'code': 400, 'message': '有效时长需在 1-168 小时之间'}), 400
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        assistant.link_expires_at = datetime.now() + timedelta(hours=hours)
        db.commit()
        logger.info(f'[App操作] 导师 {assistant.id} 链接续期 {hours} 小时')
        return jsonify({'code': 200, 'link_expires_at': assistant.link_expires_at.strftime('%m-%d %H:%M'),
                        'message': f'已复原，学生链接有效至 {assistant.link_expires_at.strftime("%m-%d %H:%M")}'})
    finally:
        db.close()


@app.route('/api/app/v1/conversation/<int:conv_id>/notify', methods=['POST'])
@app_token_required
def api_app_conversation_notify(conv_id):
    """给单个学生会话发一条系统提示（如收到举手后回「收到啦」）。校验归属当前老师。"""
    body = request.get_json(silent=True) or {}
    text = (body.get('text') or '').strip()
    if not text:
        return jsonify({'code': 400, 'message': '内容不能为空'}), 400
    if len(text) > 200:
        return jsonify({'code': 400, 'message': '内容过长'}), 400
    db = SessionLocal()
    try:
        conv = db.query(Conversation).filter_by(id=conv_id).first()
        if not conv:
            return jsonify({'code': 404, 'message': '会话不存在'}), 404
        assistant = db.query(Task).filter_by(id=conv.assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': '会话不存在'}), 404
        db.add(Message(conversation_id=conv.id, role='system', content=f'📢 老师：{text}'))
        conv.last_active_at = datetime.now()
        db.commit()
        logger.info(f'[App单发] 会话 {conv.id}（{conv.student_name}）← {text}')
        return jsonify({'code': 200, 'message': f'已回复 {conv.student_name}'})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/preview_link')
@app_token_required
def api_app_assistant_preview_link(assistant_id):
    """App 试聊：以「教师试聊」身份预备一个预览会话，返回可在浏览器打开的自鉴权链接
    （与网页版 /assistant/<id>/try 同语义；预览会话按 TEACHER_PREVIEW_CLASS 排除统计）。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        user = db.query(User).get(request.om_app_user_id)
        teacher_label = ((user.display_name if user else None) or (user.username if user else None) or '老师')[:50]
        token = _hash_student_token(TEACHER_PREVIEW_CLASS, f'preview|{request.om_app_user_id}', None)
        conv = db.query(Conversation).filter_by(assistant_id=assistant.id, student_token=token).first()
        if not conv:
            conv = Conversation(
                assistant_id=assistant.id,
                student_class=TEACHER_PREVIEW_CLASS,
                student_name=teacher_label,
                student_token=token,
                started_at=datetime.now(),
                last_active_at=datetime.now(),
                daily_reset_date=datetime.now().date(),
            )
            db.add(conv)
            db.commit()
        url = f"{request.host_url.rstrip('/')}/chat/{assistant.task_id}/preview?k={token}"
        return jsonify({'code': 200, 'url': url})
    finally:
        db.close()


@app.route('/chat/<string:task_id>/preview')
def chat_preview_enter(task_id):
    """自鉴权预览入口：凭 preview_link 颁发的 token 设置学生 cookie 后进入聊天页（App 试聊用）。"""
    k = (request.args.get('k') or '').strip()
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(task_id=task_id, is_assistant=True).first()
        if not assistant:
            abort(404)
        conv = db.query(Conversation).filter_by(
            assistant_id=assistant.id, student_token=k, student_class=TEACHER_PREVIEW_CLASS).first()
        if not k or not conv:
            abort(403)
        resp = redirect(url_for('chat_entry', task_id=task_id))
        resp.set_cookie(_student_cookie_name(task_id), k, httponly=True, samesite='Lax')
        if _assistant_is_group(assistant):
            from urllib.parse import quote
            resp.set_cookie(_member_cookie_name(task_id), quote(conv.student_name), httponly=True, samesite='Lax')
        return resp
    finally:
        db.close()


@app.route('/api/app/v1/conversation/<int:conv_id>/messages')
@app_token_required
def api_app_conversation_messages(conv_id):
    """App 审对话：某学生会话的完整消息记录（校验归属于当前老师）。"""
    db = SessionLocal()
    try:
        conv = db.query(Conversation).filter_by(id=conv_id).first()
        if not conv:
            return jsonify({'code': 404, 'message': '会话不存在'}), 404
        assistant = db.query(Task).filter_by(id=conv.assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': '会话不存在'}), 404
        msgs = (
            db.query(Message)
            .filter_by(conversation_id=conv.id)
            .order_by(Message.id.asc())
            .all()
        )
        def _jlist(s):
            try:
                return json.loads(s) if s else []
            except (json.JSONDecodeError, TypeError):
                return []
        items = []
        for m in msgs:
            imgs = _jlist(m.image_paths)
            files = _jlist(m.file_paths)
            audios = _jlist(m.audio_paths)
            items.append({
                'role': m.role,
                'content': m.content or '',
                'triggered_keyword': m.triggered_keyword or '',
                'image_urls': ['/' + p for p in imgs],
                'audio_urls': ['/' + p for p in audios],
                'files': [{'name': _display_filename(p), 'url': '/' + p} for p in files],
                'generated_image_url': ('/' + m.generated_image_path) if m.generated_image_path else None,
                'generated_video_url': ('/' + m.generated_video_path) if m.generated_video_path else None,
                'time': m.created_at.strftime('%m-%d %H:%M') if m.created_at else '',
            })
        return jsonify({'code': 200,
                        'student_name': conv.student_name,
                        'student_class': conv.student_class,
                        'is_blocked': bool(conv.is_blocked),
                        'items': items})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/config')
@app_token_required
def api_app_assistant_config_get(assistant_id):
    """App 移动控制台 + 漫游备课：读取该导师的完整设置（开关 / 黑名单 / 提示词）。"""
    db = SessionLocal()
    try:
        a = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not a:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        return jsonify({'code': 200, 'config': {
            'title': a.title,
            'description': a.description or '',
            'system_prompt': a.system_prompt or '',
            'welcome_message': a.welcome_message or '',
            'selected_model_name': a.selected_model_name or '',
            'allow_image_input': bool(a.allow_image_input),
            'allow_file_upload': bool(a.allow_file_upload),
            'allow_voice_input': bool(a.allow_voice_input),
            'allow_tts': bool(a.allow_tts),
            'allow_image_generation': bool(a.allow_image_generation),
            'allow_video_generation': bool(a.allow_video_generation),
            'allow_call_teacher': bool(a.allow_call_teacher),
            'audio_mode': a.audio_mode or 'transcribe',
            'share_mode': a.share_mode or 'independent',
            'group_max_size': a.group_max_size or 6,
            'roster_mode': a.roster_mode or 'off',
            'max_messages_per_student_daily': a.max_messages_per_student_daily or 50,
            'blocked_keywords': _load_blocked_keywords(a.blocked_keywords),
        }})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/config', methods=['POST'])
@app_token_required
def api_app_assistant_config_set(assistant_id):
    """App：保存导师设置。只更新 body 里出现的字段（开关 / 模式 / 上限 / 黑名单 / 提示词均可单独存）。"""
    body = request.get_json(silent=True) or {}
    db = SessionLocal()
    try:
        a = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not a:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        for field in ('allow_image_input', 'allow_file_upload', 'allow_voice_input', 'allow_tts',
                      'allow_image_generation', 'allow_video_generation', 'allow_call_teacher'):
            if field in body:
                setattr(a, field, _truthy(body.get(field)))
        if 'audio_mode' in body:
            a.audio_mode = 'direct' if body.get('audio_mode') == 'direct' else 'transcribe'
        if 'share_mode' in body:
            a.share_mode = 'group' if body.get('share_mode') == 'group' else 'independent'
        if 'roster_mode' in body:
            a.roster_mode = 'strict' if body.get('roster_mode') == 'strict' else 'off'
        if 'group_max_size' in body:
            try:
                g = int(body.get('group_max_size') or 6)
            except (TypeError, ValueError):
                return jsonify({'code': 400, 'message': '小组人数上限必须是数字'}), 400
            a.group_max_size = max(2, min(20, g))
        if 'max_messages_per_student_daily' in body:
            try:
                limit = int(body.get('max_messages_per_student_daily') or 50)
            except (TypeError, ValueError):
                return jsonify({'code': 400, 'message': '每日消息上限必须是数字'}), 400
            if limit < 1 or limit > 500:
                return jsonify({'code': 400, 'message': '每日消息上限需在 1-500 之间'}), 400
            a.max_messages_per_student_daily = limit
        if 'system_prompt' in body:
            a.system_prompt = (body.get('system_prompt') or '').strip() or None
        if 'welcome_message' in body:
            a.welcome_message = (body.get('welcome_message') or '').strip() or None
        if 'blocked_keywords' in body or 'categories' in body:
            a.blocked_keywords = _mobile_blocked_keywords_json(
                a.blocked_keywords, body.get('categories') or body.get('blocked_keywords') or {})
        db.commit()
        logger.info(f'[App控制台] 导师 {a.id} 设置已更新')
        return jsonify({'code': 200, 'message': '已保存'})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/students')
@app_token_required
def api_app_assistant_students(assistant_id):
    """App：某导师的学生会话列表（含违禁触发次数与封禁状态，供逐个处置）。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        convs = (
            db.query(Conversation)
            .filter(Conversation.assistant_id == assistant.id)
            .filter(Conversation.student_class != TEACHER_PREVIEW_CLASS)
            .order_by(Conversation.last_active_at.desc())
            .all()
        )
        conv_ids = [c.id for c in convs]
        viol_counts = {}
        if conv_ids:
            rows = (
                db.query(Message.conversation_id, func.count(Message.id))
                .filter(Message.conversation_id.in_(conv_ids))
                .filter(Message.triggered_keyword.isnot(None))
                .filter(Message.triggered_keyword != '')
                .filter(Message.role == 'user')
                .group_by(Message.conversation_id)
                .all()
            )
            viol_counts = dict(rows)
        return jsonify({'code': 200, 'items': [{
            'conversation_id': c.id,
            'student_class': c.student_class,
            'student_name': c.student_name,
            'is_blocked': bool(c.is_blocked),
            'violation_count': viol_counts.get(c.id, 0),
            'message_count': c.daily_message_count or 0,
            'last_active_at': c.last_active_at.strftime('%m-%d %H:%M') if c.last_active_at else '',
        } for c in convs]})
    finally:
        db.close()


@app.route('/api/chat/<string:task_id>/call_teacher', methods=['POST'])
def api_chat_call_teacher(task_id):
    """学生「找老师」按钮：发事件给老师 App（开关 + 每生 60 秒限频）。"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']
        if not assistant.allow_call_teacher:
            return jsonify({'code': 403, 'message': '该 AI 导师未开启「找老师」功能'}), 403
        member = _get_group_member_name(task_id) if _assistant_is_group(assistant) else ''
        rate_key = (conv.id, member or conv.student_name)
        now = time.time()
        with _APP_LOCK:
            last = _CALL_TEACHER_LAST.get(rate_key, 0)
            if now - last < 60:
                return jsonify({'code': 429, 'message': '刚刚已经叫过老师啦，请稍等一下'}), 429
            _CALL_TEACHER_LAST[rate_key] = now
        # 最近一条学生消息作为上下文
        last_msg = (
            db.query(Message)
            .filter_by(conversation_id=conv.id, role='user')
            .order_by(Message.id.desc())
            .first()
        )
        _emit_app_event(assistant.user_id, 'call_teacher', {
            'assistant_id': assistant.id,
            'assistant_title': assistant.title,
            'conversation_id': conv.id,
            'student_class': conv.student_class,
            'student_name': member or conv.student_name,
            'last_question': (last_msg.content[:120] if last_msg else ''),
        })
        return jsonify({'code': 200, 'message': '老师已收到你的求助信号'})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/knowledge')
@app_token_required
def api_app_knowledge_list(assistant_id):
    """App：某导师的知识库文档列表。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        docs = (
            db.query(KnowledgeDoc)
            .filter_by(assistant_id=assistant.id)
            .order_by(KnowledgeDoc.created_at.desc())
            .all()
        )
        return jsonify({'code': 200, 'kb_enabled': bool(_get_siliconflow_key(db, assistant)), 'items': [{
            'id': d.id,
            'file_name': d.file_name,
            'chunk_count': d.chunk_count or 0,
            'created_at': d.created_at.strftime('%m-%d %H:%M') if d.created_at else '',
        } for d in docs]})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/knowledge/text', methods=['POST'])
@app_token_required
def api_app_knowledge_text(assistant_id):
    """App 拍照入知识库：端侧 OCR 后的纯文本入库（复用切片/向量化管线）。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        sf_key = _get_siliconflow_key(db, assistant)
        if not sf_key:
            return jsonify({'code': 403, 'message': '知识库需要硅基流动 API Key（请在电脑端个人设置配置）'}), 403
        body = request.get_json(silent=True) or {}
        title = (body.get('title') or '').strip()[:100] or '拍照资料'
        text = (body.get('text') or '').strip()
        if len(text) < 20:
            return jsonify({'code': 400, 'message': '识别文本太短（至少 20 字），请检查拍摄内容'}), 400
        if len(text) > KB_MAX_DOC_CHARS:
            text = text[:KB_MAX_DOC_CHARS]

        # 包装成 txt 文件走现有入库管线
        from io import BytesIO
        from werkzeug.datastructures import FileStorage
        fs = FileStorage(stream=BytesIO(text.encode('utf-8')), filename=f'{title}.txt')
        doc, err = _kb_ingest_document(db, assistant, fs, sf_key)
        if err:
            db.rollback()
            return jsonify({'code': 500, 'message': f'入库失败：{err}'}), 500
        logger.info(f'[App知识库] 文本入库 assistant={assistant.id} title={title!r} chunks={doc.chunk_count}')
        return jsonify({'code': 200, 'doc_id': doc.id, 'chunk_count': doc.chunk_count})
    finally:
        db.close()


@app.route('/api/app/v1/assistant/<int:assistant_id>/knowledge/<int:doc_id>', methods=['DELETE'])
@app_token_required
def api_app_knowledge_delete(assistant_id, doc_id):
    """App：删除知识库文档。"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=request.om_app_user_id, is_assistant=True).first()
        if not assistant:
            return jsonify({'code': 404, 'message': 'AI 导师不存在'}), 404
        doc = db.query(KnowledgeDoc).filter_by(id=doc_id, assistant_id=assistant.id).first()
        if not doc:
            return jsonify({'code': 404, 'message': '资料不存在'}), 404
        if doc.file_path and os.path.exists(doc.file_path):
            try:
                os.remove(doc.file_path)
            except OSError:
                pass
        db.delete(doc)
        db.commit()
        return jsonify({'code': 200})
    finally:
        db.close()


# ------------------------------------------------------------
# OpenMentor: 视频演示列表（飞书多维表格维护，导航栏电视按钮浮窗播放）
# ------------------------------------------------------------

DEMO_VIDEOS_BASE_TOKEN = (os.getenv('LARK_DEMO_BASE_TOKEN') or 'SYHRbP4FxayxmzsbK0wcsDMqngb').strip()
DEMO_VIDEOS_TABLE_ID = (os.getenv('LARK_DEMO_TABLE_ID') or 'tbl7K4oOGoX6TLzQ').strip()
DEMO_VIDEOS_VIEW_ID = (os.getenv('LARK_DEMO_VIEW_ID') or 'vewEjlUMcl').strip()
DEMO_VIDEOS_CACHE_PATH = os.path.join('data', 'demo_videos_cache.json')
DEMO_VIDEOS_CACHE_TTL_SECONDS = 600
_DEMO_VIDEOS_MEMORY_CACHE = {'ts': 0, 'items': None}


def _fetch_demo_videos():
    """从飞书多维表格读取视频列表（标题 + BV号），按「文本」列数字排序。"""
    tenant_token = _get_lark_tenant_access_token()
    headers = {'Authorization': f'Bearer {tenant_token}'}
    url = f'https://open.feishu.cn/open-apis/bitable/v1/apps/{DEMO_VIDEOS_BASE_TOKEN}/tables/{DEMO_VIDEOS_TABLE_ID}/records'
    params = {'page_size': 100}
    if DEMO_VIDEOS_VIEW_ID:
        params['view_id'] = DEMO_VIDEOS_VIEW_ID
    items = []
    page_token = ''
    while True:
        if page_token:
            params['page_token'] = page_token
        resp = requests.get(url, headers=headers, params=params, timeout=10)
        resp.raise_for_status()
        payload = resp.json()
        if payload.get('code') not in (0, None):
            raise RuntimeError(payload.get('msg') or '读取视频列表失败')
        data = payload.get('data') or {}
        for it in data.get('items') or []:
            fields = it.get('fields') or {}
            title = _prompt_matrix_cell_text(fields.get('标题'))
            bv = _prompt_matrix_cell_text(fields.get('BV号'))
            if not title or not bv:
                continue
            try:
                sort = float(_prompt_matrix_cell_text(fields.get('文本')) or 9999)
            except ValueError:
                sort = 9999
            items.append({'title': title, 'bv': bv, 'sort': sort})
        if not data.get('has_more'):
            break
        page_token = data.get('page_token') or ''
        if not page_token:
            break
    items.sort(key=lambda x: x['sort'])
    return [{'title': i['title'], 'bv': i['bv']} for i in items]


def _load_demo_videos():
    """内存缓存 → 飞书同步 → 文件缓存兜底（分发版无飞书凭证时只剩文件缓存/空列表）。"""
    now = time.time()
    if _DEMO_VIDEOS_MEMORY_CACHE['items'] is not None and now - _DEMO_VIDEOS_MEMORY_CACHE['ts'] < DEMO_VIDEOS_CACHE_TTL_SECONDS:
        return _DEMO_VIDEOS_MEMORY_CACHE['items']
    items = None
    if all((os.getenv(k) or '').strip() for k in ('LARK_APP_ID', 'LARK_APP_SECRET')):
        try:
            items = _fetch_demo_videos()
            try:
                os.makedirs(os.path.dirname(DEMO_VIDEOS_CACHE_PATH), exist_ok=True)
                with open(DEMO_VIDEOS_CACHE_PATH, 'w', encoding='utf-8') as f:
                    json.dump({'updated_at': datetime.now().isoformat(), 'items': items}, f, ensure_ascii=False, indent=2)
            except Exception as e:
                logger.warning(f'[视频演示] 写缓存失败: {e}')
        except Exception as e:
            logger.warning(f'[视频演示] 同步失败: {e}')
    if items is None:
        try:
            with open(DEMO_VIDEOS_CACHE_PATH, 'r', encoding='utf-8') as f:
                items = (json.load(f) or {}).get('items') or []
        except Exception:
            items = []
    _DEMO_VIDEOS_MEMORY_CACHE.update({'ts': now, 'items': items})
    return items


@app.route('/api/demo_videos')
@login_required
def api_demo_videos():
    return jsonify({'items': _load_demo_videos()})


# ------------------------------------------------------------
# OpenMentor: AI 导师知识库（RAG）管理
# ------------------------------------------------------------

@app.route('/assistant/<int:assistant_id>/knowledge/upload', methods=['POST'])
@login_required
def assistant_knowledge_upload(assistant_id):
    """老师上传知识库文档（同步处理：提取 → 切片 → 向量化 → 入库）"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))

        sf_key = _get_siliconflow_key(db, assistant)
        if not sf_key:
            flash('知识库需要硅基流动 API Key 做向量化（免费），请先在「个人设置」中配置', 'warning')
            return redirect(url_for('assistant_detail', assistant_id=assistant.id))

        files = [f for f in request.files.getlist('kb_files') if f and f.filename]
        if not files:
            flash('请选择要上传的资料文件', 'warning')
            return redirect(url_for('assistant_detail', assistant_id=assistant.id))

        ok_count = 0
        for f in files:
            doc, err = _kb_ingest_document(db, assistant, f, sf_key)
            if err:
                db.rollback()
                flash(f'「{f.filename}」入库失败：{err}', 'danger')
            else:
                ok_count += 1
                flash(f'「{f.filename}」已入库（切成 {doc.chunk_count} 个知识片段）', 'success')
        if ok_count:
            logger.info(f'[知识库] 老师 {current_user.id} 为导师 {assistant.id} 上传 {ok_count} 份资料')
        return redirect(url_for('assistant_detail', assistant_id=assistant.id))
    finally:
        db.close()


@app.route('/assistant/<int:assistant_id>/knowledge/<int:doc_id>/delete', methods=['POST'])
@login_required
def assistant_knowledge_delete(assistant_id, doc_id):
    """删除一份知识库文档（含切片与原件）"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(id=assistant_id, user_id=current_user.id, is_assistant=True).first()
        if not assistant:
            flash('AI 导师不存在或无权访问', 'danger')
            return redirect(url_for('assistant_list'))
        doc = db.query(KnowledgeDoc).filter_by(id=doc_id, assistant_id=assistant.id).first()
        if not doc:
            flash('资料不存在', 'warning')
            return redirect(url_for('assistant_detail', assistant_id=assistant.id))
        name = doc.file_name
        if doc.file_path and os.path.exists(doc.file_path):
            try:
                os.remove(doc.file_path)
            except OSError:
                pass
        db.delete(doc)  # cascade 删除 chunks
        db.commit()
        flash(f'已删除资料「{name}」', 'success')
        return redirect(url_for('assistant_detail', assistant_id=assistant.id))
    finally:
        db.close()


# ------------------------------------------------------------
# OpenMentor: AI 辅助生成系统提示词（Week 3 重点功能）
# ------------------------------------------------------------

PROMPT_OPTIMIZE_TEMPLATE = """你是一位资深的"教育 AI 提示词工程师"。老师已经在用一个 system prompt，希望你针对性地优化它。

【老师当前的 system prompt】
{current_prompt}

【老师希望优化的方向】
{optimization_goals}

请按以下原则优化：
1. **保留老师原 prompt 中能正常工作、有价值的部分**（角色身份、已有的具体规则、风格基调等），不要全盘推翻
2. 针对老师指出的优化方向做**精准修改 / 补充 / 调整**，避免画蛇添足改无关部分
3. 如果老师的优化目标本身可能违背教学法（例如"让 AI 直接给学生答案"），可以在保留原意的同时给出更教学友好的实现方式（比如改成"先引导后给答案"）
4. **输出完整的、可直接采用的新版 system prompt**（不要输出 diff 或修改对比）
5. 总长度控制在 800 字以内，简洁有力
6. 直接输出 system prompt 正文（不要"以下是优化后版本"之类的开场白，也不要 Markdown 代码块包裹）
7. 用中文撰写
"""


PROMPT_META_TEMPLATE = """你是一位资深的"教育 AI 提示词工程师"，擅长为中国中小学/高校老师设计专业、可落地的 AI 助教系统提示词（system prompt）。

【老师对 AI 助教的需求】
- 角色定位：{role}
- 目标学生：{audience}
- 期望行为：{dos}
- 禁止行为：{donts}
- 语气与风格：{tone}
{material_section}{draft_section}
请为这位老师设计一份**system prompt 草稿**，要求：
1. 开篇明确"你是…"角色身份，最好有简短的教学法定位（如苏格拉底式提问、支架式教学等）
2. 列出 3-5 条具体可执行的行为准则（如"先引导思考再讲答案"、"用比喻代替专业术语"）
3. 至少包含 1 个示范对话片段（学生提问 + 你的回应方式）
4. 给出明确的边界（不答与学习无关的内容、保护未成年人、不评论政治宗教等）
5. **如果老师上传了素材**：
   - 先在脑内提炼素材里**最关键的 3-8 条教学内容**（核心概念 / 易错点 / 关键例子 / 学习路径）
   - 把这些要点**精炼地**融入系统提示词（用一段简短列表或自然语言概述），让 AI 在回答时知道该引用什么、强调什么
   - **绝对不要**把素材原文整段粘贴进 system prompt（会让提示词冗长低效）
   - 总量上素材相关内容不超过最终 prompt 的 40%
6. 总长度控制在 800 字以内，简洁有力，可直接复制到 system prompt 字段使用
7. 直接输出 system prompt 正文（不要包含"以下是 system prompt:"之类的开场白，也不要 Markdown 代码块包裹）
8. 用中文撰写
"""


@app.route('/assistant/generate_prompt', methods=['POST'])
@login_required
def assistant_generate_prompt():
    """AI 辅助生成 system prompt：老师填问题/上传素材 → AI 出草稿"""
    db = SessionLocal()
    try:
        # 取问题表单
        role = (request.form.get('role') or '').strip()
        audience = (request.form.get('audience') or '').strip()
        dos = (request.form.get('dos') or '').strip()
        donts = (request.form.get('donts') or '').strip()
        tone = (request.form.get('tone') or '').strip()
        existing_draft = (request.form.get('existing_draft') or '').strip()

        # 至少需要 role 或上传文件之一
        files = request.files.getlist('files')
        has_input = bool(role or audience or dos or donts or tone) or bool(files and any(f and f.filename for f in files))
        if not has_input:
            return jsonify({'code': 400, 'message': '请至少填写"角色定位"或上传一份素材'}), 400

        # 处理上传文件：仅文档（pdf/docx/txt/md/csv），就地读取不落盘
        material_blocks = []
        for f in files or []:
            if not f or not f.filename:
                continue
            ext = _ext_of(f.filename)
            if ext not in ALLOWED_DOC_EXTS:
                continue
            f.seek(0, os.SEEK_END)
            size = f.tell()
            f.seek(0)
            if size <= 0 or size > MAX_UPLOAD_BYTES:
                continue
            # 临时落盘以复用 _extract_text_from_document
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=f'.{ext}', delete=False) as tmp:
                f.save(tmp.name)
                tmp_path = tmp.name
            try:
                text, truncated = _extract_text_from_document(tmp_path, ext)
                if text:
                    # 单文件上限 5000 字以免元提示过长
                    if len(text) > 5000:
                        text = text[:5000] + '...[已截断]'
                    material_blocks.append(f'### 素材：{f.filename}\n{text}')
            finally:
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

        material_section = ''
        if material_blocks:
            material_section = (
                '\n\n【老师上传的参考素材】\n'
                '（注意：以下是素材原文，**仅供你阅读理解**用，请你自己提炼要点后融入 system prompt，'
                '不要把整段原文复制进最终输出）\n\n'
                + '\n\n'.join(material_blocks) + '\n'
            )

        draft_section = ''
        if existing_draft:
            draft_section = f'\n\n【老师已有的草稿（请在此基础上完善）】\n{existing_draft}\n'

        meta_prompt = PROMPT_META_TEMPLATE.format(
            role=role or '（老师未指定，请基于素材推断）',
            audience=audience or '（老师未指定，请按通用情况设计）',
            dos=dos or '（老师未指定）',
            donts=donts or '（老师未指定）',
            tone=tone or '（老师未指定，建议温和、耐心、鼓励式）',
            material_section=material_section,
            draft_section=draft_section,
        )

        # 调用老师配置的 AI（用 call_ai_model 即可，非流式）
        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
        if not ai_config:
            return jsonify({'code': 500, 'message': '请先在「AI 配置」中选择并配置一个 AI 模型'}), 500

        # 找到当前选中模型的配置
        cur_model_cfg = next((mc for mc in ai_config.model_configs if mc.model_name == ai_config.selected_model), None)
        if not cur_model_cfg or not (cur_model_cfg.api_key or '').strip():
            if ai_config.selected_model != 'ollama':
                return jsonify({'code': 500, 'message': f'当前选中的 {ai_config.selected_model} 未配置 API Key'}), 500

        try:
            generated = call_ai_model(meta_prompt, ai_config)
        except Exception as e:
            logger.exception('调用 AI 生成提示词失败')
            return jsonify({'code': 500, 'message': f'调用 AI 失败：{e}'}), 500

        if not generated or not generated.strip():
            return jsonify({'code': 500, 'message': 'AI 返回空内容，请重试'}), 500

        # 清理常见前缀
        generated = generated.strip()
        for prefix in ('```\n', '```text\n', '```markdown\n', '```'):
            if generated.startswith(prefix):
                generated = generated[len(prefix):].lstrip()
        if generated.endswith('```'):
            generated = generated[:-3].rstrip()

        return jsonify({
            'code': 200,
            'success': True,
            'generated_prompt': generated,
            'model_used': ai_config.selected_model,
            'meta_prompt_length': len(meta_prompt),
            'material_count': len(material_blocks),
        })
    finally:
        db.close()


@app.route('/assistant/optimize_prompt', methods=['POST'])
@login_required
def assistant_optimize_prompt():
    """AI 优化已有 system prompt：老师给出当前 prompt + 优化目标 → AI 返回优化版"""
    db = SessionLocal()
    try:
        body = request.get_json(silent=True) or {}
        current_prompt = (body.get('current_prompt') or '').strip()
        optimization_goals = (body.get('optimization_goals') or '').strip()

        if not current_prompt:
            return jsonify({'code': 400, 'message': '请先填一段 system prompt 再来优化'}), 400
        if len(current_prompt) > 8000:
            return jsonify({'code': 400, 'message': 'prompt 过长（≤ 8000 字符）'}), 400
        if not optimization_goals:
            return jsonify({'code': 400, 'message': '请说明你希望如何优化（例：更适合小学生 / 增加互动性）'}), 400
        if len(optimization_goals) > 2000:
            return jsonify({'code': 400, 'message': '优化目标过长（≤ 2000 字符）'}), 400

        meta_prompt = PROMPT_OPTIMIZE_TEMPLATE.format(
            current_prompt=current_prompt,
            optimization_goals=optimization_goals,
        )

        ai_config = db.query(AIConfig).filter_by(user_id=current_user.id).first()
        if not ai_config:
            return jsonify({'code': 500, 'message': '请先在「个人设置」中配置一个 AI 模型'}), 500
        cur_model_cfg = next((mc for mc in ai_config.model_configs if mc.model_name == ai_config.selected_model), None)
        if not cur_model_cfg or not (cur_model_cfg.api_key or '').strip():
            if ai_config.selected_model != 'ollama':
                return jsonify({'code': 500, 'message': f'当前选中的 {ai_config.selected_model} 未配置 API Key'}), 500

        try:
            generated = call_ai_model(meta_prompt, ai_config)
        except Exception as e:
            logger.exception('调用 AI 优化提示词失败')
            return jsonify({'code': 500, 'message': f'调用 AI 失败：{e}'}), 500

        if not generated or not generated.strip():
            return jsonify({'code': 500, 'message': 'AI 返回空内容，请重试'}), 500

        # 清理常见前缀（同 generate_prompt 处理）
        generated = generated.strip()
        for prefix in ('```\n', '```text\n', '```markdown\n', '```'):
            if generated.startswith(prefix):
                generated = generated[len(prefix):].lstrip()
        if generated.endswith('```'):
            generated = generated[:-3].rstrip()

        return jsonify({
            'code': 200,
            'success': True,
            'optimized_prompt': generated,
            'model_used': ai_config.selected_model,
        })
    finally:
        db.close()


# ------------------------------------------------------------
# OpenMentor: 学生端聊天 + SSE 流式（Day 5-7）
# ------------------------------------------------------------

# OpenAI 兼容协议端点（用于流式调用）
MODEL_OPENAI_COMPATIBLE = {
    'deepseek':    {'url': 'https://api.deepseek.com/v1/chat/completions',                       'default_model': 'deepseek-chat'},
    'doubao':      {'url': 'https://ark.cn-beijing.volces.com/api/v3/chat/completions',           'default_model': 'doubao-seed-1-6-251015'},
    'qwen':        {'url': 'https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions', 'default_model': 'qwen-plus'},
    'glm':         {'url': 'https://open.bigmodel.cn/api/paas/v4/chat/completions',               'default_model': 'glm-4'},
    'siliconflow': {'url': 'https://api.siliconflow.cn/v1/chat/completions',                      'default_model': 'Qwen/Qwen2.5-72B-Instruct'},
}


TEACHER_PREVIEW_CLASS = '教师试聊'  # 老师试聊会话的班级标记：所有统计/审计列表均排除


def _student_cookie_name(task_id):
    return f'om_student_{task_id}'


def _member_cookie_name(task_id):
    """小组模式下记录组员本人姓名的 cookie（组共享会话 token 之外的个人身份）"""
    return f'om_member_{task_id}'


def _assistant_is_group(assistant):
    return (assistant.share_mode or 'independent') == 'group'


def _get_group_member_name(task_id):
    """读取小组模式组员姓名 cookie（URL 编码存储）"""
    from urllib.parse import unquote
    raw = request.cookies.get(_member_cookie_name(task_id)) or ''
    try:
        return unquote(raw)[:100]
    except Exception:
        return ''


def _hash_student_token(student_class, student_name, device_uuid=None):
    """生成会话 token。
    - 不带 device_uuid（旧行为）：同一身份在所有设备共享 token → 共享会话
    - 带 device_uuid（B8 设备隔离）：(device, class, name) 三元组定位 → 互不可见
    """
    import hashlib
    if device_uuid:
        seed = f'{device_uuid}|{student_class}|{student_name}'
    else:
        seed = f'{student_class}|{student_name}'
    return hashlib.sha256(seed.encode('utf-8')).hexdigest()[:32]


def _ensure_daily_quota(conv):
    """检查并刷新 conversation 的每日额度。返回 (已用条数, 上限)。"""
    today = datetime.now().date()
    if conv.daily_reset_date != today:
        conv.daily_message_count = 0
        conv.daily_reset_date = today
    return conv.daily_message_count or 0, conv.assistant.max_messages_per_student_daily or 50


def _check_blocked_keywords(text, blocked_keywords_json):
    """检查文本是否命中黑名单。命中返回 (category_label, keyword)，否则 None。"""
    if not blocked_keywords_json:
        return None
    try:
        config = json.loads(blocked_keywords_json)
    except (json.JSONDecodeError, TypeError):
        return None
    text_lower = text.lower()
    for cat_key, cat in config.items():
        if not cat.get('enabled'):
            continue
        for kw in cat.get('keywords', []):
            if kw and kw.lower() in text_lower:
                return cat.get('label', cat_key), kw
    return None


def _get_assistant_model_config(db, assistant):
    """根据 assistant.selected_model_name 获取老师配置的对应模型。
    返回 (model_key, AIModelConfig 实例) 或 (None, None)。
    """
    ai_config = db.query(AIConfig).filter_by(user_id=assistant.user_id).first()
    if not ai_config:
        return None, None
    model_key = assistant.selected_model_name or ai_config.selected_model or 'deepseek'
    for mc in ai_config.model_configs:
        if mc.model_name == model_key:
            return model_key, mc
    return model_key, None


def _build_history_for_llm(messages, system_prompt, max_history=20, supports_vision=False):
    """把 Message 记录转换为 OpenAI 格式的 messages 数组。
    - 跳过 system 角色合成提示（如关键词拦截通知）
    - 保留最近 max_history 条用户/AI 消息
    - 视觉模型场景下，把图片附件转为 image_url(content array)
    - 文件附件的提取文本已在保存时拼到 content，本函数仅展开图片
    """
    msgs = []
    if system_prompt:
        msgs.append({'role': 'system', 'content': system_prompt})
    chat_msgs = [m for m in messages if m.role in ('user', 'assistant') and not m.triggered_keyword]
    chat_msgs = chat_msgs[-max_history:]
    for m in chat_msgs:
        # 默认纯文本
        item = {'role': m.role, 'content': m.content}
        # 仅 user 消息可能带图片；assistant 消息原样
        if m.role == 'user' and supports_vision and m.image_paths:
            try:
                paths = json.loads(m.image_paths)
            except (json.JSONDecodeError, TypeError):
                paths = []
            if paths:
                content_array = [{'type': 'text', 'text': m.content or ''}]
                for p in paths:
                    data_url = _image_to_data_url(p)
                    if data_url:
                        content_array.append({'type': 'image_url', 'image_url': {'url': data_url}})
                item['content'] = content_array
        msgs.append(item)
    return msgs


def _image_to_data_url(rel_path):
    """把相对路径的图片转成 data URL（base64 内嵌）。失败返回 None。"""
    try:
        # rel_path 形如 'static/uploads/openmentor/<conv_id>/xxx.png'
        if not os.path.exists(rel_path):
            return None
        ext = os.path.splitext(rel_path)[1].lower().lstrip('.')
        mime_map = {'jpg': 'jpeg', 'jpeg': 'jpeg', 'png': 'png', 'gif': 'gif', 'webp': 'webp'}
        mime = mime_map.get(ext, 'png')
        with open(rel_path, 'rb') as f:
            data = base64.b64encode(f.read()).decode('ascii')
        return f'data:image/{mime};base64,{data}'
    except Exception as e:
        logger.warning(f'图片转 data URL 失败 ({rel_path}): {e}')
        return None


# ==================== 语音输入（ASR 转写） ====================
# 路线：学生录音 → 上传音频 → 服务器调硅基流动 SenseVoice 转文字 → 文字走现有对话链
# （不依赖多模态模型，任何对话模型都兼容；老师只需配置 siliconflow API Key）

ASR_MODEL = 'FunAudioLLM/SenseVoiceSmall'
ASR_URL = 'https://api.siliconflow.cn/v1/audio/transcriptions'

# 录音桥镜像列表：纯静态录音页（static/voice_bridge.html）部署到公网 HTTPS 后填到这里。
# 学生端在 HTTP 局域网页面点麦克风时，会弹出可达的镜像完成录音（postMessage 传回音频）。
# 列表顺序即优先级（前端探测全部镜像后，按本顺序取第一个可达的）。
# 可用环境变量 OM_VOICE_BRIDGE_URLS 覆盖（逗号分隔多个地址）。
VOICE_BRIDGE_DEFAULT_MIRRORS = [
    'https://openmentor.top/voice_bridge',
    'https://xianjianaiedu.com/tools/voice_bridge',
    'https://xianjianaiedu.top/tools/voice_bridge',
]


def _voice_bridge_urls():
    env = (os.environ.get('OM_VOICE_BRIDGE_URLS') or '').strip()
    if env:
        return [u.strip() for u in env.split(',') if u.strip()]
    return list(VOICE_BRIDGE_DEFAULT_MIRRORS)


def _get_siliconflow_key(db, assistant):
    """该老师配置的硅基流动 API Key（语音转写 + 知识库向量化共用）。未配置返回 None。"""
    ai_config = db.query(AIConfig).filter_by(user_id=assistant.user_id).first()
    if not ai_config:
        return None
    for mc in ai_config.model_configs:
        if mc.model_name == 'siliconflow' and (mc.api_key or '').strip():
            return mc.api_key.strip()
    return None


def _transcribe_audio(file_path, api_key):
    """调硅基流动 SenseVoice 把音频转成文字。失败抛异常，成功返回文本（可能为空串）。"""
    ext = os.path.splitext(file_path)[1].lower().lstrip('.')
    mime_map = {
        'wav': 'audio/wav', 'mp3': 'audio/mpeg', 'm4a': 'audio/mp4', 'aac': 'audio/aac',
        'ogg': 'audio/ogg', 'opus': 'audio/opus', 'webm': 'audio/webm', 'amr': 'audio/amr',
        'flac': 'audio/flac',
    }
    mime = mime_map.get(ext, 'application/octet-stream')
    with open(file_path, 'rb') as f:
        resp = requests.post(
            ASR_URL,
            headers={'Authorization': f'Bearer {api_key}'},
            files={'file': (os.path.basename(file_path), f, mime)},
            data={'model': ASR_MODEL},
            timeout=60,
        )
    if resp.status_code != 200:
        err_text = (resp.text or '')[:300]
        raise RuntimeError(f'语音识别接口返回 {resp.status_code}: {err_text}')
    try:
        text_out = (resp.json().get('text') or '').strip()
    except ValueError:
        raise RuntimeError('语音识别接口返回了无法解析的内容')
    logger.info(f'[ASR] {os.path.basename(file_path)} → {len(text_out)} 字: {text_out[:80]!r}')
    return text_out


# ==================== 原声直达模型（音频理解） ====================
# audio_mode='direct'：学生音频不经转写、以 input_audio 直接发给多模态模型
# （适合朗读评分/口语练习等"声音本身就是内容"的场景）。
# 实测支持：通义 qwen3-omni 系列、豆包 doubao-seed-2.0 系列；老模型不支持时 API 会明确报错。
# 有 siliconflow key 时仍后台静默转写一份：用于审计页文字显示 + 多轮历史（历史不重复塞音频，省 token）。

AUDIO_DIRECT_MAX_BYTES = 4 * 1024 * 1024  # 单条消息音频总大小上限（base64 后约 5.3MB，约 2-4 分钟）


def _assistant_audio_direct(assistant):
    return (getattr(assistant, 'audio_mode', 'transcribe') or 'transcribe') == 'direct'


def _build_input_audio_parts(paths, provider_key):
    """把音频文件转成 OpenAI 兼容的 input_audio content parts。
    通义要求 data:;base64, 前缀；豆包等用裸 base64。"""
    parts = []
    for p in paths:
        ext = _ext_of(p)
        fmt = ext if ext in ALLOWED_AUDIO_EXTS else 'wav'
        with open(p, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode('ascii')
        data = f'data:;base64,{b64}' if provider_key == 'qwen' else b64
        parts.append({'type': 'input_audio', 'input_audio': {'data': data, 'format': fmt}})
    return parts


# ==================== AI 导师知识库（RAG） ====================
# 路线：老师传资料 → 切片 → 硅基流动 bge-m3 向量化（免费）→ SQLite 存 float32 BLOB
#       学生提问 → 问题向量化 → 暴力余弦取最相关片段 → 拼进 system prompt
# 课堂规模（每导师几百片段）暴力检索足够，不引入向量数据库、不新增 pip 依赖。

KB_EMBEDDING_MODEL = 'BAAI/bge-m3'
KB_EMBEDDING_URL = 'https://api.siliconflow.cn/v1/embeddings'
KB_UPLOAD_ROOT = os.path.join('static', 'uploads', 'knowledge')
KB_MAX_DOC_CHARS = 200_000   # 单文档最多取前 20 万字
KB_CHUNK_SIZE = 600          # 单片目标长度（字符）
KB_CHUNK_OVERLAP = 100       # 超长段落硬切时的重叠
KB_TOP_K = 4                 # 每次检索取最相关的片段数
KB_MIN_SIMILARITY = 0.45     # 低于该余弦相似度的片段不注入（bge-m3 实测：相关 0.53+，无关 ≤0.43）
KB_MAX_CONTEXT_CHARS = 3000  # 注入上下文的总字数上限


def _embed_texts(texts, api_key, batch_size=16):
    """调硅基流动 bge-m3 批量向量化。返回与输入等长的向量列表，失败抛异常。"""
    vectors = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        resp = requests.post(
            KB_EMBEDDING_URL,
            headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
            json={'model': KB_EMBEDDING_MODEL, 'input': batch},
            timeout=60,
        )
        if resp.status_code != 200:
            raise RuntimeError(f'向量化接口返回 {resp.status_code}: {(resp.text or "")[:200]}')
        data = sorted(resp.json().get('data') or [], key=lambda d: d.get('index', 0))
        if len(data) != len(batch):
            raise RuntimeError('向量化接口返回数量与输入不一致')
        vectors.extend(d['embedding'] for d in data)
    return vectors


def _pack_embedding(vec):
    from array import array
    return array('f', vec).tobytes()


def _unpack_embedding(blob):
    from array import array
    a = array('f')
    a.frombytes(blob)
    return a


def _cosine_similarity(a, b):
    dot = s_a = s_b = 0.0
    for x, y in zip(a, b):
        dot += x * y
        s_a += x * x
        s_b += y * y
    if s_a <= 0 or s_b <= 0:
        return 0.0
    return dot / ((s_a ** 0.5) * (s_b ** 0.5))


def _kb_chunk_text(text, chunk_size=KB_CHUNK_SIZE, overlap=KB_CHUNK_OVERLAP):
    """按段落累积切片；超长段落滑窗硬切。返回片段列表（过滤 < 20 字的碎片）。"""
    text = re.sub(r'\r\n?', '\n', text or '')
    text = re.sub(r'[ \t]+', ' ', text)
    paras = [p.strip() for p in text.split('\n') if p.strip()]
    chunks = []
    buf = ''
    for p in paras:
        if len(buf) + len(p) + 1 <= chunk_size:
            buf = (buf + '\n' + p) if buf else p
            continue
        if buf:
            chunks.append(buf)
            buf = ''
        while len(p) > chunk_size:
            chunks.append(p[:chunk_size])
            p = p[chunk_size - overlap:]
        buf = p
    if buf:
        chunks.append(buf)
    return [c for c in chunks if len(c.strip()) >= 20]


def _kb_ingest_document(db, assistant, file_storage, api_key):
    """处理一份上传的知识库文档：落盘 → 提取文本 → 切片 → 向量化 → 入库。
    返回 (KnowledgeDoc, None) 或 (None, 错误信息)。"""
    from werkzeug.utils import secure_filename as _sec
    ext = _ext_of(file_storage.filename)
    if ext not in ALLOWED_DOC_EXTS:
        return None, f'不支持的文件类型 .{ext}（支持 PDF/Word/TXT/Markdown 等文档）'
    file_storage.seek(0, os.SEEK_END)
    size = file_storage.tell()
    file_storage.seek(0)
    if size <= 0:
        return None, '空文件'
    if size > MAX_UPLOAD_BYTES:
        return None, f'文件过大（≤ {MAX_UPLOAD_BYTES // (1024 * 1024)} MB）'

    doc_dir = os.path.join(KB_UPLOAD_ROOT, str(assistant.id))
    os.makedirs(doc_dir, exist_ok=True)
    safe_base = _sec(os.path.splitext(file_storage.filename)[0]) or 'doc'
    save_path = os.path.join(doc_dir, f'{uuid.uuid4().hex[:8]}_{safe_base}.{ext}')
    file_storage.save(save_path)

    try:
        extracted, _truncated = _extract_text_from_document(save_path, ext)
        if not extracted or len(extracted.strip()) < 20:
            raise RuntimeError('没能从文件里提取到文本（扫描版 PDF 暂不支持，请用文字版）')
        chunks = _kb_chunk_text(extracted[:KB_MAX_DOC_CHARS])
        if not chunks:
            raise RuntimeError('文本切片后为空')
        vectors = _embed_texts(chunks, api_key)
    except Exception as e:
        try:
            os.remove(save_path)
        except OSError:
            pass
        return None, str(e)

    doc = KnowledgeDoc(
        assistant_id=assistant.id,
        file_name=file_storage.filename,
        file_path=save_path.replace('\\', '/'),
        status='ready',
        chunk_count=len(chunks),
    )
    db.add(doc)
    db.flush()
    for idx, (content, vec) in enumerate(zip(chunks, vectors)):
        db.add(KnowledgeChunk(
            doc_id=doc.id,
            assistant_id=assistant.id,
            chunk_index=idx,
            content=content,
            embedding=_pack_embedding(vec),
        ))
    db.commit()
    logger.info(f'[知识库] 入库 assistant={assistant.id} file={file_storage.filename!r} chunks={len(chunks)}')
    return doc, None


def _kb_retrieve(db, assistant, query, api_key):
    """检索与学生问题最相关的知识片段，返回可拼进 system prompt 的文本块（无结果返回 ''）。"""
    query = (query or '').strip()
    if not query:
        return ''
    rows = (
        db.query(KnowledgeChunk.content, KnowledgeChunk.embedding)
        .filter(KnowledgeChunk.assistant_id == assistant.id)
        .all()
    )
    if not rows:
        return ''
    qv = _embed_texts([query[:2000]], api_key)[0]
    scored = []
    for content, blob in rows:
        if not blob:
            continue
        sim = _cosine_similarity(qv, _unpack_embedding(blob))
        if sim >= KB_MIN_SIMILARITY:
            scored.append((sim, content))
    if not scored:
        return ''
    scored.sort(key=lambda x: x[0], reverse=True)
    picked = []
    total = 0
    for sim, content in scored[:KB_TOP_K]:
        if total + len(content) > KB_MAX_CONTEXT_CHARS:
            break
        picked.append(content)
        total += len(content)
    if not picked:
        return ''
    joined = '\n---\n'.join(picked)
    return (
        '\n\n【知识库参考资料】\n'
        '以下是老师为你准备的参考资料片段（按与学生问题的相关度排序）。'
        '回答时请优先依据这些资料；资料没有覆盖的内容再用你自己的知识补充，不要编造资料里没有的细节：\n'
        + joined
    )


# ==================== AI 回复朗读（TTS） ====================
# 主力 Edge TTS（免费，微软在线音色）；前端在服务端失败时自动降级浏览器 speechSynthesis。
# 音频只作为性能缓存（内容哈希命名，不进数据库、不算正式数据），启动时清理过期文件。

TTS_CACHE_DIR = os.path.join('static', 'tts_cache')
TTS_CACHE_MAX_AGE_DAYS = 7
TTS_MAX_CHARS = 1000
TTS_DEFAULT_VOICE = 'zh-CN-XiaoxiaoNeural'


def _tts_available():
    try:
        import edge_tts  # noqa: F401
        return True
    except ImportError:
        return False


# emoji / 颜文字符号区段（朗读前剔除，否则引擎会把表情念出来）
_TTS_EMOJI_RE = re.compile(
    '['
    '\U0001F000-\U0001FAFF'   # 表情、符号、交通工具、补充表情等
    '\U00002600-\U000027BF'   # 杂项符号 + 装饰符号（☀✨✅❌ 等）
    '\U0001F1E6-\U0001F1FF'   # 区域旗帜
    '\U00002B00-\U00002BFF'   # 星形等（⭐）
    '\U0000FE00-\U0000FE0F'   # 变体选择符
    '\U0000200D'              # 零宽连接符（组合 emoji）
    '\U00002049\U0000203C'    # ⁉ ‼
    '\U00003030\U0000303D'    # 〰 〽
    '\U000024C2\U0001F004'    # Ⓜ 🀄
    ']+'
)


def _tts_clean_text(text):
    """把 markdown 回复洗成适合朗读的纯文本（去代码块/公式/链接/格式符号/emoji，超长截断）。"""
    text = text or ''
    text = re.sub(r'```.*?```', '。这里有一段代码，请看文字。', text, flags=re.S)
    text = re.sub(r'\$\$.*?\$\$', '。这里有一个公式，请看文字。', text, flags=re.S)
    text = re.sub(r'\$[^$\n]{1,80}\$', '（公式）', text)
    text = re.sub(r'`([^`]*)`', r'\1', text)
    text = re.sub(r'!\[[^\]]*\]\([^)]*\)', '', text)          # 图片
    text = re.sub(r'\[([^\]]*)\]\([^)]*\)', r'\1', text)      # 链接保留文字
    text = re.sub(r'^#{1,6}\s*', '', text, flags=re.M)         # 标题井号
    text = re.sub(r'^\s*[-+*]\s+', '', text, flags=re.M)       # 列表符号
    text = re.sub(r'^\s*>\s?', '', text, flags=re.M)           # 引用
    text = re.sub(r'[*_~|#]', '', text)                        # 余下格式符号
    text = _TTS_EMOJI_RE.sub('', text)                         # emoji 表情
    text = re.sub(r'\n{2,}', '\n', text).strip()
    if len(text) > TTS_MAX_CHARS:
        text = text[:TTS_MAX_CHARS] + '。内容比较长，剩下的部分请阅读文字哦。'
    return text


def _tts_synthesize_edge(text, voice=TTS_DEFAULT_VOICE):
    """调 Edge TTS 合成 mp3 字节。失败抛异常。"""
    import asyncio
    import edge_tts

    async def _run():
        communicate = edge_tts.Communicate(text, voice)
        data = b''
        async for chunk in communicate.stream():
            if chunk.get('type') == 'audio':
                data += chunk['data']
        return data

    audio = asyncio.run(_run())
    if len(audio) < 1024:
        raise RuntimeError('TTS 返回音频过小，疑似合成失败')
    return audio


def _tts_cleanup_cache():
    """清理超过 TTS_CACHE_MAX_AGE_DAYS 的缓存音频（启动时调用一次）。"""
    try:
        if not os.path.isdir(TTS_CACHE_DIR):
            return
        cutoff = time.time() - TTS_CACHE_MAX_AGE_DAYS * 86400
        removed = 0
        for name in os.listdir(TTS_CACHE_DIR):
            p = os.path.join(TTS_CACHE_DIR, name)
            if os.path.isfile(p) and os.path.getmtime(p) < cutoff:
                os.remove(p)
                removed += 1
        if removed:
            logger.info(f'[TTS] 清理过期缓存音频 {removed} 个')
    except Exception as e:
        logger.warning(f'[TTS] 缓存清理失败: {e}')


_tts_cleanup_cache()


def _stream_openai_compatible(url, headers, payload):
    """以 OpenAI 兼容协议发起流式请求；逐 chunk yield 内容。"""
    # 诊断日志：记录消息结构（隐藏 base64 内容）
    try:
        debug_msgs = []
        for m in payload.get('messages', []):
            content = m.get('content')
            if isinstance(content, list):
                parts = []
                for c in content:
                    if c.get('type') == 'image_url':
                        url_str = c.get('image_url', {}).get('url', '')
                        if url_str.startswith('data:'):
                            parts.append(f'<image_data_url len={len(url_str)}>')
                        else:
                            parts.append(f'<image_url {url_str[:60]}>')
                    elif c.get('type') == 'text':
                        text = c.get('text', '')
                        parts.append(f'<text:{text[:50]}{"..." if len(text)>50 else ""}>')
                    else:
                        parts.append(f'<{c.get("type", "?")}>')
                debug_msgs.append(f'{m["role"]}=[{", ".join(parts)}]')
            else:
                txt = str(content or '')
                debug_msgs.append(f'{m["role"]}={txt[:80]}{"..." if len(txt)>80 else ""}')
        logger.info(f'[LLM 请求] model={payload.get("model")}, url={url}, messages: {" | ".join(debug_msgs)}')
    except Exception:
        pass

    try:
        with requests.post(url, headers=headers, json=payload, stream=True, timeout=180) as resp:
            if resp.status_code != 200:
                err_text = resp.text[:500] if resp.text else ''
                raise RuntimeError(f'API 返回 {resp.status_code}: {err_text}')
            for raw_line in resp.iter_lines():
                if not raw_line:
                    continue
                line = raw_line.decode('utf-8', errors='ignore')
                if line.startswith('data: '):
                    data_str = line[6:].strip()
                    if data_str == '[DONE]':
                        break
                    try:
                        data = json.loads(data_str)
                    except json.JSONDecodeError:
                        continue
                    choices = data.get('choices', [])
                    if not choices:
                        continue
                    delta = choices[0].get('delta', {})
                    chunk = delta.get('content')
                    if chunk:
                        yield chunk
    except requests.exceptions.RequestException as e:
        raise RuntimeError(f'网络请求异常: {e}')


def _stream_assistant_response(model_key, model_cfg, messages, extra_payload=None):
    """分发到对应模型的流式接口。yield 文本 chunk。extra_payload 合并进请求体（如 qwen-omni 的 modalities）。"""
    if model_key in MODEL_OPENAI_COMPATIBLE:
        info = MODEL_OPENAI_COMPATIBLE[model_key]
        url = info['url']
        api_key = (model_cfg.api_key or '').strip() if model_cfg else ''
        if not api_key:
            raise RuntimeError(f'未配置 {model_key} 的 API Key，请联系老师在「个人设置」中补充。')
        # 所有 OpenAI 兼容模型都允许通过 extra_settings 覆盖模型名（视觉模型/推理模型/自定义）
        custom_model = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        model_name = custom_model if custom_model else info['default_model']
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}',
        }
        payload = {
            'model': model_name,
            'messages': messages,
            'temperature': 0.7,
            'stream': True,
        }
        if extra_payload:
            payload.update(extra_payload)
        yield from _stream_openai_compatible(url, headers, payload)
    elif model_key == 'cherry':
        api_key = (model_cfg.api_key or '').strip() if model_cfg else ''
        if not api_key:
            raise RuntimeError('未配置 Cherry Studio 企业版的 API Key，请联系老师在「个人设置」中补充。')
        url = _cherry_chat_url(model_cfg.api_url if model_cfg else '')
        custom_model = (model_cfg.extra_settings or '').strip() if model_cfg else ''
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}',
        }
        payload = {
            'model': custom_model or CHERRY_DEFAULT_MODEL,
            'messages': messages,
            'temperature': 0.7,
            'stream': True,
        }
        if extra_payload:
            payload.update(extra_payload)
        yield from _stream_openai_compatible(url, headers, payload)
    elif model_key == 'ollama':
        api_url = (model_cfg.api_url if model_cfg else 'http://localhost:11434') or 'http://localhost:11434'
        if not api_url.startswith('http'):
            api_url = 'http://' + api_url
        url = f"{api_url.rstrip('/')}/v1/chat/completions"
        ollama_model = (model_cfg.extra_settings if model_cfg and model_cfg.extra_settings else 'llama3.2')
        headers = {'Content-Type': 'application/json'}
        payload = {
            'model': ollama_model,
            'messages': messages,
            'temperature': 0.7,
            'stream': True,
        }
        if extra_payload:
            payload.update(extra_payload)
        yield from _stream_openai_compatible(url, headers, payload)
    else:
        raise RuntimeError(f'暂不支持模型: {model_key}')


def _sse_event(event_type, data_dict):
    """构造一条 SSE 事件字符串"""
    return f"event: {event_type}\ndata: {json.dumps(data_dict, ensure_ascii=False)}\n\n"


def _check_assistant_accessible(assistant):
    """返回 (ok, status_code, reason)。"""
    if not assistant:
        return False, 404, '未找到该 AI 导师，请检查链接是否正确。'
    if assistant.status == 'disabled':
        return False, 403, '该 AI 导师已被老师暂时关闭，请稍后再试或联系老师。'
    if assistant.link_expires_at and assistant.link_expires_at < datetime.now():
        return False, 410, '该 AI 导师的访问链接已过期，请联系老师重新生成。'
    return True, 200, None


# ==================== 班级花名册（B4）====================

ROSTER_HEADER_ALIASES = {
    'class_name': ['班级', 'class', '班级名称', 'class_name', '班', '所在班级'],
    'student_name': ['姓名', '名字', 'name', 'student_name', '学生姓名', 'student'],
    'student_no': ['学号', 'no', 'id', 'student_no', 'student_id', '编号'],
    'notes': ['备注', 'note', 'notes', 'remark', 'remarks', '说明'],
}


def _normalize_text(v):
    if v is None:
        return ''
    s = str(v).strip()
    return ' '.join(s.split())  # 折叠多余空白


def _detect_roster_columns(header_row):
    """根据表头自动定位列索引；返回 dict {field: col_index 或 None}"""
    cols = {k: None for k in ROSTER_HEADER_ALIASES}
    for idx, cell in enumerate(header_row):
        key = _normalize_text(cell).lower()
        if not key:
            continue
        for field, aliases in ROSTER_HEADER_ALIASES.items():
            if cols[field] is not None:
                continue
            for alias in aliases:
                if key == alias.lower():
                    cols[field] = idx
                    break
    return cols


def _parse_roster_file(file_storage):
    """解析上传的花名册文件 (.xlsx / .xls / .csv)，返回 (entries, error)
    每个 entry 是 dict: {class_name, student_name, student_no, notes}
    """
    filename = (file_storage.filename or '').lower()
    if not filename:
        return None, '请选择要上传的文件'

    try:
        if filename.endswith('.xlsx'):
            from openpyxl import load_workbook
            wb = load_workbook(file_storage, read_only=True, data_only=True)
            ws = wb.active
            rows = [list(r) for r in ws.iter_rows(values_only=True)]
        elif filename.endswith('.xls'):
            try:
                import xlrd
            except ImportError:
                return None, '当前环境未安装 xlrd，无法解析 .xls；请另存为 .xlsx 后再试。'
            data = file_storage.read()
            book = xlrd.open_workbook(file_contents=data)
            sheet = book.sheet_by_index(0)
            rows = [sheet.row_values(i) for i in range(sheet.nrows)]
        elif filename.endswith('.csv'):
            import csv
            import io
            raw = file_storage.read()
            text_data = None
            for enc in ('utf-8-sig', 'utf-8', 'gbk', 'gb18030'):
                try:
                    text_data = raw.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue
            if text_data is None:
                return None, 'CSV 编码无法识别，请另存为 UTF-8 后再试。'
            reader = csv.reader(io.StringIO(text_data))
            rows = list(reader)
        else:
            return None, '只支持 .xlsx / .xls / .csv 三种格式'
    except Exception as e:
        logger.exception('[花名册] 解析文件失败')
        return None, f'解析失败：{e}'

    # 跳过空行
    rows = [r for r in rows if any(_normalize_text(c) for c in r)]
    if len(rows) < 2:
        return None, '文件内容为空或仅有表头'

    cols = _detect_roster_columns(rows[0])
    if cols['class_name'] is None or cols['student_name'] is None:
        return None, '未识别到「班级」和「姓名」两列，请使用模板或确保表头包含这两个字段'

    entries = []
    seen = set()
    for row in rows[1:]:
        def get(field):
            idx = cols[field]
            if idx is None or idx >= len(row):
                return ''
            return _normalize_text(row[idx])
        class_name = get('class_name')
        student_name = get('student_name')
        if not class_name or not student_name:
            continue
        if len(class_name) > 100 or len(student_name) > 100:
            continue
        key = (class_name, student_name)
        if key in seen:
            continue
        seen.add(key)
        entries.append({
            'class_name': class_name,
            'student_name': student_name,
            'student_no': get('student_no')[:50] or None,
            'notes': get('notes')[:200] or None,
        })

    if not entries:
        return None, '没有解析到有效的学生记录（需至少包含「班级 + 姓名」两列）'
    return entries, None


def _parse_class_options_text(text):
    """解析自由填写模式班级选项的文本输入。
    支持两种风格：
      高一: 1班, 2班, 3班
      高二: 1班, 2班
    或直接列班级（不带年级）：
      1班
      2班
      3班
    返回 [{'grade': str|None, 'class_name': str}, ...]
    """
    result = []
    seen = set()
    if not text:
        return result
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        # 找 grade : classes
        grade = None
        rest = line
        for sep in (':', '：'):
            if sep in line:
                grade, _, rest = line.partition(sep)
                grade = grade.strip()[:50] or None
                break
        # 拆班级列表
        parts = []
        for piece in rest.replace('，', ',').replace('、', ',').split(','):
            piece = piece.strip()
            if piece:
                parts.append(piece[:100])
        if not parts and grade:
            # 只写了 "高一:" 没班级 → 跳过
            continue
        if not parts:
            # 整行就是一个班级，无年级
            parts = [line[:100]]
            grade = None
        for cn in parts:
            key = (grade or '', cn)
            if key in seen:
                continue
            seen.add(key)
            result.append({'grade': grade, 'class_name': cn})
    return result


def _serialize_class_options_text(options):
    """把数据库 [(grade, class_name), ...] 反序列化成可编辑的多行文本。
    带年级的同年级班级合并到一行；不带年级的每行一个。
    """
    by_grade = {}
    no_grade = []
    for o in options:
        if o.grade:
            by_grade.setdefault(o.grade, []).append(o.class_name)
        else:
            no_grade.append(o.class_name)
    lines = []
    for grade, classes in by_grade.items():
        lines.append(f"{grade}: {', '.join(classes)}")
    lines.extend(no_grade)
    return '\n'.join(lines)


def _load_class_options_for_user(db, user_id):
    """返回结构化数据：
    {
        'has_options': bool,
        'has_grade_dim': bool,            # 是否启用年级层
        'flat_classes': [str, ...],       # 不分年级时直接用
        'classes_by_grade': {grade: [classes]},
        'grades': [str, ...],
    }
    """
    options = db.query(ClassOption).filter_by(user_id=user_id).order_by(
        ClassOption.sort_order, ClassOption.id
    ).all()
    if not options:
        return {
            'has_options': False,
            'has_grade_dim': False,
            'flat_classes': [],
            'classes_by_grade': {},
            'grades': [],
        }
    has_grade_dim = any((o.grade or '').strip() for o in options)
    classes_by_grade = {}
    flat_classes = []
    for o in options:
        if o.grade:
            classes_by_grade.setdefault(o.grade, []).append(o.class_name)
        else:
            classes_by_grade.setdefault('', []).append(o.class_name)
            flat_classes.append(o.class_name)
    grades = [g for g in classes_by_grade.keys() if g]
    return {
        'has_options': True,
        'has_grade_dim': has_grade_dim,
        'flat_classes': flat_classes,
        'classes_by_grade': classes_by_grade,
        'grades': grades,
    }


@app.route('/roster')
@login_required
def roster_index():
    db = SessionLocal()
    try:
        entries = db.query(RosterEntry).filter_by(user_id=current_user.id).order_by(
            RosterEntry.class_name, RosterEntry.student_name
        ).all()
        # 按班级分组
        grouped = {}
        for e in entries:
            grouped.setdefault(e.class_name, []).append(e)
        # 统计：用过该班级的 AI 导师数（roster_mode=strict）
        strict_assistants = db.query(Task).filter_by(
            user_id=current_user.id, is_assistant=True, roster_mode='strict'
        ).count()
        # B9: 自由填写模式班级选项
        class_options = db.query(ClassOption).filter_by(user_id=current_user.id).order_by(
            ClassOption.sort_order, ClassOption.id
        ).all()
        class_options_text = _serialize_class_options_text(class_options)
        return render_template(
            'roster.html',
            grouped=grouped,
            total=len(entries),
            strict_assistants=strict_assistants,
            class_options=class_options,
            class_options_text=class_options_text,
            class_options_count=len(class_options),
        )
    finally:
        db.close()


@app.route('/roster/class_options/save', methods=['POST'])
@login_required
def class_options_save():
    """B9: 替换式保存班级选项（自由填写模式用）"""
    text = request.form.get('class_options_text', '')
    parsed = _parse_class_options_text(text)
    db = SessionLocal()
    try:
        # 替换：先清空再插入，保证旧条目不残留
        db.query(ClassOption).filter_by(user_id=current_user.id).delete()
        for i, item in enumerate(parsed):
            db.add(ClassOption(
                user_id=current_user.id,
                grade=item['grade'],
                class_name=item['class_name'],
                sort_order=i,
            ))
        db.commit()
        if parsed:
            flash(f'已保存 {len(parsed)} 个班级选项（自由填写模式生效）', 'success')
        else:
            flash('已清空班级选项；自由填写模式将回退为学生自行输入班级', 'info')
    finally:
        db.close()
    return redirect(url_for('roster_index') + '#class-options')


@app.route('/roster/class_options/clear', methods=['POST'])
@login_required
def class_options_clear():
    db = SessionLocal()
    try:
        deleted = db.query(ClassOption).filter_by(user_id=current_user.id).delete()
        db.commit()
        flash(f'已清空班级选项（共 {deleted} 条）', 'success')
    finally:
        db.close()
    return redirect(url_for('roster_index') + '#class-options')


@app.route('/roster/import', methods=['POST'])
@login_required
def roster_import():
    file = request.files.get('file')
    if not file:
        flash('请选择要上传的文件', 'warning')
        return redirect(url_for('roster_index'))
    entries, err = _parse_roster_file(file)
    if err:
        flash(f'导入失败：{err}', 'danger')
        return redirect(url_for('roster_index'))

    mode = request.form.get('import_mode', 'merge')  # merge / replace
    db = SessionLocal()
    try:
        if mode == 'replace':
            db.query(RosterEntry).filter_by(user_id=current_user.id).delete()
            db.commit()

        added, skipped = 0, 0
        for ent in entries:
            existing = db.query(RosterEntry).filter_by(
                user_id=current_user.id,
                class_name=ent['class_name'],
                student_name=ent['student_name'],
            ).first()
            if existing:
                # 更新可选字段
                if ent['student_no']:
                    existing.student_no = ent['student_no']
                if ent['notes']:
                    existing.notes = ent['notes']
                skipped += 1
            else:
                db.add(RosterEntry(
                    user_id=current_user.id,
                    class_name=ent['class_name'],
                    student_name=ent['student_name'],
                    student_no=ent['student_no'],
                    notes=ent['notes'],
                ))
                added += 1
        db.commit()
        flash(f'导入完成：新增 {added} 名学生{("，更新 " + str(skipped) + " 名学生信息") if skipped else ""}。', 'success')
    finally:
        db.close()
    return redirect(url_for('roster_index'))


@app.route('/roster/add', methods=['POST'])
@login_required
def roster_add():
    class_name = _normalize_text(request.form.get('class_name'))
    student_name = _normalize_text(request.form.get('student_name'))
    student_no = _normalize_text(request.form.get('student_no'))[:50] or None
    notes = _normalize_text(request.form.get('notes'))[:200] or None

    if not class_name or not student_name:
        flash('请填写班级和姓名', 'warning')
        return redirect(url_for('roster_index'))
    if len(class_name) > 100 or len(student_name) > 100:
        flash('班级或姓名过长（不超过 100 字符）', 'warning')
        return redirect(url_for('roster_index'))

    db = SessionLocal()
    try:
        existing = db.query(RosterEntry).filter_by(
            user_id=current_user.id, class_name=class_name, student_name=student_name
        ).first()
        if existing:
            flash(f'「{class_name} · {student_name}」已存在', 'info')
        else:
            db.add(RosterEntry(
                user_id=current_user.id,
                class_name=class_name,
                student_name=student_name,
                student_no=student_no,
                notes=notes,
            ))
            db.commit()
            flash(f'已添加 {class_name} · {student_name}', 'success')
    finally:
        db.close()
    return redirect(url_for('roster_index'))


@app.route('/roster/<int:entry_id>/delete', methods=['POST'])
@login_required
def roster_delete(entry_id):
    db = SessionLocal()
    try:
        entry = db.query(RosterEntry).filter_by(id=entry_id, user_id=current_user.id).first()
        if entry:
            db.delete(entry)
            db.commit()
            flash(f'已删除 {entry.class_name} · {entry.student_name}', 'success')
        else:
            flash('记录不存在', 'warning')
    finally:
        db.close()
    return redirect(url_for('roster_index'))


@app.route('/roster/clear_class', methods=['POST'])
@login_required
def roster_clear_class():
    class_name = _normalize_text(request.form.get('class_name'))
    if not class_name:
        flash('未指定班级', 'warning')
        return redirect(url_for('roster_index'))
    db = SessionLocal()
    try:
        deleted = db.query(RosterEntry).filter_by(
            user_id=current_user.id, class_name=class_name
        ).delete()
        db.commit()
        flash(f'已删除「{class_name}」共 {deleted} 名学生', 'success')
    finally:
        db.close()
    return redirect(url_for('roster_index'))


@app.route('/roster/clear', methods=['POST'])
@login_required
def roster_clear():
    db = SessionLocal()
    try:
        deleted = db.query(RosterEntry).filter_by(user_id=current_user.id).delete()
        db.commit()
        flash(f'已清空全部花名册（共 {deleted} 条）', 'success')
    finally:
        db.close()
    return redirect(url_for('roster_index'))


@app.route('/roster/template.xlsx')
@login_required
def roster_template():
    """下载导入模板"""
    from openpyxl import Workbook
    import io
    wb = Workbook()
    ws = wb.active
    ws.title = '花名册'
    ws.append(['班级', '姓名', '学号', '备注'])
    ws.append(['高一(3)班', '张三', '20240101', '示例-可删除'])
    ws.append(['高一(3)班', '李四', '20240102', ''])
    ws.append(['高一(4)班', '王五', '20240201', ''])
    for col, width in zip('ABCD', (16, 14, 14, 24)):
        ws.column_dimensions[col].width = width
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    from flask import send_file
    return send_file(
        buf, as_attachment=True, download_name='openmentor_roster_template.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    )


@app.route('/roster/export.xlsx')
@login_required
def roster_export():
    """导出当前花名册"""
    from openpyxl import Workbook
    import io
    db = SessionLocal()
    try:
        entries = db.query(RosterEntry).filter_by(user_id=current_user.id).order_by(
            RosterEntry.class_name, RosterEntry.student_name
        ).all()
    finally:
        db.close()
    wb = Workbook()
    ws = wb.active
    ws.title = '花名册'
    ws.append(['班级', '姓名', '学号', '备注'])
    for e in entries:
        ws.append([e.class_name, e.student_name, e.student_no or '', e.notes or ''])
    for col, width in zip('ABCD', (16, 14, 14, 24)):
        ws.column_dimensions[col].width = width
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    from flask import send_file
    return send_file(
        buf, as_attachment=True,
        download_name=f'openmentor_roster_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    )


# ==================== 学情大屏（B5）====================

@app.route('/analytics')
@login_required
def analytics_index():
    db = SessionLocal()
    try:
        assistants = db.query(Task).filter_by(
            user_id=current_user.id, is_assistant=True
        ).order_by(Task.created_at.desc()).all()
        # 全部班级（来自 conversation 实际产生过的 + 花名册预设）
        classes_in_use = {c.student_class for c in db.query(Conversation).join(
            Task, Conversation.assistant_id == Task.id
        ).filter(Task.user_id == current_user.id).all() if c.student_class}
        classes_in_roster = {r.class_name for r in db.query(RosterEntry).filter_by(
            user_id=current_user.id
        ).all()}
        classes = sorted(classes_in_use | classes_in_roster)
        return render_template(
            'analytics.html',
            assistants=[{'id': a.id, 'title': a.title} for a in assistants],
            classes=classes,
        )
    finally:
        db.close()


@app.route('/api/analytics/data')
@login_required
def api_analytics_data():
    """返回学情大屏所有图表数据"""
    from datetime import timedelta as _td
    db = SessionLocal()
    try:
        assistant_id = request.args.get('assistant_id', 'all')
        class_name = request.args.get('class_name', 'all')
        try:
            days = max(1, min(180, int(request.args.get('days') or 30)))
        except ValueError:
            days = 30

        since = datetime.now() - _td(days=days)
        # 基础查询：当前老师所有 AI 导师下符合时间窗的消息
        q = db.query(Message).join(
            Conversation, Message.conversation_id == Conversation.id
        ).join(
            Task, Conversation.assistant_id == Task.id
        ).filter(
            Task.user_id == current_user.id,
            Task.is_assistant == True,
            Message.created_at >= since,
        )

        if assistant_id != 'all':
            try:
                aid = int(assistant_id)
                q = q.filter(Task.id == aid)
            except ValueError:
                pass
        if class_name != 'all':
            q = q.filter(Conversation.student_class == class_name)

        rows = []
        # 按需取字段，避免 N+1
        for m in q.all():
            rows.append({
                'date': m.created_at.date().isoformat() if m.created_at else None,
                'hour': m.created_at.hour if m.created_at else 0,
                'role': m.role,
                'rating': m.rating,
                'triggered_keyword': m.triggered_keyword,
                'tokens_used': m.tokens_used or 0,
                'class_name': m.conversation.student_class if m.conversation else '',
                'student_name': m.conversation.student_name if m.conversation else '',
                'assistant_id': m.conversation.assistant_id if m.conversation else None,
                'assistant_title': m.conversation.assistant.title if (m.conversation and m.conversation.assistant) else '',
                'conversation_id': m.conversation_id,
            })

        # ---- 聚合 ----
        # 1) 每日消息趋势（按日期 / 按 role）
        date_keys = []
        cur = since.date()
        end = datetime.now().date()
        while cur <= end:
            date_keys.append(cur.isoformat())
            cur += _td(days=1)

        daily_user = {d: 0 for d in date_keys}
        daily_asst = {d: 0 for d in date_keys}
        for r in rows:
            if r['date'] not in daily_user:
                continue
            if r['role'] == 'user':
                daily_user[r['date']] += 1
            elif r['role'] == 'assistant':
                daily_asst[r['date']] += 1

        # 2) 班级活跃度（消息数 + 学生数）
        cls_msgs = {}
        cls_students = {}
        for r in rows:
            cn = r['class_name'] or '未填写'
            cls_msgs[cn] = cls_msgs.get(cn, 0) + 1
            cls_students.setdefault(cn, set()).add(r['student_name'])
        class_activity = sorted([
            {'class_name': k, 'messages': v, 'students': len(cls_students[k])}
            for k, v in cls_msgs.items()
        ], key=lambda x: x['messages'], reverse=True)[:15]

        # 3) 学生 Top 10（按 user 消息数）
        stu_msgs = {}
        for r in rows:
            if r['role'] != 'user':
                continue
            key = (r['class_name'] or '未填写', r['student_name'] or '匿名')
            stu_msgs[key] = stu_msgs.get(key, 0) + 1
        top_students = sorted([
            {'name': k[1], 'class_name': k[0], 'messages': v}
            for k, v in stu_msgs.items()
        ], key=lambda x: x['messages'], reverse=True)[:10]

        # 4) AI 导师消息分布
        asst_msgs = {}
        for r in rows:
            t = r['assistant_title'] or '未知'
            asst_msgs[t] = asst_msgs.get(t, 0) + 1
        assistant_distribution = [
            {'title': k, 'messages': v} for k, v in
            sorted(asst_msgs.items(), key=lambda x: x[1], reverse=True)
        ]

        # 5) 满意度
        up = sum(1 for r in rows if r['role'] == 'assistant' and r['rating'] == 1)
        down = sum(1 for r in rows if r['role'] == 'assistant' and r['rating'] == -1)
        unrated = sum(1 for r in rows if r['role'] == 'assistant' and not r['rating'])
        rated = up + down
        satisfaction_pct = round(up * 100 / rated) if rated > 0 else None

        # 6) 关键词命中 Top 10
        kw_counts = {}
        for r in rows:
            if r['triggered_keyword']:
                kw_counts[r['triggered_keyword']] = kw_counts.get(r['triggered_keyword'], 0) + 1
        top_keywords = [
            {'keyword': k, 'hits': v} for k, v in
            sorted(kw_counts.items(), key=lambda x: x[1], reverse=True)[:10]
        ]

        # 7) 24 小时活跃热度（按消息小时）
        hour_buckets = [0] * 24
        for r in rows:
            hour_buckets[r['hour']] += 1

        # 8) 概览
        unique_students = len({(r['class_name'], r['student_name']) for r in rows if r['student_name']})
        unique_classes = len({r['class_name'] for r in rows if r['class_name']})
        total_msgs = len(rows)
        total_user_msgs = sum(1 for r in rows if r['role'] == 'user')
        blocked_hits = sum(1 for r in rows if r['triggered_keyword'])
        total_tokens = sum(r['tokens_used'] for r in rows)

        return jsonify({
            'code': 200,
            'window': {'days': days, 'since': since.date().isoformat(), 'until': end.isoformat()},
            'summary': {
                'active_students': unique_students,
                'active_classes': unique_classes,
                'total_messages': total_msgs,
                'total_user_messages': total_user_msgs,
                'total_tokens': total_tokens,
                'blocked_hits': blocked_hits,
                'satisfaction_pct': satisfaction_pct,
                'thumbs_up': up,
                'thumbs_down': down,
            },
            'daily_messages': [
                {'date': d, 'user': daily_user[d], 'assistant': daily_asst[d]}
                for d in date_keys
            ],
            'class_activity': class_activity,
            'top_students': top_students,
            'assistant_distribution': assistant_distribution,
            'satisfaction': {'up': up, 'down': down, 'unrated': unrated},
            'top_keywords': top_keywords,
            'hour_heatmap': [{'hour': h, 'count': hour_buckets[h]} for h in range(24)],
        })
    finally:
        db.close()


@app.route('/chat/<string:task_id>')
def chat_entry(task_id):
    """学生端入口：未识别身份显示填写页，已识别进入聊天室"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(task_id=task_id, is_assistant=True).first()
        ok, code, reason = _check_assistant_accessible(assistant)
        if not ok:
            return render_template('chat_unavailable.html', reason=reason), code

        # 检查 cookie 是否已有身份
        cookie_name = _student_cookie_name(task_id)
        student_token = request.cookies.get(cookie_name)
        conv = None
        if student_token:
            conv = db.query(Conversation).filter_by(
                assistant_id=assistant.id, student_token=student_token
            ).first()

        is_group = _assistant_is_group(assistant)
        member_name = _get_group_member_name(task_id) if is_group else ''
        if conv and not conv.is_blocked and (not is_group or member_name):
            # 已有身份 → 进入聊天室（小组模式还需要组员姓名 cookie，缺失则重新填写）
            cap = MODEL_CAPABILITIES.get(assistant.selected_model_name or '', {})
            return render_template(
                'chat_room.html',
                assistant=assistant,
                conversation=conv,
                model_label=cap.get('label', assistant.selected_model_name or 'AI'),
                is_group=is_group,
                member_name=member_name,
            )

        # 显示身份填写页（如果是 strict 花名册模式，传入老师的花名册数据）
        roster_data = None
        if (assistant.roster_mode or 'off') == 'strict':
            entries = db.query(RosterEntry).filter_by(user_id=assistant.user_id).all()
            class_map = {}
            for e in entries:
                class_map.setdefault(e.class_name, []).append(e.student_name)
            for k in class_map:
                class_map[k].sort()
            roster_data = {
                'mode': 'strict',
                'classes': sorted(class_map.keys()),
                'students_by_class': class_map,
            }
        # B9: 自由填写模式 + 老师配置了班级选项 → 给学生提供下拉
        class_options_data = None
        if (assistant.roster_mode or 'off') == 'off':
            cfg = _load_class_options_for_user(db, assistant.user_id)
            if cfg['has_options']:
                class_options_data = cfg
        return render_template(
            'chat_identity.html',
            assistant=assistant,
            roster=roster_data,
            class_options=class_options_data,
        )
    finally:
        db.close()


@app.route('/chat/<string:task_id>/identify', methods=['POST'])
def chat_identify(task_id):
    """学生提交班级+姓名，创建会话并设置 cookie"""
    db = SessionLocal()
    try:
        assistant = db.query(Task).filter_by(task_id=task_id, is_assistant=True).first()
        ok, code, reason = _check_assistant_accessible(assistant)
        if not ok:
            return render_template('chat_unavailable.html', reason=reason), code

        student_class = (request.form.get('student_class') or '').strip()
        student_name = (request.form.get('student_name') or '').strip()
        if not student_class or not student_name or len(student_class) > 100 or len(student_name) > 100:
            flash('请填写有效的班级和姓名（各不超过 100 字符）', 'warning')
            return redirect(url_for('chat_entry', task_id=task_id))

        # 花名册严格校验：必须在老师的花名册里
        if (assistant.roster_mode or 'off') == 'strict':
            in_roster = db.query(RosterEntry).filter_by(
                user_id=assistant.user_id,
                class_name=student_class,
                student_name=student_name,
            ).first()
            if not in_roster:
                flash(f'花名册中找不到「{student_class} · {student_name}」，请确认班级和姓名都从下拉中选择。', 'warning')
                return redirect(url_for('chat_entry', task_id=task_id))
        elif (assistant.roster_mode or 'off') == 'off':
            # B9: 自由填写模式下，如果老师配置了班级选项 → 提交的 student_class 必须在选项里
            cfg = _load_class_options_for_user(db, assistant.user_id)
            if cfg['has_options']:
                allowed_classes = set()
                for grade, classes in cfg['classes_by_grade'].items():
                    for cn in classes:
                        if grade:
                            allowed_classes.add(f'{grade} {cn}')
                        else:
                            allowed_classes.add(cn)
                if student_class not in allowed_classes:
                    flash('请从老师配置的班级中选择，不要手动输入', 'warning')
                    return redirect(url_for('chat_entry', task_id=task_id))

        # 小组模式：会话归属 = 班级 + 小组名（不带设备号 → 多设备进同一会话）；
        # 学生本人姓名另存 member cookie，发言时做【姓名】前缀
        is_group = _assistant_is_group(assistant)
        group_name = (request.form.get('group_name') or '').strip()
        if is_group:
            if not group_name or len(group_name) > 50:
                flash('请填写小组名称（不超过 50 字符）', 'warning')
                return redirect(url_for('chat_entry', task_id=task_id))
            token = _hash_student_token(student_class, f'⟨小组⟩{group_name}', None)
            conv_display_name = group_name
        else:
            # B8 设备隔离：前端 localStorage 的 device_uuid 参与 token 生成；
            # 不传时回退老逻辑（兼容老链接，保留上线前的会话）
            device_uuid = (request.form.get('device_uuid') or '').strip()
            # 简单校验：长度 8-128 字符 + 仅 ASCII 字母数字和短横线，防止注入
            if device_uuid:
                import re as _re
                if not _re.fullmatch(r'[A-Za-z0-9\-]{8,128}', device_uuid):
                    device_uuid = ''
            token = _hash_student_token(student_class, student_name, device_uuid or None)
            conv_display_name = student_name

        # 找已有会话或创建新的
        conv = db.query(Conversation).filter_by(
            assistant_id=assistant.id, student_token=token
        ).first()
        if conv:
            # 更新班级姓名（如果有变化）和最近活跃时间
            conv.student_class = student_class
            conv.student_name = conv_display_name
            conv.last_active_at = datetime.now()
        else:
            conv = Conversation(
                assistant_id=assistant.id,
                student_class=student_class,
                student_name=conv_display_name,
                student_token=token,
                started_at=datetime.now(),
                last_active_at=datetime.now(),
                daily_reset_date=datetime.now().date(),
            )
            db.add(conv)
        db.commit()

        # 设置 cookie，过期时间与链接一致
        max_age = None
        if assistant.link_expires_at:
            seconds = int((assistant.link_expires_at - datetime.now()).total_seconds())
            max_age = max(60, seconds)
        resp = redirect(url_for('chat_entry', task_id=task_id))
        resp.set_cookie(
            _student_cookie_name(task_id), token,
            max_age=max_age, httponly=True, samesite='Lax'
        )
        if is_group:
            from urllib.parse import quote
            resp.set_cookie(
                _member_cookie_name(task_id), quote(student_name),
                max_age=max_age, httponly=True, samesite='Lax'
            )
        return resp
    finally:
        db.close()


@app.route('/chat/<string:task_id>/leave', methods=['POST'])
def chat_leave(task_id):
    """学生切换身份：清除 cookie 后重新填写"""
    resp = redirect(url_for('chat_entry', task_id=task_id))
    resp.delete_cookie(_student_cookie_name(task_id))
    return resp


def _resolve_conversation(db, task_id):
    """根据 task_id 和 cookie 解析出 (assistant, conversation)。
    返回 (assistant, conv, error_dict 或 None)
    """
    assistant = db.query(Task).filter_by(task_id=task_id, is_assistant=True).first()
    ok, code, reason = _check_assistant_accessible(assistant)
    if not ok:
        return assistant, None, {'code': code, 'message': reason}

    cookie_name = _student_cookie_name(task_id)
    student_token = request.cookies.get(cookie_name)
    if not student_token:
        return assistant, None, {'code': 401, 'message': '请先填写班级和姓名'}
    conv = db.query(Conversation).filter_by(
        assistant_id=assistant.id, student_token=student_token
    ).first()
    if not conv:
        return assistant, None, {'code': 401, 'message': '会话不存在，请重新进入'}
    if conv.is_blocked:
        return assistant, conv, {'code': 403, 'message': '该会话已被老师封禁'}
    return assistant, conv, None


@app.route('/api/chat/<string:task_id>/history')
def api_chat_history(task_id):
    """获取学生历史消息。
    带 ?after_id=N 时为增量模式（小组轮询用）：只返回 id > N 的消息 + 额度，省略能力等静态字段。"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']

        after_id = request.args.get('after_id', type=int)

        used, limit = _ensure_daily_quota(conv)
        db.commit()

        def _display_filename(path):
            """从 '<uuid8>_<original>' 还原 '<original>' 用于展示"""
            base = os.path.basename(path or '')
            if '_' in base:
                prefix, _, rest = base.partition('_')
                if len(prefix) == 8 and all(c in '0123456789abcdef' for c in prefix.lower()):
                    return rest
            return base

        msgs = []
        for m in conv.messages:
            if after_id is not None and m.id <= after_id:
                continue
            try:
                imgs = json.loads(m.image_paths) if m.image_paths else []
            except (json.JSONDecodeError, TypeError):
                imgs = []
            try:
                files = json.loads(m.file_paths) if m.file_paths else []
            except (json.JSONDecodeError, TypeError):
                files = []
            try:
                audios = json.loads(m.audio_paths) if m.audio_paths else []
            except (json.JSONDecodeError, TypeError):
                audios = []
            msgs.append({
                'id': m.id,
                'role': m.role,
                'content': m.content,
                'triggered_keyword': m.triggered_keyword,
                'rating': m.rating,
                'image_urls': ['/' + p for p in imgs],
                'files': [{'name': _display_filename(p), 'url': '/' + p} for p in files],
                'audio_urls': ['/' + p for p in audios],
                'generated_image_url': ('/' + m.generated_image_path) if m.generated_image_path else None,
                'generated_video_url': ('/' + m.generated_video_path) if m.generated_video_path else None,
                'created_at': m.created_at.isoformat() if m.created_at else None,
            })

        # 增量模式：轻量返回（轮询高频，省掉能力探测等静态字段）
        if after_id is not None:
            return jsonify({
                'messages': msgs,
                'daily_used': used,
                'daily_limit': limit,
            })

        # 助学多模态能力（前端用于决定是否显示按钮）
        # 图片生成能力 = 助学开关 + 提供商支持 + 老师在 image_gen_model 实际填了模型名
        provider_key = assistant.selected_model_name or ''
        cap = MODEL_CAPABILITIES.get(provider_key, {})
        ai_cfg = db.query(AIConfig).filter_by(user_id=assistant.user_id).first()
        image_gen_configured = False
        if ai_cfg and provider_key in IMAGE_GEN_PROVIDERS:
            for mc in ai_cfg.model_configs:
                if mc.model_name == provider_key and (mc.image_gen_model or '').strip() and (mc.api_key or '').strip():
                    image_gen_configured = True
                    break
        video_gen_configured = False
        if ai_cfg and provider_key in VIDEO_GEN_PROVIDERS:
            for mc in ai_cfg.model_configs:
                if mc.model_name == provider_key and (mc.video_gen_model or '').strip() and (mc.api_key or '').strip():
                    video_gen_configured = True
                    break
        capabilities = {
            'image_input': bool(assistant.allow_image_input and cap.get('image_input')),
            'file_upload': bool(assistant.allow_file_upload),
            'image_generation': bool(assistant.allow_image_generation and cap.get('image_generation') and image_gen_configured),
            'video_generation': bool(assistant.allow_video_generation and cap.get('video_generation') and video_gen_configured),
            'voice_input': bool(assistant.allow_voice_input and (_get_siliconflow_key(db, assistant) or _assistant_audio_direct(assistant))),
            'tts': bool(assistant.allow_tts),
            'call_teacher': bool(assistant.allow_call_teacher),
        }

        return jsonify({
            'messages': msgs,
            'daily_used': used,
            'daily_limit': limit,
            'student_class': conv.student_class,
            'student_name': conv.student_name,
            'welcome_message': assistant.welcome_message,
            'capabilities': capabilities,
            'voice_bridge_urls': _voice_bridge_urls(),
        })
    finally:
        db.close()


# --- 多模态附件相关 ---

OPENMENTOR_UPLOAD_ROOT = os.path.join('static', 'uploads', 'openmentor')

ALLOWED_IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp', 'tiff', 'tif', 'heic', 'heif'}
ALLOWED_AUDIO_EXTS = {'wav', 'mp3', 'm4a', 'aac', 'ogg', 'opus', 'webm', 'amr', 'flac'}
ALLOWED_DOC_EXTS = {
    'pdf', 'docx', 'doc',           # Word
    'xlsx', 'xls',                  # Excel
    'pptx', 'ppt',                  # PowerPoint
    'txt', 'md', 'markdown', 'rst', # 纯文本类
    'csv', 'tsv',                   # 表格
    'html', 'htm', 'xml',           # 标记语言
    'json', 'yaml', 'yml',          # 数据
}
MAX_UPLOAD_BYTES = 10 * 1024 * 1024  # 10 MB（HEIC 照片可能较大）
MAX_EXTRACTED_CHARS = 50000  # 文档文本最多保留 5 万字（防 token 爆炸）

# 注册 HEIF 支持（Pillow 默认不支持 HEIC）
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
    HEIC_SUPPORTED = True
except ImportError:
    HEIC_SUPPORTED = False
    logger.warning('pillow-heif 未安装，HEIC/HEIF 图片将无法处理')


def _ext_of(filename):
    return os.path.splitext(filename or '')[1].lower().lstrip('.')


def _conv_upload_dir(conv_id):
    """新写入路径：static/uploads/openmentor/<YYYY-MM-DD>/<conv_id>/
    按"自然日"分桶，避免单层文件夹太多；同一会话跨天会有多个日期子目录，互不影响。
    旧文件保留在原 static/uploads/openmentor/<conv_id>/ 不动，由 _is_safe_attachment_path 同时识别。
    """
    today = datetime.now().strftime('%Y-%m-%d')
    path = os.path.join(OPENMENTOR_UPLOAD_ROOT, today, str(conv_id))
    if not os.path.exists(path):
        os.makedirs(path, exist_ok=True)
    return path


def _is_safe_attachment_path(rel_path, conv_id):
    """校验路径属于本会话目录，防止跨会话偷文件。
    兼容两种布局：
      - 新（B15 加日期）：static/uploads/openmentor/<YYYY-MM-DD>/<conv_id>/...
      - 旧：              static/uploads/openmentor/<conv_id>/...
    """
    abs_rel = os.path.normpath(rel_path)
    if not os.path.exists(abs_rel):
        return False
    legacy_prefix = os.path.normpath(os.path.join(OPENMENTOR_UPLOAD_ROOT, str(conv_id))) + os.sep
    if abs_rel.startswith(legacy_prefix):
        return True
    # 新布局：路径形如 OPENMENTOR_UPLOAD_ROOT/<date>/<conv_id>/...
    parts = abs_rel.split(os.sep)
    root_parts = os.path.normpath(OPENMENTOR_UPLOAD_ROOT).split(os.sep)
    n = len(root_parts)
    if len(parts) < n + 2:
        return False
    if parts[:n] != root_parts:
        return False
    # parts[n] 是日期、parts[n+1] 是 conv_id
    if parts[n + 1] != str(conv_id):
        return False
    # 简单校验日期格式（10 字符 YYYY-MM-DD）
    if len(parts[n]) != 10 or parts[n][4] != '-' or parts[n][7] != '-':
        return False
    return True


# ============ B16: QuickForm 任务的附件上传 ============

QUICKFORM_UPLOAD_ROOT = os.path.join('static', 'uploads', 'quickform')
QUICKFORM_ALLOWED_EXTS = ALLOWED_IMAGE_EXTS | ALLOWED_DOC_EXTS | {
    'mp3', 'wav', 'm4a', 'aac', 'ogg',          # 音频
    'mp4', 'mov', 'avi', 'webm', 'mkv',          # 视频
    'zip', 'rar', '7z',                          # 压缩包
}
QUICKFORM_MAX_PER_FILE = 20 * 1024 * 1024        # 单文件 20 MB（比 OpenMentor 略宽，因为可能有视频）
QUICKFORM_MAX_PER_SUBMISSION = 80 * 1024 * 1024  # 单次提交总上限 80 MB


def _quickform_upload_dir(task_id_short):
    """static/uploads/quickform/<YYYY-MM-DD>/<task_id>/"""
    today = datetime.now().strftime('%Y-%m-%d')
    safe_id = re.sub(r'[^A-Za-z0-9_\-]', '', task_id_short or 'unknown')[:32] or 'unknown'
    path = os.path.join(QUICKFORM_UPLOAD_ROOT, today, safe_id)
    if not os.path.exists(path):
        os.makedirs(path, exist_ok=True)
    return path


def _save_quickform_uploads(files_storage, task_id_short):
    """把 QuickForm 表单 multipart 中的所有文件保存到磁盘。
    返回 (attachments, errors)：
      - attachments: [{'field','name','url','size','ext'}, ...]
      - errors: 友好错误消息列表（仍可继续保存其他合法文件）
    """
    from werkzeug.utils import secure_filename as _sec
    attachments = []
    errors = []
    total_bytes = 0
    for field in files_storage.keys():
        for f in files_storage.getlist(field):
            if not f or not f.filename:
                continue
            ext = _ext_of(f.filename)
            if not ext or ext not in QUICKFORM_ALLOWED_EXTS:
                errors.append(f'忽略不支持的扩展名：{f.filename}')
                continue
            f.seek(0, os.SEEK_END)
            size = f.tell()
            f.seek(0)
            if size <= 0:
                errors.append(f'忽略空文件：{f.filename}')
                continue
            if size > QUICKFORM_MAX_PER_FILE:
                errors.append(f'忽略超大文件 {f.filename}（{size // (1024*1024)} MB > {QUICKFORM_MAX_PER_FILE // (1024*1024)} MB）')
                continue
            if total_bytes + size > QUICKFORM_MAX_PER_SUBMISSION:
                errors.append(f'忽略 {f.filename}：本次提交累计已超过 {QUICKFORM_MAX_PER_SUBMISSION // (1024*1024)} MB 上限')
                continue
            original_base = os.path.splitext(f.filename)[0]
            safe_base = _sec(original_base) or 'file'  # 中文/特殊字符回落
            unique_name = f'{uuid.uuid4().hex[:8]}_{safe_base}.{ext}'
            save_dir = _quickform_upload_dir(task_id_short)
            save_path = os.path.join(save_dir, unique_name)
            try:
                f.save(save_path)
            except Exception as e:
                errors.append(f'保存 {f.filename} 失败：{e}')
                continue
            total_bytes += size
            # 计算前端可访问的 URL（相对 / 绝对都可，模板自动 prefix）
            url_path = '/' + save_path.replace(os.sep, '/')
            attachments.append({
                'field': field,
                'name': f.filename,
                'url': url_path,
                'size': size,
                'ext': ext,
                'is_image': ext in ALLOWED_IMAGE_EXTS,
            })
    return attachments, errors


@app.route('/attachment/reveal', methods=['POST'])
@login_required
def attachment_reveal():
    """B17: 在服务器本机的文件管理器中定位/打开附件（仅本机有桌面环境时有意义）。
    支持 QuickForm（static/uploads/quickform）和 OpenMentor（static/uploads/openmentor）两类附件。
    """
    import platform as _platform
    import subprocess as _sp
    rel = (request.form.get('path') or '').strip()
    if not rel:
        return jsonify({'success': False, 'message': '缺少 path'}), 400
    abs_path = os.path.normpath(rel.lstrip('/'))
    safe_roots = [os.path.normpath(QUICKFORM_UPLOAD_ROOT), os.path.normpath(OPENMENTOR_UPLOAD_ROOT)]
    if not any(abs_path == r or abs_path.startswith(r + os.sep) for r in safe_roots):
        return jsonify({'success': False, 'message': '路径不在允许范围内'}), 400
    if not os.path.exists(abs_path):
        return jsonify({'success': False, 'message': '文件不存在'}), 404
    abs_full = os.path.abspath(abs_path)
    folder = os.path.dirname(abs_full)
    sysname = _platform.system()
    try:
        if sysname == 'Darwin':
            # macOS: open -R 高亮该文件
            _sp.Popen(['open', '-R', abs_full])
        elif sysname == 'Windows':
            # Windows: explorer /select 高亮该文件
            _sp.Popen(['explorer', f'/select,{abs_full}'])
        elif sysname == 'Linux':
            # Linux: xdg-open 打开父目录（多数桌面无文件高亮 API）
            _sp.Popen(['xdg-open', folder])
        else:
            return jsonify({'success': False, 'message': f'不支持的系统：{sysname}'}), 500
        logger.info(f'[reveal] {current_user.username} from {request.remote_addr} → {abs_full}')
        return jsonify({
            'success': True,
            'message': '已在服务器本机弹出文件夹（如果你不是从服务器本机访问，可能在那台机器上弹出）',
            'folder': folder,
        })
    except FileNotFoundError as e:
        return jsonify({'success': False, 'message': f'系统命令未找到：{e}'}), 500
    except Exception as e:
        logger.exception('[reveal] 打开文件夹失败')
        return jsonify({'success': False, 'message': str(e)}), 500


def _is_safe_quickform_attachment(rel_path, task_id_short):
    """校验某个 url 路径确实属于该 task 的 quickform 附件目录"""
    abs_rel = os.path.normpath(rel_path.lstrip('/'))
    if not os.path.exists(abs_rel):
        return False
    parts = abs_rel.split(os.sep)
    root_parts = os.path.normpath(QUICKFORM_UPLOAD_ROOT).split(os.sep)
    n = len(root_parts)
    if len(parts) < n + 2 or parts[:n] != root_parts:
        return False
    # parts[n] 是日期，parts[n+1] 是 task_id
    if parts[n + 1] != task_id_short:
        return False
    if len(parts[n]) != 10 or parts[n][4] != '-' or parts[n][7] != '-':
        return False
    return True


def _extract_text_from_document(file_path, ext):
    """从文档中提取纯文本，返回 (text, truncated_bool)。覆盖：
    PDF / Word / Excel / PowerPoint / 纯文本 / CSV / HTML / JSON / YAML / Markdown
    """
    try:
        if ext == 'pdf':
            import pypdf
            reader = pypdf.PdfReader(file_path)
            text = '\n\n'.join((page.extract_text() or '') for page in reader.pages)

        elif ext == 'docx':
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            # 表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            text = '\n'.join(paragraphs)

        elif ext == 'doc':
            # 老 .doc 格式：尝试用 textutil（macOS）或返回提示
            try:
                import subprocess
                r = subprocess.run(['textutil', '-convert', 'txt', '-stdout', file_path],
                                   capture_output=True, timeout=15)
                text = r.stdout.decode('utf-8', errors='ignore') if r.returncode == 0 else ''
                if not text:
                    text = '[.doc 老版本格式，已尝试解析但未能提取文本。建议另存为 .docx 后重新上传]'
            except Exception:
                text = '[.doc 老版本格式无法解析。建议在 Word 中"另存为 .docx"格式后重新上传]'

        elif ext == 'xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets_text = []
            for sheet in wb.worksheets:
                rows_text = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(c.strip() for c in cells):
                        rows_text.append(' | '.join(cells))
                if rows_text:
                    sheets_text.append(f'## 工作表：{sheet.title}\n' + '\n'.join(rows_text))
            text = '\n\n'.join(sheets_text)

        elif ext == 'xls':
            import xlrd
            wb = xlrd.open_workbook(file_path)
            sheets_text = []
            for sheet in wb.sheets():
                rows_text = []
                for row_idx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(row_idx, c)) for c in range(sheet.ncols)]
                    if any(c.strip() for c in cells):
                        rows_text.append(' | '.join(cells))
                if rows_text:
                    sheets_text.append(f'## 工作表：{sheet.name}\n' + '\n'.join(rows_text))
            text = '\n\n'.join(sheets_text)

        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                shape_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        shape_texts.append(shape.text)
                if shape_texts:
                    slides_text.append(f'## 第 {i} 页\n' + '\n'.join(shape_texts))
            text = '\n\n'.join(slides_text)

        elif ext == 'ppt':
            text = '[.ppt 老版本格式无法解析。建议在 PowerPoint 中"另存为 .pptx"后重新上传]'

        elif ext in ('html', 'htm'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                raw = f.read()
            # 简单去标签（不依赖 beautifulsoup）
            no_script = re.sub(r'<(script|style)[^>]*>.*?</\1>', '', raw, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', no_script)
            text = re.sub(r'\s+', ' ', text).strip()

        elif ext in ('txt', 'md', 'markdown', 'rst', 'csv', 'tsv', 'xml', 'json', 'yaml', 'yml'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()

        else:
            return '', False
    except Exception as e:
        logger.warning(f'文档解析失败 ({file_path}, ext={ext}): {e}')
        return f'[文件解析失败: {e}]', False

    truncated = False
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS]
        truncated = True
    return text, truncated


@app.route('/api/chat/<string:task_id>/tts', methods=['POST'])
def api_chat_tts(task_id):
    """学生端朗读 AI 回复：洗文本 → Edge TTS 合成 → 内容哈希缓存 → 返回音频 URL。"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']
        if not assistant.allow_tts:
            return jsonify({'code': 403, 'message': '该 AI 导师未开启朗读功能'}), 403
        if not _tts_available():
            return jsonify({'code': 503, 'message': '服务器朗读引擎不可用'}), 503

        body = request.get_json(silent=True) or {}
        try:
            message_id = int(body.get('message_id') or 0)
        except (TypeError, ValueError):
            message_id = 0
        msg = db.query(Message).filter_by(id=message_id, conversation_id=conv.id, role='assistant').first()
        if not msg:
            return jsonify({'code': 404, 'message': '消息不存在'}), 404

        text = _tts_clean_text(msg.content)
        if not text:
            return jsonify({'code': 400, 'message': '这条消息没有可朗读的文字'}), 400

        import hashlib as _hashlib
        cache_key = _hashlib.md5(f'{TTS_DEFAULT_VOICE}|{text}'.encode('utf-8')).hexdigest()
        os.makedirs(TTS_CACHE_DIR, exist_ok=True)
        cache_path = os.path.join(TTS_CACHE_DIR, f'{cache_key}.mp3')
        if not os.path.exists(cache_path):
            try:
                audio = _tts_synthesize_edge(text)
            except Exception as e:
                logger.warning(f'[TTS] 合成失败: {e}')
                return jsonify({'code': 502, 'message': f'朗读合成失败：{e}'}), 502
            with open(cache_path, 'wb') as f:
                f.write(audio)
            logger.info(f'[TTS] 合成 {len(text)} 字 → {os.path.basename(cache_path)} ({len(audio) // 1024} KB)')
        return jsonify({'code': 200, 'url': '/' + cache_path.replace('\\', '/')})
    finally:
        db.close()


@app.route('/voice-bridge')
def voice_bridge():
    """录音桥页面（纯静态）。本地访问（localhost）可直接用；
    部署到公网 HTTPS 镜像后，局域网 HTTP 页面通过弹窗 + postMessage 使用同一个文件。"""
    return app.send_static_file('voice_bridge.html')


@app.route('/api/chat/<string:task_id>/upload', methods=['POST'])
def api_chat_upload(task_id):
    """学生上传图片或文档。返回附件信息供发消息时引用。"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            logger.warning(f'[上传被拒] _resolve_conversation 失败: {err}')
            return jsonify(err), err['code']

        if 'file' not in request.files:
            logger.warning(f'[上传被拒] 请求中无 file 字段；form keys = {list(request.form.keys())}, files keys = {list(request.files.keys())}')
            return jsonify({'code': 400, 'message': '请选择文件（请求中找不到 file 字段）'}), 400
        f = request.files['file']
        if not f or not f.filename:
            logger.warning(f'[上传被拒] file 字段为空或无文件名')
            return jsonify({'code': 400, 'message': '请选择文件'}), 400

        ext = _ext_of(f.filename)
        logger.info(f'[上传] task_id={task_id}, conv_id={conv.id}, filename={f.filename!r}, ext={ext!r}, content_type={f.content_type!r}')

        if ext in ALLOWED_IMAGE_EXTS:
            kind = 'image'
            # 开了视频生成时放行图片上传（图生视频要用），不要求开图片理解
            if not assistant.allow_video_generation:
                if not assistant.allow_image_input:
                    logger.warning(f'[上传被拒] 助学未开启图片上传')
                    return jsonify({'code': 403, 'message': '该 AI 导师未开启图片上传'}), 403
                cap = MODEL_CAPABILITIES.get(assistant.selected_model_name or '', {})
                if not cap.get('image_input'):
                    logger.warning(f'[上传被拒] 模型 {assistant.selected_model_name!r} 不支持视觉')
                    return jsonify({'code': 403, 'message': '当前模型不支持图片理解，请联系老师切换模型'}), 403
        elif ext in ALLOWED_DOC_EXTS:
            kind = 'document'
            if not assistant.allow_file_upload:
                logger.warning(f'[上传被拒] 助学未开启文件上传')
                return jsonify({'code': 403, 'message': '该 AI 导师未开启文件上传'}), 403
        elif ext in ALLOWED_AUDIO_EXTS:
            kind = 'audio'
            if not assistant.allow_voice_input:
                logger.warning(f'[上传被拒] 助学未开启语音输入')
                return jsonify({'code': 403, 'message': '该 AI 导师未开启语音输入'}), 403
            if not _get_siliconflow_key(db, assistant) and not _assistant_audio_direct(assistant):
                logger.warning(f'[上传被拒] 老师未配置语音识别 key 且未开原声模式')
                return jsonify({'code': 403, 'message': '老师还没配置语音识别（需在个人设置中填写硅基流动 API Key）'}), 403
        else:
            logger.warning(f'[上传被拒] 不支持的扩展名 .{ext}（filename={f.filename!r}）')
            return jsonify({'code': 400, 'message': f'不支持的文件类型 .{ext}（filename={f.filename}）'}), 400

        # 大小校验
        f.seek(0, os.SEEK_END)
        size = f.tell()
        f.seek(0)
        if size > MAX_UPLOAD_BYTES:
            logger.warning(f'[上传被拒] 文件过大: {size} bytes')
            return jsonify({'code': 413, 'message': f'文件过大（{size//1024} KB ≤ {MAX_UPLOAD_BYTES // (1024*1024)} MB）'}), 413
        if size <= 0:
            logger.warning(f'[上传被拒] 空文件')
            return jsonify({'code': 400, 'message': '空文件'}), 400

        # 安全的文件名：先 secure 原 basename，再显式拼回扩展名（解决中文文件名被 secure_filename 剥光的 bug）
        from werkzeug.utils import secure_filename as _sec
        original_base = os.path.splitext(f.filename)[0]
        safe_base = _sec(original_base) or 'file'  # 中文/特殊字符全被剥时回落 'file'

        # HEIC/HEIF（iPhone 拍照默认格式）：浏览器原生不支持显示，且 LLM API 也不接受
        # 自动转换为 JPEG
        save_ext = ext
        if ext in ('heic', 'heif'):
            save_ext = 'jpg'

        unique_name = f'{uuid.uuid4().hex[:8]}_{safe_base}.{save_ext}'
        conv_dir = _conv_upload_dir(conv.id)
        save_path = os.path.join(conv_dir, unique_name)

        if ext in ('heic', 'heif'):
            # 用 Pillow + pillow-heif 转换
            try:
                from PIL import Image
                from io import BytesIO
                buf = BytesIO(f.read())
                img = Image.open(buf)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                img.save(save_path, 'JPEG', quality=92, optimize=True)
            except Exception as e:
                logger.exception('HEIC 转 JPEG 失败')
                return jsonify({'code': 500, 'message': f'HEIC 图片转换失败：{e}（请把照片转成 JPG/PNG 后重试）'}), 500
        else:
            f.save(save_path)

        # 文档：提取文本预览
        extracted_text = None
        truncated = False
        if kind == 'document':
            extracted_text, truncated = _extract_text_from_document(save_path, ext)

        # 返回前端
        rel_path = save_path.replace('\\', '/')
        public_url = '/' + rel_path
        return jsonify({
            'code': 200,
            'kind': kind,
            'file_path': rel_path,
            'file_name': f.filename,  # 原始文件名（含中文）用于前端 chip 显示
            'saved_name': unique_name,  # 实际落盘的安全文件名
            'url': public_url,
            'size_bytes': size,
            'extracted_chars': len(extracted_text) if extracted_text else 0,
            'extracted_truncated': truncated,
            'extracted_preview': (extracted_text[:300] + '...') if extracted_text and len(extracted_text) > 300 else (extracted_text or None),
        })
    finally:
        db.close()


# --- 图片生成 ---

def _default_image_size(provider_key, image_model):
    """按 provider/model 选择合适的输出尺寸
    - 豆包 seedream 4.x / 5.x：要求 ≥ 3,686,400 像素（1920×1920），用 2048×2048
    - 豆包 seedream 3.x、Endpoint ID：1024×1024
    - 阿里通义 qwen-image* 系列：2048×2048（高分辨率）
    - 其他（GLM cogview / SiliconFlow FLUX）：1024×1024
    """
    name = (image_model or '').lower()
    if provider_key == 'doubao':
        if any(tag in name for tag in ('seedream-4', 'seedream-5', 'seedream-6', 'seedream-7')):
            return '2048x2048'
        return '1024x1024'
    if provider_key == 'qwen' and name.startswith('qwen-image'):
        return '2048x2048'
    return '1024x1024'


def _save_image_bytes(img_bytes, conv_id, ext='png'):
    """落盘并返回相对路径"""
    conv_dir = _conv_upload_dir(conv_id)
    unique_name = f'gen_{uuid.uuid4().hex[:8]}.{ext}'
    save_path = os.path.join(conv_dir, unique_name)
    with open(save_path, 'wb') as f:
        f.write(img_bytes)
    return save_path.replace('\\', '/')


def _ext_from_url(url, default='png'):
    url_low = url.split('?')[0].lower()
    for e in ('png', 'jpg', 'jpeg', 'webp', 'gif'):
        if url_low.endswith('.' + e):
            return 'jpeg' if e == 'jpg' else e
    return default


def _poll_dashscope_task(api_key, task_id, conv_id, max_wait_seconds=120, poll_interval=3):
    """轮询 DashScope 异步任务直到 SUCCEEDED，下载图片落盘。
    兼容两种结果格式：
      - text2image 系列：output.results[0].url
      - multimodal-generation 系列：output.choices[0].message.content[i].image
    """
    info = IMAGE_GEN_PROVIDERS['qwen']
    poll_url = info['task_url_template'].format(task_id=task_id)
    poll_headers = {'Authorization': f'Bearer {api_key}'}
    elapsed = 0
    while elapsed < max_wait_seconds:
        try:
            r = requests.get(poll_url, headers=poll_headers, timeout=20)
        except requests.exceptions.RequestException as e:
            return None, f'轮询异常: {e}'
        if r.status_code != 200:
            return None, f'轮询失败 {r.status_code}: {r.text[:200]}'
        try:
            pdata = r.json()
        except (ValueError, json.JSONDecodeError):
            return None, f'轮询返回非 JSON: {r.text[:200]}'
        po = pdata.get('output') or {}
        status = (po.get('task_status') or '').upper()
        if status == 'SUCCEEDED':
            # 提取图片 URL
            url = None
            results = po.get('results') or []
            if results:
                url = results[0].get('url')
            if not url:
                # multimodal 格式
                choices = po.get('choices') or []
                if choices:
                    content = (choices[0].get('message') or {}).get('content') or []
                    for c in content:
                        if isinstance(c, dict) and c.get('image'):
                            url = c['image']
                            break
            if not url:
                return None, f'任务成功但未找到图片 URL: {str(po)[:300]}'
            try:
                img_resp = requests.get(url, timeout=60)
                img_resp.raise_for_status()
                img_bytes = img_resp.content
            except Exception as e:
                return None, f'下载远程图失败: {e}'
            ext = _ext_from_url(url)
            return _save_image_bytes(img_bytes, conv_id, ext), None
        if status in ('FAILED', 'CANCELED', 'UNKNOWN'):
            err_code = po.get('code')
            err_msg = po.get('message') or '任务失败'
            return None, f'任务{status}: {err_msg} (code={err_code})'
        # PENDING / RUNNING：等待
        time.sleep(poll_interval)
        elapsed += poll_interval
    return None, f'任务超时（已等待 {max_wait_seconds} 秒），可能任务仍在排队，建议稍后重试'


def _extract_image_url_from_dashscope_output(output):
    """从 DashScope 同步返回中提取图片 url（兼容多种格式）"""
    if not isinstance(output, dict):
        return None
    # multimodal-generation 格式：output.choices[0].message.content[i].image
    choices = output.get('choices') or []
    if choices:
        content = (choices[0].get('message') or {}).get('content') or []
        for c in content:
            if isinstance(c, dict) and c.get('image'):
                return c['image']
    # text2image 格式：output.results[0].url
    results = output.get('results') or []
    if results and isinstance(results[0], dict):
        return results[0].get('url')
    return None


def _generate_image_qwen_async(model_cfg, prompt, conv_id, size='1024*1024',
                               max_wait_seconds=120, poll_interval=3):
    """阿里通义图生：按模型名前缀分发到不同 DashScope 接口
      - wanx* / wan2*：text2image/image-synthesis（input.prompt）→ 异步
      - qwen-image*：multimodal-generation/generation（input.messages）→ 同步（部分套餐不支持异步）
    """
    info = IMAGE_GEN_PROVIDERS['qwen']
    api_key = (model_cfg.api_key or '').strip()
    image_model = (model_cfg.image_gen_model or '').strip() or info['default_model']
    image_model_lower = image_model.lower()

    is_qwen_image = image_model_lower.startswith('qwen-image')

    if is_qwen_image:
        # qwen-image 系列：multimodal-generation 接口，**同步模式**（不加 X-DashScope-Async）
        submit_url = info['qwen_image_submit_url']
        submit_payload = {
            'model': image_model,
            'input': {
                'messages': [
                    {'role': 'user', 'content': [{'text': prompt}]}
                ]
            },
            'parameters': {'size': size, 'n': 1},
        }
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json',
        }
        log_path = 'multimodal-generation (sync)'
        # 同步生成可能需要 30-60 秒，timeout 给宽
        submit_timeout = max_wait_seconds
    else:
        # wanx / wan2 系列：text2image 接口，异步模式
        submit_url = info['wanx_submit_url']
        submit_payload = {
            'model': image_model,
            'input': {'prompt': prompt},
            'parameters': {'size': size, 'n': 1},
        }
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json',
            'X-DashScope-Async': 'enable',
        }
        log_path = 'text2image (async)'
        submit_timeout = 30

    logger.info(f'[图生请求] provider=qwen ({log_path}), model={image_model}, size={size}, prompt={prompt[:80]}')

    # 提交任务
    try:
        resp = requests.post(submit_url, headers=headers, json=submit_payload, timeout=submit_timeout)
    except requests.exceptions.RequestException as e:
        return None, f'网络异常: {e}'
    if resp.status_code != 200:
        # 特殊错误：套餐不支持异步 → 提示用户
        if 'asynchronous calls' in resp.text or 'AccessDenied' in resp.text:
            return None, f'当前 API 不支持异步调用（DashScope 套餐限制）。{resp.text[:200]}'
        return None, f'提交任务失败 {resp.status_code}: {resp.text[:300]}'
    try:
        data = resp.json()
    except (ValueError, json.JSONDecodeError):
        return None, f'提交返回非 JSON: {resp.text[:200]}'
    output = data.get('output') or {}

    # 同步返回（qwen-image 走同步路径）：直接从 output 取图
    sync_url = _extract_image_url_from_dashscope_output(output)
    if sync_url:
        try:
            img_resp = requests.get(sync_url, timeout=60)
            img_resp.raise_for_status()
            return _save_image_bytes(img_resp.content, conv_id, _ext_from_url(sync_url)), None
        except Exception as e:
            return None, f'下载失败: {e}'

    # 异步返回：拿 task_id 轮询
    task_id = output.get('task_id')
    if task_id:
        return _poll_dashscope_task(api_key, task_id, conv_id, max_wait_seconds, poll_interval)

    # 既无图也无 task_id：报错
    code = data.get('code') or output.get('code')
    msg = data.get('message') or output.get('message') or str(data)[:200]
    return None, f'未返回结果（code={code}, msg={msg}）'


def _generate_image(provider_key, model_cfg, prompt, conv_id, size=None):
    """调用图生 API，下载结果到本地。按 provider 分发到 OpenAI 兼容 / DashScope 异步。
    返回 (relative_path, error_or_none)"""
    if provider_key not in IMAGE_GEN_PROVIDERS:
        return None, f'{provider_key} 不支持图片生成'
    info = IMAGE_GEN_PROVIDERS[provider_key]
    if not model_cfg or not (model_cfg.api_key or '').strip():
        return None, f'未配置 {provider_key} 的 API Key'
    image_model = (model_cfg.image_gen_model or '').strip()
    if not image_model:
        return None, f'老师未在「个人设置」配置 {provider_key} 的图片生成模型名'

    if not size:
        size = _default_image_size(provider_key, image_model)

    # 通义（万相 + qwen-image 系列）走异步路径
    if info.get('mode') == 'qwen_dispatch':
        # DashScope 用 1024*1024 格式（星号），不是 1024x1024
        size_dashscope = size.replace('x', '*')
        return _generate_image_qwen_async(model_cfg, prompt, conv_id, size=size_dashscope)

    # OpenAI 兼容（同步）路径
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {model_cfg.api_key.strip()}',
    }
    payload = {
        'model': image_model,
        'prompt': prompt,
        'size': size,
        'n': 1,
    }
    gen_url = info['url']
    if provider_key == 'cherry':
        # 服务端地址可自定义；网关后面模型五花八门，size/n 交给网关默认值
        gen_url = _cherry_api_base(model_cfg.api_url) + '/images/generations'
        payload.pop('size', None)
        payload.pop('n', None)
    logger.info(f'[图生请求] provider={provider_key}, model={image_model}, size={payload.get("size", "默认")}, prompt={prompt[:80]}')

    try:
        resp = requests.post(gen_url, headers=headers, json=payload, timeout=180)
    except requests.exceptions.RequestException as e:
        return None, f'网络异常: {e}'

    if resp.status_code != 200:
        return None, f'API {resp.status_code}: {resp.text[:300]}'

    try:
        data = resp.json()
    except (ValueError, json.JSONDecodeError):
        return None, f'API 返回非 JSON: {resp.text[:200]}'

    items = data.get('data') or []
    if not items:
        return None, f'API 未返回图片: {str(data)[:200]}'
    item = items[0]

    img_bytes = None
    ext = 'png'
    if item.get('b64_json'):
        try:
            img_bytes = base64.b64decode(item['b64_json'])
        except Exception as e:
            return None, f'base64 解码失败: {e}'
    elif item.get('url'):
        # 下载远程图
        url = item['url']
        url_low = url.split('?')[0].lower()
        for e in ('png', 'jpg', 'jpeg', 'webp', 'gif'):
            if url_low.endswith('.' + e):
                ext = 'jpeg' if e == 'jpg' else e
                break
        try:
            img_resp = requests.get(url, timeout=60)
            img_resp.raise_for_status()
            img_bytes = img_resp.content
        except Exception as e:
            return None, f'下载远程图失败: {e}'
    else:
        return None, f'API 返回的图片无 url 或 b64_json: {item}'

    # 落盘
    conv_dir = _conv_upload_dir(conv_id)
    unique_name = f'gen_{uuid.uuid4().hex[:8]}.{ext}'
    save_path = os.path.join(conv_dir, unique_name)
    try:
        with open(save_path, 'wb') as f:
            f.write(img_bytes)
    except Exception as e:
        return None, f'保存图片失败: {e}'
    return save_path.replace('\\', '/'), None


@app.route('/api/chat/<string:task_id>/generate_image', methods=['POST'])
def api_chat_generate_image(task_id):
    """学生端：根据描述生成图片"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']

        # 助学是否开启图片生成
        if not assistant.allow_image_generation:
            return jsonify({'code': 403, 'message': '该 AI 导师未开启图片生成'}), 403

        provider_key = assistant.selected_model_name or ''
        cap = MODEL_CAPABILITIES.get(provider_key, {})
        if not cap.get('image_generation'):
            return jsonify({'code': 403, 'message': f'当前模型（{provider_key}）暂不支持图片生成'}), 403

        # 取请求体
        body = request.get_json(silent=True) or {}
        prompt = (body.get('prompt') or '').strip()
        if not prompt:
            return jsonify({'code': 400, 'message': '请填写图片描述'}), 400
        if len(prompt) > 1000:
            return jsonify({'code': 400, 'message': '描述过长（≤ 1000 字符）'}), 400

        # 每日额度（图生与对话共用配额）
        used, limit = _ensure_daily_quota(conv)
        if used >= limit:
            return jsonify({'code': 429, 'message': f'今日额度已用完（{limit} 条），明天再来吧'}), 429

        # 黑名单
        hit = _check_blocked_keywords(prompt, assistant.blocked_keywords)
        if hit:
            label, kw = hit
            user_msg = Message(
                conversation_id=conv.id, role='user', content=f'[生成图片] {prompt}',
                triggered_keyword=f'{label}:{kw}',
            )
            sys_msg = Message(
                conversation_id=conv.id, role='system',
                content='这个描述不适合用来生成图片哦~ 我们换个内容吧。',
                # 不在系统提示上标记 triggered_keyword，否则违禁触发会被统计成 2 次（用户消息已标记）
            )
            db.add_all([user_msg, sys_msg])
            conv.daily_message_count = (conv.daily_message_count or 0) + 1
            conv.last_active_at = datetime.now()
            db.commit()
            return jsonify({
                'code': 200,
                'blocked': True,
                'category': label,
                'notice': sys_msg.content,
                'sys_message_id': sys_msg.id,
                'daily_used': conv.daily_message_count,
                'daily_limit': limit,
            })

        # 取老师的 model_cfg
        ai_cfg = db.query(AIConfig).filter_by(user_id=assistant.user_id).first()
        model_cfg = None
        if ai_cfg:
            for mc in ai_cfg.model_configs:
                if mc.model_name == provider_key:
                    model_cfg = mc
                    break

        # 保存 user 消息（图生请求）
        user_msg = Message(
            conversation_id=conv.id, role='user',
            content=f'[生成图片] {prompt}',
        )
        db.add(user_msg)
        conv.daily_message_count = (conv.daily_message_count or 0) + 1
        conv.last_active_at = datetime.now()
        db.commit()
        user_message_id = user_msg.id

        # 调用图生
        rel_path, error = _generate_image(provider_key, model_cfg, prompt, conv.id)

        if error or not rel_path:
            # 失败：保存 system 消息
            err_msg = Message(
                conversation_id=conv.id, role='system',
                content=f'图片生成失败：{error or "未知错误"}',
            )
            db.add(err_msg)
            db.commit()
            return jsonify({
                'code': 500,
                'success': False,
                'message': error or '图片生成失败',
                'user_message_id': user_message_id,
                'sys_message_id': err_msg.id,
                'daily_used': conv.daily_message_count,
                'daily_limit': limit,
            }), 500

        # 成功：保存 assistant 消息（含生成图）
        asst_msg = Message(
            conversation_id=conv.id, role='assistant',
            content='已生成图片',
            generated_image_path=rel_path,
        )
        db.add(asst_msg)
        db.commit()

        return jsonify({
            'code': 200,
            'success': True,
            'image_url': '/' + rel_path,
            'prompt': prompt,
            # 前端用这两个推进 lastMsgId，避免 4 秒轮询把已入库消息重复渲染
            'user_message_id': user_message_id,
            'message_id': asst_msg.id,
            'daily_used': conv.daily_message_count,
            'daily_limit': limit,
        })
    finally:
        db.close()


@app.route('/api/chat/<string:task_id>/generate_video', methods=['POST'])
def api_chat_generate_video(task_id):
    """学生端：创建视频生成任务（文生视频 / 图生视频），返回任务 id 供前端轮询。"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']

        if not assistant.allow_video_generation:
            return jsonify({'code': 403, 'message': '该 AI 导师未开启视频生成'}), 403
        provider_key = assistant.selected_model_name or ''
        cap = MODEL_CAPABILITIES.get(provider_key, {})
        if not cap.get('video_generation') or provider_key not in VIDEO_GEN_PROVIDERS:
            return jsonify({'code': 403, 'message': f'当前模型（{provider_key}）暂不支持视频生成'}), 403

        body = request.get_json(silent=True) or {}
        prompt = (body.get('prompt') or '').strip()
        mode = body.get('mode') if body.get('mode') in ('text', 'image') else 'text'
        image_path = (body.get('image_path') or '').strip()
        if not prompt:
            return jsonify({'code': 400, 'message': '请填写视频描述'}), 400
        if len(prompt) > 800:
            return jsonify({'code': 400, 'message': '描述过长（≤ 800 字符）'}), 400
        if mode == 'image':
            if not image_path or not _is_safe_attachment_path(image_path, conv.id) or _ext_of(image_path) not in ALLOWED_IMAGE_EXTS:
                return jsonify({'code': 400, 'message': '图生视频需要先上传一张图片'}), 400

        used, limit = _ensure_daily_quota(conv)
        if used >= limit:
            return jsonify({'code': 429, 'message': f'今日额度已用完（{limit} 条），明天再来吧'}), 429

        # 黑名单
        hit = _check_blocked_keywords(prompt, assistant.blocked_keywords)
        mode_label = '图生视频' if mode == 'image' else '生成视频'
        if hit:
            label, kw = hit
            user_msg = Message(
                conversation_id=conv.id, role='user', content=f'[{mode_label}] {prompt}',
                triggered_keyword=f'{label}:{kw}',
            )
            sys_msg = Message(
                conversation_id=conv.id, role='system',
                content='这个描述不适合用来生成视频哦~ 我们换个内容吧。',
                # 不在系统提示上标记 triggered_keyword，避免统计成 2 次（用户消息已标记）
            )
            db.add_all([user_msg, sys_msg])
            conv.daily_message_count = (conv.daily_message_count or 0) + 1
            conv.last_active_at = datetime.now()
            db.commit()
            return jsonify({
                'code': 200, 'blocked': True, 'category': label, 'notice': sys_msg.content,
                'sys_message_id': sys_msg.id,
                'daily_used': conv.daily_message_count, 'daily_limit': limit,
            })

        # 老师的 doubao 配置（key + 视频模型名）
        ai_cfg = db.query(AIConfig).filter_by(user_id=assistant.user_id).first()
        model_cfg = None
        if ai_cfg:
            for mc in ai_cfg.model_configs:
                if mc.model_name == provider_key:
                    model_cfg = mc
                    break
        if not model_cfg or not (model_cfg.api_key or '').strip():
            return jsonify({'code': 500, 'message': '老师未配置该提供商的 API Key'}), 500
        video_model = (model_cfg.video_gen_model or '').strip()
        if not video_model:
            return jsonify({'code': 500, 'message': '老师未在「个人设置」配置视频生成模型名'}), 500

        # 组装 Ark 异步任务请求
        info = VIDEO_GEN_PROVIDERS[provider_key]
        content = [{'type': 'text', 'text': f"{prompt} {info['default_params']}"}]
        if mode == 'image':
            data_url = _image_to_data_url(image_path)
            if not data_url:
                return jsonify({'code': 400, 'message': '图片读取失败，请重新上传'}), 400
            content.append({'type': 'image_url', 'image_url': {'url': data_url}})

        try:
            resp = requests.post(
                info['create_url'],
                headers={'Authorization': f"Bearer {model_cfg.api_key.strip()}", 'Content-Type': 'application/json'},
                json={'model': video_model, 'content': content},
                timeout=30,
            )
            payload = resp.json()
        except Exception as e:
            return jsonify({'code': 502, 'message': f'创建视频任务失败：{e}'}), 502
        if resp.status_code != 200 or not payload.get('id'):
            err_text = (payload.get('error') or {}).get('message') or str(payload)[:200]
            return jsonify({'code': 502, 'message': f'创建视频任务失败：{err_text}'}), 502
        gen_task_id = payload['id']

        # 保存 user 消息 + 计额度
        user_msg = Message(conversation_id=conv.id, role='user', content=f'[{mode_label}] {prompt}')
        db.add(user_msg)
        conv.daily_message_count = (conv.daily_message_count or 0) + 1
        conv.last_active_at = datetime.now()
        db.commit()

        _VIDEO_TASKS[gen_task_id] = {'conv_id': conv.id, 'prompt': prompt, 'done': False, 'result': None}
        logger.info(f'[视频生成] 任务创建 {gen_task_id} mode={mode} conv={conv.id} model={video_model}')
        return jsonify({
            'code': 200, 'task_id': gen_task_id,
            'user_message_id': user_msg.id,
            'daily_used': conv.daily_message_count, 'daily_limit': limit,
        })
    finally:
        db.close()


@app.route('/api/chat/<string:task_id>/video_task/<string:gen_task_id>')
def api_chat_video_task(task_id, gen_task_id):
    """轮询视频生成任务状态；成功时下载视频落盘并存为 assistant 消息。"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']
        info = _VIDEO_TASKS.get(gen_task_id)
        if not info or info['conv_id'] != conv.id:
            return jsonify({'code': 404, 'status': 'expired', 'message': '任务不存在或已过期，请重新生成'}), 404
        if info['done']:
            return jsonify(info['result'])

        provider_key = assistant.selected_model_name or ''
        prov = VIDEO_GEN_PROVIDERS.get(provider_key)
        ai_cfg = db.query(AIConfig).filter_by(user_id=assistant.user_id).first()
        model_cfg = None
        if ai_cfg:
            for mc in ai_cfg.model_configs:
                if mc.model_name == provider_key:
                    model_cfg = mc
                    break
        if not prov or not model_cfg:
            return jsonify({'code': 500, 'status': 'failed', 'message': '视频生成配置缺失'}), 500

        try:
            resp = requests.get(
                prov['task_url_template'].format(task_id=gen_task_id),
                headers={'Authorization': f"Bearer {model_cfg.api_key.strip()}"},
                timeout=15,
            )
            payload = resp.json()
        except Exception as e:
            return jsonify({'code': 200, 'status': 'running', 'message': f'查询暂时失败，继续等待：{e}'})

        status = payload.get('status')
        if status in ('queued', 'running'):
            return jsonify({'code': 200, 'status': 'running'})

        if status == 'succeeded':
            video_url = ((payload.get('content') or {}).get('video_url') or '').strip()
            rel_path = None
            error = None
            if not video_url:
                error = '任务成功但未返回视频地址'
            else:
                try:
                    v = requests.get(video_url, timeout=120)
                    v.raise_for_status()
                    conv_dir = _conv_upload_dir(conv.id)
                    unique_name = f'genv_{uuid.uuid4().hex[:8]}.mp4'
                    save_path = os.path.join(conv_dir, unique_name)
                    with open(save_path, 'wb') as f:
                        f.write(v.content)
                    rel_path = save_path.replace('\\', '/')
                except Exception as e:
                    error = f'视频下载失败：{e}'
            if error:
                result = {'code': 500, 'status': 'failed', 'message': error}
                sys_msg = Message(conversation_id=conv.id, role='system', content=f'视频生成失败：{error}')
                db.add(sys_msg)
                db.commit()
            else:
                asst_msg = Message(
                    conversation_id=conv.id, role='assistant',
                    content='已生成视频', generated_video_path=rel_path,
                )
                db.add(asst_msg)
                db.commit()
                result = {
                    'code': 200, 'status': 'succeeded',
                    'video_url': '/' + rel_path, 'message_id': asst_msg.id,
                }
                logger.info(f'[视频生成] 完成 {gen_task_id} → {rel_path} ({len(v.content) // 1024} KB)')
            info['done'] = True
            info['result'] = result
            return jsonify(result)

        # failed / cancelled
        err_text = ((payload.get('error') or {}).get('message') or status or '未知错误')
        result = {'code': 500, 'status': 'failed', 'message': f'视频生成失败：{err_text}'}
        sys_msg = Message(conversation_id=conv.id, role='system', content=f'视频生成失败：{err_text}')
        db.add(sys_msg)
        db.commit()
        info['done'] = True
        info['result'] = result
        return jsonify(result)
    finally:
        db.close()


def _validate_and_load_attachments(conv, image_paths, file_paths):
    """校验前端传来的附件路径（必须是本会话目录），返回 (valid_image_paths, valid_file_paths, file_text_blocks)"""
    valid_images = []
    for p in (image_paths or []):
        if not isinstance(p, str):
            logger.warning(f'[附件丢弃] 非字符串: {p!r}')
            continue
        if not _is_safe_attachment_path(p, conv.id):
            logger.warning(f'[附件丢弃] 路径越权或文件不存在: {p}')
            continue
        if _ext_of(p) not in ALLOWED_IMAGE_EXTS:
            logger.warning(f'[附件丢弃] 图片扩展名非法: {p} (ext={_ext_of(p)!r})')
            continue
        valid_images.append(p)

    valid_files = []
    file_text_blocks = []
    for p in (file_paths or []):
        if not isinstance(p, str):
            logger.warning(f'[附件丢弃] 非字符串: {p!r}')
            continue
        if not _is_safe_attachment_path(p, conv.id):
            logger.warning(f'[附件丢弃] 路径越权或文件不存在: {p}')
            continue
        ext = _ext_of(p)
        if ext not in ALLOWED_DOC_EXTS:
            logger.warning(f'[附件丢弃] 文档扩展名非法: {p} (ext={ext!r})')
            continue
        valid_files.append(p)
        text, truncated = _extract_text_from_document(p, ext)
        file_name = os.path.basename(p)
        block = f'\n\n[附件「{file_name}」内容{"（已截断到前 5 万字）" if truncated else ""}]\n{text}'
        file_text_blocks.append(block)

    return valid_images, valid_files, file_text_blocks


@app.route('/api/chat/<string:task_id>/message', methods=['POST'])
def api_chat_message(task_id):
    """学生发送消息，SSE 流式返回 AI 回复"""
    db = SessionLocal()
    captured = {}
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']

        # 读取消息内容
        body = request.get_json(silent=True) or {}
        user_text = (body.get('content') or '').strip()
        image_paths_in = body.get('image_paths') or []
        file_paths_in = body.get('file_paths') or []
        audio_paths_in = body.get('audio_paths') or []

        # 小组模式：ask_ai=False 表示组内讨论（只落库不调模型）；独立模式恒为问 AI
        is_group = _assistant_is_group(assistant)
        member_name = _get_group_member_name(task_id) if is_group else ''
        ask_ai = body.get('ask_ai')
        ask_ai = True if ask_ai is None else bool(ask_ai)
        if not is_group:
            ask_ai = True

        # 至少有文本或附件
        if not user_text and not image_paths_in and not file_paths_in and not audio_paths_in:
            return jsonify({'code': 400, 'message': '消息内容不能为空'}), 400
        if len(user_text) > 4000:
            return jsonify({'code': 400, 'message': '消息过长（≤ 4000 字符）'}), 400

        # 每日额度检查（小组讨论消息不耗额度，只有问 AI 才检查/扣减）
        used, limit = _ensure_daily_quota(conv)
        if ask_ai and used >= limit:
            return jsonify({'code': 429, 'message': f'今日额度已用完（{limit} 条），明天再来吧'}), 429

        # 语音附件：先校验路径，再转写成文字（转写文本与打字文本同样参与黑名单检查）
        valid_audios = []
        transcript_text = ''
        for p in audio_paths_in:
            if not isinstance(p, str) or not _is_safe_attachment_path(p, conv.id):
                logger.warning(f'[语音丢弃] 路径越权或文件不存在: {p!r}')
                continue
            if _ext_of(p) not in ALLOWED_AUDIO_EXTS:
                logger.warning(f'[语音丢弃] 音频扩展名非法: {p}')
                continue
            valid_audios.append(p)
        # 原声直达模式：音频直接发给多模态模型；转写降级为"审计/历史副产品"（可缺省、失败不阻断）
        direct_audio = _assistant_audio_direct(assistant) and bool(valid_audios)
        if direct_audio:
            total_bytes = sum(os.path.getsize(p) for p in valid_audios if os.path.exists(p))
            if total_bytes > AUDIO_DIRECT_MAX_BYTES:
                return jsonify({'code': 413, 'message': f'音频太大（{total_bytes // (1024 * 1024)}MB），原声模式请控制在约 2 分钟以内'}), 413
        if valid_audios:
            asr_key = _get_siliconflow_key(db, assistant)
            if not asr_key and not direct_audio:
                return jsonify({'code': 403, 'message': '老师还没配置语音识别（需在个人设置中填写硅基流动 API Key）'}), 403
            transcripts = []
            if asr_key:
                for p in valid_audios:
                    try:
                        t = _transcribe_audio(p, asr_key)
                    except Exception as e:
                        if direct_audio:
                            logger.warning(f'[原声模式] 静默转写失败（不影响发送）: {e}')
                            break
                        logger.exception(f'语音转写失败: {p}')
                        return jsonify({'code': 502, 'message': f'语音识别失败，请重试：{e}'}), 502
                    if t:
                        transcripts.append(t)
            transcript_text = ' '.join(transcripts).strip()
            if not transcript_text and not direct_audio and not user_text and not image_paths_in and not file_paths_in:
                return jsonify({'code': 400, 'message': '没有听清你说的话，请离麦克风近一点再试一次'}), 400

        # 黑名单检查（打字文本 + 语音转写文本）
        check_text = ' '.join(x for x in (user_text, transcript_text) if x)
        if check_text:
            hit = _check_blocked_keywords(check_text, assistant.blocked_keywords)
            if hit:
                label, kw = hit
                blocked_content = (f'【{member_name}】' if is_group and member_name else '') + check_text
                user_msg = Message(
                    conversation_id=conv.id, role='user', content=blocked_content,
                    audio_paths=json.dumps(valid_audios, ensure_ascii=False) if valid_audios else None,
                    triggered_keyword=f'{label}:{kw}',
                )
                sys_msg = Message(
                    conversation_id=conv.id, role='system',
                    content='这个问题不适合在这里讨论哦~ 我们换个话题继续学习吧。',
                    # 不在系统提示上标记 triggered_keyword，避免统计成 2 次（用户消息已标记）
                )
                db.add_all([user_msg, sys_msg])
                conv.daily_message_count = (conv.daily_message_count or 0) + 1
                conv.last_active_at = datetime.now()
                db.commit()
                _emit_app_event(assistant.user_id, 'blocked_keyword', {
                    'assistant_id': assistant.id, 'assistant_title': assistant.title,
                    'conversation_id': conv.id,
                    'student_class': conv.student_class, 'student_name': member_name or conv.student_name,
                    'category': label, 'keyword': kw, 'content': check_text[:120],
                })
                return jsonify({
                    'code': 200,
                    'blocked': True,
                    'category': label,
                    'notice': sys_msg.content,
                    # 前端用它推进 lastMsgId，否则 4 秒轮询会把已入库的这两条重新渲染一遍（重复气泡）
                    'sys_message_id': sys_msg.id,
                    'daily_used': conv.daily_message_count,
                    'daily_limit': limit,
                })

        # 校验附件
        valid_images, valid_files, file_text_blocks = _validate_and_load_attachments(conv, image_paths_in, file_paths_in)

        # 把语音转写文本和文档提取文本拼到消息内容（让模型读到学生说的话和文件正文）
        text_parts = [x for x in (user_text, transcript_text) if x]
        if direct_audio and not text_parts:
            merged_content = '🎤 语音消息'
        else:
            merged_content = '\n'.join(text_parts) or '请帮我看下这些附件。'
        if file_text_blocks:
            merged_content += ''.join(file_text_blocks)
        # 小组模式：消息前缀标记发言人（AI 和组员都能看到是谁说的）
        if is_group and member_name:
            merged_content = f'【{member_name}】{merged_content}'

        # 保存 user message
        user_msg = Message(
            conversation_id=conv.id, role='user', content=merged_content,
            image_paths=json.dumps(valid_images, ensure_ascii=False) if valid_images else None,
            file_paths=json.dumps(valid_files, ensure_ascii=False) if valid_files else None,
            audio_paths=json.dumps(valid_audios, ensure_ascii=False) if valid_audios else None,
        )
        db.add(user_msg)
        if ask_ai:
            conv.daily_message_count = (conv.daily_message_count or 0) + 1
        conv.last_active_at = datetime.now()
        db.commit()
        if ask_ai and conv.daily_message_count == max(1, int(limit * 0.8)):
            _emit_app_event(assistant.user_id, 'quota_warning', {
                'assistant_id': assistant.id, 'assistant_title': assistant.title,
                'conversation_id': conv.id,
                'student_class': conv.student_class, 'student_name': conv.student_name,
                'used': conv.daily_message_count, 'limit': limit,
            })

        # 小组讨论消息：只落库，不调模型，直接返回
        if not ask_ai:
            return jsonify({
                'code': 200,
                'discussion': True,
                'message_id': user_msg.id,
                'content': merged_content,
                'transcript': transcript_text or None,
                'daily_used': conv.daily_message_count or 0,
                'daily_limit': limit,
            })

        # 决定模型 + 是否走视觉路径（按 provider 能力判断；具体是否支持图片由老师在配置时把握）
        model_key, model_cfg = _get_assistant_model_config(db, assistant)
        if not model_key:
            return jsonify({'code': 500, 'message': '老师未配置 AI 模型，请联系老师'}), 500
        cap = MODEL_CAPABILITIES.get(model_key, {})
        supports_vision = bool(cap.get('image_input'))

        # 知识库检索（RAG）：用学生本轮问题检索相关片段，拼进本次请求的 system prompt
        # 检索失败绝不影响正常对话
        system_prompt_for_llm = assistant.system_prompt or ''
        try:
            kb_query = ' '.join(x for x in (user_text, transcript_text) if x)
            if kb_query:
                sf_key = _get_siliconflow_key(db, assistant)
                if sf_key:
                    kb_block = _kb_retrieve(db, assistant, kb_query, sf_key)
                    if kb_block:
                        system_prompt_for_llm += kb_block
                        logger.info(f'[知识库] 命中并注入 {len(kb_block)} 字参考资料 (assistant={assistant.id})')
        except Exception as e:
            logger.warning(f'[知识库] 检索失败（已跳过，不影响对话）: {e}')

        # 繁体模式：界面由出口转换层处理，AI 回复直接要求模型用繁体输出
        if _lang_is_hant():
            system_prompt_for_llm += '\n\n請始終使用繁體中文（正體中文）回答學生。'

        # 小组模式：告诉模型这是多人共享会话，消息前缀【】是发言人姓名
        if is_group:
            system_prompt_for_llm += (
                '\n\n当前是小组共享会话，多名学生在同一个对话里。每条学生消息开头【】内是发言学生的姓名。'
                '回答时自然地称呼提问学生的名字；之前其他组员的讨论内容可以作为上下文参考，'
                '如果讨论中已有不错的思路，可以先肯定再补充。'
            )

        # 构造 LLM 请求
        all_msgs = list(conv.messages)
        llm_messages = _build_history_for_llm(
            all_msgs, system_prompt_for_llm, max_history=20, supports_vision=supports_vision,
        )

        # 捕获 generator 所需数据
        # 原声直达：把本轮音频以 input_audio 附到刚保存的这条 user 消息上
        # （历史里的旧音频消息只保留转写文本，不重复传音频，省 token）
        extra_payload = None
        if direct_audio:
            audio_parts = _build_input_audio_parts(valid_audios, model_key)
            for item in reversed(llm_messages):
                if item['role'] == 'user':
                    if isinstance(item['content'], str):
                        item['content'] = [{'type': 'text', 'text': item['content']}] + audio_parts
                    else:
                        item['content'] = list(item['content']) + audio_parts
                    break
            if model_key == 'qwen':
                # qwen-omni 系列需显式声明只要文本输出
                extra_payload = {'modalities': ['text']}
            logger.info(f'[原声模式] 附加 {len(audio_parts)} 段音频直达模型 {model_key}')

        captured['conv_id'] = conv.id
        captured['model_key'] = model_key
        captured['daily_limit'] = limit
        captured['model_cfg_id'] = model_cfg.id if model_cfg else None
        captured['llm_messages'] = llm_messages
        captured['transcript'] = transcript_text
        captured['user_msg_id'] = user_msg.id
        captured['extra_payload'] = extra_payload
        captured['display_transcript'] = transcript_text or ('🎤 语音消息' if direct_audio else '')

    except Exception as e:
        logger.exception(f'chat message 处理失败: {e}')
        db.rollback()
        return jsonify({'code': 500, 'message': f'服务器内部错误: {e}'}), 500
    finally:
        db.close()

    # 构造流式响应
    def generate():
        full_text = ''
        error_msg = None
        try:
            yield _sse_event('start', {
                'daily_limit': captured['daily_limit'],
                'transcript': captured.get('display_transcript') or captured.get('transcript') or None,
                'user_message_id': captured.get('user_msg_id'),
            })

            # 重新查询 model_cfg（避免跨 session 失效）
            db2 = SessionLocal()
            try:
                model_cfg2 = (
                    db2.query(AIModelConfig).filter_by(id=captured['model_cfg_id']).first()
                    if captured['model_cfg_id'] else None
                )
                for chunk in _stream_assistant_response(captured['model_key'], model_cfg2, captured['llm_messages'], extra_payload=captured.get('extra_payload')):
                    full_text += chunk
                    yield _sse_event('delta', {'content': chunk})
            finally:
                db2.close()
        except Exception as e:
            error_msg = str(e)
            logger.exception(f'流式调用失败: {e}')

        # 保存 assistant 回复（即使中途出错，已生成的部分仍保存）
        saved_msg_id = None
        db3 = SessionLocal()
        try:
            conv3 = db3.query(Conversation).filter_by(id=captured['conv_id']).first()
            if conv3:
                if full_text:
                    asst_msg = Message(
                        conversation_id=conv3.id, role='assistant', content=full_text,
                    )
                    db3.add(asst_msg)
                    db3.flush()
                    saved_msg_id = asst_msg.id
                if error_msg and not full_text:
                    err_msg = Message(
                        conversation_id=conv3.id, role='system',
                        content=f'抱歉，AI 回复出错了：{error_msg}',
                    )
                    db3.add(err_msg)
                conv3.last_active_at = datetime.now()
                db3.commit()
        except Exception:
            logger.exception('保存 assistant 消息失败')
            db3.rollback()
        finally:
            db3.close()

        if error_msg:
            yield _sse_event('error', {'message': error_msg})
        yield _sse_event('done', {'full_text_length': len(full_text), 'message_id': saved_msg_id})

    return Response(generate(), mimetype='text/event-stream', headers={
        'Cache-Control': 'no-cache',
        'X-Accel-Buffering': 'no',  # nginx 不缓冲
    })


@app.route('/api/chat/<string:task_id>/end', methods=['POST'])
def api_chat_end(task_id):
    """sendBeacon 兜底端点：浏览器关闭时调用，更新 last_active_at"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err or not conv:
            return ('', 204)
        conv.last_active_at = datetime.now()
        db.commit()
        return ('', 204)
    finally:
        db.close()


@app.route('/api/chat/<string:task_id>/rate', methods=['POST'])
def api_chat_rate(task_id):
    """学生为某条 AI 回复评分（👍 / 👎）。同一条消息可重复点：相同则取消，不同则切换"""
    db = SessionLocal()
    try:
        assistant, conv, err = _resolve_conversation(db, task_id)
        if err:
            return jsonify(err), err['code']

        body = request.get_json(silent=True) or {}
        message_id = body.get('message_id')
        rating = body.get('rating')  # 1 / -1
        if not message_id or rating not in (1, -1):
            return jsonify({'code': 400, 'message': '参数错误'}), 400

        msg = db.query(Message).filter_by(id=int(message_id), conversation_id=conv.id).first()
        if not msg:
            return jsonify({'code': 404, 'message': '消息不存在'}), 404
        if msg.role != 'assistant':
            return jsonify({'code': 400, 'message': '只能给 AI 回复打分'}), 400

        # 重复点同一类型 → 取消；点不同类型 → 切换
        if msg.rating == rating:
            msg.rating = None
        else:
            msg.rating = rating
        db.commit()
        return jsonify({'code': 200, 'rating': msg.rating})
    finally:
        db.close()


# ============================================================
# OpenMentor 路由结束
# ============================================================


if __name__ == '__main__':
    # 创建必要的目录
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)
    if not os.path.exists('static/reports'):
        os.makedirs('static/reports')
    
    # 修复socket.getfqdn()的UnicodeDecodeError问题
    import socket
    original_getfqdn = socket.getfqdn
    def safe_getfqdn(name=''):
        try:
            return original_getfqdn(name)
        except UnicodeDecodeError:
            return name if name else 'localhost'
    socket.getfqdn = safe_getfqdn
    
    # 掌上讲台下线提醒：正常退出（Ctrl+C / kill / 关机）时给 App 推一条
    import atexit as _atexit
    import signal as _signal
    _OFFLINE_PUSH_STATE['serving'] = True
    _atexit.register(_push_server_offline)
    def _om_sigterm(_sig, _frm):
        raise SystemExit(0)
    try:
        _signal.signal(_signal.SIGTERM, _om_sigterm)
    except (ValueError, OSError):
        pass  # 非主线程或平台不支持时跳过

    # 启动应用
    _host, _port = '0.0.0.0', 5001
    if (os.getenv('OM_DEBUG') or '').strip().lower() in ('1', 'true', 'yes'):
        # 开发模式：自动重载 + 调试器（OM_DEBUG=1 python app.py）
        app.run(debug=True, host=_host, port=_port)
    else:
        # 生产默认：waitress（多线程、稳定挂机）；未安装则自动回退 Flask 开发服务器，保证无痛升级
        try:
            from waitress import serve
            logger.info(f'[OpenMentor] waitress 生产服务器启动 http://{_host}:{_port} (threads=16)')
            serve(app, host=_host, port=_port, threads=16)
        except ImportError:
            logger.warning('[OpenMentor] 未安装 waitress，回退 Flask 开发服务器（pip install waitress 可升级体验）')
            app.run(debug=False, host=_host, port=_port)
